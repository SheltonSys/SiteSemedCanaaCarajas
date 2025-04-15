from datetime import date
from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth.models import User, Group, Permission
from django.contrib.contenttypes.models import ContentType
from django.db.models import Sum, Count
from django.contrib.auth.decorators import login_required, user_passes_test
from django.shortcuts import redirect
from django.contrib.auth import logout
from django.views.decorators.csrf import csrf_exempt
from django.contrib.auth.models import User
from django.contrib import messages
from django.urls import reverse
from django.http import JsonResponse
from django.middleware.csrf import get_token
import csv
import locale
from datetime import datetime
from datetime import date
import chardet
import io
import mimetypes
from django.contrib.auth import authenticate, login
from .models import Module, UserModulePermission
from .forms import UserModulePermissionForm
from django.core.exceptions import PermissionDenied
from .models import Habilidade
from django.http import HttpResponse
from .models import Acompanhamento
import datetime
from datetime import date
from django.shortcuts import render, redirect
from django.contrib.auth import authenticate, login
from django.contrib import messages
from .models import CandidatoCurriculo  # Certifique-se de usar o modelo correto
import os
import pandas as pd
from reportlab.pdfgen import canvas
from .models import PDDE
from .models import Diretoria
from .models import MembroConselho
from .models import LivroCaixa
from .forms import LivroCaixaForm
from .models import EscrituraFiscal
from .forms import EscrituraFiscalForm
from .models import Relatorio
from .models import Escola, Membro, Caixa, Certidao
from .models import Legislacao
from .forms import LegislacaoForm  # Supondo que exista um formulário para legislação
from .models import InscricaoPlanetario  # Substitua pelo nome do modelo correto, se aplicável
from .models import Imagem, Conteudo, Evento, Usuario
from .forms import ImagemForm, ConteudoForm, EventoForm, UsuarioForm
from .models import Funcionario  # Supondo que tenha um modelo Funcionario
from .models import Bairro
from .models import Disciplina
from .models import Inscricao
from django.core.exceptions import ValidationError
from django.db import IntegrityError
from django.contrib.auth.views import LoginView
from .models import DiagnoseInicProfPort, DiagnoseAlunoMatematica, DiagnoseAlunoPortugues
from .models import Professor,Aluno,HabilidadeMatematica, HabilidadePortugues
from django.core.files.storage import FileSystemStorage
from django.core.files.storage import default_storage
from django.core.files.base import ContentFile
from django.contrib.auth import authenticate, login, logout
from django.contrib.auth import login as auth_login
from django.contrib.auth.decorators import login_required
from django.http import FileResponse
from django.conf import settings
from django.templatetags.static import static
from django.contrib.messages import get_messages
from .models import Student, UploadedFile
from django.contrib.auth.models import User  # Importe o modelo padrão do Django
from semedapp.models import IndicadoresTransporte
from .models import IndicadoresTransporte
from django.db.models import Count
from datetime import date
from django.shortcuts import render, get_object_or_404, redirect
from django.urls import reverse_lazy
from django.views.generic import ListView, CreateView, UpdateView, DeleteView
from .models import Funcionario, RegimentoCadastro, Regimento, Registro
#from .forms import FuncionarioForm, RegimentoCadastroForm, RegimentoForm, RegistroForm
from .models import RegimentoCadastro  # Substitua pelo nome do seu modelo
from .models import SemedAppRegimentoCadastro
from semedapp.models import RegimentoCadastro  # Substitua pelo caminho correto do modelo
from .models import HabilidadeProf
from .forms import UploadPlanilhaForm
from .models import HabilidadeProfAnosFinais
from .models import HabilidadeProfFinal  # Ajuste o modelo para anos finais
from .forms import CurriculoForm
from .forms import EscolaForm
from semedapp.forms import DiretorForm  # Certifique-se de que o formulário está configurado corretamente
from .models import Curriculo
from django.shortcuts import render, redirect
from django.contrib.auth import login, authenticate
from django.contrib.auth.models import User
from .models import CandidatoAutenticado
from django.contrib import messages
from django.db import IntegrityError
from .forms import FormDeAtualizacao
from semedapp.models import RoleModulePermission
from reportlab.lib.pagesizes import letter
from django.db import IntegrityError
import logging
from .models import CadastroEI, ConceitoLancado
from django.contrib.auth import get_user_model

logger = logging.getLogger(__name__)

from django.shortcuts import render, redirect
from django.contrib.auth import login, authenticate
from django.contrib.auth.models import User
from .models import CadastroCandidato
from django.contrib import messages
from .forms import TipoDemandaForm
from .models import TipoDemanda
from .models import Candidato

import matplotlib
matplotlib.use('Agg')  # Usa o backend de renderização para arquivos, sem interface gráfica


import qrcode
from io import BytesIO
from django.template.loader import render_to_string
from weasyprint import HTML
from django.views.generic.edit import DeleteView
from .forms import DiretorForm  # Certifique-se de ter um formulário para o modelo Diretor
from django.views.generic import DetailView
from .models import Diretor
from semedapp.models import Diretor
from django.shortcuts import get_object_or_404, redirect
from .forms import ProfessorForm

from django.contrib.auth.backends import ModelBackend
from django.contrib.auth import get_user_model

import xlwt
import openpyxl
from xhtml2pdf import pisa
from django.template.loader import get_template
from django.forms import inlineformset_factory
from .models import Formacao, Experiencia
from .forms import CurriculoForm, FormacaoForm, ExperienciaForm

from django.contrib.auth import get_user_model

User = get_user_model()




class Site1Backend(ModelBackend):
    def authenticate(self, request, username=None, password=None, **kwargs):
        User = get_user_model()
        try:
            user = User.objects.get(username=username, login_type='site1')
            if user.check_password(password):
                return user
        except User.DoesNotExist:
            return None



@login_required
def dashboard(request):
    return render(request, 'dashboard.html')
# *********************************************************************************************************************

def BASE(request):
    return render(request, 'base.html')
# *********************************************************************************************************************

@login_required
def dashboard_admin(request):
    user = request.user  # Obtém o usuário autenticado
    return render(request, 'dashboardadmin.html', {'user': user})
# *********************************************************************************************************************

# def login_view(request):
#     if request.method == 'POST':
#         username = request.POST.get('username')
#         password = request.POST.get('password')

#         user = authenticate(request, username=username, password=password)
#         if user is not None:
#             login(request, user)
#             return JsonResponse({
#                 'status': 'success',
#                 'message': f'Bem-vindo, {user.first_name}!',
#                 'redirect_url': '/dashboardadmin/'
#             })
#         else:
#             return JsonResponse({'status': 'error', 'message': 'Usuário ou senha inválidos.'}, status=401)

#     # Para GET, retorne o CSRF token
#     return JsonResponse({'csrf_token': get_token(request)})

from django.http import JsonResponse
from django.middleware.csrf import get_token
from django.contrib.auth import authenticate, login

def login_view(request):
    if request.method == 'POST':
        username = request.POST.get('username')
        password = request.POST.get('password')

        # Autenticar usuário
        user = authenticate(request, username=username, password=password)
        if user is not None:
            # Realizar login do usuário
            login(request, user)
            
            # Redirecionar de acordo com o tipo de usuário
            if user.is_superuser:
                redirect_url = '/modulo-pedagogico/'  # Painel do administrador
            elif hasattr(user, 'professorescola') and user.professorescola.escolas.exists():
                redirect_url = '/modulo-pedagogico/'  # Página do módulo pedagógico
            else:
                redirect_url = '/pagina-inicial/'  # Página padrão para outros usuários
            
            return JsonResponse({
                'status': 'success',
                'message': f'Bem-vindo, {user.first_name if user.first_name else user.username}!',
                'redirect_url': redirect_url
            })
        else:
            # Retornar erro se as credenciais forem inválidas
            return JsonResponse({'status': 'error', 'message': 'Usuário ou senha inválidos.'}, status=401)

    # Para requisições GET, retorne o token CSRF
    return JsonResponse({'csrf_token': get_token(request)})



# *********************************************************************************************************************

def definir_permissoes(request):
    return render(request, 'definir_permissoes.html')
# *********************************************************************************************************************

def vinculacoes(request):
    return render(request, 'vinculacoes.html')
# *********************************************************************************************************************

def controle_usuarios(request):
    users = User.objects.all()  # Ou o modelo correto de usuários
    context = {'users': users}
    return render(request, 'controle_usuarios.html', context)

# *********************************************************************************************************************

def gestao_escolar(request):
    return render(request, 'gestao_escolar.html')
# *********************************************************************************************************************

def transporte(request):
    return render(request, 'transporte.html')
# *********************************************************************************************************************

def financeiro(request):
    return render(request, 'financeiro.html')
# *********************************************************************************************************************

def contabilidade(request):
    return render(request, 'contabilidade.html')
# *********************************************************************************************************************

def soe(request):
    return render(request, 'soe.html')
# *********************************************************************************************************************

def pedagogico(request):
    return render(request, 'pedagogico.html')
# *********************************************************************************************************************

def legislacao(request):
    return render(request, 'legislacao.html')
# *********************************************************************************************************************

def logout_view(request):
    logout(request)  # Logs the user out and clears their session
    return redirect('login')  # Replace 'login' with the name of your login URL
# *********************************************************************************************************************

from django.contrib.auth import get_user_model

def controle_usuarios(request):
    User = get_user_model()  # Obtém o modelo correto (CustomUser)
    usuarios = User.objects.all()  # Consulta todos os usuários
    
    context = {
        'usuarios': usuarios,
    }
    return render(request, 'semedapp/controle_usuarios.html', context)

# *********************************************************************************************************************

from django.shortcuts import render, redirect
from django.contrib.auth import get_user_model
from django.contrib import messages

User = get_user_model()

# def adicionar_usuario(request):
#     if request.method == "POST":
#         first_name = request.POST.get("first_name")
#         last_name = request.POST.get("last_name")
#         email = request.POST.get("email")
#         username = request.POST.get("username")
#         password = request.POST.get("password")
#         is_superuser = request.POST.get("is_superuser") == "True"
#         is_staff = request.POST.get("is_staff") == "True"
#         is_active = request.POST.get("is_active") == "True"

#         try:
#             user = User.objects.create_user(
#                 first_name=first_name,
#                 last_name=last_name,
#                 email=email,
#                 username=username,
#                 password=password,
#                 is_superuser=is_superuser,
#                 is_staff=is_staff,
#                 is_active=is_active,
#             )
#             messages.success(request, "Usuário cadastrado com sucesso!")
#             return redirect("controle_usuarios")  # Substitua pelo nome correto da URL de listagem
#         except Exception as e:
#             messages.error(request, f"Erro ao cadastrar usuário: {e}")

#     return render(request, "adicionar_usuario.html")



# *********************************************************************************************************************

def transporte_dashboard(request):
    # Adicione lógica para retornar dados relevantes, se necessário
    return render(request, 'transporte_dashboard.html')
# *********************************************************************************************************************

def financeiro_dashboard(request):
    # Adicione lógica ou dados para a dashboard financeira, se necessário
    return render(request, 'financeiro_dashboard.html')
# *********************************************************************************************************************

# Indicadores de Transporte Escolar
def transporte_indicadores(request):
    return render(request, 'indicadores/transporte.html')
# *********************************************************************************************************************

def upload_arquivos(request):
    if request.method == 'POST':
        if 'file' not in request.FILES:
            messages.error(request, 'Nenhum arquivo foi enviado.')
            return redirect('upload_arquivos')

        csv_file = request.FILES['file']
        if not csv_file.name.endswith('.csv'):
            messages.error(request, 'O arquivo deve estar no formato CSV.')
            return redirect('upload_arquivos')

        try:
            # Decodifica o arquivo CSV
            file_data = csv_file.read().decode('utf-8')
            csv_data = csv.reader(file_data.splitlines(), delimiter=';')
            
            # Ignorar o cabeçalho
            header = next(csv_data, None)

            for row in csv_data:
                try:
                    # Converte o formato da data
                    data_nascimento = datetime.strptime(row[3], '%d/%m/%Y').strftime('%Y-%m-%d')

                    # Cria o objeto no banco de dados
                    IndicadoresTransporte.objects.create(
                        numero_matricula=row[0],
                        nome_completo=row[1],
                        escola=row[2],
                        data_nascimento=data_nascimento,
                        sexo=row[4],
                        turno=row[5],
                        telefone=row[6],
                        nivel_escolaridade=row[7],
                        endereco=row[8],
                        bairro=row[9],
                        cep=row[10],
                        codigo_rota=row[11],
                        email_responsavel=row[12],
                        cpf_responsavel=row[13],
                        nome_responsavel=row[14],
                        parentesco=row[15],
                        telefone_responsavel=row[16],
                        responsavel_legal=True if row[17].strip().upper() == 'SIM' else False,
                    )
                except Exception as e:
                    # Log de erro por linha
                    print(f"Erro ao processar a linha: {row}. {str(e)}")
                    messages.warning(request, f'Erro ao salvar a linha: {row}. {str(e)}')

            messages.success(request, 'Arquivo processado e dados salvos com sucesso.')
            return redirect('upload_arquivos')
        except Exception as e:
            print(f"Erro ao processar o arquivo: {str(e)}")  # Log para diagnóstico
            messages.error(request, f'Ocorreu um erro ao processar o arquivo: {str(e)}')
            return redirect('upload_arquivos')
    else:
        return render(request, 'upload_arquivos.html')



# *********************************************************************************************************************

# Gestão de Dados
def gestao_dados(request):
    # Exemplo de dados fictícios
    dados = {
        'total_escolas': 50,
        'total_alunos': 1000,
        'alunos_transportados': 300,
    }
    return render(request, 'indicadores/gestao.html', {'dados': dados})
# *********************************************************************************************************************

# Diagnosis
def diagnosis(request):
    # Dados fictícios ou lógica para a página
    dados_diagnosis = {
        'total_escolas_avaliadas': 20,
        'total_alunos_avaliados': 300,
        'medias_gerais': 7.5,
    }
    return render(request, 'pedagogico/diagnosis.html', {'dados': dados_diagnosis})
# *********************************************************************************************************************

# Orientação Educacional
def orientacao_educacional(request):
    return render(request, 'pedagogico/orientacao.html')
# *********************************************************************************************************************

# Legislação
def legislacao(request):
    # Dados fictícios ou lógica para a página
    legislacao_dados = {
        'total_leis': 15,
        'total_normas': 30,
    }
    return render(request, 'setor_pedagogico/legislacao.html', {'dados': legislacao_dados})
# *********************************************************************************************************************

# Escolas
def escolas(request):
    dados_escolas = {
        'total_escolas': 10,
        'escolas_ativas': 8,
        'escolas_inativas': 2,
    }
    return render(request, 'contabilidade/escolas.html', {'dados': dados_escolas})
# *********************************************************************************************************************

# Conselho
def conselho(request):
    dados_conselho = {
        'total_conselheiros': 5,
        'reunioes_ano': 12,
    }
    return render(request, 'contabilidade/conselho.html', {'dados': dados_conselho})
# *********************************************************************************************************************

# Contas
def contas(request):
    dados_contas = {
        'saldo_total': 50000,
        'despesas_mes': 15000,
        'receitas_mes': 20000,
    }
    return render(request, 'contabilidade/contas.html', {'dados': dados_contas})
# *********************************************************************************************************************

def dashboardadmin(request):
    modules = [
        {"name": "Gestão Escolar", "url": reverse('gestao_escolar'), "icon": "fas fa-school", "color": "primary"},
        {"name": "Transporte Escolar", "url": reverse('transporte_dashboard'), "icon": "fas fa-bus", "color": "success"},
        {"name": "Financeiro", "url": reverse('financeiro_dashboard'), "icon": "fas fa-dollar-sign", "color": "warning"},
        {"name": "Indicadores", "url": reverse('indicadores'), "icon": "fas fa-chart-bar", "color": "info"},
        {"name": "Pedagógico", "url": reverse('pedagogico'), "icon": "fas fa-book", "color": "secondary"},
    ]
    return render(request, 'dashboardadmin.html', {'modules': modules})
# ********************************************************************************************************************

def transporte_escolar(request):
    try:
        # Consulta os dados do banco de dados
        total_alunos = IndicadoresTransporte.objects.count()
        total_escolas = IndicadoresTransporte.objects.values('escola').distinct().count()
        alunos_ensino_medio = IndicadoresTransporte.objects.filter(
            nivel_escolaridade__icontains='Ensino Médio'
        ).count()
        alunos_ensino_fundamental = IndicadoresTransporte.objects.filter(
            nivel_escolaridade__icontains='Ensino Fundamental'
        ).count()

        # Prepara os dados para o gráfico de barras
        alunos_por_escola = (
            IndicadoresTransporte.objects
            .values('escola')
            .annotate(total=Count('id'))
            .order_by('-total')
        )

        escolas = [item['escola'] for item in alunos_por_escola]
        alunos = [item['total'] for item in alunos_por_escola]

        # Passa os dados para o template
        context = {
            'total_alunos': total_alunos,
            'total_escolas': total_escolas,
            'alunos_ensino_medio': alunos_ensino_medio,
            'alunos_ensino_fundamental': alunos_ensino_fundamental,
            'escolas': escolas,
            'alunos': alunos,
        }

        return render(request, 'indicadores/transporte.html', context)

    except Exception as e:
        # Em caso de erro, exibe uma página de erro ou mensagem
        return render(request, 'indicadores/erro.html', {'mensagem_erro': str(e)})


# *********************************************************************************************************************

# def upload_arquivos(request):
#     # Substitua o conteúdo abaixo com a lógica da sua visão
#     return render(request, 'indicadores/upload.html')
# *********************************************************************************************************************

def gestao_dados(request):
    # Substitua o conteúdo abaixo com a lógica necessária
    return render(request, 'indicadores/gestao.html')
# *********************************************************************************************************************

def diagnosis(request):
    # Lógica ou dados para a página de Diagnosis
    return render(request, 'diagnosis.html')
# *********************************************************************************************************************

def orientacao_educacional(request):
    # Adicione lógica para Orientação Educacional, se necessário
    return render(request, 'orientacao_educacional.html')
# *********************************************************************************************************************

def escolas(request):
    # Adicione lógica específica, se necessário
    return render(request, 'escolas.html')
# *********************************************************************************************************************

def conselho(request):
    # Adicione lógica específica para o conselho, se necessário
    return render(request, 'conselho.html')
# *********************************************************************************************************************

def contas(request):
    # Adicione lógica específica para contas, se necessário
    return render(request, 'contas.html')
# *********************************************************************************************************************

def site_curriculos(request):
    return render(request, 'site_curriculos.html')  # Use o template correto
# *********************************************************************************************************************

def cadastro_curriculo(request):
    return render(request, 'cadastro_curriculo.html')  # Certifique-se de que o template existe
# *********************************************************************************************************************

def gestao_curriculo(request):
    return render(request, 'gestao_curriculo.html')  # Certifique-se de criar o template.
# *********************************************************************************************************************

def gestao_conteudo(request):
    return render(request, 'gestao_conteudo.html')
# *********************************************************************************************************************

def index(request):
    return render(request, 'index.html')
# *********************************************************************************************************************

def cadastro_imagens(request):
    return render(request, 'cadastro_imagens.html')
# *********************************************************************************************************************

def cadastro_conteudo(request):
    return render(request, 'cadastro_conteudo.html')
# *********************************************************************************************************************

def gestao_noticias(request):
    return render(request, 'gestao_noticias.html')
# *********************************************************************************************************************

def gestao_eventos(request):
    return render(request, 'gestao_eventos.html')
# *********************************************************************************************************************

def configuracoes_site(request):
    return render(request, 'configuracoes_site.html')
# *********************************************************************************************************************

def estatisticas_site(request):
    return render(request, 'estatisticas_site.html')
# *********************************************************************************************************************

def gestao_usuarios_site(request):
    return render(request, 'gestao_usuarios_site.html')
# *********************************************************************************************************************

def gestao_processos(request):
    # Lógica da view (por exemplo, buscar processos do banco de dados)
    return render(request, 'gestao_processos.html')
# *********************************************************************************************************************

def relatorios(request):
    # Lógica da view para exibir os relatórios
    return render(request, 'relatorios.html')
# *********************************************************************************************************************

def setor_pedagogico(request):
    return render(request, 'setor_pedagogico/index.html')
# *********************************************************************************************************************

def gestao_habilidades(request):
    return render(request, 'setor_pedagogico/gestao_habilidades.html')
# *********************************************************************************************************************

def indicadores_educacionais(request):
    return render(request, 'setor_pedagogico/indicadores.html')
# *********************************************************************************************************************

def relatorios_pedagogicos(request):
    return render(request, 'setor_pedagogico/relatorios.html')
# *********************************************************************************************************************

def diagnosis_prof_portugues(request):
    return render(request, 'diagnosis/prof/portugues.html')
# *********************************************************************************************************************

def diagnosis_prof_matematica(request):
    return render(request, 'diagnosis/prof/matematica.html')
# *********************************************************************************************************************

def diagnosis_aluno_portugues(request):
    return render(request, 'diagnosis/aluno/portugues.html')
# *********************************************************************************************************************

def diagnosis_aluno_matematica(request):
    return render(request, 'diagnosis/aluno/matematica.html')
# *********************************************************************************************************************

def diagnosis_eja_site(request):
    return render(request, 'diagnosis/eja/site.html')
# *********************************************************************************************************************

def diagnosis_eja_adm(request):
    return render(request, 'diagnosis/eja/adm.html')
# *********************************************************************************************************************

def admin_eja_noticias(request):
    return render(request, 'diagnosis/eja/admin_noticias.html')
# *********************************************************************************************************************

def admin_eja_eventos(request):
    return render(request, 'diagnosis/eja/admin_eventos.html')
# *********************************************************************************************************************

def admin_eja_relatorios(request):
    return render(request, 'diagnosis/eja/admin_relatorios.html')
# *********************************************************************************************************************
def admin_site_eja(request):
    return render(request, 'diagnosis/eja/admin_dashboard.html')
# *********************************************************************************************************************

def diagnosis_emcceja(request):
    return render(request, 'diagnosis/emcceja.html')
# *********************************************************************************************************************

def indicadores_setor_pedagogico(request):
    return render(request, 'setor_pedagogico/indicadores.html')
# *********************************************************************************************************************

@login_required
def gerenciar_permissoes(request):
    users = User.objects.all()
    permissions = Permission.objects.all()
    groups = Group.objects.all()
    return render(request, 'configuracoes/gerenciar_permissoes.html', {
        'users': users,
        'permissions': permissions,
        'groups': groups,
    })

# *********************************************************************************************************************

@login_required
def add_user(request):
    if request.method == "POST":
        first_name = request.POST.get('first_name')
        email = request.POST.get('email')
        user_type = request.POST.get('user_type')

        # Criar usuário
        user = User.objects.create_user(
            username=email, 
            email=email, 
            first_name=first_name
        )

        # Atribuir tipo
        if user_type == "admin":
            user.is_superuser = True
            user.is_staff = True
        elif user_type == "staff":
            user.is_staff = True

        user.save()
        messages.success(request, "Usuário adicionado com sucesso!")
        return redirect('gerenciar_permissoes')

    return redirect('gerenciar_permissoes')
# *********************************************************************************************************************

from django.utils.timezone import now

@login_required
def edit_permissions(request, user_id):
    user = get_object_or_404(User, id=user_id)

    if request.method == "POST":
        permissions = request.POST.getlist('permissions')
        user.user_permissions.set(Permission.objects.filter(id__in=permissions))

        groups = request.POST.getlist('groups')
        user.groups.set(Group.objects.filter(id__in=groups))

        group_permissions = request.POST.getlist('group_permissions')
        group_permission_map = {}
        for group_permission in group_permissions:
            group_id, permission_id = group_permission.split('-')
            group_permission_map.setdefault(group_id, []).append(permission_id)

        for group_id, permission_ids in group_permission_map.items():
            group = Group.objects.get(id=group_id)
            group.permissions.set(Permission.objects.filter(id__in=permission_ids))

        # ✅ Correção do erro
        if not user.date_joined:
            user.date_joined = now()

        user.save()
        messages.success(request, "Permissões e grupos atualizados com sucesso.")
        return redirect('gerenciar_permissoes')

    return redirect('gerenciar_permissoes')

# *********************************************************************************************************************

@login_required
def delete_user(request, user_id):
    user = get_object_or_404(User, id=user_id)
    user.delete()
    messages.success(request, "Usuário excluído com sucesso!")
    return redirect('gerenciar_permissoes')
# *********************************************************************************************************************

def manage_permissions(request):
    if request.method == "POST":
        user_id = request.POST.get("user_id")
        user = get_object_or_404(User, id=user_id)

        group_ids = request.POST.getlist("groups")
        user.groups.set(Group.objects.filter(id__in=group_ids))

        permission_ids = request.POST.getlist("permissions")
        user.user_permissions.set(Permission.objects.filter(id__in=permission_ids))

        for group in Group.objects.all():
            group_permission_ids = request.POST.getlist(f"group_permissions_{group.id}")
            group.permissions.set(Permission.objects.filter(id__in=group_permission_ids))

        messages.success(request, "Permissões atualizadas!")
        return redirect("gerenciar_permissoes")
# *********************************************************************************************************************

def add_group(request):
    if request.method == 'POST':
        group_name = request.POST.get('group_name')
        if group_name:
            group, created = Group.objects.get_or_create(name=group_name)
            if created:
                messages.success(request, f'Grupo "{group_name}" adicionado com sucesso.')
            else:
                messages.warning(request, f'O grupo "{group_name}" já existe.')
        else:
            messages.error(request, 'O nome do grupo não pode estar vazio.')
    return redirect('gerenciar_permissoes')
# *********************************************************************************************************************

def add_permission(request):
    if request.method == 'POST':
        permission_name = request.POST.get('permission_name')
        permission_codename = request.POST.get('permission_codename')
        app_label = request.POST.get('app_label')
        model_name = request.POST.get('model_name')

        if permission_name and permission_codename and app_label and model_name:
            content_type, created = ContentType.objects.get_or_create(
                app_label=app_label, model=model_name
            )
            Permission.objects.create(
                name=permission_name,
                codename=permission_codename,
                content_type=content_type,
            )
            messages.success(request, f'Permissão "{permission_name}" adicionada com sucesso.')
        else:
            messages.error(request, 'Todos os campos são obrigatórios.')

    return redirect('gerenciar_permissoes')
# *********************************************************************************************************************

@login_required
def dashboard_view(request):
    user_modules = UserModulePermission.objects.filter(user=request.user)
    return render(request, 'dashboard.html', {'user_modules': user_modules})
# *********************************************************************************************************************

def check_module_permission(user, module_name):
    if not UserModulePermission.objects.filter(user=user, module=module_name).exists():
        raise PermissionDenied
# *********************************************************************************************************************

@login_required
def configuracao_view(request):
    if not check_module_permission(request.user, 'configuracao'):
        return redirect('no_access')  # Redireciona para página de acesso negado
    return render(request, 'configuracao.html')
# *********************************************************************************************************************
# Somente administradores podem configurar acessos
def is_admin(user):
    return user.is_staff or user.is_superuser

@user_passes_test(is_admin)
@login_required
def configure_access_view(request):
    if request.method == 'POST':
        form = UserModulePermissionForm(request.POST)
        if form.is_valid():
            form.save()
            return redirect('configure_access')
    else:
        form = UserModulePermissionForm()

    permissions = UserModulePermission.objects.all()
    return render(request, 'configure_access.html', {'form': form, 'permissions': permissions})
# *********************************************************************************************************************

def check_module_permission(user, module):
    return UserModulePermission.objects.filter(user=user, module=module).exists()
# *********************************************************************************************************************

# Exemplo de uso em uma view
@login_required
def some_module_view(request):
    if not check_module_permission(request.user, 'pedagogico'):
        return redirect('no_access')  # Redireciona se o usuário não tiver acesso
    return render(request, 'pedagogico.html')
# *********************************************************************************************************************
@login_required
def indicadores_view(request):
    if not check_module_permission(request.user, 'indicadores'):
        return redirect('no_access')
    return render(request, 'indicadores.html')

# Repita para cada módulo
# *********************************************************************************************************************

@login_required
def access_user_modules(request, user_id):
    user = get_object_or_404(User, id=user_id)
    user_modules = UserModulePermission.objects.filter(user=user)
    return render(request, 'dashboard.html', {'user': user, 'user_modules': user_modules})

# *********************************************************************************************************************

@login_required
def configure_user_modules(request, user_id):
    user = get_object_or_404(User, id=user_id)
    modules = Module.objects.all()

    if request.method == 'POST':
        selected_modules = request.POST.getlist('modules')
        # Limpa permissões antigas
        UserModulePermission.objects.filter(user=user).delete()
        # Adiciona novas permissões
        for module_id in selected_modules:
            module = get_object_or_404(Module, id=module_id)
            UserModulePermission.objects.create(user=user, module=module)
        return redirect('dashboard')  # Redireciona ao dashboard ou outra página

    user_modules = user.module_permissions.values_list('module_id', flat=True)
    return render(request, 'configure_modules.html', {
        'user': user,
        'modules': modules,
        'user_modules': user_modules,
    })

@login_required
def curriculos_view(request):
    # Substitua 'curriculos.html' pelo template correto para esta página
    return render(request, 'curriculos.html', {})
# *********************************************************************************************************************

@login_required
def modulo_view(request, slug):
    module = get_object_or_404(Module, slug=slug)
    return render(request, 'modulo.html', {'module': module})
# *********************************************************************************************************************

def update_habilidades(request):
    if request.method == 'POST':
        aluno_id = request.POST.get('aluno')
        habilidade_nome = request.POST.get('habilidade')
        # Salvar ou atualizar no banco de dados
        Habilidade.objects.create(aluno_id=aluno_id, nome=habilidade_nome)
        return redirect('gestao_habilidades')  # Redirecionar para a página atual
    return HttpResponse("Método não permitido.", status=405)
# *********************************************************************************************************************

def download_relatorio(request, tipo_relatorio):
    # Lógica para gerar ou buscar o relatório com base no tipo
    if tipo_relatorio == "aprovacao":
        # Simule o arquivo do relatório
        response = HttpResponse(content_type='application/pdf')
        response['Content-Disposition'] = 'attachment; filename="taxa_aprovacao.pdf"'
        response.write("Conteúdo do relatório de aprovação")  # Substitua por geração de PDF
    elif tipo_relatorio == "frequencia":
        response = HttpResponse(content_type='application/pdf')
        response['Content-Disposition'] = 'attachment; filename="frequencia.pdf"'
        response.write("Conteúdo do relatório de frequência")  # Substitua por geração de PDF
    elif tipo_relatorio == "desempenho":
        response = HttpResponse(content_type='application/pdf')
        response['Content-Disposition'] = 'attachment; filename="desempenho.pdf"'
        response.write("Conteúdo do relatório de desempenho")  # Substitua por geração de PDF
    else:
        return HttpResponse("Relatório não encontrado", status=404)

    return response
# *********************************************************************************************************************

def agenda_educacional(request):
    return render(request, 'setor_pedagogico/agenda_educacional.html')
# *********************************************************************************************************************

def acompanhamento_individual(request):
    return render(request, 'setor_pedagogico/acompanhamento_individual.html')
# *********************************************************************************************************************

def registrar_acompanhamento(request):
    if request.method == 'POST':
        aluno_id = request.POST.get('aluno')
        descricao = request.POST.get('descricao')

        if aluno_id and descricao:
            Acompanhamento.objects.create(aluno_id=aluno_id, descricao=descricao)
            messages.success(request, 'Acompanhamento registrado com sucesso.')
            return redirect('orientacao_educacional')  # Substitua pelo nome correto da view de redirecionamento
        else:
            messages.error(request, 'Preencha todos os campos.')

    return render(request, 'orientacao_educacional.html')
# *********************************************************************************************************************

def pdde_view(request):
    anos = [2022, 2023, 2024]  # Example data
    pdde_data = [
        {'ano': 2023, 'escola': 'Escola A', 'status': 'aprovado'},
        {'ano': 2024, 'escola': 'Escola B', 'status': 'pendente'},
    ]
    return render(request, 'contabilidade/pdde.html', {'anos': anos, 'pdde_data': pdde_data})
# *********************************************************************************************************************

from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from .models import Diretoria, EscolaPdde
from datetime import datetime

def conselho_diretoria(request):
    membros = Diretoria.objects.all().order_by("cargo")
    escolas = EscolaPdde.objects.all().order_by("nome")
    conselhos = EscolaPdde.objects.values_list("nome_conselho", flat=True).distinct()

    if request.method == "POST":
        try:
            nome = request.POST.get("nome")
            cargo = request.POST.get("cargo")
            endereco = request.POST.get("endereco")
            bairro = request.POST.get("bairro")
            telefone = request.POST.get("telefone")
            email = request.POST.get("email")
            cep = request.POST.get("cep")
            cpf = request.POST.get("cpf")
            vencimento = request.POST.get("vencimento")
            escola_id = request.POST.get("escola")
            conselho_nome = request.POST.get("conselho")

            if not all([nome, cargo, endereco, bairro, telefone, email, cep, cpf, vencimento, escola_id, conselho_nome]):
                raise ValueError("Todos os campos são obrigatórios.")

            vencimento_date = datetime.strptime(vencimento, "%Y-%m-%d").date()
            escola_obj = get_object_or_404(EscolaPdde, id=escola_id)

            Diretoria.objects.create(
                nome=nome,
                cargo=cargo,
                endereco=endereco,
                bairro=bairro,
                telefone=telefone,
                email=email,
                cep=cep,
                cpf=cpf,
                vencimento=vencimento_date,
                escola=escola_obj,
                conselho=conselho_nome
            )
            messages.success(request, "✅ Membro cadastrado com sucesso!")

        except Exception as e:
            print(f"Erro ao cadastrar membro: {e}")
            messages.error(request, "❌ Erro ao cadastrar membro. Tente novamente.")

        return redirect("conselho_diretoria")

    return render(request, "contabilidade/conselho/diretoria.html", {
        "diretoria_data": membros,
        "escolas": escolas,
        "conselhos": conselhos
    })






# *********************************************************************************************************************

def editar_diretoria(request, id):
    membro = get_object_or_404(Diretoria, id=id)
    escolas = EscolaPdde.objects.all()

    if request.method == "POST":
        try:
            membro.nome = request.POST.get("nome")
            membro.cargo = request.POST.get("cargo")
            membro.endereco = request.POST.get("endereco")
            membro.bairro = request.POST.get("bairro")
            membro.telefone = request.POST.get("telefone")
            membro.email = request.POST.get("email")
            membro.cep = request.POST.get("cep")
            membro.cpf = request.POST.get("cpf")
            vencimento = request.POST.get("vencimento")
            escola_id = request.POST.get("escola")
            conselho_nome = request.POST.get("conselho")

            if not all([membro.nome, membro.cargo, membro.endereco, membro.bairro, membro.telefone, membro.email, membro.cep, membro.cpf, vencimento, escola_id, conselho_nome]):
                messages.error(request, "Todos os campos são obrigatórios!")
                return redirect("editar_diretoria", id=id)

            membro.vencimento = datetime.strptime(vencimento, "%Y-%m-%d").date()
            membro.escola_id = escola_id
            membro.conselho = conselho_nome
            membro.save()
            messages.success(request, "✅ Membro atualizado com sucesso!")
        except Exception as e:
            print(f"Erro ao editar membro: {e}")
            messages.error(request, "❌ Erro ao atualizar o membro.")

        return redirect("conselho_diretoria")

    return render(request, "contabilidade/conselho/editar_diretoria.html", {
        "membro": membro,
        "escolas": escolas
    })



# *********************************************************************************************************************

def excluir_diretoria(request, id):
    membro = get_object_or_404(Diretoria, id=id)
    if request.method == "POST":
        membro.delete()
        messages.success(request, "✅ Membro excluído com sucesso!")
        return redirect('conselho_diretoria')

# *********************************************************************************************************************

def download_page(request):
    # Diretório onde os manuais estão armazenados
    manuals_dir = os.path.join('media', 'manuals')
    
    # Listar todos os arquivos no diretório de manuais
    if os.path.exists(manuals_dir):
        manuals = os.listdir(manuals_dir)
    else:
        manuals = []

    return render(request, 'downloads/manuals_list.html', {'manuals': manuals})
# *********************************************************************************************************************

from django.shortcuts import render
from .models import TermoDoacao  # ajuste o caminho se necessário

def pdde(request):
    termo = TermoDoacao.objects.last()  # ou use algum filtro mais específico

    context = {
        'termo': termo,
        'range_10': range(10),
        'range_3': range(3),
        'range_5': range(5),
    }

    return render(request, 'contabilidade/pdde.html', context)

# *********************************************************************************************************************

def view_pdde(request, pdde_id):
    return HttpResponse(f"Detalhes do PDDE com ID: {pdde_id}")
# *********************************************************************************************************************

def download_manual(request, file_name):
    file_path = os.path.join('media/manuals', file_name)  # Caminho relativo
    if os.path.exists(file_path):
        with open(file_path, 'rb') as fh:
            response = HttpResponse(fh.read(), content_type="application/pdf")
            response['Content-Disposition'] = f'attachment; filename={file_name}'
            return response
    else:
        return HttpResponse("Arquivo não encontrado.", status=404)
# *********************************************************************************************************************

def pdde_view(request):
    # Filtros
    ano = request.GET.get('year', None)
    status = request.GET.get('status', None)

    # Consulta com filtros
    pdde_data = PDDE.objects.all()
    if ano:
        pdde_data = pdde_data.filter(ano=ano)
    if status:
        pdde_data = pdde_data.filter(status=status)

    # Anos únicos para o filtro
    anos = PDDE.objects.values_list('ano', flat=True).distinct().order_by('-ano')

    context = {
        'pdde_data': pdde_data,
        'anos': anos,
    }
    return render(request, 'app/pdde.html', context)
# *********************************************************************************************************************

from django.contrib import messages
from django.shortcuts import render, redirect
from .models import Diretoria
from datetime import datetime

def cadastro_diretoria(request):
    if request.method == "POST":
        nome = request.POST.get("nome")
        cargo = request.POST.get("cargo")
        endereco = request.POST.get("endereco")
        bairro = request.POST.get("bairro")
        telefone = request.POST.get("telefone")
        email = request.POST.get("email")
        cep = request.POST.get("cep")
        cpf = request.POST.get("cpf")
        vencimento = request.POST.get("vencimento")

        try:
            # Validação simples de campos obrigatórios
            if not all([nome, cargo, endereco, bairro, telefone, email, cep, cpf, vencimento]):
                messages.error(request, "Todos os campos são obrigatórios!")
                return redirect("conselho_diretoria")

            # Validação da data
            try:
                data_vencimento = datetime.strptime(vencimento, "%Y-%m-%d").date()
            except ValueError:
                messages.error(request, "Data de vencimento inválida! Use o formato AAAA-MM-DD.")
                return redirect("conselho_diretoria")

            # Salvar no banco
            Diretoria.objects.create(
                nome=nome,
                cargo=cargo,
                endereco=endereco,
                bairro=bairro,
                telefone=telefone,
                email=email,
                cep=cep,
                cpf=cpf,
                vencimento=data_vencimento,
            )
            messages.success(request, "Membro cadastrado com sucesso!")
        except Exception as e:
            print(e)
            messages.error(request, "Erro ao cadastrar membro. Tente novamente.")

        return redirect("conselho_diretoria")

    return render(request, "contabilidade/conselho/diretoria.html")


# *********************************************************************************************************************

def listar_membros(request):
    """
    Exibe a lista de membros do conselho.
    """
    membros_data = MembroConselho.objects.all()
    return render(request, 'contabilidade/conselho/membros.html', {'membros_data': membros_data})
# *********************************************************************************************************************

def cadastrar_membro(request):
    escolas = EscolaPdde.objects.all().order_by("nome")

    if request.method == 'POST':
        try:
            inep = request.POST.get('inep')
            escola_nome = request.POST.get('escola')
            data_abertura = request.POST.get('data_abertura')
            data_vencimento = request.POST.get('data_vencimento')
            nome = request.POST.get('nome')
            funcao = request.POST.get('funcao')
            endereco = request.POST.get('endereco')
            bairro = request.POST.get('bairro')
            telefone = request.POST.get('telefone')
            email = request.POST.get('email')
            cep = request.POST.get('cep')
            cpf = request.POST.get('cpf')

            if not all([inep, escola_nome, data_abertura, data_vencimento, nome, funcao, endereco, bairro, telefone, email, cep, cpf]):
                messages.error(request, "❌ Todos os campos são obrigatórios.")
                return redirect('conselho_membros')

            escola_obj = EscolaPdde.objects.filter(nome=escola_nome).first()
            conselho_nome = escola_obj.nome_conselho if escola_obj else ""

            novo_membro = MembroConselho.objects.create(
                inep=inep,
                escola=escola_nome,
                conselho=conselho_nome,
                data_abertura=data_abertura,
                data_vencimento=data_vencimento,
                nome=nome,
                funcao=funcao,
                endereco=endereco,
                bairro=bairro,
                telefone=telefone,
                email=email,
                cep=cep,
                cpf=cpf,
            )

            messages.success(request, f"✅ Membro cadastrado com sucesso! Conselho: {conselho_nome}")
            return redirect('listar_membros')

        except Exception as e:
            print(f"Erro ao cadastrar membro: {e}")
            messages.error(request, f"❌ Ocorreu um erro ao cadastrar o membro: {e}")
            return redirect('conselho_membros')

    return render(request, "contabilidade/conselho/membros.html", {
        "escolas": escolas
    })


# *********************************************************************************************************************

from .models import MembroConselho

def conselho_membros(request):
    membros_data = MembroConselho.objects.all().order_by("nome")
    escolas = EscolaPdde.objects.all().order_by("nome")
    return render(request, 'contabilidade/conselho/membros.html', {
        'membros_data': membros_data,
        'escolas': escolas
    })
# *********************************************************************************************************************

from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from django.http import JsonResponse
from .models import LivroCaixa, EscolaPdde


def listar_livro_caixa(request):
    livros = LivroCaixa.objects.all().order_by('-ano_base')
    return render(request, 'contabilidade/livro_caixa/listar.html', {'livros': livros})
# *********************************************************************************************************************

def adicionar_livro_caixa(request):
    escolas = EscolaPdde.objects.all().order_by("nome")

    if request.method == "POST":
        try:
            ano_base = request.POST.get('ano_base')
            escola_id = request.POST.get('escola')
            rendimentos = float(request.POST.get('rendimentos_aplicacao') or 0)
            saldo = float(request.POST.get('saldo_anterior') or 0)
            despesas = float(request.POST.get('despesas_manutencao') or 0)

            receita_total = rendimentos + saldo
            despesa_total = despesas
            superavit_deficit = receita_total - despesa_total

            escola = get_object_or_404(EscolaPdde, id=escola_id)

            LivroCaixa.objects.create(
                ano_base=ano_base,
                escola=escola,
                conselho_escolar=escola.nome_conselho,
                cnpj=escola.cnpj,
                rendimentos_aplicacao=rendimentos,
                saldo_anterior=saldo,
                receita_total=receita_total,
                despesas_manutencao=despesas,
                despesa_total=despesa_total,
                superavit_deficit=superavit_deficit
            )
            messages.success(request, "✅ Escrituração cadastrada com sucesso!")
        except Exception as e:
            messages.error(request, f"❌ Erro ao cadastrar: {e}")

        return redirect('listar_livro_caixa')

    return render(request, 'contabilidade/adicionar_caixa.html', {'escolas': escolas})

# *********************************************************************************************************************

def editar_livro_caixa(request, id):
    livro = get_object_or_404(LivroCaixa, id=id)
    escolas = EscolaPdde.objects.all().order_by('nome')

    if request.method == "POST":
        try:
            livro.ano_base = request.POST.get('ano_base')
            escola_id = request.POST.get('escola')
            escola = get_object_or_404(EscolaPdde, id=escola_id)

            livro.escola = escola
            livro.conselho_escolar = escola.nome_conselho
            livro.cnpj = escola.cnpj
            livro.rendimentos_aplicacao = float(request.POST.get('rendimentos_aplicacao') or 0)
            livro.saldo_anterior = float(request.POST.get('saldo_anterior') or 0)
            livro.receita_total = livro.rendimentos_aplicacao + livro.saldo_anterior
            livro.despesas_manutencao = float(request.POST.get('despesas_manutencao') or 0)
            livro.despesa_total = livro.despesas_manutencao
            livro.superavit_deficit = livro.receita_total - livro.despesa_total
            livro.save()
            messages.success(request, "✅ Escrituração atualizada com sucesso!")
        except Exception as e:
            messages.error(request, f"❌ Erro ao atualizar: {e}")

        return redirect('listar_livro_caixa')

    return render(request, 'contabilidade/livro_caixa/editar.html', {
        'livro': livro,
        'escolas': escolas
    })
# *********************************************************************************************************************

def excluir_livro_caixa(request, id):
    livro = get_object_or_404(LivroCaixa, id=id)
    if request.method == "POST":
        livro.delete()
        messages.success(request, "✅ Escrituração excluída com sucesso!")
        return redirect('listar_livro_caixa')
    return render(request, 'contabilidade/livro_caixa/excluir.html', {'livro': livro})
# *********************************************************************************************************************

from django.http import JsonResponse
from .models import EscolaPdde, SemedAppEscolaPddeProgramas, Programa, Receita

def escola_info(request, escola_id):
    """
    Retorna dados da escola + programas vinculados + dados financeiros.
    """
    try:
        escola = EscolaPdde.objects.get(id=escola_id)
        vinculos = SemedAppEscolaPddeProgramas.objects.filter(escolapdde=escola).select_related('programa')
        programas = [v.programa.nome for v in vinculos]

        # Busca dados financeiros mais recentes
        receita = Receita.objects.filter(escola=escola).order_by('-data_inicio').first()

        if receita:
            rendimentos = receita.rendimento_aplicacao_custeio + receita.rendimento_aplicacao_capital
            saldo_anterior = (
                receita.saldo_anterior_custeio + receita.saldo_anterior_capital +
                receita.valor_creditado_custeio + receita.valor_creditado_capital +
                receita.recursos_proprios_custeio + receita.recursos_proprios_capital
            )
        else:
            rendimentos = 0.00
            saldo_anterior = 0.00

        data = {
            'conselho': escola.nome_conselho,
            'cnpj': escola.cnpj,
            'programas': programas,
            'rendimentos': float(rendimentos),
            'saldo_anterior': float(saldo_anterior),
        }
        return JsonResponse(data)
    except EscolaPdde.DoesNotExist:
        return JsonResponse({'error': 'Escola não encontrada.'}, status=404)



# *********************************************************************************************************************

from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from .models import MembroConselho, EscolaPdde
from datetime import datetime

def editar_membro(request, id):
    membro = get_object_or_404(MembroConselho, id=id)
    escolas = EscolaPdde.objects.all().order_by("nome")

    if request.method == "POST":
        try:
            membro.inep = request.POST.get('inep')
            membro.escola = request.POST.get('escola')
            membro.data_abertura = request.POST.get('data_abertura')
            membro.data_vencimento = request.POST.get('data_vencimento')
            membro.nome = request.POST.get('nome')
            membro.funcao = request.POST.get('funcao')
            membro.endereco = request.POST.get('endereco')
            membro.bairro = request.POST.get('bairro')
            membro.telefone = request.POST.get('telefone')
            membro.email = request.POST.get('email')
            membro.cep = request.POST.get('cep')
            membro.cpf = request.POST.get('cpf')
            membro.save()
            messages.success(request, "✅ Membro atualizado com sucesso!")
            return redirect('conselho_membros')
        except Exception as e:
            print(e)
            messages.error(request, "❌ Erro ao atualizar membro.")
            return redirect('conselho_membros')

    return render(request, 'contabilidade/conselho/editar_membro.html', {
        'membro': membro,
        'escolas': escolas
    })
# *********************************************************************************************************************

def excluir_membro(request, id):
    membro = get_object_or_404(MembroConselho, id=id)
    try:
        membro.delete()
        messages.success(request, "✅ Membro excluído com sucesso!")
    except Exception as e:
        print(e)
        messages.error(request, "❌ Erro ao excluir membro.")
    return redirect('conselho_membros')


# *********************************************************************************************************************

def adicionar_escritura_fiscal(request):
    if request.method == 'POST':
        form = EscrituraFiscalForm(request.POST)
        if form.is_valid():
            form.save()
            return redirect('listar_livro_caixa')
    else:
        form = EscrituraFiscalForm()
    return render(request, 'contabilidade/adicionar_caixa.html', {'form': form})
# *********************************************************************************************************************

def adicionar_escrituracao(request):
    if request.method == 'POST':
        # Captura os dados do formulário
        ano_base = request.POST.get('ano_base')
        conselho_escolar = request.POST.get('conselho_escolar')
        cnpj = request.POST.get('cnpj')
        rendimentos_aplicacao = request.POST.get('rendimentos_aplicacao')
        saldo_anterior = request.POST.get('saldo_anterior')
        despesas_manutencao = request.POST.get('despesas_manutencao')

        # Realiza cálculos
        receita_total = float(rendimentos_aplicacao) + float(saldo_anterior)
        despesa_total = float(despesas_manutencao)
        superavit_deficit = receita_total - despesa_total

        # Aqui você pode salvar os dados no banco de dados
        # Exemplo:
        # model_instance = LivroCaixa(
        #     ano_base=ano_base,
        #     conselho_escolar=conselho_escolar,
        #     cnpj=cnpj,
        #     rendimentos_aplicacao=rendimentos_aplicacao,
        #     saldo_anterior=saldo_anterior,
        #     receita_total=receita_total,
        #     despesas_manutencao=despesas_manutencao,
        #     despesa_total=despesa_total,
        #     superavit_deficit=superavit_deficit
        # )
        # model_instance.save()

        # Exibe uma mensagem de sucesso
        messages.success(request, 'Escrituração contábil adicionada com sucesso!')
        return redirect('contabilidade:caixa')  # Redirecione para a lista de registros ou página desejada

    # Renderiza o template do modal
    return render(request, 'contabilidade/adicionar_caixa.html')
# *********************************************************************************************************************

def listar_relatorios_pdde(request):
    # Filtros
    tipo_relatorio = request.GET.get('tipo_relatorio', '')
    ano = request.GET.get('ano', '')
    escola = request.GET.get('escola', '')

    # Query básica
    relatorios = Relatorio.objects.all()

    # Aplica filtros se disponíveis
    if tipo_relatorio:
        relatorios = relatorios.filter(tipo__icontains=tipo_relatorio)
    if ano:
        relatorios = relatorios.filter(ano=ano)
    if escola:
        relatorios = relatorios.filter(escola__icontains=escola)

    context = {
        'relatorios_data': relatorios,
    }
    return render(request, 'contabilidade/relatorios_pdde.html', context)
# *********************************************************************************************************************

def dashboard(request):
    context = {
        "total_escolas": Escola.objects.count(),
        "total_membros": Membro.objects.count(),
        "total_pdde": PDDE.objects.aggregate(total=Sum('valor'))['total'] or 0,
        "total_caixas": Caixa.objects.aggregate(total=Sum('saldo'))['total'] or 0,
        "total_certidoes": Certidao.objects.count(),
        "recent_certidoes": Certidao.objects.order_by('-data_emissao')[:5],
        "recent_movements": PDDE.objects.order_by('-data_movimentacao')[:5],
    }
    return render(request, 'dashboard/seppec_dashboard.html', context)
# *********************************************************************************************************************

def dashboard_view(request):
    # Lógica da dashboard (substitua conforme necessário)
    context = {
        'title': 'Dashboard SEPPEC',
        # Adicione outros dados ao contexto, se necessário
    }
    return render(request, 'dashboard.html', context)
# *********************************************************************************************************************

def listar_documentos(request):
    # Substitua pelo código real para buscar documentos
    documentos = []  # Exemplo: Documento.objects.all()
    return render(request, 'setor_pedagogico/documentos_list.html', {'documentos': documentos})
# *********************************************************************************************************************

def gerenciar_categorias(request):
    # Substitua pelo código real para buscar categorias
    categorias = []  # Exemplo: Categoria.objects.all()
    return render(request, 'setor_pedagogico/categorias_list.html', {'categorias': categorias})
# *********************************************************************************************************************

def relatorio_atualizacoes(request):
    # Exemplo: Carregar dados de atualizações recentes
    atualizacoes = []  # Exemplo: Atualizacao.objects.filter(data__gte=hoje - timedelta(days=30))
    total_atualizacoes = len(atualizacoes)
    return render(request, 'setor_pedagogico/relatorio_atualizacoes.html', {
        'atualizacoes': atualizacoes,
        'total_atualizacoes': total_atualizacoes,
    })
# *********************************************************************************************************************

def adicionar_legislacao(request):
    if request.method == 'POST':
        form = LegislacaoForm(request.POST)
        if form.is_valid():
            form.save()
            messages.success(request, 'Legislação adicionada com sucesso!')
            return redirect('legislacao')  # Redireciona para a página principal de legislação
    else:
        form = LegislacaoForm()

    return render(request, 'setor_pedagogico/adicionar_legislacao.html', {'form': form})
# *********************************************************************************************************************

def planetario(request):
    return render(request, 'planetario.html')
# *********************************************************************************************************************

def planetario_inscricao(request):
    if request.method == 'POST':
        nome = request.POST.get('nome')
        email = request.POST.get('email')
        telefone = request.POST.get('telefone')
        atividade = request.POST.get('atividade')
        mensagem = request.POST.get('mensagem')

        InscricaoPlanetario.objects.create(
            nome=nome,
            email=email,
            telefone=telefone,
            atividade=atividade,
            mensagem=mensagem
        )

        return render(request, 'planetario_sucesso.html')

    return render(request, 'planetario_inscricao.html')
# *********************************************************************************************************************

def galeria(request):
    imagens = Imagem.objects.all()
    return render(request, 'galeria.html', {'imagens': imagens})
# *********************************************************************************************************************

def galeria_view(request):
    imagens = Imagem.objects.all()  # Certifique-se de que o modelo `Imagem` está configurado corretamente
    return render(request, 'nome_do_template.html', {'imagens': imagens})
# *********************************************************************************************************************

# Gestão de Imagens
def gerenciar_imagens(request):
    imagens = Imagem.objects.all()
    return render(request, 'cadastro_imagens.html', {'imagens': imagens})
# *********************************************************************************************************************

# Gestão de Imagens
def adicionar_imagem(request):
    if request.method == 'POST':
        form = ImagemForm(request.POST, request.FILES)
        if form.is_valid():
            form.save()
            return redirect('galeria')  # Redirecione para a galeria após salvar
    else:
        form = ImagemForm()
    return render(request, 'adicionar_imagem.html', {'form': form})
# *********************************************************************************************************************

def lista_imagens(request):
    imagens = Imagem.objects.all().order_by('-data_criacao')
    return render(request, 'lista_imagens.html', {'imagens': imagens})
# *********************************************************************************************************************
# Gestão de Conteúdo
def gerenciar_conteudo(request):
    conteudos = Conteudo.objects.all()
    return render(request, 'cadastro_conteudo.html', {'conteudos': conteudos})
# *********************************************************************************************************************
# Gestão de Conteúdo
def adicionar_conteudo(request):
    if request.method == 'POST':
        form = ConteudoForm(request.POST, request.FILES)
        if form.is_valid():
            form.save()
            return redirect('gerenciar_conteudo')
    else:
        form = ConteudoForm()
    return render(request, 'adicionar_conteudo.html', {'form': form})

# *********************************************************************************************************************
# Gestão de Eventos
def gerenciar_eventos(request):
    eventos = Evento.objects.all()
    return render(request, 'gestao_eventos.html', {'eventos': eventos})
# *********************************************************************************************************************

def adicionar_evento(request):
    if request.method == 'POST':
        form = EventoForm(request.POST)
        if form.is_valid():
            form.save()
            return redirect('gerenciar_eventos')
    else:
        form = EventoForm()
    return render(request, 'adicionar_evento.html', {'form': form})
# *********************************************************************************************************************
# Gestão de Usuários
def gerenciar_usuarios(request):
    usuarios = Usuario.objects.all()
    return render(request, 'gestao_usuarios_site.html', {'usuarios': usuarios})

# *********************************************************************************************************************

def adicionar_usuario(request):
    if request.method == "POST":
        first_name = request.POST.get("first_name")
        last_name = request.POST.get("last_name")
        email = request.POST.get("email")
        username = request.POST.get("username")
        password = request.POST.get("password")
        is_superuser = request.POST.get("is_superuser") == "True"
        is_staff = request.POST.get("is_staff") == "True"
        is_active = request.POST.get("is_active") == "True"

        try:
            user = User.objects.create_user(
                first_name=first_name,
                last_name=last_name,
                email=email,
                username=username,
                password=password,
                is_superuser=is_superuser,
                is_staff=is_staff,
                is_active=is_active,
            )
            messages.success(request, "Usuário cadastrado com sucesso!")
            return redirect("controle_usuarios")  # Substitua pelo nome correto da URL de listagem
        except Exception as e:
            messages.error(request, f"Erro ao cadastrar usuário: {e}")

    return render(request, "adicionar_usuario.html")
# *********************************************************************************************************************

def site_semed(request):
    return render(request, 'sitesemed.html')
# *********************************************************************************************************************

def cadastro_imagens(request):
    return render(request, 'cadastro_imagens.html')
# *********************************************************************************************************************

def cadastro_conteudo(request):
    return render(request, 'cadastro_conteudo.html')
# *********************************************************************************************************************

def gestao_conteudo(request):
    return render(request, 'gestao_conteudo.html')
# *********************************************************************************************************************

def gestao_noticias(request):
    return render(request, 'gestao_noticias.html')
# *********************************************************************************************************************

def gestao_eventos(request):
    return render(request, 'gestao_eventos.html')
# *********************************************************************************************************************

def configuracoes_site(request):
    return render(request, 'configuracoes_site.html')
# *********************************************************************************************************************

def estatisticas_site(request):
    return render(request, 'estatisticas_site.html')
# *********************************************************************************************************************

def gestao_usuarios_site(request):
    return render(request, 'gestao_usuarios_site.html')
# *********************************************************************************************************************

# View para a apresentação
def apresentacao(request):
    return render(request, 'apresentacao.html')
# *********************************************************************************************************************

# View para as atribuições
def atribuicoes(request):
    return render(request, 'atribuicoes.html')
# *********************************************************************************************************************

# View para os contatos
def contatos(request):
    return render(request, 'contatos.html')
# *********************************************************************************************************************

def organograma(request):
    """
    Renderiza a página do Organograma.
    """
    return render(request, 'organograma.html')
# *********************************************************************************************************************

def escolas(request):
    return render(request, 'escolas.html')
# *********************************************************************************************************************

def escola(request):
    return render(request, 'escola.html')
# *********************************************************************************************************************

def centro_viver_conviver(request):
    return render(request, 'centro_viver_conviver.html')
# *********************************************************************************************************************

def centro_formacao(request):
    return render(request, 'centro_formacao.html')
# *********************************************************************************************************************

def diretoria_ensino_superior(request):
    return render(request, 'diretoria_ensino_superior.html')
# *********************************************************************************************************************

def centro_midias_educacionais(request):
    return render(request, 'centro_midias_educacionais.html')
# *********************************************************************************************************************

def conselho_municipal_educacao(request):
    return render(request, 'conselho_municipal_educacao.html')
# *********************************************************************************************************************

def cacs_fundeb(request):
    return render(request, 'cacs_fundeb.html')
# *********************************************************************************************************************

def conselho_alimentacao_escolar(request):
    return render(request, 'conselho_alimentacao_escolar.html')
# *********************************************************************************************************************

def forum_municipal_educacao(request):
    return render(request, 'forum_municipal_educacao.html')
# *********************************************************************************************************************

def fale_conosco(request):
    return render(request, 'faleconosco.html')
# *********************************************************************************************************************

def repositorio(request):
    return render(request, 'repositorio.html')
# *********************************************************************************************************************

def curriculo_ensino(request):
    return render(request, 'curriculo_ensino.html')
# *********************************************************************************************************************

def programas_projetos(request):
    return render(request, 'programas_projetos.html')
# *********************************************************************************************************************

def calendario_escolar(request):
    return render(request, 'calendario_escolar.html')
# *********************************************************************************************************************

def portal_aluno(request):
    return render(request, 'portal_aluno.html')
# *********************************************************************************************************************

def sige(request):
    return render(request, 'sige.html')
# *********************************************************************************************************************

def consulta_publica(request):
    if request.method == 'POST':
        titulo = request.POST.get('titulo')
        capitulo = request.POST.get('capitulo')
        tipo_alteracao = request.POST.get('tipo_alteracao')
        justificativa = request.POST.get('justificativa')
        nome_completo = request.POST.get('nome_completo')
        email = request.POST.get('email')
        cpf = request.POST.get('cpf')
        telefone = request.POST.get('telefone')
        cargo = request.POST.get('cargo')
        lotacao = request.POST.get('lotacao')
        observacoes_adicionais = request.POST.get('observacoes_adicionais')

        # Criação do objeto e salvamento no banco
        try:
            RegimentoCadastro.objects.create(
                titulo=titulo,
                capitulo=capitulo,
                tipo_alteracao=tipo_alteracao,
                justificativa=justificativa,
                nome_completo=nome_completo,
                email=email,
                cpf=cpf,
                telefone=telefone,
                cargo=cargo,
                lotacao=lotacao,
                observacoes_adicionais=observacoes_adicionais,
            )
            messages.success(request, "Sua sugestão foi enviada com sucesso!")
        except Exception as e:
            messages.error(request, f"Erro ao salvar a sugestão: {e}")

        return redirect('consulta_publica')  # Redirecione para evitar duplicação ao recarregar a página

    return render(request, 'consulta_publica.html')
# *********************************************************************************************************************

def jornal_escolar(request):
    return render(request, 'jornal_escolar.html')
# *********************************************************************************************************************

def noticia_unifesspa(request):
    return render(request, 'noticias/noticia_unifesspa.html')
# *********************************************************************************************************************

def noticia_eja(request):
    return render(request, 'noticias/noticia_eja.html')
# *********************************************************************************************************************

def noticia_planetario(request):
    return render(request, 'noticias/noticia_planetario.html')
# *********************************************************************************************************************

def programas_projetos(request):
    return render(request, 'programas_projetos.html')
# *********************************************************************************************************************

def portal_aluno(request):
    return render(request, 'portal_aluno.html')
# *********************************************************************************************************************

def emeb_gercino_correa(request):
    return render(request, 'escolas/emeb_gercino_correa.html')
# *********************************************************************************************************************

def emeb_luis_carlos_prestes(request):
    return render(request, 'escolas/emeb_luis_carlos_prestes.html')
# *********************************************************************************************************************

def emeb_ronilton(request):
    return render(request, 'escolas/emeb_ronilton.html')
# *********************************************************************************************************************

def agenda_view(request):
    return render(request, 'agenda.html')
# *********************************************************************************************************************

def inscricoes_view(request):
    return render(request, 'inscricoes.html')
# *********************************************************************************************************************

def certificado_view(request):
    return render(request, 'certificado.html')
# *********************************************************************************************************************

def encontro_pedagogico(request):
    return render(request, 'encontro_pedagogico.html')
# *********************************************************************************************************************

def emef_alexandro_nunes(request):
    return render(request, 'escolas/emef_alexandro_nunes.html')
# *********************************************************************************************************************

def emef_benedita_torres(request):
    return render(request, 'escolas/emef_benedita_torres.html')
# *********************************************************************************************************************

def emef_carmelo_mendes(request):
    return render(request, 'escolas/emef_carmelo_mendes.html')
# *********************************************************************************************************************

def emef_francisca_romana(request):
    return render(request, 'escolas/emef_francisca_romana.html')
# *********************************************************************************************************************

def emef_joao_nelson(request):
    return render(request, 'escolas/emef_joao_nelson.html')
# *********************************************************************************************************************

def emef_maria_lourdes(request):
    return render(request, 'escolas/emef_maria_lourdes.html')
# *********************************************************************************************************************

def sebastiao_agripino(request):
    return render(request, 'escolas/sebastiao_agripino.html')
# *********************************************************************************************************************

def adelaide_molinari(request):
    return render(request, 'escolas/adelaide_molinari.html')
# *********************************************************************************************************************

def carlos_henrique(request):
    return render(request, 'escolas/carlos_henrique.html')
# *********************************************************************************************************************

def juscelino_kubitschek(request):
    return render(request, 'escolas/juscelino_kubitschek.html')
# *********************************************************************************************************************

def emeif_raimundo_oliveira(request):
    return render(request, 'escolas/emeif_raimundo_oliveira.html')
# *********************************************************************************************************************

def emeif_magalhaes_barata(request):
    return render(request, 'escolas/emeif_magalhaes_barata.html')
# *********************************************************************************************************************

def magalhaes_barata(request):
    return render(request, 'escolas/magalhaes_barata.html')
# *********************************************************************************************************************

def tancredo_almeida(request):
    return render(request, 'escolas/tancredo_almeida.html')
# *********************************************************************************************************************

def teotonio_vilela(request):
    return render(request, 'escolas/teotonio_vilela.html')
# *********************************************************************************************************************

def detalhes_escola(request, escola_id):
    # Simulação de dados da escola. Substitua pelos dados reais do seu banco de dados.
    escolas = {
        1: {
            "nome": "EMEIF Maria de Lourdes",
            "descricao": "Uma escola comprometida com o futuro e com a qualidade da educação para todos.",
            "endereco": "Rua Itamarati, S/Nº – Novo Horizonte, Canaã dos Carajás-PA",
            "imagem": "assets/dist/img/escola_1.jpg",
        },
        2: {
            "nome": "EMEIF Juscelino Kubitschek",
            "descricao": "Escola de excelência que transforma vidas através da educação.",
            "endereco": "Avenida Brasil, 123 – Centro, Canaã dos Carajás-PA",
            "imagem": "assets/dist/img/escola_2.jpg",
        },
        3: {
            "nome": "EMEIF Teotônio Vilela",
            "descricao": "Onde o aprendizado se encontra com a dedicação.",
            "endereco": "Rua Amazonas, 456 – Bairro Amazonas, Canaã dos Carajás-PA",
            "imagem": "assets/dist/img/escola_3.jpg",
        },
    }

    escola = escolas.get(escola_id, None)

    if not escola:
        return render(request, "404.html", status=404)  # Exibe página de erro se a escola não existir

    return render(request, "escolas/detalhes_escola.html", {"escola": escola})

# *********************************************************************************************************************

def adicionar_escola(request):
    if request.method == "POST":
        # Aqui você adicionará a lógica para salvar os dados da nova escola.
        nome = request.POST.get("nome")
        descricao = request.POST.get("descricao")
        endereco = request.POST.get("endereco")
        imagem = request.POST.get("imagem")  # Isso dependerá do upload de arquivos

        # Simulação de salvamento (substitua pela lógica de banco de dados)
        print(f"Escola adicionada: {nome}, {descricao}, {endereco}, {imagem}")

        return redirect("escolas")  # Redireciona para a página de escolas

    return render(request, "escolas/adicionar_escola.html")
# *********************************************************************************************************************

def site_pedagogico(request):
    """
    View para a página 'Site Pedagógico'.
    """
    return render(request, 'setor_pedagogico/site_pedagogico.html')
# *********************************************************************************************************************

def verificar_cpf_ajax(request):
    if request.method == 'POST':
        cpf = request.POST.get('cpf')
        try:
            # Fetch user by CPF and check if they are a superuser
            user = User.objects.get(cpf=cpf)
            if user.is_superuser:
                return JsonResponse({'status': 'success', 'redirect_url': '/admin-login/'})
            else:
                return JsonResponse({'status': 'error', 'message': 'CPF válido, mas o usuário não é um administrador.'})
        except User.DoesNotExist:
            return JsonResponse({'status': 'error', 'message': 'CPF não encontrado.'})
    
    return JsonResponse({'status': 'error', 'message': 'Método de requisição inválido.'})
# *********************************************************************************************************************

def admin_login_view(request):
    if request.method == 'POST':
        cpf = request.POST.get('username')  # Assuming CPF is entered in 'username' field
        password = request.POST.get('password')

        # Authenticate user by CPF and password
        user = authenticate(request, username=cpf, password=password)

        if user is not None:
            # Check if the authenticated user is a superuser
            if user.is_superuser:
                login(request, user)
                # Redirect to the admin dashboard if the user is a superuser
                return redirect('regimento_list')
            else:
                messages.error(request, 'CPF válido, mas você não é um administrador.')
        else:
            messages.error(request, 'CPF ou senha inválidos.')

    # Render the admin login template for GET requests or if authentication fails
    return render(request, 'registration/admin_login.html')
# *********************************************************************************************************************

def adm_site_pedagogico(request):
    return render(request, 'setor_pedagogico/adm_site_pedagogico.html')
# *********************************************************************************************************************

def upload_arquivo(request):
    if request.method == 'POST' and request.FILES.get('arquivo'):
        arquivo_csv = request.FILES['arquivo']
        
        try:
            # Ler o arquivo CSV usando Pandas
            df = pd.read_csv(arquivo_csv)

            # Verifica se o dataframe foi lido corretamente
            if df.empty:
                messages.error(request, "O arquivo está vazio ou não pôde ser lido.")
                return redirect('admin_dashboard')

            print(f"DataFrame lido com {len(df)} linhas.")
            
            registros_salvos = 0
            registros_pulados = 0

            # Iterar sobre as linhas do dataframe e salvar no banco de dados
            for index, row in df.iterrows():
                print(f"Processando linha {index + 1}: {row.to_dict()}")

                # Verifica se o CPF já existe no banco
                if Funcionario.objects.filter(cpf=row['CPF']).exists():
                    print(f"CPF {row['CPF']} já existe. Pulando registro.")
                    registros_pulados += 1
                    continue

                # Criar o novo funcionário
                funcionario = Funcionario(
                    nome_completo=row['NOME COMPLETO'],
                    rg=row['RG'],
                    cpf=row['CPF'],
                    telefone=row['TELEFONE'],
                    email=row['E-MAIL'],
                    cargo=row['CARGO'],
                    lotacao=row['LOTAÇÃO'],
                )

                # Salvando no banco
                funcionario.save()
                registros_salvos += 1
                print(f"Funcionario {funcionario.nome_completo} salvo com sucesso.")

            messages.success(request, f'Arquivo enviado com sucesso! {registros_salvos} registros salvos, {registros_pulados} registros pulados.')
        except Exception as e:
            print(f"Erro ao processar o arquivo: {str(e)}")
            messages.error(request, f'Ocorreu um erro ao processar o arquivo: {e}')

        return redirect('admin_dashboard')

    return render(request, 'admin_dashboard.html')
# *********************************************************************************************************************

def cadastro_usuario(request):
    if request.method == 'POST':
        # Coletando os dados do formulário
        nome_completo = request.POST.get('nome_completo')
        cpf = request.POST.get('cpf')
        email = request.POST.get('email')
        senha = request.POST.get('senha')
        confirmar_senha = request.POST.get('confirmar_senha')
        data_nascimento = request.POST.get('data_nascimento')

        errors = {}

        # Tratamento para CPF duplicado
        if Usuario.objects.filter(cpf=cpf).exists():
            errors['cpf_error'] = f"Este CPF ({cpf}) já está cadastrado."

        # Tratamento para email duplicado
        if Usuario.objects.filter(email=email).exists():
            errors['email_error'] = f"Este email ({email}) já está cadastrado."

        # Tratamento para senhas que não coincidem
        if senha != confirmar_senha:
            errors['password_error'] = "As senhas não coincidem."

        # Se houver erros, retorna a página de erro com todos os detalhes
        if errors:
            return render(request, 'error_page.html', errors)

        # Coletando outros campos para Inscricao
        responsavel_legal = request.POST.get('nome_responsavel')  # Corrigido para corresponder ao nome do campo no HTML
        tipo_responsavel = request.POST.get('tipo_responsavel')  # Corrigido para corresponder ao nome do campo no HTML
        cpf_responsavel = request.POST.get('cpf_responsavel')
        rg_responsavel = request.POST.get('rg_responsavel')
        telefone = request.POST.get('telefone')
        telefone_secundario = request.POST.get('telefone_2')
        endereco = request.POST.get('endereco')
        bairro_id = request.POST.get('bairro')  # Assuming you get 'bairro' ID from a form
        bairro_instance = Bairro.objects.get(id=bairro_id)  # Fetch the Bairro instance by ID
        cidade = request.POST.get('cidade')
        ponto_referencia = request.POST.get('ponto_referencia')
        possui_necessidade_especial = request.POST.get('possui_necessidade_especial')
        tipo_necessidade_especial = request.POST.get('necessidade_especial_detalhe')
        turno_disponivel = request.POST.get('turno_disponivel')
        etapa_pretendida = request.POST.get('etapa_pretendida')
        prova_todas_disciplinas = request.POST.get('prova_todas_disciplinas')
        fez_exame_supletivo = request.POST.get('fez_exame_supletivo')
        # Fixing the variable names to match the database columns
        local_exame = request.POST.get('local_prova')  # This should match the DB field 'local_exame'
        escola = request.POST.get('escola_2024')  # This should match the DB field 'escola'

        # Verifying passwords
        if senha != confirmar_senha:
            messages.error(request, "As senhas não coincidem!")
            return render(request, 'cadastro_usuario.html')

        # Check if CPF already exists
        if Usuario.objects.filter(cpf=cpf).exists():
            messages.error(request, "Este CPF já está cadastrado.")
            return render(request, 'cadastro_usuario.html')

        try:
            # Create user with encrypted password
            usuario = Usuario.objects.create_user(
                username=email,
                nome_completo=nome_completo,
                cpf=cpf,
                email=email,
                password=senha,
                data_nascimento=data_nascimento,
            )

            # Create the corresponding Inscricao and store it in a variable
            inscricao = Inscricao.objects.create(
                candidato=usuario,
                responsavel_legal=responsavel_legal,
                tipo_responsavel=tipo_responsavel,
                cpf_responsavel=cpf_responsavel,
                rg_responsavel=rg_responsavel,
                telefone=telefone,
                telefone_secundario=telefone_secundario,
                endereco=endereco,
                bairro=bairro_instance,  # Assign the Bairro instance
                cidade=cidade,
                ponto_referencia=ponto_referencia,
                necessidade_especial=(possui_necessidade_especial == 'Sim'),
                tipo_necessidade_especial=tipo_necessidade_especial,
                turno_disponivel=turno_disponivel,
                etapa_pretendida=etapa_pretendida,
                prova_todas_disciplinas=prova_todas_disciplinas,
                exame_supletivo=fez_exame_supletivo,
                local_exame=local_exame,  # Now this matches the database column
                escola=escola,  # Now this matches the database column
            )

            # Handle disciplines based on "prova_todas_disciplinas"
            if prova_todas_disciplinas == 'Sim':
                todas_disciplinas = Disciplina.objects.filter(nome__in=[
                    'Matemática', 'Ciências', 'Arte', 'Educação Física',
                    'História', 'Geografia', 'Língua Portuguesa', 'Inglês'
                ])
                inscricao.disciplinas.set(todas_disciplinas)
            else:
                # Adiciona apenas as disciplinas selecionadas pelo usuário
                disciplinas_selecionadas = request.POST.getlist('disciplinas')
                inscricao.disciplinas.set(disciplinas_selecionadas)

            # Log the user in and redirect to the candidate area
            login(request, usuario)
            return redirect('area_do_candidato')

        except IntegrityError as e:
            # Tratamento para outros erros do banco de dados
            errors['db_error'] = "Erro no banco de dados: " + str(e)
            return render(request, 'error_page.html', errors)

        except ValidationError as e:
            # Tratamento para outros erros de validação
            errors['validation_error'] = "Erro de validação: " + str(e)
            return render(request, 'error_page.html', errors)

    # Render the registration page if not a POST request
    return render(request, 'cadastro_usuario.html')

#**********************************************************************************************************

def home(request):
    return render(request, 'home.html')  # Substitua 'home.html' pelo caminho correto do seu template.
#**********************************************************************************************************

def eja_cadastro(request):
    return render(request, 'eja_cadastro.html')  # Substitua pelo caminho correto do seu template.
#**********************************************************************************************************

class CustomLoginView(LoginView):
    template_name = 'registration/login.html'  # Substitua pelo caminho correto do template
    redirect_authenticated_user = True
#**********************************************************************************************************

def logout_confirm(request):
    return render(request, 'logout_confirm.html')

def area_do_candidato(request):
    return render(request, 'area_do_candidato.html')

def get_diagnose_data(request):
    turmas = [
        '401', '403', '404', '406', '408', '409', '410', '413', '414', '415', '417',
        '421', '423', '426', '428', '429', '430', '431', '432', '433', '434', '435',
        '436', '437', '438', '439', '441', '442', '447', '451', '471'
    ]
    
    items = []
    for item in DiagnoseInicProfPort.objects.all():
        responses = {f"professor_{turma}": getattr(item, f"professor_{turma}", "N") for turma in turmas}
        items.append({
            "numero": item.item,
            "habilidade": item.habilidade,
            "descricao_habilidade": item.descricao_habilidade,
            "respostas": responses
        })
    
    return JsonResponse({"items": items, "turmas": turmas})

def get_diagnose_data_inic_alunos_matematica(request):
    # Lista de séries que serão usadas para os alunos
    series = ['3º', '4º', '5º', '6º']

    # Lista para armazenar os itens e suas respostas associadas
    items = []

    # Recupera todos os itens da tabela DiagnoseAlunoMatematica
    dados_db = DiagnoseAlunoMatematica.objects.all()

    # Itera sobre cada item recuperado do banco de dados
    for item in dados_db:
        # Adiciona os dados de cada item (série, habilidade, acertos e erros) na lista
        items.append({
            "serie": item.serie,  # Série do aluno
            "habilidade": item.habilidade,  # Habilidade correspondente
            "acerto": item.acerto,  # Total de acertos
            "erro": item.erro  # Total de erros
        })

    # Retorna os dados em formato JSON
    return JsonResponse({"items": items, "series": series})

def aluno_matematica_view(request):
    # Inicialização das variáveis
    series = ['3º', '4º', '5º', '6º']
    serie_3_corretas = serie_4_corretas = serie_5_corretas = serie_6_corretas = 0
    serie_3_erradas = serie_4_erradas = serie_5_erradas = serie_6_erradas = 0
    serie_3_total = serie_4_total = serie_5_total = serie_6_total = 0
    total_corretas = total_erradas = total_respostas = 0
    itens = []

    # Buscar dados do banco
    itens_db = DiagnoseAlunoMatematica.objects.all()

    for item in itens_db:
        try:
            # Convertendo para float se for uma string antes de arredondar
            acerto = round(float(item.acerto), 2)  # Arredondamento para duas casas decimais
            erro = round(float(item.erro), 2)  # Arredondamento para duas casas decimais
        except ValueError:
            # Se não for possível converter, ignorar a linha ou lidar com o erro conforme necessário
            continue

        total_corretas += acerto
        total_erradas += erro
        total_respostas += 1

        # Acumula os dados por série
        if item.serie == '3º':
            serie_3_corretas += acerto
            serie_3_erradas += erro
            serie_3_total += 1
        elif item.serie == '4º':
            serie_4_corretas += acerto
            serie_4_erradas += erro
            serie_4_total += 1
        elif item.serie == '5º':
            serie_5_corretas += acerto
            serie_5_erradas += erro
            serie_5_total += 1
        elif item.serie == '6º':
            serie_6_corretas += acerto
            serie_6_erradas += erro
            serie_6_total += 1

        itens.append({
            'serie': item.serie,
            'habilidade': item.habilidade,
            'acerto': acerto,
            'erro': erro
        })

    # Cálculo dos percentuais por série
    serie_3_percentual_acertos = round((serie_3_corretas / serie_3_total) * 100, 2) if serie_3_total > 0 else 0
    serie_3_percentual_erros = round((serie_3_erradas / serie_3_total) * 100, 2) if serie_3_total > 0 else 0
    serie_4_percentual_acertos = round((serie_4_corretas / serie_4_total) * 100, 2) if serie_4_total > 0 else 0
    serie_4_percentual_erros = round((serie_4_erradas / serie_4_total) * 100, 2) if serie_4_total > 0 else 0
    serie_5_percentual_acertos = round((serie_5_corretas / serie_5_total) * 100, 2) if serie_5_total > 0 else 0
    serie_5_percentual_erros = round((serie_5_erradas / serie_5_total) * 100, 2) if serie_5_total > 0 else 0
    serie_6_percentual_acertos = round((serie_6_corretas / serie_6_total) * 100, 2) if serie_6_total > 0 else 0
    serie_6_percentual_erros = round((serie_6_erradas / serie_6_total) * 100, 2) if serie_6_total > 0 else 0

    percentual_acertos = round((total_corretas / total_respostas) * 100, 2) if total_respostas > 0 else 0
    percentual_erros = round((total_erradas / total_respostas) * 100, 2) if total_respostas > 0 else 0

    uploaded_file_url = None

    if request.method == 'POST':
        if 'upload_planilha' in request.POST and request.FILES.get('planilha'):
            planilha = request.FILES['planilha']

            fs = FileSystemStorage()
            filename = fs.save(planilha.name, planilha)
            uploaded_file_url = fs.url(filename)

            try:
                # Carrega a planilha com pandas
                df = pd.read_excel(fs.path(filename))

                # Limpa a tabela antes de adicionar novos dados
                DiagnoseAlunoMatematica.objects.all().delete()

                # Itera sobre as linhas da planilha e processa cada uma
                for index, row in df.iterrows():
                    try:
                        serie = row['AI']
                        habilidade = row['Habilidade']
                        acerto = float(str(row['Acerto']).replace(",", "."))
                        erro = float(str(row['Erro']).replace(",", "."))

                        # Validação dos dados antes de salvar
                        if not (0 <= acerto <= 1 and 0 <= erro <= 1):
                            raise ValueError(f"Valores inválidos de acerto/erro na linha {index + 1}")

                        DiagnoseAlunoMatematica.objects.create(
                            serie=serie,
                            habilidade=habilidade,
                            acerto=acerto,
                            erro=erro
                        )
                    except Exception as e:
                        messages.error(request, f"Erro ao processar a linha {index + 1}: {e}")

                messages.success(request, "Planilha carregada e processada com sucesso!")
            except Exception as e:
                messages.error(request, f"Erro ao processar o arquivo: {e}")

            return redirect('aluno_matematica_view')

    return render(request, 'lingua_matematica_aluno.html', {
        'series': series,
        'itens': itens,
        'total_corretas': total_corretas,
        'total_erradas': total_erradas,
        'percentual_acertos': percentual_acertos,
        'percentual_erros': percentual_erros,
        'serie_3_corretas': serie_3_corretas,
        'serie_3_erradas': serie_3_erradas,
        'serie_3_percentual_acertos': serie_3_percentual_acertos,
        'serie_3_percentual_erros': serie_3_percentual_erros,
        'serie_4_corretas': serie_4_corretas,
        'serie_4_erradas': serie_4_erradas,
        'serie_4_percentual_acertos': serie_4_percentual_acertos,
        'serie_4_percentual_erros': serie_4_percentual_erros,
        'serie_5_corretas': serie_5_corretas,
        'serie_5_erradas': serie_5_erradas,
        'serie_5_percentual_acertos': serie_5_percentual_acertos,
        'serie_5_percentual_erros': serie_5_percentual_erros,
        'serie_6_corretas': serie_6_corretas,
        'serie_6_erradas': serie_6_erradas,
        'serie_6_percentual_acertos': serie_6_percentual_acertos,
        'serie_6_percentual_erros': serie_6_percentual_erros,
        'uploaded_file_url': uploaded_file_url
    })


################################################################################################################

def upload_excel_aluno_matematica(request):
    if request.method == "POST" and request.FILES.get("planilha"):
        planilha = request.FILES['planilha']
        fs = FileSystemStorage()
        filename = fs.save(planilha.name, planilha)
        file_path = fs.path(filename)

        # Mapeamento das séries em texto para números
        series_mapping = {
            '3º ano': 3,
            '4º ano': 4,
            '5º ano': 5,
            '6º ano': 6,
            # Adicione mais séries conforme necessário
        }

        try:
            # Verificar se o tipo de arquivo é Excel
            mime_type, _ = mimetypes.guess_type(file_path)
            if mime_type in ['application/vnd.ms-excel', 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet']:
                df = pd.read_excel(file_path, engine='openpyxl')

                # Limpar dados existentes
                DiagnoseAlunoMatematica.objects.all().delete()

                # Processar cada linha da planilha
                for index, row in df.iterrows():
                    try:
                        # Ajuste os nomes das colunas com os corretos
                        serie_texto = row.get('AI')  # Supondo que a coluna 'AI' contém a série em texto
                        habilidade = row.get('Habilidade')
                        acerto = float(str(row.get('Acerto')).replace(",", ".")) if not pd.isna(row.get('Acerto')) else 0
                        erro = float(str(row.get('Erro')).replace(",", ".")) if not pd.isna(row.get('Erro')) else 0

                        # Converter a série de texto para número
                        serie = series_mapping.get(serie_texto)
                        if serie is None:
                            raise ValueError(f"Série inválida '{serie_texto}' na linha {index + 1}")

                        # Validação dos valores de acerto/erro
                        if not (0 <= acerto <= 1 and 0 <= erro <= 1):
                            raise ValueError(f"Valores inválidos de acerto/erro na linha {index + 1}")

                        # Salvar no banco de dados
                        DiagnoseAlunoMatematica.objects.create(
                            serie=serie,  # Salvar a série convertida
                            habilidade=habilidade,
                            acerto=acerto,
                            erro=erro
                        )
                    except Exception as e:
                        messages.error(request, f"Erro ao processar a linha {index + 1}: {e}")
                        continue

                messages.success(request, "Planilha carregada e processada com sucesso!")
            else:
                raise ValueError("Formato de arquivo inválido.")

        except Exception as e:
            messages.error(request, f"Erro ao processar o arquivo: {e}")

        return redirect('aluno_matematica_view')

    return render(request, 'upload_page_matematica.html')

################################################################################################################

def get_diagnose_data_inic_alunos_port(request):
    # Lista de séries que serão usadas para os alunos
    series = ['3º', '4º', '5º', '6º']

    # Lista para armazenar os itens e suas respostas associadas
    items = []

    # Recupera todos os itens da tabela DiagnoseAlunoPortugues (ou o modelo equivalente)
    dados_db = DiagnoseAlunoPortugues.objects.all()

    # Itera sobre cada item recuperado do banco de dados
    for item in dados_db:
        # Adiciona os dados de cada item (série, habilidade, acertos e erros) na lista
        items.append({
            "serie": item.serie if item.serie in series else "N/A",  # Verifica se a série está na lista
            "habilidade": item.habilidade if item.habilidade else "N/A",  # Verifica se a habilidade está presente
            "acerto": round(float(item.acerto or 0), 2),  # Formata os acertos e trata valores nulos
            "erro": round(float(item.erro or 0), 2),  # Formata os erros e trata valores nulos
        })

    # Retorna os dados em formato JSON
    return JsonResponse({"items": items, "series": series})
################################################################################################################

def habilidades_portugues_view(request):
    habilidades = HabilidadePortugues.objects.all()  # Buscando todos os dados
    return render(request, 'habilidades_portugues.html', {'habilidades': habilidades})
################################################################################################################

def habilidades_matematica_view(request):
    # Lógica para processar e renderizar as habilidades de Matemática
    habilidades = HabilidadeMatematica.objects.all()  # Buscando todos os dados
    return render(request, 'habilidades_matematica.html', {'habilidades': habilidades})
################################################################################################################

def habilidades_portugues_view(request):
    return render(request, 'diagnosis/eja/habilidades_portugues.html')
################################################################################################################

# View para habilidades de Matemática
def habilidades_matematica_view(request):
    # Dados simulados (você pode buscar isso do banco de dados)
    habilidades = [
        {'serie': '5º Ano', 'topico': 'Números e Operações', 'habilidade': 'H1', 'descricao': 'Resolver problemas envolvendo adição e subtração.'},
        {'serie': '5º Ano', 'topico': 'Geometria', 'habilidade': 'H2', 'descricao': 'Identificar formas geométricas em objetos do cotidiano.'},
        {'serie': '9º Ano', 'topico': 'Álgebra', 'habilidade': 'H3', 'descricao': 'Resolver equações de primeiro grau com uma incógnita.'},
        {'serie': '9º Ano', 'topico': 'Estatística', 'habilidade': 'H4', 'descricao': 'Interpretar gráficos e tabelas.'},
    ]
    return render(request, 'diagnosis/eja/habilidades_matematica.html', {'habilidades': habilidades})
################################################################################################################

# View para a página de Pesquisa
def habilidades_pesquisa_view(request):
    query = request.GET.get('habilidade', '')

    # Simulação de dados (substitua por consulta ao banco de dados)
    todas_habilidades = [
        {'serie': '5º Ano', 'nome': 'Leitura e Interpretação', 'descricao': 'Desenvolver a habilidade de interpretação textual.'},
        {'serie': '9º Ano', 'nome': 'Álgebra Básica', 'descricao': 'Resolver equações lineares simples.'},
    ]

    habilidades = [h for h in todas_habilidades if query.lower() in h['nome'].lower()] if query else todas_habilidades

    return render(request, 'diagnosis/eja/habilidades_pesquisa.html', {'habilidades': habilidades})

def lingua_portuguesa_alunos_iniciais(request):
    return render(request, 'diagnoses/lingua_portuguesa_alunos_iniciais.html')
###***********************************************************************************************************************

def lingua_portuguesa_professores_iniciais(request):
    return render(request, 'diagnoses/lingua_portuguesa_prof.html')
###***********************************************************************************************************************

def lingua_matematica_alunos_iniciais(request):
    return render(request, 'diagnoses/lingua_matematica_alunos_iniciais.html')
###***********************************************************************************************************************

def lingua_matematica_professores_iniciais(request):
    return render(request, 'diagnoses/lingua_matematica_professores_iniciais.html')
###***********************************************************************************************************************

def lingua_portuguesa_alunos_finais(request):
    return render(request, 'diagnoses/lingua_portuguesa_alunos_finais.html')
###***********************************************************************************************************************

def lingua_portuguesa_professores_finais(request):
    return render(request, 'diagnoses/lingua_portuguesa_prof.html')  # Alterado para o template correto
###***********************************************************************************************************************

def lingua_matematica_alunos_finais(request):
    return render(request, 'diagnoses/lingua_matematica_alunos_finais.html')
###***********************************************************************************************************************

def lingua_matematica_professores_finais(request):
    return render(request, 'diagnoses/lingua_matematica_professores_finais.html')
###***********************************************************************************************************************

def lingua_portuguesa_prof_view(request):
    # Definindo as turmas
    turmas = ['401', '403', '404', '406', '408', '409', '410', '413', '414', '415', '417', '421', '423', '426', 
              '428', '429', '430', '431', '432', '433', '434', '435', '436', '437', '438', '439', '441', '442', 
              '447', '451', '471']

    # Definindo os itens e habilidades
    itens = [
        {'numero': 1, 'habilidade': 'D12', 'respostas': {'401': 1, '403': 1, '404': 1, '406': 1}},
        {'numero': 2, 'habilidade': 'D03', 'respostas': {'401': 0, '403': 1, '404': 1, '406': 1}},
        {'numero': 3, 'habilidade': 'D09', 'respostas': {'401': 0, '403': 0, '404': 1, '406': 0}},
        {'numero': 4, 'habilidade': 'D01', 'respostas': {'401': 1, '403': 1, '404': 1, '406': 1}},
        {'numero': 5, 'habilidade': 'D09', 'respostas': {'401': 1, '403': 1, '404': 1, '406': 1}},
        {'numero': 6, 'habilidade': 'D06', 'respostas': {'401': 1, '403': 0, '404': 1, '406': 1}},
        {'numero': 7, 'habilidade': 'D15', 'respostas': {'401': 1, '403': 1, '404': 0, '406': 1}},
        {'numero': 8, 'habilidade': 'D02', 'respostas': {'401': 1, '403': 1, '404': 0, '406': 1}},
        {'numero': 9, 'habilidade': 'D03', 'respostas': {'401': 1, '403': 1, '404': 1, '406': 0}},
        {'numero': 10, 'habilidade': 'D05', 'respostas': {'401': 1, '403': 1, '404': 1, '406': 1}},
        {'numero': 11, 'habilidade': 'D04', 'respostas': {'401': 1, '403': 1, '404': 1, '406': 1}},
        {'numero': 12, 'habilidade': 'D12', 'respostas': {'401': 1, '403': 1, '404': 1, '406': 1}},
        {'numero': 13, 'habilidade': 'D05', 'respostas': {'401': 1, '403': 1, '404': 0, '406': 0}},
        {'numero': 14, 'habilidade': 'D13', 'respostas': {'401': 1, '403': 1, '404': 1, '406': 1}},
        {'numero': 15, 'habilidade': 'D10', 'respostas': {'401': 1, '403': 1, '404': 1, '406': 1}},
        {'numero': 16, 'habilidade': 'D13', 'respostas': {'401': 1, '403': 1, '404': 1, '406': 1}},
        {'numero': 17, 'habilidade': 'D07', 'respostas': {'401': 0, '403': 1, '404': 1, '406': 0}},
        {'numero': 18, 'habilidade': 'D14', 'respostas': {'401': 1, '403': 1, '404': 1, '406': 1}},
        {'numero': 19, 'habilidade': 'D06', 'respostas': {'401': 1, '403': 1, '404': 1, '406': 1}},
        {'numero': 20, 'habilidade': 'D07', 'respostas': {'401': 0, '403': 1, '404': 1, '406': 0}},
        {'numero': 21, 'habilidade': 'D14', 'respostas': {'401': 1, '403': 1, '404': 1, '406': 1}},
        {'numero': 22, 'habilidade': 'D10', 'respostas': {'401': 1, '403': 1, '404': 1, '406': 1}},
        {'numero': 23, 'habilidade': 'D13', 'respostas': {'401': 1, '403': 1, '404': 1, '406': 1}},
        {'numero': 24, 'habilidade': 'D15', 'respostas': {'401': 0, '403': 1, '404': 1, '406': 1}},
        {'numero': 25, 'habilidade': 'D04', 'respostas': {'401': 1, '403': 1, '404': 1, '406': 1}},
        {'numero': 26, 'habilidade': 'D02', 'respostas': {'401': 1, '403': 1, '404': 1, '406': 1}},
    ]

    itens = []
    total_corretas = 0
    total_erradas = 0

    # Carregando os dados salvos no banco de dados
    itens_db = DiagnoseInicProfPort.objects.all()

    for item in itens_db:
        respostas = {}
        for turma in turmas:
            resposta = getattr(item, f'professor_{turma}', 'N')
            respostas[turma] = resposta
            if resposta == 'S':
                total_corretas += 1
            else:
                total_erradas += 1
        itens.append({'numero': item.item, 'habilidade': item.habilidade, 'respostas': respostas})

    uploaded_file_url = None

    if request.method == 'POST':
        if 'upload_planilha' in request.POST and request.FILES.get('planilha'):
            planilha = request.FILES['planilha']
            fs = FileSystemStorage()
            filename = fs.save(planilha.name, planilha)
            uploaded_file_url = fs.url(filename)
            
            # Processar a planilha usando pandas
            df = pd.read_excel(fs.path(filename))

            for index, row in df.iterrows():
                item_num = row['Item']
                habilidade = row['Habilidade']
                
                # Criar ou atualizar o registro no banco de dados
                diagnose, created = DiagnoseInicProfPort.objects.get_or_create(
                    item=item_num,
                    habilidade=habilidade
                )

                for turma in turmas:
                    resposta = row.get(str(turma), 0)
                    # Salva a resposta no banco de dados como 'S' ou 'N'
                    setattr(diagnose, f'professor_{turma}', 'S' if resposta == 1 else 'N')
                
                diagnose.save()

            # Redirecionar de volta para a página de exibição após o salvamento
            return redirect('lingua_portuguesa_prof_view')

        if 'salvar_respostas' in request.POST:
            for item in itens:
                diagnose = DiagnoseInicProfPort.objects.get(item=item['numero'])
                for turma in turmas:
                    nova_resposta = request.POST.get(f'respostas_{item["numero"]}_{turma}', 'N')
                    setattr(diagnose, f'professor_{turma}', nova_resposta)
                diagnose.save()

            # Redirecionar após salvar
            return redirect('lingua_portuguesa_prof_view')

    return render(request, 'diagnoses/lingua_portuguesa_prof.html', {
        'turmas': turmas,
        'itens': itens,
        'total_corretas': total_corretas,
        'total_erradas': total_erradas,
        'uploaded_file_url': uploaded_file_url
    })
###***********************************************************************************************************************

def lingua_matematica_professores_finais(request):
    return render(request, 'diagnoses/lingua_matematica_professores_finais.html')
###***********************************************************************************************************************

def support(request):
    return render(request, 'webapp/support.html')
###***********************************************************************************************************************

def configuracoes_view(request):
    return render(request, 'webapp/configuracoes.html')
###***********************************************************************************************************************

def alterar_senha_view(request):
    if request.method == 'POST':
        # Lógica de alteração de senha
        pass  # Coloque aqui a lógica para alterar a senha
    return render(request, 'webapp/alterar_senha.html')
###***********************************************************************************************************************

def atualizar_preferencias_view(request):
    if request.method == 'POST':
        # Lógica para atualizar as preferências de notificação do usuário
        # Exemplo: request.user.profile.email_notifications = request.POST.get('email_notifications')
        pass  # Coloque aqui a lógica para salvar as preferências de notificação
    return redirect('configuracoes')  # Redireciona de volta para a página de configurações após salvar
###***********************************************************************************************************************

@login_required
def upload_files(request):
    if request.method == 'POST':
        csv_file = request.FILES['file']

        # Detect encoding
        raw_data = csv_file.read()
        result = chardet.detect(raw_data)
        encoding = result['encoding']
        csv_file.seek(0)  # Reset file pointer

        try:
            # Read CSV file with detected encoding and ensure the first row is treated as header
            df = pd.read_csv(csv_file, encoding=encoding, delimiter=';', header=0, on_bad_lines='skip')

            expected_headers = [
                'numero_matricula', 'nome_completo', 'escola', 'data_nascimento', 'sexo', 'turno',
                'telefone', 'nivel_escolaridade', 'endereco', 'bairro', 'cep', 'codigo_rota',
                'email_responsavel', 'cpf_responsavel', 'nome_responsavel', 'parentesco',
                'telefone_responsavel', 'responsavel_legal'
            ]

            # Verify that the CSV file contains the expected headers
            csv_headers = list(df.columns)
            if csv_headers != expected_headers:
                missing_headers = set(expected_headers) - set(csv_headers)
                extra_headers = set(csv_headers) - set(expected_headers)
                error_message = f"Os cabeçalhos dos arquivos CSV não correspondem aos cabeçalhos esperados.\n"
                if missing_headers:
                    error_message += f"Cabeçalhos ausentes: {', '.join(missing_headers)}\n"
                if extra_headers:
                    error_message += f"Cabeçalhos extras: {', '.join(extra_headers)}"
                raise ValueError(error_message)

            # Convert date format
            df['data_nascimento'] = pd.to_datetime(df['data_nascimento'], format='%d/%m/%Y', errors='coerce')

            students_inserted = 0
            students_existing = []

            for _, row in df.iterrows():
                if pd.isna(row['data_nascimento']):
                    raise ValueError(f"Formato de data inválido para linha  : {row.to_dict()}")

                numero_matricula = row['numero_matricula']
                if not Student.objects.filter(numero_matricula=numero_matricula).exists():
                    Student.objects.create(
                        numero_matricula=numero_matricula,
                        nome_completo=row['nome_completo'],
                        escola=row['escola'],
                        data_nascimento=row['data_nascimento'].strftime('%Y-%m-%d'),
                        sexo=row['sexo'],
                        turno=row['turno'],
                        telefone=row['telefone'],
                        nivel_escolaridade=row['nivel_escolaridade'],
                        endereco=row['endereco'],
                        bairro=row['bairro'],
                        cep=row['cep'],
                        codigo_rota=row['codigo_rota'],
                        email_responsavel=row['email_responsavel'],
                        cpf_responsavel=row['cpf_responsavel'],
                        nome_responsavel=row['nome_responsavel'],
                        parentesco=row['parentesco'],
                        telefone_responsavel=row['telefone_responsavel'],
                        responsavel_legal=row['responsavel_legal'],
                    )
                    students_inserted += 1
                else:
                    students_existing.append(row)

            # Generate file with existing students
            if students_existing:
                output = io.StringIO()
                writer = csv.DictWriter(output, fieldnames=expected_headers)
                writer.writeheader()
                writer.writerows(students_existing)
                csv_content = output.getvalue()
                output.close()

                file_name = 'students_existing.csv'
                file_path = default_storage.save(file_name, ContentFile(csv_content))

                students_existing_url = default_storage.url(file_path)
            else:
                students_existing_url = None

            messages.success(request, f'{students_inserted} alunos foram inseridos com sucesso.')
            if students_existing:
                messages.warning(request, f'{len(students_existing)} alunos já existiam no banco de dados.')

            return render(request, 'webapp/upload_files.html', {'students_existing_url': students_existing_url})

        except Exception as e:
            messages.error(request, f'Erro ao processar arquivo CSV: {str(e)}')

    return render(request, 'webapp/upload_files.html')
###***********************************************************************************************************************
def dashboard_transporte(request):
    total_alunos = IndicadoresTransporte.objects.count()
    total_escolas = IndicadoresTransporte.objects.values('escola').distinct().count()
    alunos_ensino_medio = IndicadoresTransporte.objects.filter(
        nivel_escolaridade__icontains='Ensino Médio'
    ).count()
    alunos_ensino_fundamental = IndicadoresTransporte.objects.filter(
        nivel_escolaridade__icontains='Ensino Fundamental'
    ).count()

    alunos_por_escola = (
        IndicadoresTransporte.objects
        .values('escola')
        .annotate(total=Count('id'))
        .order_by('-total')
    )
    escolas = [item['escola'] for item in alunos_por_escola]
    alunos = [item['total'] for item in alunos_por_escola]

    context = {
        'total_alunos': total_alunos,
        'total_escolas': total_escolas,
        'alunos_ensino_medio': alunos_ensino_medio,
        'alunos_ensino_fundamental': alunos_ensino_fundamental,
        'escolas': escolas,
        'alunos': alunos,
    }
    return render(request, 'indicadores/transporte.html', context)

###***********************************************************************************************************************

def transporte_dashboard(request):
    total_alunos = IndicadoresTransporte.objects.count()
    total_escolas = IndicadoresTransporte.objects.values('escola').distinct().count()
    alunos_ensino_medio = IndicadoresTransporte.objects.filter(nivel_escolaridade__icontains='Ensino Médio').count()
    alunos_ensino_fundamental = IndicadoresTransporte.objects.filter(nivel_escolaridade__icontains='Ensino Fundamental').count()

    alunos_por_escola = IndicadoresTransporte.objects.values('escola').annotate(total=Count('id')).order_by('-total')

    context = {
        'total_alunos': total_alunos,
        'total_escolas': total_escolas,
        'alunos_ensino_medio': alunos_ensino_medio,
        'alunos_ensino_fundamental': alunos_ensino_fundamental,
        'alunos_por_escola': list(alunos_por_escola),
    }
    return render(request, 'indicadores/transporte.html', context)
###***********************************************************************************************************************

def transporte_dados_api(request):
    alunos_por_escola = IndicadoresTransporte.objects.values('escola').annotate(total=Count('id')).order_by('-total')
    data = {
        'total_alunos': IndicadoresTransporte.objects.count(),
        'total_escolas': IndicadoresTransporte.objects.values('escola').distinct().count(),
        'alunos_ensino_medio': IndicadoresTransporte.objects.filter(nivel_escolaridade__icontains='Ensino Médio').count(),
        'alunos_ensino_fundamental': IndicadoresTransporte.objects.filter(nivel_escolaridade__icontains='Ensino Fundamental').count(),
        'alunos_por_escola': list(alunos_por_escola),
    }
    return JsonResponse(data)
###***********************************************************************************************************************

def indicadores_canaa(request):
    total_escolas = IndicadoresTransporte.objects.values('escola').distinct().count()
    taxa_alfabetizacao = calcular_taxa_alfabetizacao()
    # total_creches = Creches.objects.filter(ativa=True).count()
    # alunos_creches = Alunos.objects.filter(nivel_escolaridade__icontains='Creche').count()
    media_idade = calcular_media_idade()
    investimento_educacao = calcular_investimento_educacao(2024)
    escolas_com_transporte = IndicadoresTransporte.objects.values('escola').distinct().count()
    taxa_inclusao_digital = calcular_inclusao_digital()

    context = {
        'taxa_alfabetizacao': taxa_alfabetizacao,
        # 'total_creches': total_creches,
        # 'alunos_creches': alunos_creches,
        'media_idade': media_idade,
        'investimento_educacao': investimento_educacao,
        'escolas_com_transporte': escolas_com_transporte,
        'total_escolas': total_escolas,
        'taxa_inclusao_digital': taxa_inclusao_digital,
    }
    return render(request, 'indicadores/canaa.html', context)
###***********************************************************************************************************************

def calcular_taxa_alfabetizacao():
    # Simula o total de crianças alfabetizadas e o total de crianças na faixa etária
    total_criancas = IndicadoresTransporte.objects.filter(data_nascimento__year__gte=2010).count()
    alfabetizadas = total_criancas * 0.95  # Exemplo de taxa fixa de 95%
    return round((alfabetizadas / total_criancas) * 100, 2) if total_criancas > 0 else 0
###***********************************************************************************************************************

def calcular_media_idade():
    hoje = date.today()
    alunos = IndicadoresTransporte.objects.all()
    total_idades = sum([hoje.year - aluno.data_nascimento.year for aluno in alunos])
    total_alunos = alunos.count()
    return round(total_idades / total_alunos, 1) if total_alunos > 0 else 0
###***********************************************************************************************************************

def calcular_investimento_educacao(ano):
    # Simula um investimento baseado em dados fictícios
    investimento = {
        2023: 12000000,
        2024: 15000000,
    }
    return investimento.get(ano, 0)
###***********************************************************************************************************************

def calcular_inclusao_digital():
    escolas_totais = IndicadoresTransporte.objects.values('escola').distinct().count()
    escolas_com_internet = int(escolas_totais * 0.85)  # Exemplo: 85% possuem internet
    return round((escolas_com_internet / escolas_totais) * 100, 2) if escolas_totais > 0 else 0
###***********************************************************************************************************************

# class Creches(models.Model):
#     nome = models.CharField(max_length=255)
#     ativa = models.BooleanField(default=True)


# class Alunos(models.Model):
#     nome_completo = models.CharField(max_length=255)
#     data_nascimento = models.DateField()
#     nivel_escolaridade = models.CharField(max_length=100)
###***********************************************************************************************************************
def gestao_dados(request):
    # Filtragem
    escola_selecionada = request.GET.get('escola', '')
    bairro_selecionado = request.GET.get('bairro', '')
    turno_selecionado = request.GET.get('turno', '')
    sexo_selecionado = request.GET.get('sexo', '')

    # Obter dados únicos para os filtros
    escolas = IndicadoresTransporte.objects.values_list('escola', flat=True).distinct()
    bairros = IndicadoresTransporte.objects.values_list('bairro', flat=True).distinct()

    # Aplicar filtros na consulta
    alunos = IndicadoresTransporte.objects.all()
    if escola_selecionada:
        alunos = alunos.filter(escola=escola_selecionada)
    if bairro_selecionado:
        alunos = alunos.filter(bairro=bairro_selecionado)
    if turno_selecionado:
        alunos = alunos.filter(turno=turno_selecionado)
    if sexo_selecionado:
        alunos = alunos.filter(sexo=sexo_selecionado)

    # Contexto para o template
    context = {
        'alunos': alunos,
        'escolas': escolas,
        'bairros': bairros,
        'escola_selecionada': escola_selecionada,
        'bairro_selecionado': bairro_selecionado,
        'turno_selecionado': turno_selecionado,
        'sexo_selecionado': sexo_selecionado,
    }

    return render(request, 'gestao_dados.html', context)
###***********************************************************************************************************************

def planejamento(request):
    return render(request, 'seppec/planejamento.html')
###***********************************************************************************************************************

def execucao(request):
    return render(request, 'seppec/execucao.html')
###***********************************************************************************************************************

def controle(request):
    return render(request, 'seppec/controle.html')
###***********************************************************************************************************************

def comunicacao(request):
    return render(request, 'seppec/comunicacao.html')
###***********************************************************************************************************************

def quem_somos(request):
    return render(request, 'setor_pedagogico/quem_somos.html')
###***********************************************************************************************************************

def missao_visao(request):
    return render(request, 'setor_pedagogico/missao_visao.html')
###***********************************************************************************************************************

def equipe(request):
    return render(request, 'setor_pedagogico/equipe.html')
###***********************************************************************************************************************

def projetos_especiais(request):
    return render(request, 'setor_pedagogico/projetos_especiais.html')
###***********************************************************************************************************************

def contato_pedagogico(request):
    return render(request, 'setor_pedagogico/contato_pedagogico.html')
###***********************************************************************************************************************

def painel_administrativo(request):
    # Busca todos os registros da tabela
    artigos = RegimentoCadastro.objects.all()  
    return render(request, 'painel_administrativo.html', {'artigos': artigos})
###***********************************************************************************************************************

def admin_dashboard(request):
    # Contar apenas artigos cujo ID seja maior que um valor (exemplo: 10)
    total_artigos = SemedAppRegimentoCadastro.objects.filter(id__gt=10).count()
    
    # Filtrar artigos específicos (caso deseje enviar dados detalhados ao template)
    artigos = SemedAppRegimentoCadastro.objects.filter(id__gte=1).values(
        'id', 'titulo', 'capitulo', 'tipo_alteracao', 'nome_completo',
        'email', 'cpf', 'telefone', 'cargo', 'lotacao', 'data_submissao', 'status'
    )
    
    return render(request, 'adm_site_pedagogico.html', {
        'total_artigos': total_artigos,
        'artigos': list(artigos)  # Enviar artigos ao template (se necessário)
    })
###***********************************************************************************************************************

def get_articles_data(request):
    artigos = SemedAppRegimentoCadastro.objects.all().values(
        'id', 'titulo', 'capitulo', 'tipo_alteracao', 'nome_completo',
        'email', 'cpf', 'telefone', 'cargo', 'lotacao', 'data_submissao', 'status'
    )
    return JsonResponse({'data': list(artigos)})

###***********************************************************************************************************************

def site_pedagogico_view(request):
    # Renderiza a página do Site Pedagógico
    return render(request, 'site_pedagogico.html')
###***********************************************************************************************************************


# def index(request):
#     return render(request, 'banco_curriculos/index.html')
###***********************************************************************************************************************

def vagas(request):
    return render(request, 'banco_curriculos/vagas.html')
###***********************************************************************************************************************

def curriculos(request):
    return render(request, 'banco_curriculos/curriculos.html')
###***********************************************************************************************************************

from django.shortcuts import render, redirect
from semedapp.models import CadastroCandidato
from django.contrib.auth.hashers import check_password
from django.contrib import messages

def login_curriculo(request):
    if request.method == "POST":
        email = request.POST.get("email")
        senha = request.POST.get("password")

        try:
            # Busca o candidato pelo email
            candidato = CadastroCandidato.objects.get(email=email)

            # Verifica a senha
            if check_password(senha, candidato.senha):
                # Login bem-sucedido - cria a sessão do candidato
                request.session['candidato_id'] = candidato.id
                messages.success(request, "Login realizado com sucesso!")
                return redirect("area_candidato")  # Substitua pela URL da área do candidato
            else:
                # Senha incorreta
                messages.error(request, "E-mail ou senha inválidos.")
        except CadastroCandidato.DoesNotExist:
            # Candidato não encontrado
            messages.error(request, "E-mail ou senha inválidos.")

    return render(request, "banco_curriculos/login.html")


###***********************************************************************************************************************

def registrar(request):
    return render(request, 'banco_curriculos/registrar.html')
###***********************************************************************************************************************

def registrar_view(request):
    if request.method == "POST":
        # Lógica de registro
        pass
    return render(request, 'banco_curriculos/registrar.html')
###***********************************************************************************************************************

def login_curriculo_view(request):
    return render(request, 'banco_curriculos/login.html')
###***********************************************************************************************************************

# Página Home para banco de currículos
def banco_curriculos_home(request):
    return render(request, 'banco_curriculos/index.html')  # Template específico
###***********************************************************************************************************************

def upload_planilha_prof(request):#####PRIMEIRO
    if request.method == 'POST':
        form = UploadPlanilhaForm(request.POST, request.FILES)
        if form.is_valid():
            planilha = request.FILES['planilha']
            # Ler o arquivo CSV
            df = pd.read_csv(planilha)

            # Limpa os dados existentes para evitar duplicações
            HabilidadeProf.objects.all().delete()

            # Processar os dados da planilha
            for _, row in df.iterrows():
                HabilidadeProf.objects.create(
                    item=row['Item'],
                    habilidade=row['Habilidade'],
                    turma_401=row.get('401', 0),
                    turma_403=row.get('403', 0),
                    turma_404=row.get('404', 0),
                    turma_406=row.get('406', 0),
                    turma_408=row.get('408', 0),
                    turma_409=row.get('409', 0),
                    turma_410=row.get('410', 0),
                    turma_413=row.get('413', 0),
                    turma_414=row.get('414', 0),
                    turma_415=row.get('415', 0),
                    turma_417=row.get('417', 0),
                    turma_421=row.get('421', 0),
                    turma_423=row.get('423', 0),
                    turma_426=row.get('426', 0),
                    turma_428=row.get('428', 0),
                    turma_429=row.get('429', 0),
                    turma_430=row.get('430', 0),
                    turma_431=row.get('431', 0),
                    turma_432=row.get('432', 0),
                    turma_433=row.get('433', 0),
                    turma_434=row.get('434', 0),
                    turma_435=row.get('435', 0),
                    turma_436=row.get('436', 0),
                    turma_437=row.get('437', 0),
                    turma_438=row.get('438', 0),
                    turma_439=row.get('439', 0),
                    turma_441=row.get('441', 0),
                    turma_442=row.get('442', 0),
                    turma_447=row.get('447', 0),
                    turma_451=row.get('451', 0),
                    turma_471=row.get('471', 0),
                )
            return redirect('listar_habilidades_prof')  # Redireciona para a página de listagem
    else:
        form = UploadPlanilhaForm()

    return render(request, 'upload_planilha_prof.html', {'form': form})

###***********************************************************************************************************************

def listar_habilidades_prof(request):#####PRIMEIRO
    habilidades = HabilidadeProf.objects.all()
    total_acertos = sum(h.calcular_acertos() for h in habilidades)
    total_erros = sum(h.calcular_erros() for h in habilidades)
    
    return render(request, 'listar_habilidades_prof.html', {
        'habilidades': habilidades,
        'total_acertos': total_acertos,
        'total_erros': total_erros,
    })

###***********************************************************************************************************************

def habilidade_prof_view(request):
    habilidades = HabilidadeProf.objects.all()
    
    # Calcular totais
    total_acertos = sum(h.calcular_acertos() for h in habilidades)
    total_erros = sum(h.calcular_erros() for h in habilidades)
    total_respostas = total_acertos + total_erros

    percentual_acertos = (
        round((total_acertos / total_respostas) * 100, 2) if total_respostas > 0 else 0
    )
    percentual_erros = (
        round((total_erros / total_respostas) * 100, 2) if total_respostas > 0 else 0
    )
    
    # Lista de turmas (para garantir que o template receba a estrutura correta)
    turmas = [
        '401', '403', '404', '406', '408', '409', '410', '413', '414', '415',
        '417', '421', '423', '426', '428', '429', '430', '432', '433', '434',
        '435', '436', '437', '438', '439', '441', '442', '447', '451', '471'
    ]

    # Renderizar o template com os dados
    return render(request, 'diagnosis/prof/portugues.html', {
        'habilidades': habilidades,
        'total_corretas': total_acertos,
        'total_erradas': total_erros,
        'percentual_acertos': percentual_acertos,
        'percentual_erros': percentual_erros,
        'turmas': turmas,
    })

###***********************************************************************************************************************

def carregar_dados_habilidades(request):#####PRIMEIRO
    habilidades = HabilidadeProf.objects.all()
    data = []

    # Inicializar variáveis para os totais
    total_acertos = 0
    total_erros = 0

    for h in habilidades:
        # Calcular acertos e erros para cada habilidade
        acertos = h.calcular_acertos()
        erros = h.calcular_erros()

        # Atualizar os totais
        total_acertos += acertos
        total_erros += erros

        # Adicionar habilidade e valores ao JSON
        data.append({
            'item': h.item,
            'habilidade': h.habilidade,
            'turmas': {
                '401': h.turma_401,
                '403': h.turma_403,
                '404': h.turma_404,
                '406': h.turma_406,
                '408': h.turma_408,
                '409': h.turma_409,
                '410': h.turma_410,
                '413': h.turma_413,
                '414': h.turma_414,
                '415': h.turma_415,
                '417': h.turma_417,
                '421': h.turma_421,
                '423': h.turma_423,
                '426': h.turma_426,
                '428': h.turma_428,
                '429': h.turma_429,
                '430': h.turma_430,
                '432': h.turma_432,
                '433': h.turma_433,
                '434': h.turma_434,
                '435': h.turma_435,
                '436': h.turma_436,
                '437': h.turma_437,
                '438': h.turma_438,
                '439': h.turma_439,
                '441': h.turma_441,
                '442': h.turma_442,
                '447': h.turma_447,
                '451': h.turma_451,
                '471': h.turma_471,
            },
            'acertos': acertos,
            'erros': erros,
        })

    # Calcular percentuais gerais
    total_respostas = total_acertos + total_erros
    percentual_acertos = round((total_acertos / total_respostas) * 100, 2) if total_respostas > 0 else 0
    percentual_erros = round((total_erros / total_respostas) * 100, 2) if total_respostas > 0 else 0

    # Retornar os dados com os totais e percentuais
    return JsonResponse({
        'habilidades': data,
        'total_acertos': total_acertos,
        'total_erros': total_erros,
        'percentual_acertos': percentual_acertos,
        'percentual_erros': percentual_erros,
    })
###***********************************************************************************************************************

def carregar_habilidades(request):#####PRIMEIRO
    habilidades = HabilidadeProf.objects.all()
    data = []

    for h in habilidades:
        data.append({
            'item': h.item,
            'habilidade': h.habilidade,
            'turmas': {
                '401': h.turma_401,
                '403': h.turma_403,
                '404': h.turma_404,
                '406': h.turma_406,
                '408': h.turma_408,
                '409': h.turma_409,
                '410': h.turma_410,
                '413': h.turma_413,
                '414': h.turma_414,
                '415': h.turma_415,
                '417': h.turma_417,
                '421': h.turma_421,
                '423': h.turma_423,
                '426': h.turma_426,
                '428': h.turma_428,
                '429': h.turma_429,
                '430': h.turma_430,
                '432': h.turma_432,
                '433': h.turma_433,
                '434': h.turma_434,
                '435': h.turma_435,
                '436': h.turma_436,
                '437': h.turma_437,
                '438': h.turma_438,
                '439': h.turma_439,
                '441': h.turma_441,
                '442': h.turma_442,
                '447': h.turma_447,
                '451': h.turma_451,
                '471': h.turma_471,
            }
        })

    return JsonResponse({'habilidades': data})
###***********************************************************************************************************************

def habilidade_prof_anos_finais_view(request):#####SEGUNDO
    if request.method == 'POST' and 'planilha' in request.FILES:
        planilha = request.FILES['planilha']

        try:
            df = pd.read_csv(planilha, encoding='utf-8')
            df.columns = df.columns.str.strip()

            if 'Item' not in df.columns or 'Habilidade' not in df.columns:
                return render(request, 'diagnosis/prof/anos_finais.html', {
                    'error_message': 'As colunas "Item" e "Habilidade" são obrigatórias.'
                })

            # Limpar dados antigos e salvar novos
            HabilidadeProfAnosFinais.objects.all().delete()
            for _, row in df.iterrows():
                HabilidadeProfAnosFinais.objects.create(
                    item=row['Item'],
                    habilidade=row['Habilidade'],
                    turma_101=row.get('101', 0),
                    turma_102=row.get('102', 0),
                    # Adicione todos os campos conforme necessário
                    turma_171=row.get('171', 0),
                )

        except Exception as e:
            return render(request, 'diagnosis/prof/anos_finais.html', {
                'error_message': f'Erro ao processar a planilha: {str(e)}'
            })

    # Buscar dados do banco
    habilidades = HabilidadeProfAnosFinais.objects.all()

    # Calcular totais para os cards
    total_acertos = sum(h.calcular_acertos() for h in habilidades)
    total_erros = sum(h.calcular_erros() for h in habilidades)
    total_respostas = total_acertos + total_erros
    percentual_acertos = (
        round((total_acertos / total_respostas) * 100, 2) if total_respostas > 0 else 0
    )
    percentual_erros = (
        round((total_erros / total_respostas) * 100, 2) if total_respostas > 0 else 0
    )

    return render(request, 'diagnosis/prof/anos_finais.html', {
        'habilidades': habilidades,
        'total_corretas': total_acertos,
        'total_erradas': total_erros,
        'percentual_acertos': percentual_acertos,
        'percentual_erros': percentual_erros,
    })
###***********************************************************************************************************************

from .models import HabilidadeProfFinal  # Certifique-se de que o modelo está correto
import pandas as pd

def portugues_final_view(request):#####SEGUNDO
    if request.method == 'POST' and 'planilha' in request.FILES:
        planilha = request.FILES['planilha']

        try:
            # Ler a planilha enviada
            df = pd.read_csv(planilha, encoding='utf-8')
            df.columns = df.columns.str.strip()  # Remove espaços extras das colunas

            # Verificar as colunas obrigatórias
            if 'Item' not in df.columns or 'Habilidade' not in df.columns:
                return render(request, 'diagnosis/prof/portugues_finais.html', {
                    'error_message': 'As colunas "Item" e "Habilidade" são obrigatórias na planilha.'
                })

            # Limpar dados antigos e salvar os novos
            HabilidadeProfFinal.objects.all().delete()
            for _, row in df.iterrows():
                HabilidadeProfFinal.objects.create(
                    item=row['Item'],
                    habilidade=row['Habilidade'],
                    turma_101=row.get('101', 0),
                    turma_102=row.get('102', 0),
                    turma_103=row.get('103', 0),
                    turma_104=row.get('104', 0),
                    turma_105=row.get('105', 0),
                    turma_106=row.get('106', 0),
                    turma_107=row.get('107', 0),
                    turma_109=row.get('109', 0),
                    turma_110=row.get('110', 0),
                    turma_112=row.get('112', 0),
                    turma_114=row.get('114', 0),
                    turma_117=row.get('117', 0),
                    turma_119=row.get('119', 0),
                    turma_120=row.get('120', 0),
                    turma_121=row.get('121', 0),
                    turma_124=row.get('124', 0),
                    turma_126=row.get('126', 0),
                    turma_128=row.get('128', 0),
                    turma_129=row.get('129', 0),
                    turma_130=row.get('130', 0),
                    turma_131=row.get('131', 0),
                    turma_134=row.get('134', 0),
                    turma_135=row.get('135', 0),
                    turma_137=row.get('137', 0),
                    turma_138=row.get('138', 0),
                    turma_139=row.get('139', 0),
                    turma_140=row.get('140', 0),
                    turma_142=row.get('142', 0),
                    turma_143=row.get('143', 0),
                    turma_144=row.get('144', 0),
                    turma_145=row.get('145', 0),
                    turma_146=row.get('146', 0),
                    turma_147=row.get('147', 0),
                    turma_171=row.get('171', 0),
                )
        except Exception as e:
            return render(request, 'diagnosis/prof/portugues_finais.html', {
                'error_message': f'Erro ao processar a planilha: {str(e)}'
            })

        return redirect('diagnosis_prof_portugues_finais')

    # Buscar os dados do banco
    habilidades = HabilidadeProfFinal.objects.all()
    total_acertos = sum(h.calcular_acertos() for h in habilidades)
    total_erros = sum(h.calcular_erros() for h in habilidades)
    percentual_acertos = round((total_acertos / (total_acertos + total_erros)) * 100, 2) if (total_acertos + total_erros) > 0 else 0
    percentual_erros = round((total_erros / (total_acertos + total_erros)) * 100, 2) if (total_acertos + total_erros) > 0 else 0

    return render(request, 'diagnosis/prof/portugues_finais.html', {
        'habilidades': habilidades,
        'total_corretas': total_acertos,
        'total_erradas': total_erros,
        'percentual_acertos': percentual_acertos,
        'percentual_erros': percentual_erros,
    })

###***********************************************************************************************************************

def carregar_habilidades_final(request):#####SEGUNDO
    habilidades = HabilidadeProfFinal.objects.all()
    data = []

    for h in habilidades:
        data.append({
            'item': h.item,
            'habilidade': h.habilidade,
            'turmas': {
                '101': h.turma_101,
                '102': h.turma_102,
                '103': h.turma_103,
                '104': h.turma_104,
                '105': h.turma_105,
                '106': h.turma_106,
                '107': h.turma_107,
                '109': h.turma_109,
                '110': h.turma_110,
                '112': h.turma_112,
                '114': h.turma_114,
                '117': h.turma_117,
                '119': h.turma_119,
                '120': h.turma_120,
                '121': h.turma_121,
                '124': h.turma_124,
                '126': h.turma_126,
                '128': h.turma_128,
                '129': h.turma_129,
                '130': h.turma_130,
                '131': h.turma_131,
                '134': h.turma_134,
                '135': h.turma_135,
                '137': h.turma_137,
                '138': h.turma_138,
                '139': h.turma_139,
                '140': h.turma_140,
                '142': h.turma_142,
                '143': h.turma_143,
                '144': h.turma_144,
                '145': h.turma_145,
                '146': h.turma_146,
                '147': h.turma_147,
                '171': h.turma_171,
            }
        })

    return JsonResponse({'habilidades': data})
###***********************************************************************************************************************

def carregar_habilidades(request, tipo):
    # Escolha do modelo com base no tipo
    modelo = HabilidadeProf if tipo == 'inicial' else HabilidadeProfFinal
    habilidades = modelo.objects.all()
    
    # Formatar os dados corretamente
    data = []
    for h in habilidades:
        turmas = {
            '401': h.turma_401,
            '403': h.turma_403,
            '404': h.turma_404,
            '406': h.turma_406,
            '408': h.turma_408,
            '409': h.turma_409,
            '410': h.turma_410,
            '413': h.turma_413,
            '414': h.turma_414,
            '415': h.turma_415,
            '417': h.turma_417,
            '421': h.turma_421,
            '423': h.turma_423,
            '426': h.turma_426,
            '428': h.turma_428,
            '429': h.turma_429,
            '430': h.turma_430,
            '432': h.turma_432,
            '433': h.turma_433,
            '434': h.turma_434,
            '435': h.turma_435,
            '436': h.turma_436,
            '437': h.turma_437,
            '438': h.turma_438,
            '439': h.turma_439,
            '441': h.turma_441,
            '442': h.turma_442,
            '447': h.turma_447,
            '451': h.turma_451,
            '471': h.turma_471,
        }
        data.append({
            'item': h.item,
            'habilidade': h.habilidade,
            'turmas': turmas,
        })
    
    # Retornar os dados no formato esperado pelo JavaScript
    return JsonResponse({'habilidades': data})

###***********************************************************************************************************************
###*****BANCO DE CURRÍCULOS
###***********************************************************************************************************************

def registrar_curriculo(request):
    if request.method == "POST":
        nome = request.POST.get("nome_completo")
        email = request.POST.get("email")
        cpf = request.POST.get("cpf")
        senha = request.POST.get("senha")

        CadastroCandidato.objects.create(
            nome_completo=nome,
            email=email,
            cpf=cpf,
            senha=senha  # Passe em texto plano; será criptografada no método save
        )
        messages.success(request, "Cadastro realizado com sucesso!")
        return redirect("login_candidato")

    return render(request, "banco_curriculos/registrar.html")





###***********************************************************************************************************************

def listar_curriculos(request):
    curriculos = Curriculo.objects.all().order_by('-data_cadastro')
    return render(request, 'curriculos.html', {'curriculos': curriculos})
###***********************************************************************************************************************

from django.http import HttpResponse
from .models import Diretor
from weasyprint import HTML
from django.template.loader import render_to_string

def imprimir_curriculo(request, id):
    try:
        diretor = Diretor.objects.get(id=id)
        html_string = render_to_string('semedapp/imprimir_curriculo.html', {'diretor': diretor})
        pdf = HTML(string=html_string).write_pdf()
        response = HttpResponse(pdf, content_type='application/pdf')
        response['Content-Disposition'] = f'inline; filename="curriculo_{diretor.nome_completo}.pdf"'
        return response
    except Diretor.DoesNotExist:
        return HttpResponse("Currículo não encontrado", status=404)

###***********************************************************************************************************************



FormacaoFormSet = inlineformset_factory(Curriculo, Formacao, form=FormacaoForm, extra=1, can_delete=True)
ExperienciaFormSet = inlineformset_factory(Curriculo, Experiencia, form=ExperienciaForm, extra=1, can_delete=True)

def cadastrar_curriculo(request):
    if request.method == 'POST':
        curriculo_form = CurriculoForm(request.POST, request.FILES)
        formacao_formset = FormacaoFormSet(request.POST, instance=Curriculo())
        experiencia_formset = ExperienciaFormSet(request.POST, instance=Curriculo())

        if curriculo_form.is_valid() and formacao_formset.is_valid() and experiencia_formset.is_valid():
            curriculo = curriculo_form.save()
            formacao_formset.instance = curriculo
            experiencia_formset.instance = curriculo
            formacao_formset.save()
            experiencia_formset.save()
            return redirect('curriculo_sucesso')
    else:
        curriculo_form = CurriculoForm()
        formacao_formset = FormacaoFormSet(instance=Curriculo())
        experiencia_formset = ExperienciaFormSet(instance=Curriculo())

    return render(request, 'cadastrar_curriculo.html', {
        'curriculo_form': curriculo_form,
        'formacao_formset': formacao_formset,
        'experiencia_formset': experiencia_formset,
    })



class CurriculoCreateView(CreateView):
    model = Curriculo
    template_name = 'curriculo_create.html'
    fields = ['nome_completo', 'cpf', 'email', 'telefone', 'data_nascimento']  # Campos específicos
    success_url = '/curriculos/'  # Ajuste o caminho para a listagem de currículos


def curriculo_sucesso(request):
    return render(request, 'curriculo_sucesso.html')



# def cadastrar_curriculo(request):
#     return render(request, 'banco_curriculos/registrar.html')


###***********************************************************************************************************************
###*****CADASTRO DE DEMANDAS
###***********************************************************************************************************************

def tipo_demandas(request):
    return render(request, 'demandas/tipo_demandas.html')
###***********************************************************************************************************************

def cadastro_demandas(request):
    return render(request, 'demandas/cadastro_demandas.html')
###***********************************************************************************************************************

def gestao_demandas(request):
    return render(request, 'demandas/gestao_demandas.html')
###***********************************************************************************************************************

def relatorios_demandas(request):
    return render(request, 'demandas/relatorios_demandas.html')
###***********************************************************************************************************************

def criar_demanda(request):
    if request.method == 'POST':
        form = TipoDemandaForm(request.POST)
        if form.is_valid():
            form.save()
            messages.success(request, 'Demanda cadastrada com sucesso!')
            return redirect('criar_demanda')  # Redireciona para evitar reenvio duplicado
        else:
            messages.error(request, 'Erro ao cadastrar demanda. Verifique os dados.')
    else:
        form = TipoDemandaForm()

    # Buscando todas as demandas
    demandas = TipoDemanda.objects.all()
    return render(request, 'demandas/cadastro_demandas.html', {'form': form, 'demandas': demandas})
###***********************************************************************************************************************

def excluir_demanda(request, id):
    demanda = get_object_or_404(TipoDemanda, id=id)
    demanda.delete()
    messages.success(request, 'Demanda excluída com sucesso!')
    return redirect('criar_demanda')  # Redireciona para a página principal
###***********************************************************************************************************************

def editar_demanda(request, id):
    # Recupera a demanda ou retorna 404 se não for encontrada
    demanda = get_object_or_404(TipoDemanda, id=id)
    
    if request.method == 'POST':
        # Atualiza os dados da demanda com os valores do formulário
        form = TipoDemandaForm(request.POST, instance=demanda)
        if form.is_valid():
            form.save()  # Salva as alterações no banco de dados
            messages.success(request, 'Demanda atualizada com sucesso!')
            return redirect('criar_demanda')  # Redireciona para a página principal ou tabela
        else:
            messages.error(request, 'Erro ao atualizar demanda. Verifique os dados.')
    else:
        # Carrega o formulário com os dados da demanda
        form = TipoDemandaForm(instance=demanda)
    
    return render(request, 'demandas/editar_demanda.html', {'form': form, 'demanda': demanda})
###***********************************************************************************************************************

def gestao_demandas(request):
    demandas = TipoDemanda.objects.all()
    total_demandas = demandas.count()
    concluidas = demandas.filter(status="Finalizado").count()
    em_andamento = demandas.filter(status="Em andamento").count()
    pendentes = demandas.filter(status="Pendente").count()
    
    context = {
        'demandas': demandas,
        'total_demandas': total_demandas,
        'concluidas': concluidas,
        'em_andamento': em_andamento,
        'pendentes': pendentes,
    }
    return render(request, 'demandas/gestao_demandas.html', context)
###***********************************************************************************************************************

def relatorios_demandas(request):
    demandas = TipoDemanda.objects.all()

    # Filtros
    status = request.GET.get('status')
    prioridade = request.GET.get('prioridade')
    date_start = request.GET.get('date_start')
    date_end = request.GET.get('date_end')

    if status:
        demandas = demandas.filter(status=status)
    if prioridade:
        demandas = demandas.filter(prioridade=prioridade)
    if date_start and date_end:
        demandas = demandas.filter(data_cadastro__range=[date_start, date_end])

    context = {
        'demandas': demandas,
        'total_demandas': demandas.count(),
        'concluidas': demandas.filter(status="Finalizado").count(),
        'em_andamento': demandas.filter(status="Em andamento").count(),
        'pendentes': demandas.filter(status="Pendente").count(),
    }
    return render(request, 'demandas/relatorios_demandas.html', context)
###***********************************************************************************************************************

from django.http import JsonResponse

def listar_demandas(request):
    demandas = TipoDemanda.objects.values(
        "id",
        "data_cadastro",
        "data_entrega",
        "descricao",
        "responsavel",
        "destinatario",
        "prioridade",
        "status"
    )
    return JsonResponse(list(demandas), safe=False)
###***********************************************************************************************************************

def criar_tipo_demanda(request):
    if request.method == 'POST':
        form = TipoDemandaForm(request.POST)
        if form.is_valid():
            form.save()
            messages.success(request, 'Tipo de demanda criado com sucesso!')
            return redirect('tipo_demandas')
        else:
            messages.error(request, 'Erro ao criar tipo de demanda. Verifique os dados.')
    else:
        form = TipoDemandaForm()
    
    return render(request, 'demandas/criar_tipo_demanda.html', {'form': form})


from django.contrib.auth.decorators import login_required, user_passes_test

# Função para verificar se o usuário pertence ao grupo "Cadastro de Demandas"
def is_demanda_user(user):
    return user.groups.filter(name="Cadastro de Demandas").exists()



def tipo_demandas(request):
    return render(request, 'demandas/tipo_demandas.html')


def cadastro_demandas(request):
    return render(request, 'demandas/cadastro_demandas.html')


def criar_demanda(request):
    if request.method == 'POST':
        form = TipoDemandaForm(request.POST)
        if form.is_valid():
            form.save()
            messages.success(request, 'Demanda cadastrada com sucesso!')
            return redirect('criar_demanda')
        else:
            messages.error(request, 'Erro ao cadastrar demanda. Verifique os dados.')
    else:
        form = TipoDemandaForm()

    demandas = TipoDemanda.objects.all()
    return render(request, 'demandas/cadastro_demandas.html', {'form': form, 'demandas': demandas})





def gestao_demandas(request):
    demandas = TipoDemanda.objects.all()
    total_demandas = demandas.count()
    concluidas = demandas.filter(status="Finalizado").count()
    em_andamento = demandas.filter(status="Em andamento").count()
    pendentes = demandas.filter(status="Pendente").count()

    context = {
        'demandas': demandas,
        'total_demandas': total_demandas,
        'concluidas': concluidas,
        'em_andamento': em_andamento,
        'pendentes': pendentes,
    }
    return render(request, 'demandas/gestao_demandas.html', context)


from django.contrib.auth.decorators import user_passes_test

# Função para verificar se o usuário está no grupo "Acesso Total Demandas"
def is_acesso_total_demandas(user):
    return user.groups.filter(name="Acesso Total Demandas").exists() or user.is_superuser


def tipo_demandas(request):
    return render(request, 'demandas/tipo_demandas.html')







###***********************************************************************************************************************

def gerar_relatorio_pdf(request):
    # Filtrar as demandas com base nos filtros aplicados (opcional)
    status = request.GET.get('status', '')
    prioridade = request.GET.get('prioridade', '')
    date_start = request.GET.get('date_start', '')
    date_end = request.GET.get('date_end', '')

    demandas = TipoDemanda.objects.all()

    if status:
        demandas = demandas.filter(status=status)
    if prioridade:
        demandas = demandas.filter(prioridade=prioridade)
    if date_start and date_end:
        demandas = demandas.filter(data_cadastro__range=[date_start, date_end])

    # Renderizar o template para o PDF
    template_path = 'demandas/relatorio_pdf.html'
    context = {'demandas': demandas}
    template = get_template(template_path)
    html = template.render(context)

    # Criar o PDF
    response = HttpResponse(content_type='application/pdf')
    response['Content-Disposition'] = 'attachment; filename="relatorio_demandas.pdf"'
    pisa_status = pisa.CreatePDF(html, dest=response)

    # Verificar se houve erros
    if pisa_status.err:
        return HttpResponse('Erro ao gerar o PDF', status=500)

    return response
###***********************************************************************************************************************

def editar_escola(request, id):
    escola = get_object_or_404(Escola, id=id)
    if request.method == "POST":
        # Processa o formulário aqui
        pass
    return render(request, 'editar_escola.html', {'escola': escola})
###***********************************************************************************************************************

def adicionar_escola(request):
    if request.method == 'POST':
        form = EscolaForm(request.POST, request.FILES)
        if form.is_valid():
            form.save()
            return redirect('listar_escolas')  # Certifique-se de ajustar para a URL correta
    else:
        form = EscolaForm()

    return render(request, 'escolas/adicionar_escola.html', {'form': form})
###***********************************************************************************************************************

def cadastrar_professor(request):
    if request.method == 'POST':
        form = ProfessorForm(request.POST, request.FILES)  # Inclua `request.FILES` para arquivos.
        if form.is_valid():
            form.save()
            return redirect('listar_professores')  # Redireciona para a página de listagem.
    else:
        form = ProfessorForm()

    return render(request, 'cadastrar_professor.html', {'form': form})
###***********************************************************************************************************************

def sucesso(request):
    return render(request, 'sucesso.html')
###***********************************************************************************************************************

def curriculo_sucesso(request):
    return render(request, 'semedapp/curriculo_sucesso.html')
###***********************************************************************************************************************

def listar_professores(request):
    professores = Professor.objects.all()
    return render(request, 'semedapp/listar_professores.html', {'professores': professores})
###***********************************************************************************************************************

def editar_professor(request, professor_id):
    professor = get_object_or_404(Professor, id=professor_id)
    if request.method == 'POST':
        form = ProfessorForm(request.POST, request.FILES, instance=professor)
        if form.is_valid():
            form.save()
            return redirect('listar_professores')
    else:
        form = ProfessorForm(instance=professor)
    return render(request, 'semedapp/editar_professor.html', {'form': form, 'professor': professor})
###***********************************************************************************************************************

def excluir_professor(request, professor_id):
    professor = get_object_or_404(Professor, id=professor_id)
    professor.delete()
    messages.success(request, "Professor excluído com sucesso.")
    return redirect('listar_professores')
###***********************************************************************************************************************

# View para relatórios de professores
def relatorios_professores(request):
    professores = Professor.objects.all()
    # Lógica para gerar relatórios (simplesmente listando os professores por enquanto)
    return render(request, 'semedapp/relatorios_professores.html', {'professores': professores})
###***********************************************************************************************************************

def criar_professor(request):
    if request.method == 'POST':
        # Extração de dados do formulário
        nome_completo = request.POST.get('nome_completo')
        cpf = request.POST.get('cpf')
        rg = request.POST.get('rg', None)
        email = request.POST.get('email')
        telefone = request.POST.get('telefone')
        endereco = request.POST.get('endereco')
        bairro = request.POST.get('bairro', None)
        cidade = request.POST.get('cidade', 'Canaã dos Carajás')  # Valor padrão
        cep = request.POST.get('cep', None)
        estado_civil = request.POST.get('estado_civil')
        sexo = request.POST.get('sexo')
        data_nascimento = request.POST.get('data_nascimento', None)

        # Formação Acadêmica
        formacao_academica = request.POST.get('formacao_academica', None)
        curso = request.POST.get('curso', None)
        instituicao = request.POST.get('instituicao', None)
        ano_conclusao = request.POST.get('ano_conclusao', None)

        # Experiência Profissional
        experiencia_profissional = request.POST.get('experiencia_profissional', None)

        # Arquivos
        foto = request.FILES.get('foto', None)
        curriculo_pdf = request.FILES.get('arquivo_curriculo', None)
        certificados_pdf = request.FILES.get('arquivo_certificados', None)

        # Salvando os dados no banco de dados
        professor = Professor(
            nome_completo=nome_completo,
            cpf=cpf,
            rg=rg,
            email=email,
            telefone=telefone,
            endereco=endereco,
            bairro=bairro,
            cidade=cidade,
            cep=cep,
            estado_civil=estado_civil,
            sexo=sexo,
            data_nascimento=data_nascimento,
            formacao_academica=formacao_academica,
            curso=curso,
            instituicao=instituicao,
            ano_conclusao=ano_conclusao,
            experiencia_profissional=experiencia_profissional,
            foto=foto,
            curriculo_pdf=curriculo_pdf,
            certificados_pdf=certificados_pdf,
        )
        professor.save()  # Salva no banco de dados

        return redirect('listar_professores')  # Redirecione após salvar

    return render(request, 'semedapp/criar_professor.html')
###***********************************************************************************************************************
from django.db.models import Count
from django.db.models.functions import Lower
from django.shortcuts import render
from .models import Diretor
import unicodedata
from django.db.models import Value


# Função para normalizar texto e evitar erro com NoneType
def normalizar_texto(valor):
    return valor.strip().title() if valor else ""

def listar_diretores(request):
    # Obtém todos os diretores
    diretores = Diretor.objects.all()

    # 📌 Obtém as formações acadêmicas disponíveis sem valores duplicados ou vazios
    formacoes_disponiveis = sorted(
        set(
            normalizar_texto(diretor.formacao_academica)
            for diretor in diretores.exclude(formacao_academica__isnull=True)
            .exclude(formacao_academica__exact='')
        )
    )

    # 📌 Obtém os filtros da URL
    search_query = request.GET.get('search', '').strip()
    cidade_filtro = request.GET.get('cidade', '').strip()
    cargo_filtro = request.GET.get('cargo', '').strip()

    # 📌 Obtém cidades e cargos distintos do banco, removendo valores vazios e ordenando
    cidades_disponiveis = sorted(
        set(
            normalizar_texto(cidade)
            for cidade in Diretor.objects.exclude(cidade__isnull=True)
            .exclude(cidade__exact='')
            .values_list('cidade', flat=True)
            .distinct()
        )
    )

    cargos_disponiveis = sorted(
        set(
            normalizar_texto(cargo)
            for cargo in Diretor.objects.exclude(cargo__isnull=True)
            .exclude(cargo__exact='')
            .values_list('cargo', flat=True)
            .distinct()
        )
    )

    # 📌 Aplicação dos filtros
    if search_query:
        diretores = diretores.filter(nome_completo__icontains=search_query) | diretores.filter(cpf__icontains=search_query)

    if cidade_filtro:
        diretores = diretores.filter(cidade__iexact=cidade_filtro)

    if cargo_filtro:
        diretores = diretores.filter(cargo__iexact=cargo_filtro)

    # 📌 Total de currículos filtrados
    total_curriculos = diretores.count()

    # 📌 Total de formações acadêmicas distintas
    total_formacoes = len(formacoes_disponiveis)

    # 📌 Contagem de cidades únicas atendidas
    total_cidades = len(cidades_disponiveis)

    # Contexto para o template
    context = {
        'diretores': diretores,
        'formacoes_disponiveis': formacoes_disponiveis,  # Lista única de formações acadêmicas
        'total_curriculos': total_curriculos,  # Total de currículos filtrados
        'search_query': search_query,  # Mantém a pesquisa após o filtro
        'cidade_filtro': cidade_filtro,  # Mantém a cidade selecionada
        'cargo_filtro': cargo_filtro,  # Mantém o cargo selecionado
        'cidades_disponiveis': cidades_disponiveis,  # Lista única de cidades ordenadas
        'cargos_disponiveis': cargos_disponiveis,  # Lista única de cargos ordenados
        'total_formacoes': total_formacoes,  # Total de formações acadêmicas
        'total_cidades': total_cidades,  # Total de cidades atendidas
    }

    return render(request, 'semedapp/listar_diretores.html', context)



###***********************************************************************************************************************
# View para cadastrar diretor


from .forms import DiretorForm, ExperienciaFormSet, FormacaoFormSet, CertificadoFormSet


def cadastrar_diretor(request):
    if request.method == 'POST':
        # Captura os campos do formulário
        nome_completo = request.POST.get('nome_completo')
        cpf = request.POST.get('cpf')
        rg = request.POST.get('rg')
        email = request.POST.get('email')
        telefone = request.POST.get('telefone')
        endereco = request.POST.get('endereco')
        bairro = request.POST.get('bairro')
        cidade = request.POST.get('cidade')
        cep = request.POST.get('cep')
        estado_civil = request.POST.get('estado_civil')
        # Ajusta o estado civil para capitalização correta
        if estado_civil:
            estado_civil = estado_civil.capitalize()
        sexo = request.POST.get('sexo')
        data_nascimento = request.POST.get('data_nascimento') or None
        formacao_academica = request.POST.get('formacao_academica')
        curso = request.POST.get('curso')
        instituicao = request.POST.get('instituicao')
        ano_conclusao = request.POST.get('ano_conclusao')
        experiencia_profissional = request.POST.get('experiencia_profissional')

        # Novos campos de experiência profissional
        empresa = request.POST.get('empresa')
        cargo = request.POST.get('cargo')
        data_inicio = request.POST.get('data_inicio') or None
        data_fim = request.POST.get('data_fim') or None

        # Processa arquivos enviados
        foto = request.FILES.get('foto')
        curriculo_pdf = request.FILES.get('curriculo_pdf')
        certificados_pdf = request.FILES.get('certificados_pdf')

        # Verifica se o CPF ou e-mail já estão cadastrados
        if Diretor.objects.filter(cpf=cpf).exists():
            messages.error(request, 'O CPF informado já está cadastrado.')
            return redirect('cadastrar_diretor')
        
        if Diretor.objects.filter(email=email).exists():
            messages.error(request, 'O e-mail informado já está cadastrado.')
            return redirect('cadastrar_diretor')

        # Cria a instância do modelo e salva no banco de dados
        try:
            diretor = Diretor(
                nome_completo=nome_completo,
                cpf=cpf,
                rg=rg,
                email=email,
                telefone=telefone,
                endereco=endereco,
                bairro=bairro,
                cidade=cidade,
                cep=cep,
                estado_civil=estado_civil,
                sexo=sexo,
                data_nascimento=data_nascimento,
                formacao_academica=formacao_academica,
                curso=curso,
                instituicao=instituicao,
                ano_conclusao=ano_conclusao,
                experiencia_profissional=experiencia_profissional,
                empresa=empresa,
                cargo=cargo,
                data_inicio=data_inicio,
                data_fim=data_fim,
                foto=foto,
                curriculo_pdf=curriculo_pdf,
                certificados_pdf=certificados_pdf,
            )
            diretor.save()
            messages.success(request, f"Currículo de {nome_completo} cadastrado com sucesso!")
            return redirect('cadastrar_diretor')
        except Exception as e:
            messages.error(request, f"Erro ao cadastrar o currículo: {str(e)}")
            return redirect('cadastrar_diretor')

    return render(request, 'semedapp/cadastrar_diretor.html')

###***********************************************************************************************************************

def visualizar_diretor(request, id):
    diretor = get_object_or_404(Diretor, id=id)
    return render(request, 'semedapp/visualizar_diretor.html', {'diretor': diretor})
###***********************************************************************************************************************

class VisualizarDiretorView(DetailView):
    model = Diretor
    template_name = 'semedapp/visualizar_diretor.html'
    context_object_name = 'diretor'
###***********************************************************************************************************************

def editar_diretor(request, id):
    diretor = get_object_or_404(Diretor, id=id)
    
    if request.method == 'POST':
        form = DiretorForm(request.POST, instance=diretor)
        if form.is_valid():
            form.save()
            return redirect('listar_diretores')  # Ajuste o nome da URL para redirecionar após salvar
    else:
        form = DiretorForm(instance=diretor)
    
    return render(request, 'semedapp/editar_diretor.html', {'form': form, 'diretor': diretor})
###***********************************************************************************************************************

def excluir_diretor(request, id):
    diretor = get_object_or_404(Diretor, id=id)
    diretor.delete()
    return redirect('listar_diretores')  # Ajuste o nome para a URL de listagem
###***********************************************************************************************************************

class DiretorDeleteView(DeleteView):
    model = Diretor
    success_url = reverse_lazy('listar_diretores')
    template_name = 'semedapp/confirmar_exclusao.html'
###***********************************************************************************************************************

# def imprimir_curriculo(request, diretor_id):
#     try:
#         diretor = Diretor.objects.get(id=diretor_id)
#     except Diretor.DoesNotExist:
#         return HttpResponse("Diretor não encontrado", status=404)

#     html_string = render_to_string('semedapp/imprimir_curriculo.html', {'diretor': diretor})

#     pdf = HTML(string=html_string).write_pdf()

#     response = HttpResponse(pdf, content_type='application/pdf')
#     response['Content-Disposition'] = f'inline; filename="curriculo_{diretor.nome_completo}.pdf"'
#     return response

###***********************************************************************************************************************

import qrcode
from io import BytesIO
from django.shortcuts import get_object_or_404
from django.core.files.base import ContentFile
from django.http import HttpResponse


def gerar_qrcode(request, diretor_id):
    diretor = get_object_or_404(Diretor, id=diretor_id)
    data = f"https://example.com/curriculo/{diretor.id}"  # URL personalized for the director

    qr = qrcode.QRCode(
        version=1,
        error_correction=qrcode.constants.ERROR_CORRECT_L,
        box_size=10,
        border=4,
    )
    qr.add_data(data)
    qr.make(fit=True)

    img = qr.make_image(fill_color="black", back_color="white")

    # Save QR Code to the diretor's folder
    buffer = BytesIO()
    img.save(buffer, format="PNG")
    buffer.seek(0)
    filename = f"qrcode_{diretor.id}.png"
    diretor.qr_code.save(filename, ContentFile(buffer.read()))
    buffer.close()

    return HttpResponse(diretor.qr_code.url)


###***********************************************************************************************************************

from django.shortcuts import render, redirect
from django.contrib import messages
from django.contrib.auth.models import User
from django.contrib.auth import authenticate, login
from .models import CadastroCandidato

def registrar_candidato(request):
    if request.method == "POST":
        nome_completo = request.POST.get("nome_completo")
        email = request.POST.get("email")
        cpf = request.POST.get("cpf")
        senha = request.POST.get("password")
        confirm_senha = request.POST.get("confirm_password")

        # Verifica se as senhas coincidem
        if senha != confirm_senha:
            messages.error(request, "As senhas não coincidem!")
            return redirect("registrar_candidato")

        try:
            # Cria o usuário no sistema
            user = User.objects.create_user(username=email, email=email, password=senha)
            user.first_name = nome_completo
            user.save()

            # Cria o perfil de candidato
            candidato = CadastroCandidato.objects.create(
                user=user,
                nome_completo=nome_completo,
                cpf=cpf,
            )

            # Autentica e loga o usuário automaticamente
            user = authenticate(request, username=email, password=senha)
            if user is not None:
                login(request, user)
                messages.success(request, "Cadastro realizado com sucesso!")
                return redirect("area_candidato")

        except Exception as e:
            messages.error(request, f"Erro ao cadastrar: {str(e)}")
            return redirect("registrar_candidato")

    return render(request, "banco_curriculos/registrar.html")


###***********************************************************************************************************************

from django.shortcuts import render, redirect
from django.contrib import messages
from semedapp.models import CadastroCandidato

def login_candidato(request):
    if request.method == "POST":
        email = request.POST.get("email")
        senha = request.POST.get("password")

        candidato = CadastroCandidato.objects.filter(email=email).first()
        if candidato and candidato.check_senha(senha):
            # Salve o ID e o nome na sessão para identificar o usuário logado
            request.session["candidato_id"] = candidato.id
            request.session["candidato_nome"] = candidato.nome_completo
            messages.success(request, f"Bem-vindo, {candidato.nome_completo}!")
            return redirect("area_candidato")  # Redirecione para a área do candidato
        else:
            messages.error(request, "E-mail ou senha inválidos.")
    
    return render(request, "banco_curriculos/login.html")




###***********************************************************************************************************************

@login_required
def area_administrativa(request):
    return render(request, "semedapp/area_administrativa.html", {"user": request.user})


###***********************************************************************************************************************
def exportar_excel(request):
    # Cria um workbook e uma planilha
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Diretores"

    # Cabeçalhos da tabela
    ws.append(["Nome", "CPF", "RG", "Email", "Telefone", "Endereço", "Bairro", "Cidade", "Nascimento"])

    # Adiciona os dados dos diretores
    diretores = Diretor.objects.all()
    for diretor in diretores:
        ws.append([
            diretor.nome_completo,
            diretor.cpf,
            diretor.rg,
            diretor.email,
            diretor.telefone,
            diretor.endereco,
            diretor.bairro,
            diretor.cidade,
            diretor.data_nascimento.strftime("%d/%m/%Y") if diretor.data_nascimento else "",
        ])

    # Configura a resposta HTTP
    response = HttpResponse(
        content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )
    response["Content-Disposition"] = 'attachment; filename="diretores.xlsx"'
    wb.save(response)

    return response

###***********************************************************************************************************************
from reportlab.lib.pagesizes import landscape, letter
from reportlab.pdfgen import canvas
from django.http import HttpResponse
from .models import CurriculoAntigo
from datetime import datetime
import os

# ✅ Exportação para PDF (com cabeçalho, rodapé, formato paisagem)
def exportar_pdf(request):
    # Criar resposta HTTP com tipo de arquivo PDF
    response = HttpResponse(content_type='application/pdf')
    response['Content-Disposition'] = 'attachment; filename="curriculos.pdf"'

    # Criar documento PDF em formato paisagem
    p = canvas.Canvas(response, pagesize=landscape(letter))
    largura, altura = landscape(letter)

    # Caminho da imagem para o cabeçalho
    imagem_cabecalho = os.path.join("static", "img", "logo_semed.png")  # Ajuste conforme seu projeto

    def desenhar_cabecalho_rodape(canvas_obj, doc_page_num):
        """ Adiciona cabeçalho e rodapé em cada página """
        # Cabeçalho
        if os.path.exists(imagem_cabecalho):  # Verifica se a imagem existe
            canvas_obj.drawImage(imagem_cabecalho, 30, altura - 50, width=100, height=40, preserveAspectRatio=True)

        canvas_obj.setFont("Helvetica-Bold", 14)
        canvas_obj.drawString(150, altura - 40, "Relatório de Currículos Antigos - SEMED Canaã dos Carajás")

        # Rodapé
        canvas_obj.setFont("Helvetica", 10)
        data_atual = datetime.now().strftime('%d/%m/%Y %H:%M')
        canvas_obj.drawString(30, 20, f"Gerado em: {data_atual}")
        canvas_obj.drawString(largura - 100, 20, f"Página {doc_page_num}")

    # Cabeçalhos da tabela
    p.setFont("Helvetica-Bold", 10)
    y = altura - 80  # Posição vertical inicial para o cabeçalho da tabela
    colunas = ["CPF", "Nome", "Email", "Cargo", "Telefone"]  # Mudamos "Data Nascimento" para "Telefone"
    espacamento = [50, 150, 350, 500, 650]

    # Criar numeração de páginas
    pagina_atual = 1
    desenhar_cabecalho_rodape(p, pagina_atual)

    for i, coluna in enumerate(colunas):
        p.drawString(espacamento[i], y, coluna)

    # Linhas da tabela
    p.setFont("Helvetica", 9)
    y -= 20

    curriculos = CurriculoAntigo.objects.all().values("cpf", "nome", "email", "cargo", "fone1")  # Agora pegamos "fone1"

    for curriculo in curriculos:
        cpf = str(curriculo["cpf"]) if curriculo["cpf"] else ""
        nome = str(curriculo["nome"]) if curriculo["nome"] else ""
        email = str(curriculo["email"]) if curriculo["email"] else ""
        cargo = str(curriculo["cargo"]) if curriculo["cargo"] else "Sem Cargo"
        telefone = str(curriculo["fone1"]) if curriculo["fone1"] else "Sem telefone"  # Pegamos telefone em vez de data_nascimento

        p.drawString(50, y, cpf)
        p.drawString(150, y, nome)
        p.drawString(350, y, email)
        p.drawString(500, y, cargo)
        p.drawString(650, y, telefone)  # Alterado para telefone
        y -= 15  # Move a linha para baixo

        # Evita que o conteúdo saia da página
        if y < 50:
            p.showPage()
            pagina_atual += 1
            desenhar_cabecalho_rodape(p, pagina_atual)
            p.setFont("Helvetica", 9)
            y = altura - 80  # Reinicia a posição vertical para a nova página

    p.save()
    return response



###***********************************************************************************************************************
def exportar_csv(request):
    response = HttpResponse(content_type="text/csv")
    response["Content-Disposition"] = 'attachment; filename="diretores.csv"'

    writer = csv.writer(response)
    writer.writerow(["Nome", "CPF", "RG", "Email", "Telefone", "Endereço", "Bairro", "Cidade", "Nascimento"])

    diretores = Diretor.objects.all()
    for diretor in diretores:
        writer.writerow([
            diretor.nome_completo,
            diretor.cpf,
            diretor.rg,
            diretor.email,
            diretor.telefone,
            diretor.endereco,
            diretor.bairro,
            diretor.cidade,
            diretor.data_nascimento.strftime("%d/%m/%Y") if diretor.data_nascimento else "",
        ])

    return response

###***********************************************************************************************************************



import pandas as pd
from django.http import HttpResponse
from django.utils.timezone import is_aware, localtime
from datetime import datetime
from .models import CurriculoAntigo

def exportar_xls(request):
    # Selecionar apenas as colunas desejadas
    colunas = [
        "cpf", "nome", "email", "senha", "data_nascimento", "sexo", "naturalidade", "naturalidade_uf",
        "pne", "pne_detalhe", "curriculo_lattes", "logradouro", "numero", "bairro", "cep", "complemento",
        "uf", "municipio", "fone1", "fone2", "formacao_nivel", "formacao_instituicao", "formacao_curso",
        "formacao_situacao", "formacao_inicio", "formacao_conclusao", "data_cadastro", "data_update", "ip", "cargo"
    ]

    # Recuperar os dados e criar um DataFrame
    curriculos = list(CurriculoAntigo.objects.values(*colunas))
    df = pd.DataFrame(curriculos)

    # ✅ Corrigir os campos de data
    campos_de_data = ["data_nascimento", "formacao_inicio", "formacao_conclusao", "data_cadastro", "data_update"]
    for campo in campos_de_data:
        if campo in df.columns:
            df[campo] = df[campo].apply(lambda x: formatar_data(x))

    # Criar resposta HTTP com o arquivo Excel
    response = HttpResponse(content_type='application/vnd.ms-excel')
    response['Content-Disposition'] = 'attachment; filename="curriculos.xls"'

    df.to_excel(response, index=False)

    return response

# ✅ Função auxiliar para formatar datas corretamente
def formatar_data(valor):
    if pd.isna(valor):  # Verifica se o valor está vazio (NaN)
        return ""
    elif isinstance(valor, datetime):  # Se for datetime, converte para naive e formata
        if is_aware(valor):
            valor = localtime(valor)  # Remove timezone
        return valor.strftime('%d/%m/%Y')
    elif isinstance(valor, str):  # Se for string, retorna sem modificação
        return valor
    else:  # Para valores do tipo date (sem timezone)
        return valor.strftime('%d/%m/%Y')


###***********************************************************************************************************************
# Página inicial do Portal da Transparência
def transparencia_index(request):
    return render(request, 'transparencia/index.html')
###***********************************************************************************************************************
# Página de Folha de Pagamento
def transparencia_folha_pagamento(request):
    # Lógica da folha de pagamento
    return render(request, 'transparencia/folha_pagamento.html')
###***********************************************************************************************************************

# Página de Alimentação
def transparencia_alimentacao(request):
    # Lógica para carregar os dados de alimentação
    return render(request, 'transparencia/alimentacao.html')
###***********************************************************************************************************************

# Página de Bens de Consumo
def transparencia_bens_consumo(request):
    # Lógica para carregar os dados de bens e consumo
    return render(request, 'transparencia/bens_consumo.html')
###***********************************************************************************************************************

# Página de Outros Gastos
def outros_gastos(request):
    return render(request, 'transparencia/outros_gastos.html')
###***********************************************************************************************************************

def privacy_policy(request):
    return render(request, 'semedapp/privacy_policy.html')
###***********************************************************************************************************************

def contact_us(request):
    return render(request, 'semedapp/contact_us.html')
###***********************************************************************************************************************

def help(request):
    return render(request, 'semedapp/help.html')
###***********************************************************************************************************************

def pesquisar_curriculos(request):
    return render(request, 'semedapp/pesquisar_curriculos.html')
###***********************************************************************************************************************
from django.shortcuts import get_object_or_404, render, redirect
from django.contrib import messages
from semedapp.models import FormacaoAcademica
from .forms import FormacaoAcademicaForm


def cadastrar_formacao(request, diretor_id):
    diretor = get_object_or_404(Diretor, id=diretor_id)

    if request.method == 'POST':
        nivel = request.POST.get('nivel')
        instituicao = request.POST.get('instituicao')
        curso = request.POST.get('curso')
        ano_conclusao = request.POST.get('ano_conclusao')

        try:
            FormacaoAcademica.objects.create(
                diretor=diretor,
                nivel=nivel,
                instituicao=instituicao,
                curso=curso,
                ano_conclusao=ano_conclusao,
            )
            messages.success(request, f"Formação adicionada com sucesso para o diretor {diretor.nome_completo}!")
            return redirect('detalhes_diretor', diretor_id=diretor.id)
        except Exception as e:
            messages.error(request, f"Erro ao cadastrar a formação: {str(e)}")
            return redirect('cadastrar_formacao', diretor_id=diretor.id)

    return render(request, 'semedapp/cadastrar_formacao.html', {'diretor': diretor})
###***********************************************************************************************************************

from django.shortcuts import render

def registrar_novo_curriculo(request):
    return render(request, 'banco_curriculos/registrar.html')
###***********************************************************************************************************************

from django.shortcuts import render, redirect
from django.contrib import messages
from .models import CadastroCandidato
from django.contrib.auth.hashers import make_password
from django.db import IntegrityError

def cadastrar_candidato(request):
    if request.method == "POST":
        nome_completo = request.POST.get("nome_completo")
        email = request.POST.get("email")
        cpf = request.POST.get("cpf")
        senha = request.POST.get("password")
        confirm_senha = request.POST.get("confirm_password")

        if senha != confirm_senha:
            messages.error(request, "As senhas não coincidem!")
            return redirect("registrar_candidato")

        try:
            # Cria o candidato no banco de dados
            candidato = CadastroCandidato.objects.create(
                nome_completo=nome_completo,
                email=email,
                cpf=cpf,
                senha=make_password(senha),  # Criptografa a senha
            )

            # Salva o ID do candidato na sessão para login automático
            request.session["candidato_id"] = candidato.id
            messages.success(request, f"Bem-vindo, {candidato.nome_completo}!")
            return redirect("area_candidato")

        except IntegrityError as e:
            # Mensagem de erro no caso de duplicação de email ou CPF
            messages.error(request, "Já existe um candidato com este e-mail ou CPF.")
            print(f"Erro de integridade: {e}")

        except Exception as e:
            # Tratamento de outros erros
            messages.error(request, "Erro ao cadastrar candidato. Tente novamente.")
            print(f"Erro ao cadastrar candidato: {e}")

    return render(request, "banco_curriculos/registrar.html")
###***********************************************************************************************************************

from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from .models import CadastroCandidato, Diretor

def area_candidato(request):
    """Exibe a área do candidato e busca seu currículo na tabela Diretor"""

    # 🔍 Obtém o ID do candidato da sessão
    candidato_id = request.session.get('candidato_id')

    if not candidato_id:
        messages.error(request, "Você precisa estar logado para acessar esta página.")
        return redirect("login_candidato")

    # 🔍 Busca o candidato e verifica se existe
    candidato = get_object_or_404(CadastroCandidato, id=candidato_id)

    # 🔍 Busca o diretor associado pelo CPF do candidato
    diretor = Diretor.objects.filter(cpf=candidato.cpf).first()

    return render(request, 'banco_curriculos/area_candidato.html', {
        'candidato': candidato,
        'diretor': diretor  # Passa os dados do diretor para o template
    })



###***********************************************************************************************************************

@login_required
def adicionar_curriculo(request):
    """View para adicionar um novo currículo."""
    if request.method == 'POST':
        form = CurriculoForm(request.POST, request.FILES)
        if form.is_valid():
            curriculo = form.save(commit=False)
            curriculo.usuario = request.user  # Associa o currículo ao usuário autenticado
            curriculo.save()
            messages.success(request, 'Currículo cadastrado com sucesso!')
            return redirect('area_candidato')
    else:
        form = CurriculoForm()

    return render(request, 'banco_curriculos/adicionar_curriculo.html', {'form': form})
###***********************************************************************************************************************

from django.shortcuts import redirect
from django.contrib.auth import logout
from django.contrib import messages

def logout_candidato(request):
    # Destroi a sessão
    logout(request)
    # Adiciona uma mensagem de sucesso
    messages.success(request, "Você saiu com sucesso.")
    # Redireciona para a página de login
    return redirect("login_candidato")
###***********************************************************************************************************************

from django.shortcuts import render
from django.contrib.auth.decorators import login_required

@login_required
def minhas_inscricoes(request):
    # Obtém as inscrições relacionadas ao usuário logado
    inscricoes = request.user.inscricoes_usuario.all()  # Use o related_name definido no modelo

    # Renderiza a página com uma mensagem amigável caso não existam inscrições
    contexto = {
        'inscricoes': inscricoes,
        'mensagem': 'Nenhuma inscrição encontrada!' if not inscricoes else None
    }
    return render(request, 'banco_curriculos/minhas_inscricoes.html', contexto)
###***********************************************************************************************************************

import qrcode
from django.core.files.base import ContentFile

def gerar_qr_code(diretor):
    data = f"Currículo de {diretor.nome_completo}"
    qr = qrcode.make(data)

    # Salvar o QR Code na pasta correta
    qr_code_path = f"diretores_qrcodes/qr_code_{diretor.id}.png"
    with open(f"{settings.MEDIA_ROOT}/{qr_code_path}", "wb") as qr_file:
        qr.save(qr_file)

    diretor.qr_code.save(qr_code_path, ContentFile(qr_file.read()), save=True)
###***********************************************************************************************************************

from django.http import JsonResponse
from semedapp.models import Diretor

def verificar_cpf(request):
    cpf = request.GET.get('cpf')
    if Diretor.objects.filter(cpf=cpf).exists():
        return JsonResponse({'success': True})
    return JsonResponse({'success': False, 'message': 'CPF não encontrado no sistema!'})
###***********************************************************************************************************************

## aque a impressão de currículo para a ára addministrativa

# def imprimir_curriculo(request, cpf):
#     try:

#         diretor = Diretor.objects.get(cpf=cpf)
        

#         html_string = render_to_string('semedapp/imprimir_curriculo.html', {'diretor': diretor})
        

#         pdf = HTML(string=html_string).write_pdf()
        

#         response = HttpResponse(pdf, content_type='application/pdf')
#         response['Content-Disposition'] = f'inline; filename="curriculo_{diretor.nome_completo}.pdf"'
#         return response

#     except Diretor.DoesNotExist:

#         return render(request, 'semedapp/imprimir_curriculo.html', {
#             'alerta': f'Currículo não encontrado para o CPF: {cpf}.'
#         })


###***********************************************************************************************************************   

from django import forms
from .models import Candidato

class FormDeAtualizacao(forms.ModelForm):
    class Meta:
        model = Candidato
        fields = ['nome_completo', 'email', 'cpf']  # Inclua os campos que deseja permitir edição
###***********************************************************************************************************************

from django.shortcuts import redirect

def editar_perfil(request):
    if request.method == "POST":
        cpf = request.POST.get("cpf", "").strip()

        candidato = get_object_or_404(CadastroCandidato, cpf=cpf)

        candidato.nome_completo = request.POST.get("nome_completo")
        candidato.email = request.POST.get("email")

        nova_senha = request.POST.get("senha")
        if nova_senha:
            candidato.senha = nova_senha

        candidato.save()

        # ✅ Redireciona para a área do candidato após salvar
        return redirect("area_candidato")  # Certifique-se de que esta URL está correta no seu `urls.py`

    return JsonResponse({"error": True, "message": "Método inválido"}, status=400)
###***********************************************************************************************************************

# from django.contrib.auth.models import User
# from .models import Candidato

# # Exemplo ao criar um novo candidato
# def criar_candidato():
#     user = User.objects.create_user(username='usuario', password='senha123')
#     candidato = Candidato.objects.create(
#         user=user,
#         nome_completo='João da Silva',
#         cpf='123.456.789-00',
#         email='joao@example.com',
#     )
###***********************************************************************************************************************

from django.shortcuts import redirect
from functools import wraps

def module_permission_required(module_name):
    def decorator(view_func):
        @wraps(view_func)
        def _wrapped_view(request, *args, **kwargs):
            if not request.user.is_authenticated:
                return redirect("login")
            
            # Verificar se o usuário é administrador ou tem permissão para o módulo
            if (
                request.user.role == "Administrador"
                or RoleModulePermission.objects.filter(
                    role=request.user.role,
                    module__module_name=module_name,
                ).exists()
            ):
                return view_func(request, *args, **kwargs)
            return redirect("acesso_negado")
        return _wrapped_view
    return decorator
###***********************************************************************************************************************

@module_permission_required("Cadastro Demanda")
def cadastro_demanda_view(request):
    return render(request, "modulos/cadastro_demanda.html")
###***********************************************************************************************************************

@module_permission_required("Indicadores")
def indicadores_view(request):
    return render(request, "modulos/indicadores.html")
###***********************************************************************************************************************

from django.http import JsonResponse
from django.shortcuts import get_object_or_404
from semedapp.models import CadastroCandidato, Diretor

def visualizar_curriculo(request):
    candidato_id = request.session.get("candidato_id")

    if not candidato_id:
        return JsonResponse({"success": False, "message": "Usuário não autenticado."})

    candidato = get_object_or_404(CadastroCandidato, id=candidato_id)
    diretor = Diretor.objects.filter(cpf=candidato.cpf).first()

    if not diretor:
        return JsonResponse({"success": False, "message": "Nenhum currículo encontrado."})

    # Retorna os dados do candidato no formato JSON
    return JsonResponse({
        "success": True,
        "diretor": {
            "nome_completo": diretor.nome_completo,
            "cpf": diretor.cpf,
            "rg": diretor.rg,
            "email": diretor.email,
            "telefone": diretor.telefone,
            "endereco": diretor.endereco,
            "bairro": diretor.bairro,
            "cidade": diretor.cidade,
            "cep": diretor.cep,
            "data_nascimento": diretor.data_nascimento.strftime("%d/%m/%Y"),
            "sexo": diretor.sexo,
            "estado_civil": diretor.estado_civil,
            "formacao_academica": diretor.formacao_academica,
            "curso": diretor.curso,
            "instituicao": diretor.instituicao,
            "ano_conclusao": diretor.ano_conclusao,
            "empresa": diretor.empresa,
            "cargo": diretor.cargo,
            "data_inicio": diretor.data_inicio.strftime("%d/%m/%Y") if diretor.data_inicio else None,
            "data_fim": diretor.data_fim.strftime("%d/%m/%Y") if diretor.data_fim else None,
            "foto": diretor.foto.url if diretor.foto else None,
            "curriculo_pdf": diretor.curriculo_pdf.url if diretor.curriculo_pdf else None,
        }
    })
###***********************************************************************************************************************

from django.shortcuts import render, get_object_or_404, redirect
from semedapp.models import Diretor

def editar_curriculo(request, candidato_id):
    candidato = get_object_or_404(Diretor, id=candidato_id)

    if request.method == 'POST':
        candidato.nome_completo = request.POST.get('nome_completo')
        candidato.cpf = request.POST.get('cpf')
        candidato.rg = request.POST.get('rg')
        candidato.email = request.POST.get('email')
        candidato.telefone = request.POST.get('telefone')
        candidato.endereco = request.POST.get('endereco')
        candidato.bairro = request.POST.get('bairro')
        candidato.cidade = request.POST.get('cidade')
        candidato.cep = request.POST.get('cep')
        candidato.data_nascimento = request.POST.get('data_nascimento')
        candidato.formacao_academica = request.POST.get('formacao_academica')
        candidato.curso = request.POST.get('curso')
        candidato.instituicao = request.POST.get('instituicao')
        candidato.ano_conclusao = request.POST.get('ano_conclusao')
        candidato.experiencia_profissional = request.POST.get('experiencia_profissional')
        candidato.estado_civil = request.POST.get('estado_civil')
        candidato.sexo = request.POST.get('sexo')
        candidato.data_cadastro = request.POST.get('data_cadastro')
        candidato.empresa = request.POST.get('empresa')
        candidato.cargo = request.POST.get('cargo')
        candidato.data_inicio = request.POST.get('data_inicio')
        candidato.data_fim = request.POST.get('data_fim')

        if 'foto' in request.FILES:
            candidato.foto = request.FILES['foto']
        if 'curriculo_pdf' in request.FILES:
            candidato.curriculo_pdf = request.FILES['curriculo_pdf']
        if 'certificados_pdf' in request.FILES:
            candidato.certificados_pdf = request.FILES['certificados_pdf']

        candidato.save()
        return redirect('area_candidato')  # Redireciona para a área do candidato

    return render(request, 'banco_curriculos/editar_curriculo.html', {'candidato': candidato})
###***********************************************************************************************************************

from django.http import JsonResponse
from django.core.mail import send_mail
from django.conf import settings
from django.contrib.auth.models import User

def recuperar_senha(request):
    if request.method == "POST":
        email = request.POST.get("email")

        # Verifica se o e-mail existe no banco
        try:
            user = User.objects.get(email=email)
        except User.DoesNotExist:
            return JsonResponse({"success": False, "message": "E-mail não encontrado."}, status=400)

        # Enviar e-mail de recuperação
        try:
            send_mail(
                "Recuperação de Senha",
                "Clique no link para redefinir sua senha: http://127.0.0.1:8000/banco-curriculos/resetar-senha/",
                settings.DEFAULT_FROM_EMAIL,
                [email],
                fail_silently=False,
            )
            return JsonResponse({"success": True, "message": "E-mail enviado com sucesso!"})

        except Exception as e:
            return JsonResponse({"success": False, "message": f"Erro ao enviar e-mail: {str(e)}"}, status=500)

    return JsonResponse({"success": False, "message": "Método inválido"}, status=400)
###***********************************************************************************************************************

from django.shortcuts import render, redirect
from django.contrib.auth.tokens import default_token_generator
from django.contrib.auth import get_user_model
from django.utils.http import urlsafe_base64_decode
from django.contrib.auth.forms import SetPasswordForm

User = get_user_model()

def resetar_senha(request, uidb64, token):
    try:
        uid = urlsafe_base64_decode(uidb64).decode()
        user = User.objects.get(pk=uid)

        if not default_token_generator.check_token(user, token):
            return render(request, "error.html", {"mensagem": "Token inválido ou expirado."})  # Alterado para error.html

    except (TypeError, ValueError, OverflowError, User.DoesNotExist):
        return render(request, "error.html", {"mensagem": "Usuário inválido."})  # Alterado para error.html

    if request.method == "POST":
        form = SetPasswordForm(user, request.POST)
        if form.is_valid():
            form.save()
            return redirect("login_candidato")

    else:
        form = SetPasswordForm(user)

    return render(request, "resetar_senha.html", {"form": form})
###***********************************************************************************************************************

from django.shortcuts import render, redirect
from django.core.mail import send_mail
from django.contrib import messages
from django.urls import reverse
import uuid
from django.http import JsonResponse
from django.contrib.auth.models import User  # Certifique-se de importar o modelo correto

def enviar_link_recuperacao(request):
    if request.method == "POST":
        email = request.POST.get("email")  # Captura o e-mail do formulário
        print(f"🔍 Verificando e-mail: {email}")  # DEBUG: Veja se o e-mail está sendo capturado corretamente

        try:
            user = User.objects.get(email=email)  # Busca o usuário no banco de dados
            token = str(uuid.uuid4())  # Gera um token único
            reset_url = request.build_absolute_uri(reverse("resetar_senha", args=[token]))

            # Enviar o e-mail
            send_mail(
                "Recuperação de Senha",
                f"Olá, {user.username}. Clique no link abaixo para redefinir sua senha:\n{reset_url}",
                "tecnologiasyscorp@outlook.com",  # E-mail de envio
                [email],  # E-mail do destinatário
                fail_silently=False,
            )

            print(f"✅ E-mail enviado com sucesso para {email}")  # DEBUG: Mostra que o e-mail foi enviado
            return JsonResponse({"success": True, "message": "E-mail enviado com sucesso!"})

        except User.DoesNotExist:
            print(f"❌ E-mail não encontrado: {email}")  # DEBUG: Mostra o e-mail que não foi encontrado
            return JsonResponse({"success": False, "message": f"E-mail '{email}' não encontrado no sistema."})

    return JsonResponse({"success": False, "message": "Requisição inválida."})
###***********************************************************************************************************************

import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart

# Configurações do servidor SMTP da Microsoft
EMAIL_HOST = "smtp.office365.com"
EMAIL_PORT = 587
EMAIL_HOST_USER = "tecnologiasyscorp@outlook.com"
EMAIL_HOST_PASSWORD = "@Drik16091985"

def enviar_email(destinatario, assunto, mensagem):
    try:
        # Criar o e-mail
        msg = MIMEMultipart()
        msg["From"] = EMAIL_HOST_USER
        msg["To"] = destinatario
        msg["Subject"] = assunto

        # Corpo do e-mail
        msg.attach(MIMEText(mensagem, "plain"))

        # Conectar ao servidor SMTP
        server = smtplib.SMTP(EMAIL_HOST, EMAIL_PORT)
        server.starttls()  # Ativar TLS
        server.login(EMAIL_HOST_USER, EMAIL_HOST_PASSWORD)  # Autenticar
        server.sendmail(EMAIL_HOST_USER, destinatario, msg.as_string())  # Enviar e-mail
        server.quit()

        print(f"✅ E-mail enviado com sucesso para {destinatario}")

    except Exception as e:
        print(f"❌ Erro ao enviar e-mail: {e}")

# Teste
if __name__ == "__main__":
    enviar_email(
        "destinatario@email.com",
        "Teste de Envio via Outlook",
        "Olá! Este é um teste de envio de e-mail pelo Python usando Outlook."
    )
###***********************************************************************************************************************

from django.shortcuts import render, redirect
from django.contrib.auth.hashers import make_password
from semedapp.models import CadastroCandidato

def resetar_senha(request, candidato_id):
    candidato = get_object_or_404(CadastroCandidato, id=candidato_id)

    if request.method == "POST":
        nova_senha = request.POST.get("nova_senha")
        candidato.senha = make_password(nova_senha)
        candidato.save()

        return render(request, "banco_curriculos/resetar_senha_sucesso.html")

    return render(request, "banco_curriculos/resetar_senha.html", {"candidato": candidato})
###***********************************************************************************************************************

from django.core.mail import send_mail
from django.http import JsonResponse

def enviar_email_recuperacao(request):
    """Função para enviar o link de recuperação de senha"""
    
    if request.method == "POST":
        email_destino = request.POST.get("email")

        if not email_destino:
            return JsonResponse({"success": False, "message": "E-mail inválido!"})

        assunto = "Recuperação de Senha"
        mensagem = f"Olá, clique no link para redefinir sua senha: https://seusite.com/redefinir-senha"

        try:
            send_mail(
                assunto,
                mensagem,
                "tecnologiasyscorp@outlook.com",  # Remetente
                [email_destino],  # Destinatário
                fail_silently=False,
            )
            return JsonResponse({"success": True, "message": "E-mail enviado com sucesso!"})
        except Exception as e:
            return JsonResponse({"success": False, "message": f"Erro ao enviar: {str(e)}"})

    return JsonResponse({"success": False, "message": "Método inválido"})
###***********************************************************************************************************************

from django.shortcuts import render
from .models import CurriculoAntigo  # Ajuste o nome do modelo conforme necessário

def listar_diretores_antigos(request):
    # Buscando todos os currículos antigos no banco de dados
    curriculos = CurriculoAntigo.objects.all()

    # Passando os dados para o template
    return render(request, 'banco_curriculos/listar_diretores_antigos.html', {'curriculos': curriculos})
###***********************************************************************************************************************

from django.shortcuts import render
from .models import CurriculoAntigo
from django.db.models import Q
from django.utils.timezone import now
from django.core.paginator import Paginator

def listar_curriculos_antigos(request):
    # Obtém os parâmetros da requisição GET
    search_query = request.GET.get('search', '').strip()
    filter_cargo = request.GET.get('cargo', '').strip()
    page_number = request.GET.get('page', 1)

    # Captura a data atual no formato desejado
    data_atual = now().strftime('%d/%m/%Y')

    # Filtrar os currículos pelo nome, CPF ou cargo
    curriculos = CurriculoAntigo.objects.all()

    if search_query:
        curriculos = curriculos.filter(
            Q(nome__icontains=search_query) | 
            Q(cpf__icontains=search_query)
        )

    if filter_cargo:
        curriculos = curriculos.filter(cargo__iexact=filter_cargo)

    # Paginação - 200 registros por página
    paginator = Paginator(curriculos, 200)
    page_obj = paginator.get_page(page_number)

    # Obter lista de cargos únicos para o dropdown
    cargos_disponiveis = CurriculoAntigo.objects.exclude(cargo__isnull=True).exclude(cargo="").values_list('cargo', flat=True).distinct()

    context = {
        'curriculos': page_obj,  # Currículos paginados
        'total_curriculos': curriculos.count(),
        'cargos_disponiveis': sorted(set(c.strip() for c in cargos_disponiveis)),  # Remove espaços extras
        'search_query': search_query,  # Para manter o valor no campo de busca
        'filter_cargo': filter_cargo,  # Para manter o valor no campo de filtro
        'page_obj': page_obj  # Objeto de paginação para usar no template
    }
    return render(request, 'banco_curriculos/listar_curriculos_antigos.html', context)

###***********************************************************************************************************************

def visualizar_curriculo_antigo(request, id):
    curriculo = get_object_or_404(CurriculoAntigo, id=id)
    return render(request, 'banco_curriculos/visualizar_curriculo_antigo.html', {'curriculo': curriculo})
###***********************************************************************************************************************

# Editar curriculo antigo
def editar_curriculo_antigo(request, id):
    curriculo = get_object_or_404(CurriculoAntigo, id=id)
    if request.method == 'POST':
        # Processar o formulário de edição
        pass  # Aqui você pode implementar a lógica para editar
    return render(request, 'editar_curriculo_antigo.html', {'curriculo': curriculo})
###***********************************************************************************************************************

# Excluir curriculo antigo
def excluir_curriculo_antigo(request, id):
    curriculo = get_object_or_404(CurriculoAntigo, id=id)
    if request.method == 'POST':
        curriculo.delete()
        return redirect('listar_curriculos_antigos')  # Redireciona de volta para a lista
    return render(request, 'confirmar_exclusao.html', {'curriculo': curriculo})
###***********************************************************************************************************************

from django.http import JsonResponse
from django.shortcuts import get_object_or_404
from .models import CurriculoAntigo

def imprimir_curriculo_antigo(request, cpf):
    try:
        curriculo = get_object_or_404(CurriculoAntigo, cpf=cpf)
        # Retornar a URL para imprimir caso o currículo exista
        return JsonResponse({"success": True, "url": f"/banco-curriculos/imprimir/{cpf}/"})
    except:
        # Retorna uma resposta JSON com a mensagem de erro
        return JsonResponse({"success": False, "error": "Nenhum currículo encontrado para este CPF."})
###***********************************************************************************************************************

from django.shortcuts import render

def educacao_infantil(request):
    return render(request, 'educacao_infantil.html')  # Sem a pasta "webapp"
###***********************************************************************************************************************

from django.shortcuts import render

def educacao_infantil(request):
    return render(request, 'webapp/educacao_infantil.html')
###***********************************************************************************************************************

from django.shortcuts import render
from django.core.paginator import Paginator
from .models import CadastroEI, Escolas
from django.contrib.auth.decorators import login_required

from django.contrib.auth.models import Group

@login_required
def cadastro_escola(request):
    nome_escola = request.GET.get('nome_escola', '')
    ano = request.GET.get('ano', '')
    modalidade = request.GET.get('modalidade', '')
    turma = request.GET.get('turma', '')

    # Lista de usernames que têm acesso total, além do superuser
    usernames_com_acesso_total = [
        "Valder", "Tatiane", "Edna", "Simone"
    ]

    if request.user.is_superuser or request.user.username in usernames_com_acesso_total:
        escolas = CadastroEI.objects.all()
    else:
        escolas_vinculadas = Escolas.objects.filter(coordenador=request.user)

        if escolas_vinculadas.exists():
            escolas = CadastroEI.objects.filter(unidade_ensino__in=escolas_vinculadas.values_list('nome', flat=True))
        else:
            escolas = CadastroEI.objects.filter(professor=request.user)

    # Filtros
    if nome_escola:
        escolas = escolas.filter(unidade_ensino__icontains=nome_escola)
    if ano:
        escolas = escolas.filter(ano=ano)
    if modalidade:
        escolas = escolas.filter(modalidade__icontains=modalidade)
    if turma:
        escolas = escolas.filter(turma__nome__icontains=turma)

    escolas = escolas.select_related("turma")
    total_turmas = escolas.values("unidade_ensino", "turma__nome").distinct().count()
    total_alunos = escolas.count()
    total_escolas = escolas.values("unidade_ensino").distinct().count()
    total_avaliados = escolas.filter(avaliado="SIM").count()

    paginator = Paginator(escolas.exclude(avaliado="SIM"), 50)
    page_number = request.GET.get('page')
    escolas_paginadas = paginator.get_page(page_number)

    escolas_nomes = escolas.order_by("unidade_ensino").values_list("unidade_ensino", flat=True).distinct()
    turmas = escolas.values_list("turma__nome", flat=True).distinct()

    query_params = request.GET.copy()
    query_params.pop('page', None)
    query_string = query_params.urlencode()

    is_coordenador = Escolas.objects.filter(coordenador=request.user).exists()

    context = {
        'escolas': escolas_paginadas,
        'escolas_avaliadas': escolas.filter(avaliado="SIM"),
        'escolas_nomes': escolas_nomes,
        'turmas': turmas,
        'filtros': {
            'nome_escola': nome_escola,
            'ano': ano,
            'modalidade': modalidade,
            'turma': turma,
        },
        'total_alunos': total_alunos,
        'total_turmas': total_turmas,
        'total_escolas': total_escolas,
        'total_avaliados': total_avaliados,
        'query_string': query_string,
        'range_matematica': range(1, 11),
        'range_linguagem': range(11, 21),
        'is_coordenador': is_coordenador
    }

    return render(request, 'webapp/cadastro_escola.html', context)


###***********************************************************************************************************************

from django.http import JsonResponse
from .models import CadastroEI

def get_turmas(request):
    escola_nome = request.GET.get('escola', '').strip()
    
    if escola_nome:
        turmas = CadastroEI.objects.filter(unidade_ensino=escola_nome).values_list("turma", flat=True).distinct()
        return JsonResponse(list(turmas), safe=False)
    
    return JsonResponse([], safe=False)  # Retorna lista vazia se não encontrar
###***********************************************************************************************************************

from django.shortcuts import render
from django.http import JsonResponse
from .models import CadastroEI

def cadastro_turma(request):
    # Captura os filtros da URL
    nome_escola = request.GET.get('nome_escola', '')
    turma = request.GET.get('turma', '')

    # Filtra as turmas conforme os parâmetros fornecidos
    turmas_cadastradas = CadastroEI.objects.all()

    if nome_escola:
        turmas_cadastradas = turmas_cadastradas.filter(unidade_ensino=nome_escola)

    if turma:
        turmas_cadastradas = turmas_cadastradas.filter(turma=turma)

    # Obter todas as escolas únicas
    escolas_nomes = CadastroEI.objects.values_list("unidade_ensino", flat=True).distinct()

    # Obter todas as turmas únicas associadas às escolas
    turmas = CadastroEI.objects.filter(unidade_ensino=nome_escola).values_list("turma", flat=True).distinct() if nome_escola else []

    context = {
        'turmas_cadastradas': turmas_cadastradas,
        'escolas': escolas_nomes,
        'turmas': turmas,
        'filtros': {
            'nome_escola': nome_escola,
            'turma': turma,
        }
    }
    
    return render(request, 'webapp/cadastro_turma.html', context)

# Endpoint para buscar as turmas de uma escola selecionada via AJAX
def get_turmas(request):
    escola_nome = request.GET.get('escola', '')
    turmas = CadastroEI.objects.filter(unidade_ensino=escola_nome).values_list("turma", flat=True).distinct()
    return JsonResponse(list(turmas), safe=False)

###***********************************************************************************************************************

def cadastro_alunos(request):
    return render(request, 'webapp/cadastro_alunos.html')
###***********************************************************************************************************************

def corrigir_avaliacoes(request):
    return render(request, 'webapp/corrigir_avaliacoes.html')
###***********************************************************************************************************************

import csv
import io
from datetime import datetime
from django.shortcuts import render
from django.contrib import messages
from .models import CadastroEI
from .forms import UploadCSVForm

def upload_cadastro_ei(request):
    if request.method == "POST":
        form = UploadCSVForm(request.POST, request.FILES)
        if form.is_valid():
            csv_file = request.FILES["file"]

            if not csv_file.name.endswith(".csv"):
                messages.error(request, "Por favor, envie um arquivo CSV válido.")
                return render(request, "webapp/upload_cadastro_ei.html", {"form": form})

            try:
                data_set = csv_file.read().decode("utf-8")
                io_string = io.StringIO(data_set)
                reader = csv.DictReader(io_string, delimiter=";")

                erros = []
                linhas_processadas = 0

                for index, row in enumerate(reader, start=2):
                    try:
                        # Corrigir o formato da data para aceitar %d/%m/%Y
                        data_nascimento = datetime.strptime(row["data_nascimento"], "%d/%m/%Y").strftime("%Y-%m-%d")

                        CadastroEI.objects.create(
                            unidade_ensino=row["unidade_ensino"],
                            ano=row["ano"],
                            modalidade=row["modalidade"],
                            formato_letivo=row["formato_letivo"],
                            turma=row["turma"],
                            cpf=row["cpf"],
                            pessoa_nome=row["pessoa_nome"],
                            data_nascimento=data_nascimento,  # Data corrigida para o formato esperado no banco
                            idade=int(row["idade"]) if row["idade"] else None,
                        )
                        linhas_processadas += 1

                    except ValueError as ve:
                        erros.append(f"Linha {index}: Erro de conversão de data ({ve})")
                    except Exception as e:
                        erros.append(f"Linha {index}: Erro ({e})")

                if erros:
                    messages.warning(request, f"Foram encontrados {len(erros)} erros. Veja abaixo:")
                    for erro in erros[:10]:  # Mostra os primeiros 10 erros
                        messages.warning(request, erro)
                else:
                    messages.success(request, f"Importação concluída! {linhas_processadas} registros inseridos.")

                return render(request, "webapp/upload_cadastro_ei.html", {"form": form})

            except Exception as e:
                messages.error(request, f"Erro ao processar o arquivo: {e}")
                return render(request, "webapp/upload_cadastro_ei.html", {"form": form})

    else:
        form = UploadCSVForm()

    return render(request, "webapp/upload_cadastro_ei.html", {"form": form})

###***********************************************************************************************************************

from django.shortcuts import get_object_or_404, redirect
from django.urls import reverse
from .models import CadastroEI

def excluir_escola(request, id_matricula):
    escola = get_object_or_404(CadastroEI, id_matricula=id_matricula)
    escola.delete()
    return redirect(reverse("cadastro_escola"))  # Redireciona de volta para a lista
###***********************************************************************************************************************

import csv
import io
from django.shortcuts import render
from django.contrib import messages
from .models import CadastroEscola  # Ajuste para o seu modelo

def processar_csv(request):
    if request.method == "POST" and request.FILES.get("csv_file"):
        csv_file = request.FILES["csv_file"]

        # Verifica se é um arquivo CSV válido
        if not csv_file.name.endswith('.csv'):
            messages.error(request, "Por favor, envie um arquivo CSV válido.")
            return render(request, "webapp/cadastro_escola.html")

        # Abrindo o CSV corretamente, considerando o encoding correto
        try:
            decoded_file = csv_file.read().decode("ISO-8859-1")  # Alternativa a UTF-8
        except UnicodeDecodeError:
            messages.error(request, "Erro ao processar o arquivo: encoding inválido.")
            return render(request, "webapp/cadastro_escola.html")

        # Lendo o CSV e ignorando erros
        csv_reader = csv.reader(io.StringIO(decoded_file), delimiter=";")
        next(csv_reader, None)  # Pular cabeçalho

        total_registros = 0
        registros_cadastrados = 0
        erros = []

        for row in csv_reader:
            total_registros += 1
            try:
                # Pegando os valores do CSV
                id_matricula = row[0]
                unidade_ensino = row[1]
                ano = row[2]
                modalidade = row[3]
                formato_letivo = row[4]
                turma = row[5]
                cpf = row[6]
                pessoa_nome = row[7]
                data_nascimento = row[8]
                idade = row[9]

                # Verifica se o CPF já existe
                if CadastroEscola.objects.filter(cpf=cpf).exists():
                    erros.append(f"⚠ CPF {cpf} já cadastrado, pulando...")
                    continue  # Pula para o próximo registro

                # Criar o novo registro
                CadastroEscola.objects.create(
                    id_matricula=id_matricula,
                    unidade_ensino=unidade_ensino,
                    ano=ano,
                    modalidade=modalidade,
                    formato_letivo=formato_letivo,
                    turma=turma,
                    cpf=cpf,
                    pessoa_nome=pessoa_nome,
                    data_nascimento=data_nascimento,
                    idade=idade,
                )
                registros_cadastrados += 1

            except Exception as e:
                erros.append(f"⚠ Erro ao processar linha {total_registros}: {str(e)}")

        # Exibir mensagens de sucesso e erro
        messages.success(request, f"{registros_cadastrados} registros cadastrados com sucesso!")
        if erros:
            for erro in erros[:5]:  # Mostra apenas os primeiros 5 erros
                messages.warning(request, erro)

        return render(request, "webapp/cadastro_escola.html")

    return render(request, "webapp/cadastro_escola.html")
###***********************************************************************************************************************

from django.shortcuts import render, get_object_or_404, redirect
from .models import CadastroEI

def editar_escola(request, id_matricula):
    escola = get_object_or_404(CadastroEI, id_matricula=id_matricula)

    if request.method == "POST":
        escola.unidade_ensino = request.POST["unidade_ensino"]
        escola.ano = request.POST["ano"]
        escola.modalidade = request.POST["modalidade"]
        escola.turma = request.POST["turma"]
        escola.pessoa_nome = request.POST["pessoa_nome"]
        escola.idade = request.POST["idade"]
        escola.save()
        
        return redirect("cadastro_escola")  # Redireciona para a listagem de escolas após salvar

    return render(request, "webapp/editar_escola.html", {"escola": escola})
###***********************************************************************************************************************

from django.shortcuts import render, get_object_or_404, redirect
from django.http import JsonResponse
from .models import CadastroEI, NotaAluno
from django.contrib import messages

def selecionar_turma(request):
    """ Página inicial para selecionar uma escola e turma """
    escolas = CadastroEI.objects.values_list("unidade_ensino", flat=True).distinct()

    filtros = {
        "nome_escola": request.GET.get("nome_escola", ""),
        "turma": request.GET.get("turma", ""),
        "ano": request.GET.get("ano", ""),
    }

    alunos = CadastroEI.objects.all()

    if filtros["nome_escola"]:
        alunos = alunos.filter(unidade_ensino=filtros["nome_escola"])
    
    if filtros["turma"]:
        alunos = alunos.filter(turma=filtros["turma"])

    if filtros["ano"]:
        alunos = alunos.filter(ano=filtros["ano"])

    turmas = CadastroEI.objects.filter(unidade_ensino=filtros["nome_escola"]).values_list("turma", flat=True).distinct()

    return render(request, 'webapp/selecionar_turma.html', {
        "escolas": escolas,
        "turmas": turmas,
        "filtros": filtros,
        "alunos": alunos
    })
###***********************************************************************************************************************

from django.http import JsonResponse
from .models import CadastroEI

def buscar_turmas(request):
    """ Retorna as turmas de uma escola via AJAX """
    escola_nome = request.GET.get("escola", "")
    
    if not escola_nome:
        return JsonResponse({"error": "Escola não fornecida."}, status=400)
    
    try:
        turmas = CadastroEI.objects.filter(unidade_ensino=escola_nome).values_list("turma", flat=True).distinct()
        if turmas.exists():
            return JsonResponse(list(turmas), safe=False)
        else:
            return JsonResponse({"message": "Nenhuma turma encontrada para esta escola."}, status=404)
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)

###***********************************************************************************************************************

def listar_alunos(request):
    """ Exibe os alunos da turma selecionada e permite inserir notas """
    escola = request.GET.get("escola", "")
    turma = request.GET.get("turma", "")

    alunos = CadastroEI.objects.filter(unidade_ensino=escola, turma=turma)

    return render(request, 'webapp/lancar_notas.html', {
        "alunos": alunos,
        "escola": escola,
        "turma": turma,
    })
###***********************************************************************************************************************

def salvar_notas(request):
    """ Processa o salvamento das notas no banco """
    if request.method == "POST":
        for key, value in request.POST.items():
            if key.startswith("nota_"):  # Verifica os campos de nota
                matricula_id = key.split("_")[1]  # Extrai o ID do aluno
                aluno = get_object_or_404(CadastroEI, id_matricula=matricula_id)
                disciplina = request.POST.get(f"disciplina_{matricula_id}")
                bimestre = request.POST.get(f"bimestre_{matricula_id}")
                
                # Salva a nota
                NotaAluno.objects.create(
                    aluno=aluno,
                    disciplina=disciplina,
                    nota=float(value),
                    bimestre=int(bimestre)
                )

        messages.success(request, "Notas lançadas com sucesso!")
        return redirect("selecionar_turma")
###***********************************************************************************************************************

from django.shortcuts import render, get_object_or_404, redirect
from .models import CadastroEI
from django.contrib import messages

def editar_turma(request, id_matricula):  # Alterado o nome do argumento para id_matricula
    turma = get_object_or_404(CadastroEI, id_matricula=id_matricula)  # Substituído id por id_matricula

    if request.method == "POST":
        turma.turma = request.POST.get("nome_turma")
        turma.save()
        messages.success(request, "Turma atualizada com sucesso!")
        return redirect("cadastro_turma")  # Redireciona após a edição

    return render(request, "webapp/editar_turma.html", {"turma": turma})

###***********************************************************************************************************************

from django.shortcuts import render, get_object_or_404
from .models import CadastroEI

def visualizar_avaliacao(request, id_matricula):
    escola = get_object_or_404(CadastroEI, id_matricula=id_matricula)
    
    return render(request, "webapp/visualizar_avaliacao.html", {"escola": escola})
###***********************************************************************************************************************

from django.shortcuts import render, get_object_or_404, redirect
from .models import CadastroEI

def salvar_avaliacao(request, id_matricula):
    if request.method == "POST":
        escola = get_object_or_404(CadastroEI, id_matricula=id_matricula)
        
        # Atualiza as 20 questões na tabela
        for i in range(1, 11):
            questao_matematica = request.POST.get(f"questao_matematica_{i}")
            setattr(escola, f"questao_matematica_{i}", questao_matematica)
        
        for i in range(11, 21):
            questao_linguagem = request.POST.get(f"questao_linguagem_{i}")
            setattr(escola, f"questao_linguagem_{i}", questao_linguagem)
        
        # Atualiza a coluna 'avaliado' para 'SIM'
        escola.avaliado = "SIM"
        
        escola.save()  # Salva todas as alterações
        
        # Redireciona para a página de cadastro de escolas
        return redirect('cadastro_escola')

###***********************************************************************************************************************

from django.shortcuts import get_object_or_404, redirect
from .models import CadastroEI
from django.contrib import messages

def excluir_turma(request, id_matricula):
    turma = get_object_or_404(CadastroEI, id_matricula=id_matricula)
    turma.delete()
    messages.success(request, "Turma excluída com sucesso!")
    return redirect('cadastro_turma')
###***********************************************************************************************************************

from django.shortcuts import render
from .models import Aluno, Escola, Resposta

def resumo_dashboard(request):
    total_respostas = Resposta.objects.count()
    total_conceitos = Resposta.objects.values('conceito').distinct().count()
    total_escolas = Escola.objects.count()

    respostas_por_turma = (
        Resposta.objects.values('turma__nome')
        .annotate(total=Count('id'))
        .order_by('turma__nome')
    )
    respostas_por_turma_dict = {r['turma__nome']: r['total'] for r in respostas_por_turma}

    conceitos = (
        Resposta.objects.values('conceito')
        .annotate(total=Count('id'))
        .order_by('conceito')
    )
    conceitos_dict = {c['conceito']: c['total'] for c in conceitos}

    respostas_por_escola = (
        Resposta.objects.values('escola__nome')
        .annotate(total=Count('id'))
        .order_by('escola__nome')
    )
    respostas_por_escola_dict = {e['escola__nome']: e['total'] for e in respostas_por_escola}

    context = {
        'total_respostas': total_respostas,
        'total_conceitos': total_conceitos,
        'total_escolas': total_escolas,
        'respostas_por_turma': respostas_por_turma_dict,
        'conceitos': conceitos_dict,
        'respostas_por_escola': respostas_por_escola_dict,
    }
    return render(request, 'webapp/resumo_dashboard.html', context)
###***********************************************************************************************************************

from django.shortcuts import render, get_object_or_404, redirect
from .models import Aluno, Resposta
from .forms import RespostaForm
from django.contrib import messages

def lancar_conceito(request):
    if request.method == 'POST':
        form = RespostaForm(request.POST)
        if form.is_valid():
            form.save()
            messages.success(request, "Conceito lançado com sucesso!")
            return redirect('lancar_conceito')
    else:
        form = RespostaForm()
    
    respostas = Resposta.objects.select_related('aluno').order_by('-data_resposta')
    
    context = {
        'form': form,
        'respostas': respostas,
    }
    return render(request, 'lancar_conceito.html', context)
###***********************************************************************************************************************
from django.contrib.auth.decorators import login_required
from django.shortcuts import render
from .models import CadastroEI, Escolas

@login_required
def gestao_alunos(request):
    usuario_logado = request.user
    nome_escola = request.GET.get('nome_escola', '')
    ano = request.GET.get('ano', '')
    modalidade = request.GET.get('modalidade', '')
    turma = request.GET.get('turma', '')

    usernames_com_acesso_total = [
        "Valder", "Tatiane", "Edna", "Simone"
    ]

    if usuario_logado.is_superuser or usuario_logado.username in usernames_com_acesso_total:
        registros = CadastroEI.objects.all()
    else:
        escolas_vinculadas = Escolas.objects.filter(coordenador=usuario_logado)
        if escolas_vinculadas.exists():
            registros = CadastroEI.objects.filter(unidade_ensino__in=escolas_vinculadas.values_list('nome', flat=True))
        else:
            registros = CadastroEI.objects.filter(professor=usuario_logado)

    # Filtros
    if nome_escola:
        registros = registros.filter(unidade_ensino__icontains=nome_escola)
    if ano:
        registros = registros.filter(ano=ano)
    if modalidade:
        registros = registros.filter(modalidade__icontains=modalidade)
    if turma:
        registros = registros.filter(turma__nome__icontains=turma)

    registros = registros.select_related("turma")
    turmas = registros.order_by("turma__nome").values_list("turma__nome", flat=True).distinct()
    total_turmas = registros.values("unidade_ensino", "turma__nome").distinct().count()
    escolas_nomes = registros.order_by("unidade_ensino").values_list("unidade_ensino", flat=True).distinct()
    total_alunos = registros.count()
    total_escolas = registros.values("unidade_ensino").distinct().count()

    registros_com_questoes = []
    for registro in registros:
        questoes_matematica = [getattr(registro, f"questao_matematica_{i}", "BRANCO") for i in range(1, 11)]
        questoes_linguagem = [getattr(registro, f"questao_linguagem_{i}", "BRANCO") for i in range(11, 21)]
        questoes_excluidas = {11, 13}

        pontuacao_matematica = sum(2 if r == "CERTO" else 1 if r == "PARCIAL" else 0 for r in questoes_matematica)
        pontuacao_linguagem = sum(
            2 if r == "CERTO" else 1 if r == "PARCIAL" else 0
            for i, r in enumerate(questoes_linguagem, start=11) if i not in questoes_excluidas
        )

        desempenho_matematica = (pontuacao_matematica / 20) * 100
        desempenho_linguagem = (pontuacao_linguagem / 16) * 100

        def avaliar(d):
            return "Excelente" if d == 100 else "Bom" if d >= 70 else "Regular" if d >= 50 else "INSUFICIENTE"

        registros_com_questoes.append({
            "registro": registro,
            "questoes_matematica": questoes_matematica,
            "questoes_linguagem": questoes_linguagem,
            "avaliacao_matematica": avaliar(desempenho_matematica),
            "avaliacao_linguagem": avaliar(desempenho_linguagem),
            "desempenho_matematica": f"{desempenho_matematica:.2f}%",
            "desempenho_linguagem": f"{desempenho_linguagem:.2f}%"
        })

    context = {
        'registros_com_questoes': registros_com_questoes,
        'escolas_nomes': escolas_nomes,
        'turmas': turmas,
        'total_alunos': total_alunos,
        'total_turmas': total_turmas,
        'total_escolas': total_escolas,
        'filtros': {
            'nome_escola': nome_escola,
            'ano': ano,
            'modalidade': modalidade,
            'turma': turma,
        }
    }

    return render(request, 'semedapp/gestao_alunos.html', context)






###***********************************************************************************************************************

from django.shortcuts import render
from django.contrib.auth.decorators import login_required
from .models import CadastroEI, Escolas

@login_required
def gestao_relatorios(request):
    usuario_logado = request.user

    # Lista de usernames com acesso total à rede
    usernames_com_acesso_total = [
        "Valder", "Tatiane", "Edna", "Simone"
    ]

    nome_escola = request.GET.get('nome_escola', '')
    turma = request.GET.get('turma', '')
    ano = request.GET.get('ano', '')
    modalidade = request.GET.get('modalidade', '')

    # Verifica acesso total
    if usuario_logado.is_superuser or usuario_logado.username in usernames_com_acesso_total:
        registros = CadastroEI.objects.all()
    else:
        escolas_vinculadas = Escolas.objects.filter(coordenador=usuario_logado)
        if escolas_vinculadas.exists():
            registros = CadastroEI.objects.filter(unidade_ensino__in=escolas_vinculadas.values_list('nome', flat=True))
        else:
            registros = CadastroEI.objects.filter(professor=usuario_logado)

    if nome_escola:
        registros = registros.filter(unidade_ensino__icontains=nome_escola)
    if ano:
        registros = registros.filter(ano=ano)
    if modalidade:
        registros = registros.filter(modalidade__icontains=modalidade)
    if turma:
        registros = registros.filter(turma__nome__icontains=turma)

    escolas = registros.order_by("unidade_ensino").values_list("unidade_ensino", flat=True).distinct()
    turmas = registros.order_by("turma__nome").values_list("turma__nome", flat=True).distinct()

    total_alunos_avaliados = registros.count()
    total_pontos_geral = 0

    pontuacao = {'CERTO': 2, 'PARCIAL': 1, 'ERRADO': 0, 'BRANCO': 0, 'CRIANÇA COM LAUDO': 1.5}
    total_excelente = total_bom = total_regular = total_necessita_melhorar = total_insuficiente = 0

    relatorios = []
    for registro in registros:
        questoes_matematica = [getattr(registro, f"questao_matematica_{i}", "BRANCO") for i in range(1, 11)]
        questoes_linguagem = [getattr(registro, f"questao_linguagem_{i}", "BRANCO") for i in range(11, 21)]
        questoes_excluidas = {11, 13}

        pontuacao_matematica = sum(pontuacao.get(r, 0) for r in questoes_matematica)
        pontuacao_linguagem = sum(pontuacao.get(r, 0) for i, r in enumerate(questoes_linguagem, start=11) if i not in questoes_excluidas)

        max_pontuacao_matematica = 20
        max_pontuacao_linguagem = 16
        max_pontos = max_pontuacao_matematica + max_pontuacao_linguagem

        total_pontos = pontuacao_matematica + pontuacao_linguagem
        total_pontos_geral += total_pontos
        desempenho_geral = (total_pontos / max_pontos) * 100 if max_pontos > 0 else 0

        if desempenho_geral == 100:
            avaliacao_geral = "Excelente"
            total_excelente += 1
        elif desempenho_geral >= 70:
            avaliacao_geral = "Bom"
            total_bom += 1
        elif desempenho_geral >= 50:
            avaliacao_geral = "Regular"
            total_regular += 1
        else:
            avaliacao_geral = "INSUFICIENTE"
            total_necessita_melhorar += 1
            total_insuficiente += 1

        relatorios.append({
            'registro': registro,
            'desempenho': f"{desempenho_geral:.2f}%",
            'avaliacao': avaliacao_geral,
            'questoes_matematica': questoes_matematica,
            'questoes_linguagem': questoes_linguagem,
        })

    media_desempenho = (total_pontos_geral / (total_alunos_avaliados * max_pontos) * 100) if total_alunos_avaliados else 0
    media_matematica = (sum(pontuacao.get(q, 0) for registro in registros for q in questoes_matematica) / (total_alunos_avaliados * max_pontuacao_matematica) * 100) if total_alunos_avaliados else 0
    media_linguagem = (sum(pontuacao.get(q, 0) for registro in registros for i, q in enumerate(questoes_linguagem, start=11) if i not in questoes_excluidas) / (total_alunos_avaliados * max_pontuacao_linguagem) * 100) if total_alunos_avaliados else 0

    total_turmas = registros.values("turma").distinct().count()
    maior_desempenho_qtd = sum(1 for r in relatorios if float(r['desempenho'].replace('%', '')) == 100)
    menor_desempenho_qtd = sum(1 for r in relatorios if float(r['desempenho'].replace('%', '')) < 50)

    total_faltantes_matematica = registros.filter(**{f"questao_matematica_{i}": "FALTOU" for i in range(1, 11)}).count()
    total_faltantes_linguagem = registros.filter(**{f"questao_linguagem_{i}": "FALTOU" for i in range(11, 21)}).count()

    context = {
        'maior_desempenho_qtd': maior_desempenho_qtd,
        'menor_desempenho_qtd': menor_desempenho_qtd,
        'total_turmas': total_turmas,
        'total_faltantes_matematica': total_faltantes_matematica,
        'total_faltantes_linguagem': total_faltantes_linguagem,
        'relatorios': relatorios,
        'escolas': escolas,
        'turmas': turmas,
        'total_alunos_avaliados': total_alunos_avaliados,
        'media_desempenho': f"{media_desempenho:.2f}%",
        'media_matematica': f"{media_matematica:.2f}%",
        'media_linguagem': f"{media_linguagem:.2f}%",
        'total_excelente': total_excelente,
        'total_bom': total_bom,
        'total_regular': total_regular,
        'total_necessita_melhorar': total_necessita_melhorar,
        'total_insuficiente': total_insuficiente,
        'filtros': {
            'nome_escola': nome_escola,
            'turma': turma,
            'ano': ano,
            'modalidade': modalidade,
        },
    }

    return render(request, 'webapp/gestao_relatorios.html', context)








###***********************************************************************************************************************

from django.http import HttpResponse
from django.template.loader import get_template
from django.templatetags.static import static
from xhtml2pdf import pisa
from .models import CadastroEI

def gerar_pdf_relatorio(request):
    nome_escola = request.GET.get('nome_escola', '')
    turma = request.GET.get('turma', '')
    ano = request.GET.get('ano', '')
    modalidade = request.GET.get('modalidade', '')

    registros = CadastroEI.objects.all()
    if nome_escola:
        registros = registros.filter(unidade_ensino__icontains=nome_escola)
    if turma:
        registros = registros.filter(turma__nome__icontains=turma)  # ✅ Correto

    if ano:
        registros = registros.filter(ano=ano)
    if modalidade:
        registros = registros.filter(modalidade__icontains=modalidade)

    relatorios = []
    for registro in registros:
        questoes_matematica = [getattr(registro, f"questao_matematica_{i}", "BRANCO") for i in range(1, 11)]
        questoes_linguagem = [getattr(registro, f"questao_linguagem_{i}", "BRANCO") for i in range(11, 21)]
        desempenho = sum(2 for resposta in questoes_matematica + questoes_linguagem if resposta == 'CERTO')

        relatorios.append({
            'registro': registro,
            'desempenho': desempenho,
        })

    # Caminhos absolutos para as imagens
    cabecalho_superior = request.build_absolute_uri(static('assets/dist/img/cabecalho_superior.png'))
    cabecalho_inferior = request.build_absolute_uri(static('assets/dist/img/cabecalho_inferior.png'))

    template_path = 'webapp/relatorio_pdf.html'
    context = {
        'relatorios': relatorios,
        'cabecalho_superior': cabecalho_superior,
        'cabecalho_inferior': cabecalho_inferior,
        'filtros': {
            'nome_escola': nome_escola,
            'turma': turma,
            'ano': ano,
            'modalidade': modalidade,
        }
    }

    response = HttpResponse(content_type='application/pdf')
    response['Content-Disposition'] = 'attachment; filename="relatorio_avaliacoes.pdf"'
    template = get_template(template_path)
    html = template.render(context)

    pisa_status = pisa.CreatePDF(html, dest=response)
    if pisa_status.err:
        return HttpResponse('Erro ao gerar PDF <pre>' + html + '</pre>')
    
    return response

###***********************************************************************************************************************

def get_turmas(request):
    """View para buscar turmas dinamicamente com base na escola selecionada."""
    escola = request.GET.get('escola', '')
    if escola:
        turmas = CadastroEI.objects.filter(unidade_ensino__icontains=escola).values_list('turma', flat=True).distinct()
    else:
        turmas = []
    return JsonResponse(list(turmas), safe=False)
###***********************************************************************************************************************

# views.py
from django.shortcuts import render, redirect
from django.db.models import Q
from .models import CadastroEI, ConceitoLancado

def lancamento_conceitos(request):
    escola = request.GET.get('escola', '')
    turma = request.GET.get('turma', '')
    ano = request.GET.get('ano', '')

    # Alunos sem conceito
    alunos_sem_conceito = CadastroEI.objects.exclude(
        pessoa_nome__in=ConceitoLancado.objects.values_list('aluno', flat=True)
    )
    
    # Aplicação de filtros
    if escola:
        alunos_sem_conceito = alunos_sem_conceito.filter(unidade_ensino__icontains=escola)
    if turma:
        alunos_sem_conceito = alunos_sem_conceito.filter(turma__icontains=turma)
    if ano:
        alunos_sem_conceito = alunos_sem_conceito.filter(ano=ano)

    # Alunos com conceito já lançado
    conceitos_lancados = ConceitoLancado.objects.all()

    # Contagem para os cards
    total_alunos = CadastroEI.objects.count()
    total_conceitos_lancados = conceitos_lancados.count()
    total_alunos_sem_conceito = alunos_sem_conceito.count()

    if request.method == 'POST':
        for aluno in alunos_sem_conceito:
            conceito_matematica = request.POST.get(f'conceito_matematica_{aluno.id_matricula}')
            conceito_linguagem = request.POST.get(f'conceito_linguagem_{aluno.id_matricula}')
            
            if conceito_matematica or conceito_linguagem:
                # Verifica se já existe conceito lançado para evitar duplicação
                if not ConceitoLancado.objects.filter(aluno=aluno.pessoa_nome, ano=aluno.ano, turma=aluno.turma).exists():
                    ConceitoLancado.objects.create(
                        aluno=aluno.pessoa_nome,
                        turma=aluno.turma,
                        ano=aluno.ano,
                        modalidade=aluno.modalidade,
                        conceito_matematica=conceito_matematica,
                        conceito_linguagem=conceito_linguagem
                    )
        return redirect('lancamento_conceitos')

    context = {
        'alunos_sem_conceito': alunos_sem_conceito,
        'conceitos_lancados': conceitos_lancados,
        'escolas': CadastroEI.objects.values_list('unidade_ensino', flat=True).distinct(),
        'turmas': CadastroEI.objects.values_list('turma', flat=True).distinct(),
        'total_alunos': total_alunos,
        'total_conceitos_lancados': total_conceitos_lancados,
        'total_alunos_sem_conceito': total_alunos_sem_conceito,
        'filtros': {
            'escola': escola,
            'turma': turma,
            'ano': ano,
        }
    }
    return render(request, 'webapp/lancamento_conceitos.html', context)
###***********************************************************************************************************************

# from django.templatetags.static import static
# from django.http import HttpResponse
# from weasyprint import HTML

# def gerar_pdf(request):

#     cabecalho_superior = request.build_absolute_uri(static('assets/dist/img/cabecalho_superior.png'))
#     cabecalho_inferior = request.build_absolute_uri(static('assets/dist/img/cabecalho_inferior.png'))

#     context = {
#         'relatorios': relatorios,
#         'cabecalho_superior': cabecalho_superior,
#         'cabecalho_inferior': cabecalho_inferior,
#     }

#     html_string = render_to_string('webapp/relatorio_pdf.html', context)
#     response = HttpResponse(content_type='application/pdf')
#     HTML(string=html_string).write_pdf(response)
#     return response
###***********************************************************************************************************************

from django.shortcuts import render
from .models import SynapticAtendimento
from django.db.models import Count, Avg, Q

def soe_gestao(request):
    total_estudantes = SynapticAtendimento.objects.values('nome_unidade_ensino').distinct().count()
    total_escolas = SynapticAtendimento.objects.values('nome_unidade_ensino').distinct().count()
    atendimentos_mensais = SynapticAtendimento.objects.filter(
        registro__month=1  # Ajuste para obter a média real por mês
    ).count()

    indicadores = [
        {
            "nome": "Acompanhamento Psicoemocional",
            "meta": "100% dos estudantes",
            "resultado": "85%",
            "status": "Em Progresso",
            "badge_class": "warning"
        },
        {
            "nome": "Redução de Conflitos Escolares",
            "meta": "20% de redução",
            "resultado": "15%",
            "status": "Concluído",
            "badge_class": "success"
        },
        {
            "nome": "Orientação Vocacional",
            "meta": "100% dos alunos do 9º ano",
            "resultado": "95%",
            "status": "Quase lá",
            "badge_class": "info"
        }
    ]

    equipe_gestao = [
        {"nome": "Janayna1", "cargo": "Coordenadora Pedagógica"},
        {"nome": "Janayna2", "cargo": "Psicólogo Escolar"},
        {"nome": "Janayna3", "cargo": "Assistente Social"},
        {"nome": "Janayna4", "cargo": "Orientador Educacional"}
    ]

    context = {
        "total_estudantes": total_estudantes,
        "total_escolas": total_escolas,
        "atendimentos_mensais": atendimentos_mensais,
        "indicadores": indicadores,
        "equipe_gestao": equipe_gestao,
    }

    return render(request, "webapp/soe_gestao.html", context)

###***********************************************************************************************************************

from django.db.models import Sum

def soe_servicos(request):
    # Consultando os dados de cada tabela, limitando a 10 registros
    synaptic = SynapticAtendimento.objects.all()[:10]
    sige = SIGEAtendimento.objects.all()[:10]
    autokee = AutoKeeAtendimento.objects.all()[:10]
    orientadores = OrientadoresAtendimento.objects.all()[:10]
    soe = AtendimentoSOE.objects.all()[:10]

    # Debugging: verificar contagem no terminal
    print(f"Synaptic: {synaptic.count()} registros encontrados.")
    print(f"SIGE: {sige.count()} registros encontrados.")
    print(f"AutoKee: {autokee.count()} registros encontrados.")
    print(f"Orientadores: {orientadores.count()} registros encontrados.")

    # Somatórios gerais de todas as tabelas
    total_ocorrencias = SynapticAtendimento.objects.count() + SIGEAtendimento.objects.count() + \
                        AutoKeeAtendimento.objects.count() + OrientadoresAtendimento.objects.count() + \
                        AtendimentoSOE.objects.count()

    # Exemplo de cálculo de pendentes usando um campo existente em OrientadoresAtendimento
    pendentes_total = (
        SynapticAtendimento.objects.filter(status_descricao="Pendente").count() +
        SIGEAtendimento.objects.filter(status_descricao="Pendente").count() +
        AutoKeeAtendimento.objects.filter(status_descricao="Pendente").count() +
        AtendimentoSOE.objects.filter(status_descricao="Pendente").count() +
        (OrientadoresAtendimento.objects.aggregate(total=Sum('ameaca'))['total'] or 0)  # Exemplo usando o campo 'ameaca'
    )

    abertas_total = (
        SynapticAtendimento.objects.filter(status_descricao="Aberta").count() +
        SIGEAtendimento.objects.filter(status_descricao="Aberta").count() +
        AutoKeeAtendimento.objects.filter(status_descricao="Aberta").count() +
        AtendimentoSOE.objects.filter(status_descricao="Aberta").count() +
        (OrientadoresAtendimento.objects.aggregate(total=Sum('bullying'))['total'] or 0)  # Usando 'bullying' como exemplo
    )

    resolvidas_total = (
        SynapticAtendimento.objects.filter(status_descricao="Resolvida").count() +
        SIGEAtendimento.objects.filter(status_descricao="Resolvida").count() +
        AutoKeeAtendimento.objects.filter(status_descricao="Resolvida").count() +
        AtendimentoSOE.objects.filter(status_descricao="Resolvida").count() +
        (OrientadoresAtendimento.objects.aggregate(total=Sum('uso_alcool_drogas'))['total'] or 0)  # 'uso_alcool_drogas'
    )

    context = {
        'total_ocorrencias': total_ocorrencias,
        'pendentes_total': pendentes_total,
        'abertas_total': abertas_total,
        'resolvidas_total': resolvidas_total,
        'synaptic': synaptic,
        'sige': sige,
        'autokee': autokee,
        'orientadores': orientadores,
        'soe': soe,
    }

    return render(request, 'webapp/servicos_soe.html', context)
###***********************************************************************************************************************

def soe_uploads(request):
    return render(request, 'webapp/soe_uploads.html')
###***********************************************************************************************************************

def soe_relatorios(request):
    return render(request, 'webapp/soe_relatorios.html')
###***********************************************************************************************************************
from django.shortcuts import render
from .models import AtendimentoSOE  # Assumindo que o modelo já foi criado

def servicos_soe(request):
    atendimentos = AtendimentoSOE.objects.all().order_by('-registro')  # Ordenação pela data do registro (mais recente primeiro)
    return render(request, 'webapp/servicos_soe.html', {'atendimentos': atendimentos})
###***********************************************************************************************************************

# from django.shortcuts import render
# from semedapp.models import SynapticAtendimento, SIGEAtendimento, AutoKeeAtendimento, OrientadoresAtendimento

# def soe_servicos(request):
#     # Consultando os dados de cada tabela e limitando a 10 registros
#     synaptic = SynapticAtendimento.objects.all()[:10]
#     sige = SIGEAtendimento.objects.all()[:10]
#     autokee = AutoKeeAtendimento.objects.all()[:10]
#     orientadores = OrientadoresAtendimento.objects.all()[:10]

#     # Somatórios gerais de todas as tabelas
#     total_ocorrencias = (
#         SynapticAtendimento.objects.count() +
#         SIGEAtendimento.objects.count() +
#         AutoKeeAtendimento.objects.count() +
#         OrientadoresAtendimento.objects.count()
#     )
#     pendentes_total = (
#         SynapticAtendimento.objects.filter(status_descricao="Pendente").count() +
#         SIGEAtendimento.objects.filter(status_descricao="Pendente").count() +
#         AutoKeeAtendimento.objects.filter(status_descricao="Pendente").count() +
#         OrientadoresAtendimento.objects.filter(status_descricao="Pendente").count()
#     )
#     abertas_total = (
#         SynapticAtendimento.objects.filter(status_descricao="Aberta").count() +
#         SIGEAtendimento.objects.filter(status_descricao="Aberta").count() +
#         AutoKeeAtendimento.objects.filter(status_descricao="Aberta").count() +
#         OrientadoresAtendimento.objects.filter(status_descricao="Aberta").count()
#     )
#     resolvidas_total = (
#         SynapticAtendimento.objects.filter(status_descricao="Resolvida").count() +
#         SIGEAtendimento.objects.filter(status_descricao="Resolvida").count() +
#         AutoKeeAtendimento.objects.filter(status_descricao="Resolvida").count() +
#         OrientadoresAtendimento.objects.filter(status_descricao="Resolvida").count()
#     )

#     context = {
#         'total_ocorrencias': total_ocorrencias,
#         'pendentes_total': pendentes_total,
#         'abertas_total': abertas_total,
#         'resolvidas_total': resolvidas_total,
#         'synaptic': synaptic,
#         'sige': sige,
#         'autokee': autokee,
#         'orientadores': orientadores,
#     }

#     return render(request, 'webapp/servicos_soe.html', context)
###***********************************************************************************************************************

from django.shortcuts import render
from .models import SynapticAtendimento

def dados_synaptic(request):
    dados = SynapticAtendimento.objects.all().order_by('-registro')
    return render(request, 'webapp/dados_synaptic.html', {'dados': dados})
###***********************************************************************************************************************

from django.shortcuts import render

def upload_dados(request):
    # Lógica para upload de dados
    return render(request, 'semedapp/upload_dados.html')
###***********************************************************************************************************************

from django.shortcuts import render
from .models import SIGEAtendimento
from django.core.paginator import Paginator
from django.db.models import Count
from django.http import HttpResponse
from django.template.loader import render_to_string
import weasyprint

def ocorrencias_sige(request):
    registros = SIGEAtendimento.objects.all()

    # Filtros de busca
    unidade = request.GET.get('unidade')
    classificacao = request.GET.get('classificacao')
    status_descricao = request.GET.get('status')

    if unidade:
        registros = registros.filter(nome_unidade_ensino__icontains=unidade)
    if classificacao:
        registros = registros.filter(classificacao_nome__icontains=classificacao)
    if status_descricao:
        registros = registros.filter(status_descricao__icontains=status_descricao)

    # Paginação
    paginator = Paginator(registros, 10)
    page_number = request.GET.get('page')
    page_obj = paginator.get_page(page_number)

    # Contagem de status
    status_count = registros.values('status_descricao').annotate(count=Count('id')).order_by('-count')

    # Total de ocorrências
    total_ocorrencias = registros.count()

    context = {
        'page_obj': page_obj,
        'status_count': status_count,
        'total_ocorrencias': total_ocorrencias,
        'classificacoes_list': registros.values_list('classificacao_nome', flat=True).distinct(),
        'unidades_list': registros.values_list('nome_unidade_ensino', flat=True).distinct(),
        'status_list': registros.values_list('status_descricao', flat=True).distinct(),
    }

    return render(request, 'webapp/ocorrencias_sige.html', context)
###***********************************************************************************************************************

def ocorrencias_sige_pdf(request):
    registros = SIGEAtendimento.objects.all()

    # Filtragem opcional para o PDF
    unidade = request.GET.get('unidade')
    classificacao = request.GET.get('classificacao')
    status_descricao = request.GET.get('status')

    if unidade:
        registros = registros.filter(nome_unidade_ensino__icontains=unidade)
    if classificacao:
        registros = registros.filter(classificacao_nome__icontains=classificacao)
    if status_descricao:
        registros = registros.filter(status_descricao__icontains=status_descricao)

    # Renderiza o template HTML como string
    html_string = render_to_string('webapp/ocorrencias_sige_pdf.html', {
        'registros': registros,
        'total_ocorrencias': registros.count(),
    })

    # Gera o PDF usando WeasyPrint
    response = HttpResponse(content_type='application/pdf')
    response['Content-Disposition'] = 'inline; filename="relatorio_ocorrencias_sige.pdf"'
    weasyprint.HTML(string=html_string).write_pdf(response)

    return response
###***********************************************************************************************************************

from django.shortcuts import render
from django.core.paginator import Paginator
from .models import SynapticAtendimento
from django.db.models import Count

def ocorrencias_synaptic(request):
    registros = SynapticAtendimento.objects.all()

    # Filtros
    classificacao_filtrada = request.GET.get('classificacao')
    status_filtrado = request.GET.get('status')
    unidade_filtrada = request.GET.get('unidade')

    if classificacao_filtrada:
        registros = registros.filter(classificacao_nome=classificacao_filtrada)
    if status_filtrado:
        registros = registros.filter(status_descricao=status_filtrado)
    if unidade_filtrada:
        registros = registros.filter(nome_unidade_ensino=unidade_filtrada)

    # Paginação (100 registros por página)
    paginator = Paginator(registros, 100)
    page_number = request.GET.get('page')
    page_obj = paginator.get_page(page_number)

    # Dados para os filtros
    classificacoes_list = SynapticAtendimento.objects.values_list('classificacao_nome', flat=True).distinct()
    status_list = SynapticAtendimento.objects.values_list('status_descricao', flat=True).distinct()
    unidades_list = SynapticAtendimento.objects.values_list('nome_unidade_ensino', flat=True).distinct()

    # Dados estatísticos para os cards
    total_ocorrencias = registros.count()
    status_count = registros.values('status_descricao').annotate(count=Count('status_descricao'))
    classificacoes_count = registros.values('classificacao_nome').annotate(count=Count('classificacao_nome'))
    ultimo_registro = registros.order_by('-registro').first()

    context = {
        'page_obj': page_obj,
        'total_ocorrencias': total_ocorrencias,
        'status_count': status_count,
        'classificacoes_count': classificacoes_count,
        'ultimo_registro': ultimo_registro.registro if ultimo_registro else None,
        'classificacoes_list': classificacoes_list,
        'status_list': status_list,
        'unidades_list': unidades_list,
    }

    return render(request, 'webapp/ocorrencias_synaptic.html', context)
###***********************************************************************************************************************

from django.http import HttpResponse
from django.template.loader import get_template
from xhtml2pdf import pisa
from .models import SynapticAtendimento
import io
from django.db.models import Count

def ocorrencias_synaptic_pdf(request):
    registros = SynapticAtendimento.objects.all()

    # Aplicando filtros conforme parâmetros GET
    classificacao_filtrada = request.GET.get('classificacao')
    status_filtrado = request.GET.get('status')
    unidade_filtrada = request.GET.get('unidade')

    if classificacao_filtrada:
        registros = registros.filter(classificacao_nome=classificacao_filtrada)
    if status_filtrado:
        registros = registros.filter(status_descricao=status_filtrado)
    if unidade_filtrada:
        registros = registros.filter(nome_unidade_ensino=unidade_filtrada)

    # Dados estatísticos
    total_ocorrencias = registros.count()

    context = {
        'registros': registros,
        'total_ocorrencias': total_ocorrencias,
    }

    # Gerando o PDF
    template_path = 'webapp/ocorrencias_synaptic_pdf.html'
    response = HttpResponse(content_type='application/pdf')
    response['Content-Disposition'] = 'inline; filename="ocorrencias_synaptic.pdf"'

    template = get_template(template_path)
    html = template.render(context)
    pisa_status = pisa.CreatePDF(io.BytesIO(html.encode('UTF-8')), dest=response, encoding='UTF-8')

    if pisa_status.err:
        return HttpResponse('Erro ao gerar PDF', status=500)
    return response
###***********************************************************************************************************************

from django.shortcuts import render
from .models import SIGEAtendimento, SynapticAtendimento, AutoKeeAtendimento
from django.db.models import Count

def soe_principal(request):
    # Estatísticas gerais combinadas para SIGE, Synaptic e AutoKee
    total_sige_atendimentos = SIGEAtendimento.objects.count()
    total_synaptic_atendimentos = SynapticAtendimento.objects.count()
    total_autokee_atendimentos = AutoKeeAtendimento.objects.count()
    
    total_atendimentos = total_sige_atendimentos + total_synaptic_atendimentos + total_autokee_atendimentos
    total_escolas = SIGEAtendimento.objects.values('nome_unidade_ensino').distinct().count()
    
    # Estimativa de atendimentos mensais
    atendimentos_mensais = total_atendimentos // 12

    # Eventos Recentes combinando SIGE, Synaptic e AutoKee (limite de 5 registros)
    eventos_recentes_sige = SIGEAtendimento.objects.order_by('-registro')[:2]
    eventos_recentes_synaptic = SynapticAtendimento.objects.order_by('-registro')[:2]
    eventos_recentes_autokee = AutoKeeAtendimento.objects.order_by('-registro')[:2]
    eventos_recentes = list(eventos_recentes_sige) + list(eventos_recentes_synaptic) + list(eventos_recentes_autokee)
    
    # Dados recentes de AutoKee para a tabela
    registros = AutoKeeAtendimento.objects.order_by('-registro')[:10]

    context = {
        'total_atendimentos': total_atendimentos,
        'total_escolas': total_escolas,
        'atendimentos_mensais': atendimentos_mensais,
        'eventos_recentes': eventos_recentes,
        'registros': registros,
    }

    return render(request, 'webapp/soe_principal.html', context)
###***********************************************************************************************************************

from django.shortcuts import render
from .models import AutoKeeAtendimento
from django.core.paginator import Paginator
from django.db.models import Count

def ocorrencias_autokee(request):
    registros = AutoKeeAtendimento.objects.all()

    # Filtros de busca
    unidade = request.GET.get('unidade')
    classificacao = request.GET.get('classificacao')
    status_descricao = request.GET.get('status')

    if unidade:
        registros = registros.filter(nome_unidade_ensino__icontains=unidade)
    if classificacao:
        registros = registros.filter(classificacao_nome__icontains=classificacao)
    if status_descricao:
        registros = registros.filter(status_descricao__icontains=status_descricao)

    # Paginação
    paginator = Paginator(registros, 10)  # 10 registros por página
    page_number = request.GET.get('page')
    page_obj = paginator.get_page(page_number)

    # Contagem de status
    status_count = registros.values('status_descricao').annotate(count=Count('id')).order_by('-count')

    context = {
        'page_obj': page_obj,
        'status_count': status_count,
        'total_ocorrencias': registros.count(),
        'classificacoes_list': registros.values_list('classificacao_nome', flat=True).distinct(),
        'unidades_list': registros.values_list('nome_unidade_ensino', flat=True).distinct(),
        'status_list': registros.values_list('status_descricao', flat=True).distinct(),
    }

    return render(request, 'webapp/ocorrencias_autokee.html', context)
###***********************************************************************************************************************

from django.shortcuts import render
from .models import OrientadoresAtendimento

def ocorrencias_orientadores(request):
    registros = OrientadoresAtendimento.objects.all()[:100]  # Limite de 10 registros

    context = {
        'registros': registros,
        'total_ocorrencias': registros.count(),
    }
    return render(request, 'webapp/ocorrencias_orientadores.html', context)
###***********************************************************************************************************************

from django.shortcuts import render
from django.db.models import Sum
from .models import OrientadoresAtendimento

def orientadores_ocorrencias(request):
    # Recuperando os 10 primeiros registros para exibição
    registros = OrientadoresAtendimento.objects.all()[:100]

    # Calculando os totais para exibição nos indicadores
    total_ocorrencias = registros.aggregate(total=Sum('total_geral'))['total'] or 0
    total_violencia_sexual = registros.aggregate(total=Sum('violencia_sexual'))['total'] or 0
    total_uso_alcool_drogas = registros.aggregate(total=Sum('uso_alcool_drogas'))['total'] or 0

    context = {
        'registros': registros,
        'total_ocorrencias': total_ocorrencias,
        'total_violencia_sexual': total_violencia_sexual,
        'total_uso_alcool_drogas': total_uso_alcool_drogas,
    }

    return render(request, 'webapp/orientadores_ocorrencias.html', context)
###***********************************************************************************************************************

# views.py
from django.shortcuts import get_object_or_404, redirect, render
from django.contrib.auth.models import User
from django.contrib import messages

def editar_usuario(request, user_id):
    usuario = get_object_or_404(User, id=user_id)
    
    if request.method == 'POST':
        usuario.first_name = request.POST.get('first_name')
        usuario.last_name = request.POST.get('last_name')
        usuario.email = request.POST.get('email')
        usuario.save()
        messages.success(request, 'Usuário atualizado com sucesso!')
        return redirect('controle_usuarios')

    context = {'usuario': usuario}
    return render(request, 'semedapp/editar_usuario.html', context)
###***********************************************************************************************************************

# views.py
from django.shortcuts import get_object_or_404, redirect
from django.contrib.auth.models import User
from django.contrib import messages

def excluir_usuario(request, user_id):
    usuario = get_object_or_404(User, id=user_id)
    usuario.delete()
    messages.success(request, 'Usuário excluído com sucesso!')
    return redirect('controle_usuarios')
###***********************************************************************************************************************

# views.py
# from django.shortcuts import render, redirect
# from .forms import ProfessorForm, CoordenadorForm

from django.shortcuts import render, redirect
from .models import ProfessorEI, CoordenadorEI  # Certifique-se de importar os modelos corretos

def adicionar_professor(request):
    if request.method == "POST":
        unidade_ensino = request.POST.get('unidade_ensino')
        ano = request.POST.get('ano')
        modalidade = request.POST.get('modalidade')
        formato_letivo = request.POST.get('formato_letivo')
        turma = request.POST.get('turma')
        nome_professor = request.POST.get('nome_professor')
        cpf_professor = request.POST.get('cpf_professor')
        email_professor = request.POST.get('email_professor')

        ProfessorEI.objects.create(
            unidade_ensino=unidade_ensino,
            ano=ano,
            modalidade=modalidade,
            formato_letivo=formato_letivo,
            turma=turma,
            nome_professor=nome_professor,
            cpf_professor=cpf_professor,
            email_professor=email_professor
        )
        return redirect('controle_usuarios')

    return render(request, 'semedapp/adicionar_professor.html')
###***********************************************************************************************************************

def adicionar_coordenador(request):
    if request.method == "POST":
        unidade_ensino = request.POST.get('unidade_ensino')
        ano = request.POST.get('ano')
        modalidade = request.POST.get('modalidade')
        formato_letivo = request.POST.get('formato_letivo')
        nome_coordenadora = request.POST.get('nome_Coordenadora')
        cpf_professor = request.POST.get('cpf_professor')
        email_coordenadora = request.POST.get('email_Coordenadora')

        CoordenadorEI.objects.create(
            unidade_ensino=unidade_ensino,
            ano=ano,
            modalidade=modalidade,
            formato_letivo=formato_letivo,
            nome_Coordenadora=nome_coordenadora,
            cpf_professor=cpf_professor,
            email_Coordenadora=email_coordenadora
        )
        return redirect('controle_usuarios')

    return render(request, 'semedapp/adicionar_coordenador.html')
###***********************************************************************************************************************

def imprimir_curriculo_por_cpf(request, cpf):
    try:
        diretor = Diretor.objects.get(cpf=cpf)
        html_string = render_to_string('semedapp/imprimir_curriculo.html', {'diretor': diretor})
        pdf = HTML(string=html_string).write_pdf()
        response = HttpResponse(pdf, content_type='application/pdf')
        response['Content-Disposition'] = f'inline; filename="curriculo_{diretor.nome_completo}.pdf"'
        return response
    except Diretor.DoesNotExist:
        return HttpResponse("Currículo não encontrado", status=404)
###***********************************************************************************************************************

def imprimir_curriculo_por_diretor(request, diretor_id):
    try:
        diretor = Diretor.objects.get(id=diretor_id)
        html_string = render_to_string('semedapp/imprimir_curriculo.html', {'diretor': diretor})
        pdf = HTML(string=html_string).write_pdf()
        response = HttpResponse(pdf, content_type='application/pdf')
        response['Content-Disposition'] = f'inline; filename="curriculo_{diretor.nome_completo}.pdf"'
        return response
    except Diretor.DoesNotExist:
        return HttpResponse("Currículo não encontrado", status=404)
###***********************************************************************************************************************

def imprimir_curriculo_por_candidato(request, candidato_id):
    try:
        candidato = CadastroCandidato.objects.get(id=candidato_id)
        html_string = render_to_string('banco_curriculos/curriculo_pdf.html', {'candidato': candidato})
        pdf = HTML(string=html_string).write_pdf()
        response = HttpResponse(pdf, content_type='application/pdf')
        response['Content-Disposition'] = f'inline; filename="curriculo_{candidato.nome_completo}.pdf"'
        return response
    except CadastroCandidato.DoesNotExist:
        return HttpResponse("Currículo não encontrado para o candidato.", status=404)
###***********************************************************************************************************************

def imprimir_curriculo_diretor(request, diretor_id):
    try:
        diretor = Diretor.objects.get(id=diretor_id)
        html_string = render_to_string('semedapp/imprimir_curriculo.html', {'diretor': diretor})
        pdf = HTML(string=html_string).write_pdf()
        response = HttpResponse(pdf, content_type='application/pdf')
        response['Content-Disposition'] = f'inline; filename="curriculo_{diretor.nome_completo}.pdf"'
        return response
    except Diretor.DoesNotExist:
        return HttpResponse("Currículo não encontrado", status=404)
###***********************************************************************************************************************

from django.contrib.auth import authenticate, login, logout
from django.shortcuts import render, redirect
from django.contrib import messages
from django.contrib.auth.decorators import login_required
from .models import CustomUserProf
from .forms import CustomUserProfCreationForm

def registrar_professor(request):
    if request.method == "POST":
        form = CustomUserProfCreationForm(request.POST)
        if form.is_valid():
            form.save()
            messages.success(request, "Cadastro realizado com sucesso!")
            return redirect("login_prof")
    else:
        form = CustomUserProfCreationForm()
    return render(request, "semedapp/registrar_professor.html", {"form": form})
###***********************************************************************************************************************

def login_prof(request):
    if request.method == 'POST':
        username = request.POST.get('username')
        password = request.POST.get('password')
        
        user = authenticate(request, username=username, password=password)
        
        if user is not None:
            if user.is_professor:
                login(request, user)
                return redirect('modulo_pedagogico')  # Redireciona para o Módulo Pedagógico
            else:
                messages.error(request, "Você não tem permissão para acessar esta área.")
        else:
            messages.error(request, "Nome de usuário ou senha incorretos.")
    
    return render(request, 'semedapp/login_prof.html')
###***********************************************************************************************************************

from django.shortcuts import render
from django.contrib.auth.decorators import login_required
from .models import CadastroEI, Escolas

@login_required
def modulo_pedagogico(request):
    usuario_logado = request.user
    turma_nome = request.GET.get('turma', '')

    # Lista de usernames com acesso total à rede
    usernames_com_acesso_total = [
        "Valder", "Tatiane", "Edna", "Simone"
    ]

    # Permissão total se for superuser ou nome na lista
    if usuario_logado.is_superuser or usuario_logado.username in usernames_com_acesso_total:
        alunos = CadastroEI.objects.all()
    else:
        escolas_vinculadas = Escolas.objects.filter(coordenador=usuario_logado)
        if escolas_vinculadas.exists():
            alunos = CadastroEI.objects.filter(unidade_ensino__in=escolas_vinculadas.values_list('nome', flat=True))
        else:
            alunos = CadastroEI.objects.filter(professor=usuario_logado)

    # Aplicar filtro de turma (se for texto simples, use icontains)
    if turma_nome:
        alunos = alunos.filter(turma__icontains=turma_nome)

    # Listar turmas únicas (ordenadas)
    turmas_disponiveis = alunos.order_by("turma").values_list("turma", flat=True).distinct()

    context = {
        'alunos': alunos,
        'turma_nome': turma_nome,
        'turmas_disponiveis': turmas_disponiveis,
    }

    return render(request, 'webapp/modulo_pedagogico.html', context)


###***********************************************************************************************************************

def logout_prof(request):
    logout(request)
    messages.success(request, "Você saiu com sucesso.")
    return redirect('login_prof')

# from django.contrib.auth.decorators import login_required
# from django.http import HttpResponseForbidden

# @login_required
# def pagina_professor(request):
#     if not request.user.is_authenticated:
#         return HttpResponseForbidden("Acesso negado.")
#     return render(request, 'semedapp/pagina_professor.html')
###***********************************************************************************************************************

def login_prof_view(request):
    print("View login_prof_view foi chamada")
    
    # Obter as unidades distintas
    unidades = CadastroEI.objects.values_list('unidade_ensino', flat=True).distinct()
    
    # Obter os nomes das turmas distintas usando o campo `nome` da relação de turma
    turmas = CadastroEI.objects.order_by('turma__nome').values_list('turma__nome', flat=True).distinct()
    
    print("Unidades:", list(unidades))
    print("Turmas:", list(turmas))
    
    context = {
        'unidades': list(unidades),  # Converter para lista explícita para evitar erros no template
        'turmas': list(turmas),
    }
    return render(request, 'semedapp/login_prof.html', context)
###***********************************************************************************************************************

from django.http import JsonResponse
from .models import CadastroEI

def carregar_unidades_turmas(request):
    unidades = CadastroEI.objects.values_list('unidade_ensino', flat=True).distinct()
    turmas = CadastroEI.objects.order_by('turma__nome').values_list('turma__nome', flat=True).distinct()
    
    data = {
        "unidades": list(unidades),
        "turmas": list(turmas)
    }
    return JsonResponse(data)
###***********************************************************************************************************************

from django.contrib.auth.decorators import login_required
@login_required
def your_view(request):
    is_professor = request.user.groups.filter(name="Professor").exists()
    
    context = {
        'is_professor': is_professor,
    }
    return render(request, 'your_template.html', context)
###***********************************************************************************************************************
# user_groups.py
from django import template

register = template.Library()

@register.filter
def has_group(user, group_name):
    return user.groups.filter(name=group_name).exists()
###***********************************************************************************************************************

from django.shortcuts import render, get_object_or_404, redirect
from .models import Aluno  # Certifique-se de importar o modelo correto
from .forms import AlunoForm  # Certifique-se de criar um formulário correspondente

def editar_aluno(request, id):
    aluno = get_object_or_404(Aluno, id=id)  # Busca o aluno pelo ID ou retorna 404 se não existir
    
    if request.method == "POST":
        form = AlunoForm(request.POST, instance=aluno)
        if form.is_valid():
            form.save()
            return redirect('semedapp:modulo_pedagogico')  # Redireciona para a lista de alunos ou dashboard
    else:
        form = AlunoForm(instance=aluno)
    
    return render(request, 'webapp/editar_aluno.html', {'form': form, 'aluno': aluno})
###***********************************************************************************************************************

from django.shortcuts import get_object_or_404, redirect
from .models import Aluno  # substitua pelo nome correto do seu modelo

def avaliar_aluno(request, id):
    aluno = get_object_or_404(Aluno, id=id)
    aluno.avaliado = "SIM"  # atualize o campo para indicar que foi avaliado
    aluno.save()
    return redirect('modulo_pedagogico')  # substitua pelo nome correto da sua URL de redirecionamento
###***********************************************************************************************************************

from django.http import JsonResponse
from .models import Escola, Turma  # Certifique-se de que esses modelos existem

def carregar_turmas_por_escola(request):
    escola_nome = request.GET.get('escola')
    if escola_nome:
        turmas = Turma.objects.filter(escola__nome=escola_nome).values_list('nome', flat=True).distinct()
        return JsonResponse({'turmas': list(turmas)})
    return JsonResponse({'turmas': []})
###***********************************************************************************************************************

from django.shortcuts import render
from semedapp.decorators import module_required

@module_required('pedagogico')
def minha_view_pedagogica(request):
    return render(request, "pedagogico.html")
###***********************************************************************************************************************

from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from .models import Diretor, CadastroCandidato

def editar_curriculo(request, diretor_id):
    # Buscar o currículo na tabela `semedapp_diretor` usando o ID
    diretor = get_object_or_404(Diretor, id=diretor_id)

    # Buscar o candidato pelo CPF na tabela `semedapp_cadastrocandidato`
    candidato = CadastroCandidato.objects.filter(cpf=diretor.cpf).first()

    if request.method == "POST":
        try:
            # Capturar os valores do formulário e atualizar o objeto `diretor`
            diretor.nome_completo = request.POST.get("nome_completo", diretor.nome_completo)
            diretor.cpf = request.POST.get("cpf", diretor.cpf)
            diretor.rg = request.POST.get("rg", diretor.rg)
            diretor.email = request.POST.get("email", diretor.email)
            diretor.telefone = request.POST.get("telefone", diretor.telefone)
            diretor.endereco = request.POST.get("endereco", diretor.endereco)
            diretor.bairro = request.POST.get("bairro", diretor.bairro)
            diretor.cidade = request.POST.get("cidade", diretor.cidade)
            diretor.cep = request.POST.get("cep", diretor.cep)
            diretor.estado_civil = request.POST.get("estado_civil", diretor.estado_civil)
            diretor.sexo = request.POST.get("sexo", diretor.sexo)
            diretor.data_nascimento = request.POST.get("data_nascimento") or diretor.data_nascimento
            diretor.formacao_academica = request.POST.get("formacao_academica", diretor.formacao_academica)
            diretor.curso = request.POST.get("curso", diretor.curso)
            diretor.instituicao = request.POST.get("instituicao", diretor.instituicao)
            diretor.ano_conclusao = request.POST.get("ano_conclusao") or diretor.ano_conclusao
            diretor.experiencia_profissional = request.POST.get("experiencia_profissional", diretor.experiencia_profissional)
            diretor.empresa = request.POST.get("empresa", diretor.empresa)
            diretor.cargo = request.POST.get("cargo", diretor.cargo)  # Atualiza o cargo
            diretor.data_inicio = request.POST.get("data_inicio") or diretor.data_inicio
            diretor.data_fim = request.POST.get("data_fim") or diretor.data_fim

            # Processar arquivos enviados (caso haja upload de novos arquivos)
            if "foto" in request.FILES:
                diretor.foto = request.FILES["foto"]
            if "curriculo_pdf" in request.FILES:
                diretor.curriculo_pdf = request.FILES["curriculo_pdf"]
            if "certificados_pdf" in request.FILES:
                diretor.certificados_pdf = request.FILES["certificados_pdf"]

            # **Salvar no banco de dados**
            diretor.save()

            # Mensagem de sucesso
            messages.success(request, "Currículo atualizado com sucesso!")
            return redirect("area_candidato")

        except Exception as e:
            # Mensagem de erro e print do erro no console
            print(f"Erro ao atualizar currículo: {e}")
            messages.error(request, f"Erro ao atualizar currículo: {str(e)}")
            return redirect("editar_curriculo", diretor_id=diretor.id)

    return render(request, "banco_curriculos/editar_curriculo.html", {
        "candidato": candidato,
        "diretor": diretor,
    })
###***********************************************************************************************************************


from django.shortcuts import get_object_or_404
from django.http import HttpResponse
from django.template.loader import render_to_string
from weasyprint import HTML
from .models import Diretor

def imprimir_curriculo_alternativo(request, diretor_id):
    # 🔍 Busca o diretor pelo ID
    diretor = get_object_or_404(Diretor, id=diretor_id)

    # 🔍 Renderiza o HTML com os dados do diretor
    html_content = render_to_string("banco_curriculos/curriculo_pdf_alternativo.html", {"diretor": diretor})

    # 🔍 Gera o PDF com WeasyPrint
    pdf = HTML(string=html_content).write_pdf()

    # 🔍 Retorna o PDF como resposta
    response = HttpResponse(pdf, content_type="application/pdf")
    response["Content-Disposition"] = f'attachment; filename="curriculo_alternativo_{diretor_id}.pdf"'

    return response
###***********************************************************************************************************************

###****************************************SEPECC*************************************************************************
###****************************************SEPECC*************************************************************************
###****************************************SEPECC*************************************************************************
from django.shortcuts import render, redirect
from django.contrib import messages
from django.views import View
from .models import EscolaPdde, Programa

class EscolaPddeCreateView(View):
    template_name = 'semedapp/escola_pdde_form.html'

    def get(self, request):
        """Renderiza o formulário vazio com a lista de programas disponíveis."""
        programas = Programa.objects.all()
        return render(request, self.template_name, {
            "programas": programas,
        })

    def post(self, request):
        """Recebe os dados do formulário e salva no banco de dados."""
        try:
            # Captura os dados do formulário
            nome = request.POST.get('nome')
            endereco = request.POST.get('endereco')
            zona = request.POST.get("zona")
            ensino = request.POST.get("ensino")
            cep = request.POST.get('cep')
            bairro = request.POST.get('bairro')
            cidade = request.POST.get('cidade', "Canaã dos Carajás")
            uf = request.POST.get('uf', "PA")
            tipo = request.POST.get('tipo')
            dependencia_administrativa = request.POST.get('dependencia_administrativa')
            codigo_inep = request.POST.get('codigo_inep')
            cnpj = request.POST.get('cnpj')
            nome_conselho = request.POST.get('nome_conselho')

            # Captura a data de validade da procuração (pode ser None)
            validade_procuracao = request.POST.get('validade_procuracao', None)

            # Captura os números (com valores padrão se estiverem vazios)
            quantidade_salas = int(request.POST.get('quantidade_salas', 0))
            quantidade_turmas = int(request.POST.get('quantidade_turmas', 0))
            quantidade_professores = int(request.POST.get('quantidade_professores', 0))
            quantidade_alunos = int(request.POST.get('quantidade_alunos', 0))

            # Captura o arquivo de procuração (pode ser None)
            procuracao = request.FILES.get('procuracao', None)

            # Verificação de campos obrigatórios
            if not nome or not codigo_inep or not cnpj:
                messages.error(request, "Os campos Nome, Código INEP e CNPJ são obrigatórios.")
                return render(request, self.template_name, {
                    "programas": Programa.objects.all(),
                    "nome": nome, "endereco": endereco, "zona": zona, "ensino": ensino,
                    "cep": cep, "bairro": bairro, "cidade": cidade, "uf": uf, "tipo": tipo,
                    "dependencia_administrativa": dependencia_administrativa, "codigo_inep": codigo_inep,
                    "cnpj": cnpj, "nome_conselho": nome_conselho, "validade_procuracao": validade_procuracao,
                    "quantidade_salas": quantidade_salas, "quantidade_turmas": quantidade_turmas,
                    "quantidade_professores": quantidade_professores, "quantidade_alunos": quantidade_alunos
                })

            # Criar a instância da escola
            escola = EscolaPdde.objects.create(
                nome=nome,
                endereco=endereco,
                zona=zona,
                ensino=ensino,
                cep=cep,
                bairro=bairro,
                cidade=cidade,
                uf=uf,
                tipo=tipo,
                dependencia_administrativa=dependencia_administrativa,
                codigo_inep=codigo_inep,
                cnpj=cnpj,
                nome_conselho=nome_conselho,
                procuracao=procuracao if procuracao else None,  # Garante que pode ser None
                validade_procuracao=validade_procuracao if validade_procuracao else None,  # Garante que pode ser None
                quantidade_salas=quantidade_salas,
                quantidade_turmas=quantidade_turmas,
                quantidade_professores=quantidade_professores,
                quantidade_alunos=quantidade_alunos
            )

            # Captura os programas selecionados
            programas_selecionados = request.POST.getlist("programas")
            if programas_selecionados:
                escola.programas.set(Programa.objects.filter(id__in=programas_selecionados))

            messages.success(request, "Escola cadastrada com sucesso!")
            return redirect('lista_escolas_pdde')

        except Exception as e:
            messages.error(request, f"Ocorreu um erro ao salvar a escola: {str(e)}")
            return render(request, self.template_name, {
                "programas": Programa.objects.all(),
            })
###***********************************************************************************************************************

from django.shortcuts import render
from django.views import View
from .models import EscolaPdde

class EscolaPddeListView(View):
    template_name = 'semedapp/escola_pdde_list.html'  # Certifique-se de que o template correto está sendo usado

    def get(self, request):
        escolas = EscolaPdde.objects.all()
        total_escolas = escolas.count()
        total_urbanas = escolas.filter(zona="Urbana").count()
        total_rurais = escolas.filter(zona="Rural").count()
        
        # Contagem por nível de ensino
        total_ensino_fundamental = escolas.filter(ensino__icontains="Ensino Fundamental").count()
        total_ensino_medio = escolas.filter(ensino__icontains="Ensino Médio").count()
        total_educacao_infantil = escolas.filter(ensino__icontains="Educação Infantil").count()

        context = {
            'escolas': escolas,
            'total_escolas': total_escolas,
            'total_urbanas': total_urbanas,
            'total_rurais': total_rurais,
            'total_ensino_fundamental': total_ensino_fundamental,
            'total_ensino_medio': total_ensino_medio,
            'total_educacao_infantil': total_educacao_infantil
        }
        return render(request, self.template_name, context)
###***********************************************************************************************************************

from django.shortcuts import get_object_or_404, redirect
from django.views import View
from django.contrib import messages
from .models import EscolaPdde

class EscolaPddeDeleteView(View):
    def get(self, request, pk):
        escola = get_object_or_404(EscolaPdde, pk=pk)
        escola.delete()
        messages.success(request, "Escola excluída com sucesso!")
        return redirect("lista_escolas_pdde")  # Redireciona para a lista após excluir
###***********************************************************************************************************************

from django.shortcuts import render, get_object_or_404
from django.views import View
from .models import EscolaPdde, Pagamento

class EscolaPddeDetailView(View):
    template_name = "semedapp/escola_pdde_detail.html"

    def get(self, request, pk):
        # Obtém a escola específica
        escola = get_object_or_404(EscolaPdde, pk=pk)
        
        # Obtém os pagamentos relacionados à escola
        pagamentos = Pagamento.objects.filter(escola=escola)

        # Calcula o total dos pagamentos
        total_pagamentos = sum(pagamento.valor for pagamento in pagamentos)

        context = {
            "escola": escola,
            "pagamentos": pagamentos,
            "total_pagamentos": total_pagamentos
        }
        return render(request, self.template_name, context)
###***********************************************************************************************************************

from django.shortcuts import render, get_object_or_404, redirect
from django.contrib import messages
from django.views import View
from .models import EscolaPdde, Programa
from .forms import EscolaPddeForm

class EscolaPddeUpdateView(View):
    template_name = "semedapp/escola_pdde_edit.html"

    def get(self, request, escola_id):
        escola = get_object_or_404(EscolaPdde, id=escola_id)
        form = EscolaPddeForm(instance=escola)
        return render(request, self.template_name, {
            "form": form,
            "escola": escola
        })

    def post(self, request, escola_id):
        escola = get_object_or_404(EscolaPdde, id=escola_id)
        form = EscolaPddeForm(request.POST, request.FILES, instance=escola)

        if form.is_valid():
            # Captura a instância sem salvar ainda
            escola = form.save(commit=False)

            # Verifica se há um novo arquivo de procuração antes de sobrescrever
            if "procuracao" in request.FILES:
                escola.procuracao = request.FILES["procuracao"]

            # Salva a instância da escola
            escola.save()

            # Atualiza os programas vinculados
            form.save_m2m()

            messages.success(request, "Escola atualizada com sucesso!")
            return redirect("lista_escolas_pdde")
        else:
            messages.error(request, "Erro ao atualizar a escola. Verifique os dados!")

        return render(request, self.template_name, {
            "form": form,
            "escola": escola
        })
###***********************************************************************************************************************

from django.http import JsonResponse
from django.shortcuts import get_object_or_404
from .models import EscolaPdde, ReceitaDespesa

def get_escola_pdde(request, pk):
    """
    Retorna os dados da escola PDDE em formato JSON para preencher os campos automaticamente.
    """
    escola = get_object_or_404(EscolaPdde, pk=pk)

    data = {
        "id": escola.id,
        "nome_conselho": escola.nome_conselho if escola.nome_conselho else "Não informado",
        "cnpj": escola.cnpj,
        "endereco": escola.endereco,
        "bairro": escola.bairro,
        "cidade": escola.cidade,
        "uf": escola.uf,
    }

    return JsonResponse(data)
###***********************************************************************************************************************

from django.http import JsonResponse
from django.db.models import Sum, Min, Max
from .models import Receita, Pagamento

def get_receita_despesa(request, escola_id):
    try:
        # Busca o primeiro objeto Receita relacionado à escola
        receita_obj = Receita.objects.filter(escola_id=escola_id).first()

        if not receita_obj:
            return JsonResponse({"erro": "Nenhuma receita encontrada para esta escola."}, status=404)

        # Agregados para somar valores financeiros
        receita = Receita.objects.filter(escola_id=escola_id).aggregate(
            saldo_anterior_custeio=Sum('saldo_anterior_custeio', default=0),
            saldo_anterior_capital=Sum('saldo_anterior_capital', default=0),
            valor_creditado_custeio=Sum('valor_creditado_custeio', default=0),
            valor_creditado_capital=Sum('valor_creditado_capital', default=0),
            recursos_proprios_custeio=Sum('recursos_proprios_custeio', default=0),
            recursos_proprios_capital=Sum('recursos_proprios_capital', default=0),
            rendimento_aplicacao_custeio=Sum('rendimento_aplicacao_custeio', default=0),
            rendimento_aplicacao_capital=Sum('rendimento_aplicacao_capital', default=0),
            devolucao_fnde_custeio=Sum('devolucao_fnde_custeio', default=0),
            devolucao_fnde_capital=Sum('devolucao_fnde_capital', default=0),
            saldo_devolvido_custeio=Sum('saldo_devolvido_custeio', default=0),
            saldo_devolvido_capital=Sum('saldo_devolvido_capital', default=0),
            valor_despesa_realizada_custeio=Sum('valor_despesa_realizada_custeio', default=0),
            valor_despesa_realizada_capital=Sum('valor_despesa_realizada_capital', default=0),
            saldo_reprogramar_custeio=Sum('saldo_reprogramar_custeio', default=0),
            saldo_reprogramar_capital=Sum('saldo_reprogramar_capital', default=0),
            data_inicio=Min('data_inicio'),
            data_fim=Max('data_fim')
        )

        total_pagamentos = Pagamento.objects.filter(escola_id=escola_id).aggregate(
            total_pagamentos=Sum('valor', default=0)
        )['total_pagamentos'] or 0

        # Calculando os totais corretamente
        valor_total_receita_custeio = (
            receita['saldo_anterior_custeio'] +
            receita['valor_creditado_custeio'] +
            receita['recursos_proprios_custeio'] +
            receita['rendimento_aplicacao_custeio']
        )
        valor_total_receita_capital = (
            receita['saldo_anterior_capital'] +
            receita['valor_creditado_capital'] +
            receita['recursos_proprios_capital'] +
            receita['rendimento_aplicacao_capital']
        )

        total_receita = valor_total_receita_custeio + valor_total_receita_capital
        total_despesa = (
            receita['devolucao_fnde_custeio'] +
            receita['devolucao_fnde_capital'] +
            receita['saldo_devolvido_custeio'] +
            receita['saldo_devolvido_capital'] +
            total_pagamentos
        )
        saldo_reprogramar = total_receita - total_despesa

        data_inicio_formatada = receita['data_inicio'].strftime("%Y-%m-%d") if receita['data_inicio'] else ""
        data_fim_formatada = receita['data_fim'].strftime("%Y-%m-%d") if receita['data_fim'] else ""

        data = {
            "data_inicio": data_inicio_formatada,
            "data_fim": data_fim_formatada,
            "saldo_anterior_custeio": f"{receita['saldo_anterior_custeio']:.2f}".replace(".", ","),
            "saldo_anterior_capital": f"{receita['saldo_anterior_capital']:.2f}".replace(".", ","),
            "valor_creditado_custeio": f"{receita['valor_creditado_custeio']:.2f}".replace(".", ","),
            "valor_creditado_capital": f"{receita['valor_creditado_capital']:.2f}".replace(".", ","),
            "recursos_proprios_custeio": f"{receita['recursos_proprios_custeio']:.2f}".replace(".", ","),
            "recursos_proprios_capital": f"{receita['recursos_proprios_capital']:.2f}".replace(".", ","),
            "rendimento_aplicacao_custeio": f"{receita['rendimento_aplicacao_custeio']:.2f}".replace(".", ","),
            "rendimento_aplicacao_capital": f"{receita['rendimento_aplicacao_capital']:.2f}".replace(".", ","),
            "devolucao_fnde_custeio": f"{receita['devolucao_fnde_custeio']:.2f}".replace(".", ","),
            "devolucao_fnde_capital": f"{receita['devolucao_fnde_capital']:.2f}".replace(".", ","),
            "saldo_devolvido_custeio": f"{receita['saldo_devolvido_custeio']:.2f}".replace(".", ","),
            "saldo_devolvido_capital": f"{receita['saldo_devolvido_capital']:.2f}".replace(".", ","),
            "valor_total_receita_custeio": f"{valor_total_receita_custeio:.2f}".replace(".", ","),
            "valor_total_receita_capital": f"{valor_total_receita_capital:.2f}".replace(".", ","),
            "valor_despesa_realizada_custeio": f"{receita['valor_despesa_realizada_custeio']:.2f}".replace(".", ","),
            "valor_despesa_realizada_capital": f"{receita['valor_despesa_realizada_capital']:.2f}".replace(".", ","),
            "saldo_reprogramar_custeio": f"{receita['saldo_reprogramar_custeio']:.2f}".replace(".", ","),
            "saldo_reprogramar_capital": f"{receita['saldo_reprogramar_capital']:.2f}".replace(".", ","),
            "total_receita": f"{total_receita:.2f}".replace(".", ","),
            "total_despesa": f"{total_despesa:.2f}".replace(".", ","),
            "saldo_reprogramar": f"{saldo_reprogramar:.2f}".replace(".", ","),
            "escolas_atendidas": receita_obj.escolas_atendidas or 0
        }

        return JsonResponse(data)

    except Exception as e:
        return JsonResponse({"erro": f"Erro ao buscar dados financeiros: {str(e)}"}, status=500)


    
###***********************************************************************************************************************

from django.http import JsonResponse
from django.shortcuts import get_object_or_404
from .models import EscolaPdde  # Certifique-se de importar corretamente o modelo

def get_escola_pdde_modal(request, pk):
    """
    Retorna os dados da escola PDDE em formato JSON para preencher a modal via AJAX.
    """
    escola = get_object_or_404(EscolaPdde, pk=pk)

    data = {
        "id": escola.id,
        "nome": escola.nome,
        "cnpj": escola.cnpj,
        "endereco": escola.endereco,
        "bairro": escola.bairro,
        "cidade": escola.cidade,
        "uf": escola.uf,
    }

    return JsonResponse(data)
###***********************************************************************************************************************

from django.shortcuts import render
from django.views import View
from .models import EscolaPdde

class EscolaPddeModalView(View):
    template_name = "contabilidade/pdde.html"

    def get(self, request):
        """
        Renderiza a página da modal carregando todas as escolas disponíveis no banco de dados.
        """
        escolas = EscolaPdde.objects.all().order_by("nome")  # Busca todas as escolas
        context = {"escolas": escolas, "escola_selecionada": None}  # Nenhuma escola pré-selecionada
        return render(request, self.template_name, context)

    def post(self, request):
        """
        Processa o formulário e retorna os dados da escola selecionada para preencher os campos da modal.
        """
        escola_id = request.POST.get("escola")  # Obtém o ID da escola selecionada
        escola_selecionada = EscolaPdde.objects.filter(id=escola_id).first() if escola_id else None

        escolas = EscolaPdde.objects.all().order_by("nome")  # Lista todas as escolas

        context = {
            "escolas": escolas,
            "escola_selecionada": escola_selecionada,  # Define a escola selecionada, se houver
        }
        return render(request, self.template_name, context)
###***********************************************************************************************************************

from django.http import JsonResponse
from django.shortcuts import render
from .models import Pagamento, EscolaPdde


def lancamento_pagamento(request):
    escolas = EscolaPdde.objects.all().order_by("nome")  # Busca todas as escolas ordenadas
    return render(request, "lancar_pagamento.html", {"escolas": escolas})
###***********************************************************************************************************************

from django.http import JsonResponse
from .models import Pagamento

def get_pagamentos(request, escola_id):
    try:
        #print(f"🔹 Escola ID recebido: {escola_id}")  # Para depuração no terminal
        pagamentos = Pagamento.objects.filter(escola_id=escola_id).order_by("-data_pagamento")

        if not pagamentos.exists():
            return JsonResponse({"pagamentos": []}, safe=False)  # Retorna JSON válido

        lista_pagamentos = []
        for pagamento in pagamentos:
            try:
                lista_pagamentos.append({
                    "id": pagamento.id,
                    "nome_favorecido": pagamento.nome_favorecido or "Não informado",
                    "cnpj_cpf": pagamento.cnpj_cpf or "N/A",
                    "tipo_bem_servico": pagamento.tipo_bem_servico or "N/A",
                    "origem": pagamento.origem if pagamento.origem else "FNDE",
                    "tipo_pagamento": pagamento.tipo_pagamento or "N/A",
                    "tipo_documento": pagamento.tipo_documento or "N/A",
                    "numero_documento": pagamento.numero_documento or "",
                    "data_documento": pagamento.data_documento.strftime("%d/%m/%Y") if pagamento.data_documento else "",
                    "data_pagamento": pagamento.data_pagamento.strftime("%d/%m/%Y") if pagamento.data_pagamento else "Não informado",
                    "numero_pagamento": pagamento.numero_documento_pagamento or "",
                    "valor": float(pagamento.valor) if pagamento.valor else 0.00,
                    "url_editar": f"/editar-pagamento/{pagamento.id}/",
                })
            except Exception as e:
                print(f"❌ Erro ao processar pagamento {pagamento.id}: {e}")  # Log de erro específico
        #print(f"🔹 Pagamentos enviados: {lista_pagamentos}")  # Log para verificar resposta
        return JsonResponse({"pagamentos": lista_pagamentos}, safe=False)
    except Exception as e:
        #print(f"❌ Erro geral na view get_pagamentos: {e}")  # Mostra erro no console do servidor
        return JsonResponse({"erro": f"Erro ao buscar pagamentos: {str(e)}"}, status=500)
###***********************************************************************************************************************

from django.shortcuts import render
from .models import EscolaPdde

def pddereceita_despesa_view(request):
    escolas = EscolaPdde.objects.all()  # Buscar todas as escolas no banco de dados
    context = {"escolas": escolas}
    return render(request, "contabilidade/pddereceitadespesa.html", context)
###***********************************************************************************************************************

from django.shortcuts import render
from .models import EscolaPdde

def pdde_list(request):
    pdde_data = EscolaPdde.objects.all().order_by("-id")  # Busca todas as escolas cadastradas no PDDE
    return render(request, "seu_template.html", {"pdde_data": pdde_data})
###***********************************************************************************************************************

from django.shortcuts import render, get_object_or_404
from django.http import JsonResponse
from django.db.models import Sum
from .models import ReceitaDespesa, EscolaPdde, Pagamento, Programa, SemedAppEscolaPddeProgramas

def pddereceita_despesa(request, escola_id):
    """
    View que retorna os dados financeiros da escola, incluindo receitas, despesas, pagamentos e programas vinculados.
    """
    try:
        # 🔹 Obtém a escola
        escola = get_object_or_404(EscolaPdde, id=escola_id)

        # 🔹 Buscar IDs dos programas vinculados
        programas_ids = SemedAppEscolaPddeProgramas.objects.filter(escolapdde_id=escola_id).values_list('programa_id', flat=True)

        # 🔹 Buscar nomes dos programas vinculados
        programas_vinculados = Programa.objects.filter(id__in=programas_ids).values_list('nome', flat=True)

        # 🔹 Obtém as receitas e despesas da escola (somando os valores)
        receita_despesa = ReceitaDespesa.objects.filter(escola=escola).aggregate(
            saldo_anterior_custeio=Sum("saldo_anterior_custeio", default=0),
            saldo_anterior_capital=Sum("saldo_anterior_capital", default=0),
            valor_creditado_custeio=Sum("valor_creditado_custeio", default=0),
            valor_creditado_capital=Sum("valor_creditado_capital", default=0),

            # ✅ Separação dos Recursos Próprios
            recursos_proprios_custeio=Sum("recursos_proprios_custeio", default=0),
            recursos_proprios_capital=Sum("recursos_proprios_capital", default=0),

            rendimento_aplicacao_custeio=Sum("rendimento_aplicacao_custeio", default=0),
            rendimento_aplicacao_capital=Sum("rendimento_aplicacao_capital", default=0),
            devolucao_fnde_custeio=Sum("devolucao_fnde_custeio", default=0),
            devolucao_fnde_capital=Sum("devolucao_fnde_capital", default=0),
            saldo_devolvido_custeio=Sum("saldo_devolvido_custeio", default=0),
            saldo_devolvido_capital=Sum("saldo_devolvido_capital", default=0),
            valor_despesa_realizada_custeio=Sum("valor_despesa_realizada_custeio", default=0),
            valor_despesa_realizada_capital=Sum("valor_despesa_realizada_capital", default=0),
            saldo_reprogramar_custeio=Sum("saldo_reprogramar_custeio", default=0),
            saldo_reprogramar_capital=Sum("saldo_reprogramar_capital", default=0),
            escolas_atendidas=Sum("escolas_atendidas", default=0)
        )

        # 🔹 Obtém os pagamentos realizados pela escola
        pagamentos = Pagamento.objects.filter(escola=escola).values(
            "id",
            "nome_favorecido",
            "cnpj_cpf",
            "tipo_pagamento",
            "tipo_bem_servico",
            "tipo_documento",
            "numero_documento",
            "data_documento",
            "numero_pagamento",
            "data_pagamento",
            "origem",
            "valor"
        )

        # Adiciona o campo "url_editar" dinamicamente
        for pagamento in pagamentos:
            pagamento["url_editar"] = f"/editar-pagamento/{pagamento['id']}/"


        # 🔹 Calcula o total de receitas e despesas
        total_receita = sum([
            receita_despesa["saldo_anterior_custeio"],
            receita_despesa["saldo_anterior_capital"],
            receita_despesa["valor_creditado_custeio"],
            receita_despesa["valor_creditado_capital"],
            receita_despesa["recursos_proprios_custeio"],  # ✅ Novo Campo
            receita_despesa["recursos_proprios_capital"],  # ✅ Novo Campo
            receita_despesa["rendimento_aplicacao_custeio"],
            receita_despesa["rendimento_aplicacao_capital"]
        ])

        total_despesa = sum([
            receita_despesa["devolucao_fnde_custeio"],
            receita_despesa["devolucao_fnde_capital"],
            receita_despesa["saldo_devolvido_custeio"],
            receita_despesa["saldo_devolvido_capital"],
            receita_despesa["valor_despesa_realizada_custeio"],
            receita_despesa["valor_despesa_realizada_capital"]
        ])

        saldo_reprogramar = total_receita - total_despesa
        receita_liquida = total_receita - total_despesa

        # 🔹 Prepara os dados para JSONResponse
        data = {
            "escola": escola.nome,
            "programas_vinculados": list(programas_vinculados),
            "saldo_anterior_custeio": str(receita_despesa["saldo_anterior_custeio"]),
            "saldo_anterior_capital": str(receita_despesa["saldo_anterior_capital"]),
            "valor_creditado_custeio": str(receita_despesa["valor_creditado_custeio"]),
            "valor_creditado_capital": str(receita_despesa["valor_creditado_capital"]),

            # ✅ Incluindo os novos campos
            "recursos_proprios_custeio": str(receita_despesa["recursos_proprios_custeio"]),
            "recursos_proprios_capital": str(receita_despesa["recursos_proprios_capital"]),

            "rendimento_aplicacao_custeio": str(receita_despesa["rendimento_aplicacao_custeio"]),
            "rendimento_aplicacao_capital": str(receita_despesa["rendimento_aplicacao_capital"]),
            "devolucao_fnde_custeio": str(receita_despesa["devolucao_fnde_custeio"]),
            "devolucao_fnde_capital": str(receita_despesa["devolucao_fnde_capital"]),
            "saldo_devolvido_custeio": str(receita_despesa["saldo_devolvido_custeio"]),
            "saldo_devolvido_capital": str(receita_despesa["saldo_devolvido_capital"]),
            "valor_despesa_realizada_custeio": str(receita_despesa["valor_despesa_realizada_custeio"]),
            "valor_despesa_realizada_capital": str(receita_despesa["valor_despesa_realizada_capital"]),
            "saldo_reprogramar_custeio": str(receita_despesa["saldo_reprogramar_custeio"]),
            "saldo_reprogramar_capital": str(receita_despesa["saldo_reprogramar_capital"]),
            "escolas_atendidas": str(receita_despesa["escolas_atendidas"]),
            "total_receita": str(total_receita),
            "total_despesa": str(total_despesa),
            "saldo_reprogramar": str(saldo_reprogramar),
            "receita_liquida": str(receita_liquida),
            "pagamentos": list(pagamentos),
        }

        return JsonResponse(data)

    except Exception as e:
        return JsonResponse({"erro": f"Erro ao buscar dados financeiros: {str(e)}"}, status=500)
###***********************************************************************************************************************

import datetime
from django.shortcuts import render, get_object_or_404, redirect
from django.contrib import messages
from .models import Receita, EscolaPdde, Programa

def pddelancar_receita(request):
    escolas = EscolaPdde.objects.all().order_by("nome")
    receitas = None  
    receita = None  

    if request.method == "POST":
        try:
            escola_id = request.POST.get("escola")
            programa_id = request.POST.get("programa")

            if not escola_id or not programa_id:
                messages.error(request, "Escola e Programa são obrigatórios.")
                return redirect("pddereceita_despesa")

            escola = get_object_or_404(EscolaPdde, id=escola_id)
            programa = get_object_or_404(Programa, id=programa_id)

            data_inicio = request.POST.get("data_inicio")
            data_fim = request.POST.get("data_fim") or "2025-12-31"

            def parse_decimal(value):
                return float(value.replace(",", ".")) if value else 0.0
            
            saldo_anterior_custeio = parse_decimal(request.POST.get("saldo_anterior_custeio", "0"))
            saldo_anterior_capital = parse_decimal(request.POST.get("saldo_anterior_capital", "0"))
            valor_creditado_custeio = parse_decimal(request.POST.get("valor_creditado_custeio", "0"))
            valor_creditado_capital = parse_decimal(request.POST.get("valor_creditado_capital", "0"))

            # ✅ Ajustando os nomes dos campos corretamente
            recursos_proprios_custeio = parse_decimal(request.POST.get("recursos_proprios_custeio", "0"))
            recursos_proprios_capital = parse_decimal(request.POST.get("recursos_proprios_capital", "0"))

            rendimento_aplicacao_custeio = parse_decimal(request.POST.get("rendimento_aplicacao_custeio", "0"))
            rendimento_aplicacao_capital = parse_decimal(request.POST.get("rendimento_aplicacao_capital", "0"))
            devolucao_fnde_custeio = parse_decimal(request.POST.get("devolucao_fnde_custeio", "0"))
            devolucao_fnde_capital = parse_decimal(request.POST.get("devolucao_fnde_capital", "0"))

            saldo_devolvido_custeio = parse_decimal(request.POST.get("saldo_devolvido_custeio", "0"))
            saldo_devolvido_capital = parse_decimal(request.POST.get("saldo_devolvido_capital", "0"))

            valor_despesa_realizada_custeio = parse_decimal(request.POST.get("valor_despesa_realizada_custeio", "0"))
            valor_despesa_realizada_capital = parse_decimal(request.POST.get("valor_despesa_realizada_capital", "0"))

            saldo_reprogramar_custeio = parse_decimal(request.POST.get("saldo_reprogramar_custeio", "0"))
            saldo_reprogramar_capital = parse_decimal(request.POST.get("saldo_reprogramar_capital", "0"))

            escolas_atendidas = int(request.POST.get("escolas_atendidas", "1"))

            valor_total_receita_custeio = saldo_anterior_custeio + valor_creditado_custeio + rendimento_aplicacao_custeio
            valor_total_receita_capital = saldo_anterior_capital + valor_creditado_capital + rendimento_aplicacao_capital + recursos_proprios_capital

            receita, created = Receita.objects.update_or_create(
                escola=escola,
                programa=programa,
                defaults={
                    "data_inicio": data_inicio,
                    "data_fim": data_fim,
                    "saldo_anterior_custeio": saldo_anterior_custeio,
                    "saldo_anterior_capital": saldo_anterior_capital,
                    "valor_creditado_custeio": valor_creditado_custeio,
                    "valor_creditado_capital": valor_creditado_capital,
                    "recursos_proprios_custeio": recursos_proprios_custeio,
                    "recursos_proprios_capital": recursos_proprios_capital,
                    "rendimento_aplicacao_custeio": rendimento_aplicacao_custeio,
                    "rendimento_aplicacao_capital": rendimento_aplicacao_capital,
                    "devolucao_fnde_custeio": devolucao_fnde_custeio,
                    "devolucao_fnde_capital": devolucao_fnde_capital,
                    "saldo_devolvido_custeio": saldo_devolvido_custeio,
                    "saldo_devolvido_capital": saldo_devolvido_capital,
                    "valor_despesa_realizada_custeio": valor_despesa_realizada_custeio,
                    "valor_despesa_realizada_capital": valor_despesa_realizada_capital,
                    "saldo_reprogramar_custeio": saldo_reprogramar_custeio,
                    "saldo_reprogramar_capital": saldo_reprogramar_capital,
                    "valor_total_receita_custeio": valor_total_receita_custeio,
                    "valor_total_receita_capital": valor_total_receita_capital,
                    "escolas_atendidas": escolas_atendidas,
                }
            )

            messages.success(request, "Receita lançada com sucesso!")
            return redirect("pddereceita_despesa")

        except Exception as e:
            messages.error(request, f"Erro ao salvar: {str(e)}")
            return render(request, "erro.html", {"error_message": f"Erro ao salvar: {str(e)}"})

    return render(request, "pdde/lancar_receita.html", {"escolas": escolas})

###***********************************************************************************************************************

def pddelancar_despesa(request):
    return render(request, "pdde/lancar_despesa.html")
###***********************************************************************************************************************

from django.shortcuts import render, redirect
from django.contrib import messages
from django.db import transaction
from decimal import Decimal
from django.utils.dateparse import parse_date
from .models import Pagamento, EscolaPdde

def pddelancar_pagamento(request):
    escolas = EscolaPdde.objects.all()
    pagamentos = Pagamento.objects.all().order_by("-data_pagamento")  # Pagamentos mais recentes primeiro

    if request.method == "POST":
        try:
            with transaction.atomic():
                # 📌 Pegando dados do formulário
                escola_id = request.POST.get("escola")
                programa = request.POST.get("programa")
                nome_favorecido = request.POST.get("nome_favorecido")
                cnpj_cpf = request.POST.get("cnpj_cpf")
                tipo_pagamento = request.POST.get("tipo_pagamento")
                tipo_bem_servico = request.POST.get("tipo_bem_servico")

                tipo_documento = request.POST.get("tipo_documento")
                numero_documento = request.POST.get("numero_documento")
                data_documento = request.POST.get("data_documento")

                tipo_pagamento_efetuado = request.POST.get("tipo_pagamento_efetuado")
                numero_documento_pagamento = request.POST.get("numero_documento_pagamento")
                data_pagamento = request.POST.get("data_pagamento")

                valor = request.POST.get("valor")
                exercicio = request.POST.get("exercicio", "2025")

                # ✅ Validação de campos obrigatórios
                campos_obrigatorios = [
                    escola_id, programa, nome_favorecido, cnpj_cpf, tipo_pagamento, tipo_bem_servico,
                    tipo_documento, numero_documento, data_documento,
                    tipo_pagamento_efetuado, numero_documento_pagamento, data_pagamento, valor
                ]

                if not all(campos_obrigatorios):
                    messages.error(request, "⚠️ Todos os campos são obrigatórios. Preencha corretamente.")
                    return redirect("pddelancar_pagamento")

                # 🔍 Verifica se a escola existe
                escola = EscolaPdde.objects.get(id=escola_id)

                # 📅 Converte datas
                data_documento = parse_date(data_documento)
                data_pagamento = parse_date(data_pagamento)

                if not data_documento or not data_pagamento:
                    messages.error(request, "⚠️ Datas inválidas. Verifique os campos de data.")
                    return redirect("pddelancar_pagamento")

                # 💰 Converte valor para Decimal
                valor = Decimal(valor.replace(",", ".").strip())

                # 🚀 Cria o pagamento
                Pagamento.objects.create(
                    escola=escola,
                    programa=programa,
                    nome_favorecido=nome_favorecido,
                    cnpj_cpf=cnpj_cpf,
                    tipo_pagamento=tipo_pagamento,
                    tipo_bem_servico=tipo_bem_servico,
                    tipo_documento=tipo_documento,
                    numero_documento=numero_documento,
                    data_documento=data_documento,
                    tipo_pagamento_efetuado=tipo_pagamento_efetuado,
                    numero_documento_pagamento=numero_documento_pagamento,
                    data_pagamento=data_pagamento,
                    valor=valor,
                    exercicio=int(exercicio),
                )

                messages.success(request, "✅ Pagamento lançado com sucesso!")
                return redirect("pddelancar_pagamento")

        except EscolaPdde.DoesNotExist:
            messages.error(request, "❌ Escola não encontrada.")
            return redirect("pddelancar_pagamento")

        except Exception as e:
            messages.error(request, f"❌ Erro ao lançar pagamento: {str(e)}")
            return redirect("pddelancar_pagamento")

    return render(request, "pdde/lancar_pagamento.html", {"escolas": escolas, "pagamentos": pagamentos})




###***********************************************************************************************************************

from django.http import JsonResponse
from decimal import Decimal
from django.db.models import Sum
from .models import ReceitaDespesa, EscolaPdde, Pagamento, Programa

def get_dados_financeiros(request, escola_id):
    """
    Retorna os dados financeiros da escola selecionada,
    incluindo total de despesas realizadas, saldo disponível,
    lista de programas e pagamentos detalhados.
    """
    try:
        # 🔹 Verifica se a escola existe
        escola = EscolaPdde.objects.get(id=escola_id)

        # 🔹 Busca todas as receitas e despesas associadas à escola
        receita_despesa = ReceitaDespesa.objects.filter(escola=escola).first()

        # 🔹 Se não houver registros, retorna valores padrão (evita erro 404)
        if not receita_despesa:
            return JsonResponse({
                "escola": escola.nome,
                "saldo_reprogramado_custeio": "0.00",
                "saldo_reprogramado_capital": "0.00",
                "total_pagamentos_custeio": "0.00",
                "total_pagamentos_capital": "0.00",
                "saldo_disponivel_custeio": "0.00",
                "saldo_disponivel_capital": "0.00",
                "total_despesa_realizada": "0.00",
                "saldo_disponivel": "0.00",
                "programas": [],
                "pagamentos": [],
            })

        # 🔹 Calcula total de despesas realizadas (somando todas as despesas registradas)
        total_despesa_realizada = ReceitaDespesa.objects.filter(escola=escola).aggregate(
            Sum("total_despesa")
        )["total_despesa__sum"] or Decimal("0.00")

        # 🔹 Calcula saldo disponível (soma do saldo reprogramado)
        saldo_disponivel = ReceitaDespesa.objects.filter(escola=escola).aggregate(
            Sum("saldo_reprogramar")
        )["saldo_reprogramar__sum"] or Decimal("0.00")

        # 🔹 Pagamentos Custeio e Capital
        pagamentos_custeio = Pagamento.objects.filter(escola=escola, natureza="Custeio").aggregate(
            total=Sum("valor")
        )["total"] or 0

        pagamentos_capital = Pagamento.objects.filter(escola=escola, natureza="Capital").aggregate(
            total=Sum("valor")
        )["total"] or 0

        # 🔹 Saldo disponível após pagamentos
        saldo_disponivel_custeio = receita_despesa.saldo_reprogramado_custeio - pagamentos_custeio
        saldo_disponivel_capital = receita_despesa.saldo_reprogramado_capital - pagamentos_capital

        # 🔹 Busca os programas vinculados à escola na tabela ReceitaDespesa
        programas_ids = ReceitaDespesa.objects.filter(escola=escola).values_list("programa", flat=True).distinct()
        programas_nomes = [Programa.objects.get(id=prog_id).nome for prog_id in programas_ids if Programa.objects.filter(id=prog_id).exists()]

        # 🔹 Lista de pagamentos detalhados
        pagamentos_detalhados = Pagamento.objects.filter(escola=escola).values(
            "id", "nome_favorecido", "cnpj_cpf", "tipo_bem_servico",
            "origem", "natureza", "documento", "data_pagamento", "valor"
        )

        # 🔹 Criar resposta JSON com todos os dados necessários
        data = {
            "escola": escola.nome,
            "saldo_reprogramado_custeio": f"{receita_despesa.saldo_reprogramado_custeio:.2f}",
            "saldo_reprogramado_capital": f"{receita_despesa.saldo_reprogramado_capital:.2f}",
            "total_pagamentos_custeio": f"{pagamentos_custeio:.2f}",
            "total_pagamentos_capital": f"{pagamentos_capital:.2f}",
            "saldo_disponivel_custeio": f"{saldo_disponivel_custeio:.2f}",
            "saldo_disponivel_capital": f"{saldo_disponivel_capital:.2f}",
            "total_despesa_realizada": f"{total_despesa_realizada:.2f}",
            "saldo_disponivel": f"{saldo_disponivel:.2f}",
            "programas": programas_nomes,  # ✅ Lista de programas vinculados
            "pagamentos": list(pagamentos_detalhados),  # ✅ Lista de pagamentos detalhados
        }

        return JsonResponse(data)

    except EscolaPdde.DoesNotExist:
        return JsonResponse({"erro": "Escola não encontrada."}, status=404)
    except Exception as e:
        return JsonResponse({"erro": f"Erro ao buscar dados financeiros: {str(e)}"}, status=500)


###***********************************************************************************************************************

from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
import json
from .models import ReceitaDespesa

@csrf_exempt
def atualizar_receita_despesa(request):
    if request.method == "POST":
        try:
            data = json.loads(request.body)
            escola_id = data.get("escola_id")

            # Buscar ou criar o registro da escola
            receita_despesa, created = ReceitaDespesa.objects.get_or_create(escola_id=escola_id)

            # Atualizar valores recebidos
            receita_despesa.saldo_anterior_custeio = data.get("saldo_anterior_custeio", receita_despesa.saldo_anterior_custeio)
            receita_despesa.saldo_anterior_capital = data.get("saldo_anterior_capital", receita_despesa.saldo_anterior_capital)
            receita_despesa.valor_creditado_custeio = data.get("valor_creditado_custeio", receita_despesa.valor_creditado_custeio)
            receita_despesa.valor_creditado_capital = data.get("valor_creditado_capital", receita_despesa.valor_creditado_capital)
            receita_despesa.recursos_proprios = data.get("recursos_proprios", receita_despesa.recursos_proprios)
            receita_despesa.rendimento_aplicacao = data.get("rendimento_aplicacao", receita_despesa.rendimento_aplicacao)
            receita_despesa.devolucao_fnde = data.get("devolucao_fnde", receita_despesa.devolucao_fnde)
            receita_despesa.saldo_devolvido = data.get("saldo_devolvido", receita_despesa.saldo_devolvido)

            # Calcular total receita
            receita_despesa.total_receita = (
                receita_despesa.saldo_anterior_custeio +
                receita_despesa.saldo_anterior_capital +
                receita_despesa.valor_creditado_custeio +
                receita_despesa.valor_creditado_capital +
                receita_despesa.recursos_proprios +
                receita_despesa.rendimento_aplicacao
            )

            # Calcular total despesa
            receita_despesa.total_despesa = receita_despesa.devolucao_fnde + receita_despesa.saldo_devolvido

            # Calcular receita líquida
            receita_despesa.receita_liquida = receita_despesa.total_receita - receita_despesa.total_despesa

            receita_despesa.save()  # Salvar no banco

            return JsonResponse({
                "mensagem": "Valores atualizados com sucesso",
                "total_receita": receita_despesa.total_receita,
                "total_despesa": receita_despesa.total_despesa,
                "receita_liquida": receita_despesa.receita_liquida,
            })

        except Exception as e:
            return JsonResponse({"erro": str(e)}, status=500)
###***********************************************************************************************************************

from django.http import JsonResponse
from .models import EscolaPrograma, Programa

def get_programas(request, escola_id):
    try:
        programas = EscolaPrograma.objects.filter(escola_id=escola_id).select_related("programa")

        if not programas.exists():
            return JsonResponse({"programas": []})  # Retorna uma lista vazia se a escola não tiver programas

        data = {
            "programas": [{"id": p.programa.id, "nome": p.programa.nome} for p in programas]
        }
        return JsonResponse(data)

    except Exception as e:
        return JsonResponse({"erro": f"Erro ao buscar programas: {str(e)}"}, status=500)
###***********************************************************************************************************************

from django.http import JsonResponse
from .models import Pagamento

def get_pagamentos_pdde(request, escola_id):
    """
    Retorna os pagamentos do programa PDDE para uma escola específica.
    """
    try:
        # Filtra os pagamentos apenas do programa PDDE para a escola selecionada
        pagamentos = Pagamento.objects.filter(escola_id=escola_id, programa="PDDE")

        # Retorna um JSON vazio se não houver pagamentos
        if not pagamentos.exists():
            return JsonResponse({"pagamentos": []})

        data = {
            "pagamentos": [
                {
                    "id": pagamento.id,
                    "favorecido": pagamento.nome_favorecido,
                    "cnpj_cpf": pagamento.cnpj_cpf,
                    "descricao": pagamento.tipo_bem_servico,
                    "origem": pagamento.origem,
                    "nat_despesa": pagamento.natureza,
                    "documento": pagamento.documento,
                    "data_pagamento": pagamento.data_pagamento.strftime("%d/%m/%Y"),  # Converte a data para string
                    "valor": f"{pagamento.valor:.2f}"  # Formata o valor com 2 casas decimais
                }
                for pagamento in pagamentos
            ]
        }

        return JsonResponse(data)

    except Exception as e:
        return JsonResponse({"erro": f"Erro ao buscar pagamentos PDDE: {str(e)}"}, status=500)
###***********************************************************************************************************************

from django.http import JsonResponse
from django.db.models import Sum
from .models import Pagamento, EscolaPdde, ReceitaDespesa

def get_dados_pagamentos_pdde(request, escola_id=None):
    """
    Retorna os dados completos de pagamentos dos programas vinculados a uma escola específica.
    Se nenhuma escola for selecionada, retorna todas as escolas cadastradas.
    Inclui: escolas, programas, exercício, total de despesas realizadas e saldo disponível.
    """
    try:
        # Buscar todas as escolas cadastradas para exibição no frontend
        todas_escolas = EscolaPdde.objects.values("id", "nome")

        # Se nenhuma escola for selecionada, apenas retorna a lista de escolas
        if not escola_id:
            return JsonResponse({"escolas": list(todas_escolas)})

        # Buscar a escola específica
        escola = EscolaPdde.objects.get(id=escola_id)

        # Buscar os programas vinculados à escola selecionada
        programas = ReceitaDespesa.objects.filter(escola=escola).values_list("programa", flat=True).distinct()

        # Buscar os pagamentos da escola
        pagamentos = Pagamento.objects.filter(escola_id=escola_id)

        # Calcular total de despesas realizadas
        total_despesa_realizada = pagamentos.aggregate(total=Sum("valor"))["total"] or 0.00

        # Calcular saldo disponível somando todos os programas vinculados à escola
        saldo_disponivel = ReceitaDespesa.objects.filter(escola=escola).aggregate(saldo=Sum("saldo_reprogramar"))["saldo"] or 0.00

        # Construção dos dados para resposta JSON
        data = {
            "escolas": list(todas_escolas),  # Todas as escolas para preenchimento do <select>
            "escola": escola.nome,  # Nome da escola selecionada
            "programas": list(programas),  # Lista de programas vinculados à escola
            "exercicio": escola.ano,  # Ano do exercício
            "total_despesa_realizada": f"{total_despesa_realizada:.2f}",
            "saldo_disponivel": f"{saldo_disponivel:.2f}",
            "pagamentos": [
                {
                    "id": pagamento.id,
                    "favorecido": pagamento.nome_favorecido,
                    "cnpj_cpf": pagamento.cnpj_cpf,
                    "descricao": pagamento.tipo_bem_servico,
                    "origem": pagamento.origem,
                    "nat_despesa": pagamento.natureza,
                    "documento": pagamento.documento,
                    "data_pagamento": pagamento.data_pagamento.strftime("%d/%m/%Y"),
                    "valor": f"{pagamento.valor:.2f}"
                }
                for pagamento in pagamentos
            ]
        }

        return JsonResponse(data)

    except EscolaPdde.DoesNotExist:
        return JsonResponse({"erro": "Escola não encontrada."}, status=404)
    except Exception as e:
        return JsonResponse({"erro": f"Erro ao buscar dados de pagamentos: {str(e)}"}, status=500)
###***********************************************************************************************************************

from django.http import JsonResponse
from .models import EscolaPdde

def get_escolas(request):
    """Retorna todas as escolas cadastradas no sistema em formato JSON."""
    escolas = EscolaPdde.objects.all().order_by("nome")
    data = {"escolas": [{"id": escola.id, "nome": escola.nome} for escola in escolas]}
    return JsonResponse(data)


def get_contas_por_escola(request, escola_id):
    """Retorna todas as contas bancárias vinculadas a uma escola específica."""
    contas = ContaBancaria.objects.filter(escola_id=escola_id).values("id", "banco", "agencia", "conta")
    return JsonResponse({"contas": list(contas)})
###***********************************************************************************************************************

import locale
from datetime import datetime
from django.shortcuts import get_object_or_404
from django.http import HttpResponse
from django.template.loader import render_to_string
from weasyprint import HTML
from django.db.models import Sum, Min, Max
from decimal import Decimal
from .models import Pagamento, Receita, EscolaPdde
import babel.dates
from django.utils.timezone import localtime
from django.utils.timezone import now
from babel.dates import format_date

def pdde_gerar_pdf(request):
    escola_id = request.GET.get('escola_id', "").strip()

    if escola_id.isdigit():
        escola_id = int(escola_id)
        escola = get_object_or_404(EscolaPdde, id=escola_id)
        receitas = Receita.objects.filter(escola=escola)
        receita = receitas.first()  # Obtém um único objeto ou None
    else:
        escolas = EscolaPdde.objects.all()
        if not escolas.exists():
            return HttpResponse("Nenhuma escola cadastrada no sistema.", status=404)

        receitas = Receita.objects.all()
        receita = receitas.first()  # Obtém um único objeto para o relatório geral

    pagamentos = Pagamento.objects.filter(escola=escola) if escola_id else Pagamento.objects.all()

    # 🔹 Identificação da entidade
    if receita:
        entidade = receita.escola.nome_conselho or "Não disponível"
        cnpj = receita.escola.cnpj or "Não disponível"
        municipio = receita.escola.cidade or "Não disponível"
        estado = receita.escola.uf or "Não disponível"
        periodo_inicio = receitas.aggregate(Min('data_inicio'))['data_inicio__min']
        periodo_fim = receitas.aggregate(Max('data_fim'))['data_fim__max']
        periodo = f"{periodo_inicio.strftime('%d/%m/%Y')} até {periodo_fim.strftime('%d/%m/%Y')}" if periodo_inicio and periodo_fim else "Não disponível"
    else:
        entidade = "Múltiplas entidades"
        cnpj = "N/A"
        municipio = "Canaã dos Carajás"
        estado = "PA"
        periodo = "Ano Atual"

    # 🔹 Cálculo dos totais de receitas e despesas
    total_receitas = receitas.aggregate(
        total_custeio=Sum('valor_total_receita_custeio', default=0),
        total_capital=Sum('valor_total_receita_capital', default=0)
    )
    total_receitas = (total_receitas['total_custeio'] or Decimal(0)) + (total_receitas['total_capital'] or Decimal(0))

    total_despesas = pagamentos.aggregate(Sum('valor'))['valor__sum'] or Decimal(0)
    saldo_geral = total_receitas - total_despesas

    # 🔹 Formatação correta de valores monetários
    def format_currency(value):
        return f"R$ {value:,.2f}".replace(',', 'X').replace('.', ',').replace('X', '.')

    data_atual = format_date(now(), format="d 'de' MMMM 'de' yyyy", locale="pt_BR")

    # 🔹 Contexto para o template
    context = {
        'entidade': entidade,
        'cnpj': cnpj,
        'municipio': municipio,
        'estado': estado,
        'periodo': periodo,
        'escolas': escolas if not escola_id else [escola],
        'pagamentos': pagamentos,
        'receita': receita,  # Passando um único objeto para o template
        'total_receitas': format_currency(total_receitas),
        'total_despesas': format_currency(total_despesas),
        'saldo_geral': format_currency(saldo_geral),
        'data_atual': data_atual,
    }

    html_string = render_to_string("pdde/relatorio_pdf.html", context)
    pdf_file = HTML(string=html_string).write_pdf()

    response = HttpResponse(pdf_file, content_type="application/pdf")
    response["Content-Disposition"] = f'inline; filename="relatorio_{escola.nome if escola_id else "geral"}.pdf"'

    return response
###***********************************************************************************************************************

from django.shortcuts import render, get_object_or_404, redirect
from django.contrib.auth.decorators import login_required
from django.contrib import messages
from .models import LancamentoBancario

# 🔹 Conciliação Bancária - Exibição de Páginas
def conciliacao_extrato(request):
    return render(request, 'conciliacao/conciliacao_extrato.html')
###***********************************************************************************************************************

def conciliacao_lancamento(request):
    """ Lista os lançamentos bancários """
    lancamentos = LancamentoBancario.objects.all()
    return render(request, "conciliacao/conciliacao_lancamento.html", {"lancamentos": lancamentos})
###***********************************************************************************************************************

def conciliacao_saldo(request):
    return render(request, 'conciliacao/conciliacao_saldo.html')
###***********************************************************************************************************************

from django.shortcuts import render
from .models import EscolaPdde, ContaBancaria

def conciliacao_conferencia(request):
    """Renderiza a página de conferência com as escolas e contas bancárias."""
    
    escolas = EscolaPdde.objects.all()
    contas = ContaBancaria.objects.all()

    return render(request, "conciliacao/conciliacao_conferencia.html", {
        "escolas": escolas,
        "contas": contas,
    })

###***********************************************************************************************************************

def conciliacao_relatorios(request):
    return render(request, 'conciliacao/conciliacao_relatorios.html')
###***********************************************************************************************************************

from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from django.contrib.auth.decorators import login_required
from .models import LancamentoBancario, EscolaPdde, ContaBancaria

from decimal import Decimal

@login_required
def criar_lancamento(request):
    """Cria um novo lançamento bancário associado a uma escola e conta."""
    
    if request.method == "POST":
        data = request.POST.get("data")
        descricao = request.POST.get("descricao")
        tipo = request.POST.get("tipo")
        valor = request.POST.get("valor")
        categoria = request.POST.get("categoria")
        escola_id = request.POST.get("escola")
        conta_id = request.POST.get("conta_bancaria")

        # Convertendo 'valor' para Decimal
        try:
            valor = Decimal(valor.replace(",", "."))  # Substitui vírgula por ponto antes de converter
        except (ValueError, TypeError):
            messages.error(request, "O valor informado é inválido.")
            return redirect("criar_lancamento")

        if not escola_id or not conta_id:
            messages.error(request, "Por favor, selecione uma escola e uma conta bancária.")
            return redirect("criar_lancamento")

        escola = get_object_or_404(EscolaPdde, id=escola_id)
        conta = get_object_or_404(ContaBancaria, id=conta_id)

        LancamentoBancario.objects.create(
            usuario=request.user,
            data=data,
            descricao=descricao,
            tipo=tipo,
            valor=valor,
            categoria=categoria,
            escola=escola,
            conta_bancaria=conta
        )

        messages.success(request, "Lançamento cadastrado com sucesso!")
        return redirect("listar_lancamentos")

    escolas = EscolaPdde.objects.all()
    contas = ContaBancaria.objects.all()

    return render(request, "conciliacao/criar_lancamento.html", {
        "escolas": escolas,
        "contas": contas
    })
###***********************************************************************************************************************

@login_required
def editar_lancamento(request, pk):
    """ Edita um lançamento bancário existente """
    lancamento = get_object_or_404(LancamentoBancario, id=pk)

    if request.method == "POST":
        lancamento.data = request.POST.get("data")
        lancamento.descricao = request.POST.get("descricao")
        lancamento.tipo = request.POST.get("tipo")
        lancamento.categoria = request.POST.get("categoria")
        lancamento.valor = request.POST.get("valor")
        lancamento.save()

        messages.success(request, "Lançamento atualizado com sucesso!")
        return redirect("conciliacao_lancamento")

    return render(request, "conciliacao/editar_lancamento.html", {"lancamento": lancamento})
###***********************************************************************************************************************

@login_required
def excluir_lancamento(request, pk):
    """ Exclui um lançamento bancário """
    lancamento = get_object_or_404(LancamentoBancario, id=pk)
    lancamento.delete()

    messages.success(request, "Lançamento excluído com sucesso!")
    return redirect("conciliacao_lancamento")
###***********************************************************************************************************************

from django.shortcuts import render
from .models import LancamentoBancario

def listar_lancamentos(request):
    """Lista todos os lançamentos bancários."""
    lancamentos = LancamentoBancario.objects.all()
    return render(request, "conciliacao/listar_lancamentos.html", {"lancamentos": lancamentos})
###***********************************************************************************************************************

from django.shortcuts import render
from .models import EscolaPdde  # ✅ CORRETO (ajustado para o nome exato do modelo)

def listar_pdde(request):
    """View para listar os registros da tabela EscolaPDDE"""
    pdde_data = EscolaPdde.objects.all()
    return render(request, "seu_template.html", {"pdde_data": pdde_data})
###***********************************************************************************************************************

from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from .models import ContaBancaria, EscolaPdde, Programa

def criar_conta_bancaria(request):
    """Cria uma nova conta bancária vinculada a uma escola e programa."""
    if request.method == "POST":
        escola_id = request.POST.get("escola")
        programa_id = request.POST.get("programa")
        banco = request.POST.get("banco")
        agencia = request.POST.get("agencia")
        conta = request.POST.get("conta")
        tipo_conta = request.POST.get("tipo_conta")

        if not escola_id:
            messages.error(request, "Selecione uma escola para vincular a conta bancária.")
            return redirect("nova_conta")

        escola = get_object_or_404(EscolaPdde, id=escola_id)
        programa = get_object_or_404(Programa, id=programa_id)

        # Criar a conta bancária vinculada à escola e ao programa
        ContaBancaria.objects.create(
            escola=escola,
            banco=banco,
            agencia=agencia,
            conta=conta,
            tipo_conta=tipo_conta,
            conselho=escola.nome_conselho,
            programa=programa  # Supondo que você adicionou o campo ForeignKey no modelo
        )

        messages.success(request, "Conta bancária criada com sucesso!")
        return redirect("listar_contas")

    # GET
    escolas = EscolaPdde.objects.prefetch_related("semedappescolapddeprogramas_set__programa")
    return render(request, "conciliacao/criar_conta.html", {"escolas": escolas})



###***********************************************************************************************************************

from django.shortcuts import render
from .models import ContaBancaria

def listar_contas(request):
    """Lista todas as contas bancárias cadastradas."""
    contas = ContaBancaria.objects.select_related("escola").all()
    return render(request, "conciliacao/listar_contas.html", {"contas": contas})
###***********************************************************************************************************************

from django.shortcuts import render, get_object_or_404, redirect
from django.contrib import messages
from .models import ContaBancaria
from .forms import ContaBancariaForm

def editar_conta(request, conta_id):
    conta = get_object_or_404(ContaBancaria, id=conta_id)

    if request.method == "POST":
        form = ContaBancariaForm(request.POST, instance=conta)
        if form.is_valid():
            form.save()
            messages.success(request, "Conta bancária atualizada com sucesso!")
            return redirect("listar_contas")
        else:
            messages.error(request, "Erro ao atualizar a conta bancária.")
    else:
        form = ContaBancariaForm(instance=conta)

    return render(request, "conciliacao/editar_conta.html", {"form": form, "conta": conta})


###***********************************************************************************************************************

from django.shortcuts import get_object_or_404, redirect
from django.contrib import messages
from .models import ContaBancaria

def excluir_conta(request, conta_id):
    conta = get_object_or_404(ContaBancaria, id=conta_id)

    if request.method == "POST":
        conta.delete()
        messages.success(request, "Conta bancária excluída com sucesso.")
        return redirect("listar_contas")

    return render(request, "conciliacao/confirmar_exclusao.html", {"conta": conta})
###***********************************************************************************************************************

from django.shortcuts import render
from django.contrib.auth.decorators import login_required
from django.db.models import Sum
from decimal import Decimal
from .models import EscolaPdde, ContaBancaria, LancamentoBancario, Receita

@login_required
def consultar_saldo(request):
    escolas = EscolaPdde.objects.all()
    contas = ContaBancaria.objects.all()
    saldos = []

    data_inicial = request.GET.get("data_inicial")
    data_final = request.GET.get("data_final")
    escola_id = request.GET.get("escola")
    conta_id = request.GET.get("conta_bancaria")

    if data_inicial and data_final and escola_id and conta_id:
        escola = EscolaPdde.objects.get(id=escola_id)
        conta = ContaBancaria.objects.get(id=conta_id)

        total_receitas = Receita.objects.filter(
            escola=escola,
            conta_bancaria=conta
        ).aggregate(total=Sum("valor_total_receita"))["total"] or Decimal("0.00")

        lancamentos = LancamentoBancario.objects.filter(
            escola=escola,
            conta_bancaria=conta,
            data__range=[data_inicial, data_final]
        ).order_by("data")

        saldo_atual = total_receitas

        for lancamento in lancamentos:
            valor_lancamento = lancamento.valor if lancamento.tipo == "credito" else -lancamento.valor
            saldo_atual += valor_lancamento
            if lancamento.tipo == "credito":
                saldo_atual += lancamento.valor
                tipo = "Entrada"
                cor_vermelha = False
            else:
                saldo_atual -= lancamento.valor
                tipo = "Saída"
                cor_vermelha = True

            saldos.append({
                "data": lancamento.data.strftime("%d/%m/%Y"),
                "descricao": lancamento.descricao,
                "tipo": tipo,
                "valor": abs(saldo_atual),
                "valor_lancamento": float(valor_lancamento),
                "is_saida": cor_vermelha
            })

    return render(request, "conciliacao/consulta_saldo.html", {
        "escolas": escolas,
        "contas": contas,
        "saldos": saldos
    })

###***********************************************************************************************************************

from django.http import JsonResponse
from .models import EscolaPdde

def api_escolas(request):
    """Retorna todas as escolas cadastradas em formato JSON"""
    escolas = EscolaPdde.objects.values("id", "nome")
    return JsonResponse({"escolas": list(escolas)})
###***********************************************************************************************************************

from .models import ContaBancaria

def api_contas(request, escola_id):
    """Retorna todas as contas bancárias vinculadas a uma escola específica"""
    contas = ContaBancaria.objects.filter(escola_id=escola_id).values("id", "banco", "agencia", "conta")
    return JsonResponse({"contas": list(contas)})
###***********************************************************************************************************************

from django.http import JsonResponse
from django.db.models import Sum
from decimal import Decimal
from .models import LancamentoBancario

def api_saldos(request):
    data_inicial = request.GET.get("data_inicial")
    data_final = request.GET.get("data_final")
    escola_id = request.GET.get("escola")
    conta_id = request.GET.get("conta")

    if not data_inicial or not data_final or not escola_id or not conta_id:
        return JsonResponse({"erro": "Todos os campos são obrigatórios!"}, status=400)

    # Calcula o saldo inicial antes da data_inicial
    saldo_credito = LancamentoBancario.objects.filter(
        escola_id=escola_id,
        conta_bancaria_id=conta_id,
        tipo="Entrada",  # já em português
        data__lt=data_inicial
    ).aggregate(total=Sum("valor"))["total"] or Decimal("0.00")

    saldo_debito = LancamentoBancario.objects.filter(
        escola_id=escola_id,
        conta_bancaria_id=conta_id,
        tipo="Saída",
        data__lt=data_inicial
    ).aggregate(total=Sum("valor"))["total"] or Decimal("0.00")

    saldo_atual = saldo_credito - saldo_debito

    # Busca os lançamentos dentro do intervalo
    lancamentos = LancamentoBancario.objects.filter(
        escola_id=escola_id,
        conta_bancaria_id=conta_id,
        data__range=[data_inicial, data_final]
    ).order_by("data")

    resultado = []

    for lancamento in lancamentos:
        valor = lancamento.valor or Decimal("0.00")

        if lancamento.tipo == "Entrada":
            saldo_atual += valor
            valor_lancamento = valor
        else:  # Saída
            saldo_atual -= valor
            valor_lancamento = -valor

        resultado.append({
            "data": lancamento.data.strftime("%d/%m/%Y"),
            "descricao": lancamento.descricao,
            "tipo": lancamento.tipo,  # já vem como 'Entrada' ou 'Saída'
            "tipo_raw": lancamento.tipo,  # idem
            "valor_lancamento": float(valor_lancamento),
            "saldo_acumulado": float(saldo_atual)
        })

    return JsonResponse({"saldos": resultado})




###***********************************************************************************************************************

@login_required
def conferencia_movimentacao(request):
    escolas = EscolaPdde.objects.all()
    contas = ContaBancaria.objects.all()
    conferencias = []

    data_inicial = request.GET.get("data_inicial")
    data_final = request.GET.get("data_final")
    escola_id = request.GET.get("escola")
    conta_id = request.GET.get("conta_bancaria")

    if data_inicial and data_final and escola_id and conta_id:
        escola = EscolaPdde.objects.get(id=escola_id)
        conta = ContaBancaria.objects.get(id=conta_id)

        lancamentos = LancamentoBancario.objects.filter(
            escola=escola,
            conta_bancaria=conta,
            data__range=[data_inicial, data_final]
        ).order_by("data")

        for lancamento in lancamentos:
            tipo = lancamento.tipo  # "credito" ou "debito"
            valor = lancamento.valor if tipo == "credito" else -lancamento.valor
            conferencias.append({
                "data": lancamento.data.strftime("%d/%m/%Y"),
                "descricao": lancamento.descricao,
                "escola": escola.nome,
                "conta_bancaria": f"{conta.banco} - Ag: {conta.agencia} - Conta: {conta.conta}",
                "valor_lancamento": float(valor),
                "tipo_raw": tipo,
                "status": "Entrada" if tipo == "credito" else "Saída"
            })

    return render(request, "conciliacao/conferencia_movimentacao.html", {
        "escolas": escolas,
        "contas": contas,
        "conferencias": conferencias
    })


###***********************************************************************************************************************

from django.http import JsonResponse
from django.db.models import Sum
from decimal import Decimal
from .models import EscolaPdde, ContaBancaria, LancamentoBancario

def api_conferencia(request):
    data_inicial = request.GET.get("data_inicial")
    data_final = request.GET.get("data_final")
    escola_id = request.GET.get("escola")
    conta_id = request.GET.get("conta")
    tipo_conferencia = request.GET.get("tipo")

    if not all([data_inicial, data_final, escola_id, conta_id, tipo_conferencia]):
        return JsonResponse({"error": "Parâmetros inválidos."}, status=400)

    try:
        escola = EscolaPdde.objects.get(id=escola_id)
        conta = ContaBancaria.objects.get(id=conta_id)
    except (EscolaPdde.DoesNotExist, ContaBancaria.DoesNotExist):
        return JsonResponse({"error": "Escola ou Conta não encontrada."}, status=404)

    saldo_credito = LancamentoBancario.objects.filter(
        escola=escola,
        conta_bancaria=conta,
        tipo="credito",
        data__lt=data_inicial
    ).aggregate(total=Sum("valor"))["total"] or Decimal("0.00")

    saldo_debito = LancamentoBancario.objects.filter(
        escola=escola,
        conta_bancaria=conta,
        tipo="debito",
        data__lt=data_inicial
    ).aggregate(total=Sum("valor"))["total"] or Decimal("0.00")

    saldo_atual = saldo_credito - saldo_debito

    lancamentos = LancamentoBancario.objects.filter(
        escola=escola,
        conta_bancaria=conta,
        data__range=[data_inicial, data_final]
    ).order_by("data")

    conferencias = []
    for lancamento in lancamentos:
        tipo_raw = lancamento.tipo  # ← Deve retornar 'credito' ou 'debito'
        status = "Entrada" if tipo_raw == "credito" else "Saída"

        conferencias.append({
            "data": lancamento.data.strftime("%d/%m/%Y"),
            "descricao": lancamento.descricao,
            "escola": escola.nome,
            "conta_bancaria": f"{conta.banco} - Ag: {conta.agencia} - Conta: {conta.conta}",
            "valor_lancamento": float(lancamento.valor),
            "tipo_raw": tipo_raw,
            "status": status  # Isso será usado no frontend
        })



    return JsonResponse({"conferencias": conferencias})


###***********************************************************************************************************************

from django.http import JsonResponse
from django.db.models import Sum
from .models import EscolaPdde, ContaBancaria, LancamentoBancario

def api_extrato(request):
    """Retorna os lançamentos bancários filtrados por escola e conta bancária."""
    escola_id = request.GET.get("escola")
    conta_id = request.GET.get("conta")

    if not escola_id or not conta_id:
        return JsonResponse({"error": "Parâmetros inválidos."}, status=400)

    extratos = LancamentoBancario.objects.filter(escola_id=escola_id, conta_bancaria_id=conta_id).order_by("-data")

    extrato_list = [
        {
            "data": extrato.data.strftime("%d/%m/%Y"),
            "descricao": extrato.descricao,
            "escola": extrato.escola.nome if extrato.escola else "N/A",
            "conta_bancaria": f"{extrato.conta_bancaria.banco} - Ag: {extrato.conta_bancaria.agencia} - Conta: {extrato.conta_bancaria.conta}" if extrato.conta_bancaria else "N/A",
            "tipo": extrato.tipo,
            "valor": float(extrato.valor),
        }
        for extrato in extratos
    ]

    return JsonResponse({"extratos": extrato_list})
###***********************************************************************************************************************

from django.shortcuts import render, redirect
from .forms import ProgramaForm
from django.http import JsonResponse
from .models import Programa  # Importando o modelo Programa


def listar_programas(request):
    if request.headers.get('X-Requested-With') == 'XMLHttpRequest':  
        # Certifique-se de que está buscando os programas corretamente
        programas = Programa.objects.all().values("id", "nome").order_by("nome")
        return JsonResponse({"programas": list(programas)})

    # Caso não seja uma requisição AJAX, renderiza o formulário normalmente
    programas = Programa.objects.all().order_by("nome")
    return render(request, "escola_pdde_form.html", {"programas": programas})
###***********************************************************************************************************************

def criar_programa(request):
    if request.method == "POST":
        form = ProgramaForm(request.POST)
        if form.is_valid():
            form.save()
            return redirect("listar_programas")
    else:
        form = ProgramaForm()
    
    return render(request, "semedapp/form_programa.html", {"form": form})
###***********************************************************************************************************************

from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
from .models import Programa  # Removemos EscolaPdde porque não precisamos mais dela
import json

@csrf_exempt
def cadastrar_programa(request):
    if request.method == "POST":
        try:
            data = json.loads(request.body)  # Lendo JSON do request body

            nome = data.get("nome")
            descricao = data.get("descricao", "")
            resolucao = data.get("resolucao", "")
            data_inicio = data.get("data_inicio")
            data_fim = data.get("data_fim")

            # Validação de campos obrigatórios
            if not nome or not data_inicio or not data_fim:
                return JsonResponse(
                    {"success": False, "error": "Todos os campos obrigatórios devem ser preenchidos!"},
                    status=400
                )

            # Criando o programa (sem escola)
            programa = Programa.objects.create(
                nome=nome,
                descricao=descricao,
                resolucao=resolucao,
                data_inicio=data_inicio,
                data_fim=data_fim
            )

            return JsonResponse(
                {"success": True, "id": programa.id, "nome": programa.nome},
                status=201
            )

        except json.JSONDecodeError:
            return JsonResponse(
                {"success": False, "error": "Formato JSON inválido."},
                status=400
            )
        except Exception as e:
            return JsonResponse(
                {"success": False, "error": f"Erro inesperado: {str(e)}"},
                status=500
            )

    return JsonResponse({"success": False, "error": "Método inválido!"}, status=405)
###***********************************************************************************************************************

from django.shortcuts import render
from .models import EscolaPdde  # Importa o modelo correto

def lancar_receita(request):
    escolas = EscolaPdde.objects.all().order_by("nome")
    print(escolas)  # 🔹 Isso ajudará a verificar se as escolas estão sendo carregadas
    return render(request, "pdde/lancar_receita.html", {"escolas": escolas})
###***********************************************************************************************************************

from django.http import JsonResponse
from .models import EscolaPdde

def listar_escolas(request):
    if request.headers.get("X-Requested-With") == "XMLHttpRequest":
        escolas = EscolaPdde.objects.values("id", "nome").order_by("nome")
        return JsonResponse({"escolas": list(escolas)}, safe=False)

    return JsonResponse({"error": "Requisição inválida"}, status=400)
###***********************************************************************************************************************

from django.shortcuts import render, get_object_or_404, redirect
from django.contrib import messages
from .models import EscolaPdde, Programa

def editar_escola_pdde(request, escola_id):
    """
    View para edição da escola PDDE, carregando os dados necessários e permitindo a atualização.
    """
    escola = get_object_or_404(EscolaPdde, id=escola_id)
    programas = Programa.objects.all()  # 🔹 Lista de todos os programas disponíveis

    if request.method == "POST":
        try:
            # Atualiza os dados manualmente
            escola.nome = request.POST.get("nome", escola.nome)
            escola.tipo = request.POST.get("tipo", escola.tipo)
            escola.nome_conselho = request.POST.get("nome_conselho", escola.nome_conselho)
            escola.ano = request.POST.get("ano", escola.ano)
            escola.status = request.POST.get("status", escola.status)
            escola.endereco = request.POST.get("endereco", escola.endereco)
            escola.bairro = request.POST.get("bairro", escola.bairro)
            escola.cep = request.POST.get("cep", escola.cep)
            escola.cidade = request.POST.get("cidade", escola.cidade)
            escola.uf = request.POST.get("uf", escola.uf)
            escola.dependencia_administrativa = request.POST.get("dependencia_administrativa", escola.dependencia_administrativa)
            escola.codigo_inep = request.POST.get("codigo_inep", escola.codigo_inep)
            escola.cnpj = request.POST.get("cnpj", escola.cnpj)
            escola.zona = request.POST.get("zona", escola.zona)
            escola.ensino = request.POST.get("ensino", escola.ensino)
            escola.quantidade_salas = request.POST.get("quantidade_salas", escola.quantidade_salas)
            escola.quantidade_turmas = request.POST.get("quantidade_turmas", escola.quantidade_turmas)
            escola.quantidade_professores = request.POST.get("quantidade_professores", escola.quantidade_professores)
            escola.quantidade_alunos = request.POST.get("quantidade_alunos", escola.quantidade_alunos)

            # Atualiza o programa vinculado
            programa_id = request.POST.get("programa")
            if programa_id:
                escola.programa = Programa.objects.get(id=programa_id)

            escola.save()

            messages.success(request, "Escola atualizada com sucesso!")
            return redirect("lista_escolas_pdde")

        except Exception as e:
            messages.error(request, f"Erro ao atualizar a escola: {str(e)}")

    return render(request, "semedapp/escola_pdde_edit.html", {
        "escola": escola,
        "programas": programas,
        "programa_vinculado": escola.programa if hasattr(escola, 'programa') else None,  # 🔹 Pega o programa vinculado
    })
###***********************************************************************************************************************

from django.shortcuts import render, get_object_or_404
from .models import EscolaPdde, Programa, SemedAppEscolaPddeProgramas

def detalhes_escola_pdde(request, escola_id):
    escola = get_object_or_404(EscolaPdde, id=escola_id)

    # 🔹 Buscar os programas vinculados à escola pela tabela intermediária
    programas_vinculados = Programa.objects.filter(
        id__in=SemedAppEscolaPddeProgramas.objects.filter(escolapdde_id=escola_id).values_list('programa_id', flat=True)
    )

    context = {
        'escola': escola,
        'programas_vinculados': programas_vinculados,
    }

    return render(request, 'semedapp/pddereceitadespesa.html', context)
###***********************************************************************************************************************

from django.http import JsonResponse
from .models import SemedAppEscolaPddeProgramas, Programa

def get_programas_vinculados(request, escola_id):
    """
    Retorna os programas vinculados a uma escola específica.
    """
    programas_ids = SemedAppEscolaPddeProgramas.objects.filter(escolapdde_id=escola_id).values_list('programa_id', flat=True)
    programas = Programa.objects.filter(id__in=programas_ids).values_list('nome', flat=True)

    return JsonResponse({"programas": list(programas)})
###***********************************************************************************************************************

from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from .models import EscolaPdde, Programa, SemedAppEscolaPddeProgramas
from .forms import VincularEscolaProgramaForm

def vincular_escola_programa(request):
    """
    View para vincular uma escola a um ou mais programas.
    """
    if request.method == "POST":
        form = VincularEscolaProgramaForm(request.POST)
        if form.is_valid():
            escola = form.cleaned_data["escola"]
            programas = form.cleaned_data["programas"]

            # Removendo vínculos existentes da escola
            SemedAppEscolaPddeProgramas.objects.filter(escolapdde=escola).delete()

            # Criando novos vínculos
            for programa in programas:
                SemedAppEscolaPddeProgramas.objects.create(escolapdde=escola, programa=programa)

            messages.success(request, f"Os programas foram vinculados à escola {escola.nome} com sucesso!")
            return redirect("lista_escolas_pdde")  # Redireciona para a lista de escolas vinculadas

    else:
        form = VincularEscolaProgramaForm()

    return render(request, "semedapp/vincular_escola_programa.html", {"form": form})
###***********************************************************************************************************************

from django.http import JsonResponse
from .models import Escola, Programa

def get_programas_escola(request, escola_id):
    try:
        escola = Escola.objects.get(id=escola_id)
        programas = Programa.objects.filter(escola=escola)  # Ajuste de acordo com seu modelo

        data = {"programas": list(programas.values("id", "nome"))}
        return JsonResponse(data)
    except Escola.DoesNotExist:
        return JsonResponse({"error": "Escola não encontrada"}, status=404)
###***********************************************************************************************************************

from django.http import JsonResponse
from semedapp.models import Receita

def get_programas_por_escola(request, escola_id):
    """
    Retorna uma lista de programas vinculados a uma escola específica,
    baseada na tabela semedapp_receita.
    """
    try:
        programas = Receita.objects.filter(escola_id=escola_id).values_list('programa', flat=True).distinct()
        
        # Convertendo QuerySet para uma lista
        programas_lista = list(programas)

        return JsonResponse({"programas": programas_lista})

    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)
###***********************************************************************************************************************

from django.http import JsonResponse
from .models import EscolaPdde, SemedAppEscolaPddeProgramas, Programa

def listar_programas_por_escola(request, escola_id):
    """
    Retorna os programas cadastrados para uma determinada escola na tabela `semedapp_escolapdde`.
    """
    try:
        # Verifica se a escola realmente existe
        escola = EscolaPdde.objects.filter(id=escola_id).first()
        if not escola:
            return JsonResponse({"erro": "Escola não encontrada."}, status=404)

        # Obtém os IDs dos programas vinculados à escola
        programas_ids = SemedAppEscolaPddeProgramas.objects.filter(escolapdde=escola).values_list("programa", flat=True).distinct()

        # Obtém os nomes dos programas correspondentes
        programas_nomes = list(Programa.objects.filter(id__in=programas_ids).values_list("nome", flat=True))

        return JsonResponse({"programas": programas_nomes})

    except Exception as e:
        return JsonResponse({"erro": f"Erro ao buscar programas: {str(e)}"}, status=500)
###***********************************************************************************************************************

from django.http import JsonResponse
from django.db.models import Sum
from .models import Receita, EscolaPdde  # Certifique-se de que os modelos corretos estão importados

def get_dados_receita_despesa_pdde(request, escola_id):
    """
    Retorna os dados financeiros da escola na tabela semedapp_receita.
    """
    try:
        # 🔹 Verifica se há registros para a escola na tabela semedapp_receita
        receita = Receita.objects.filter(escola_id=escola_id).first()

        if not receita:
            return JsonResponse({"erro": "Nenhuma receita ou despesa encontrada para esta escola."}, status=404)

        # 🔹 Criamos a resposta JSON com os nomes corretos dos campos
        data = {
            "saldo_anterior_custeio": str(receita.saldo_anterior_custeio),
            "saldo_anterior_capital": str(receita.saldo_anterior_capital),
            "valor_creditado_custeio": str(receita.valor_creditado_custeio),
            "valor_creditado_capital": str(receita.valor_creditado_capital),
            "recursos_proprios": str(receita.recursos_proprios),
            "rendimento_aplicacao_custeio": str(receita.rendimento_aplicacao_custeio),
            "rendimento_aplicacao_capital": str(receita.rendimento_aplicacao_capital),
            "devolucao_fnde_custeio": str(receita.devolucao_fnde_custeio),
            "devolucao_fnde_capital": str(receita.devolucao_fnde_capital),
            "valor_total_receita_custeio": str(receita.valor_total_receita_custeio),
            "valor_total_receita_capital": str(receita.valor_total_receita_capital),
            "valor_despesa_realizada_custeio": str(receita.valor_despesa_realizada_custeio),
            "valor_despesa_realizada_capital": str(receita.valor_despesa_realizada_capital),
            "saldo_reprogramar_custeio": str(receita.saldo_reprogramar_custeio),
            "saldo_reprogramar_capital": str(receita.saldo_reprogramar_capital),
            "saldo_devolvido_custeio": str(receita.saldo_devolvido_custeio),
            "saldo_devolvido_capital": str(receita.saldo_devolvido_capital),
            "escolas_atendidas": receita.escolas_atendidas,
        }
        
        return JsonResponse(data)

    except Exception as e:
        return JsonResponse({"erro": f"Erro ao buscar dados financeiros: {str(e)}"}, status=500)
###***********************************************************************************************************************    

from django.http import JsonResponse
from .models import Receita
from decimal import Decimal

def get_dados_receita(request, escola_id):
    """
    Retorna os dados financeiros da escola selecionada para o PDDE, 
    focando apenas nas receitas.
    """
    try:
        # 🔹 Busca a receita associada à escola
        receita = Receita.objects.filter(escola_id=escola_id).first()

        # 🔹 Se não houver registros, retorna valores padrão (evita erro 404)
        if not receita:
            return JsonResponse({
                "escola_id": escola_id,
                "erro": "Nenhuma receita encontrada para essa escola",
                "saldo_anterior_custeio": "0.00",
                "saldo_anterior_capital": "0.00",
                "valor_creditado_custeio": "0.00",
                "valor_creditado_capital": "0.00",
                "rendimento_aplicacao_custeio": "0.00",
                "rendimento_aplicacao_capital": "0.00",
                "valor_total_receita_custeio": "0.00",
                "valor_total_receita_capital": "0.00",
                "saldo_reprogramar_custeio": "0.00",
                "saldo_reprogramar_capital": "0.00",
                "saldo_devolvido_custeio": "0.00",
                "saldo_devolvido_capital": "0.00",
                "escolas_atendidas": 0
            })

        # 🔹 Retorna os dados da receita encontrada
        data = {
            "escola_id": receita.escola.id,
            "programa": receita.programa,
            "data_inicio": receita.data_inicio.strftime("%Y-%m-%d"),
            "data_fim": receita.data_fim.strftime("%Y-%m-%d"),
            "saldo_anterior_custeio": f"{receita.saldo_anterior_custeio:.2f}",
            "saldo_anterior_capital": f"{receita.saldo_anterior_capital:.2f}",
            "valor_creditado_custeio": f"{receita.valor_creditado_custeio:.2f}",
            "valor_creditado_capital": f"{receita.valor_creditado_capital:.2f}",
            "rendimento_aplicacao_custeio": f"{receita.rendimento_aplicacao_custeio:.2f}",
            "rendimento_aplicacao_capital": f"{receita.rendimento_aplicacao_capital:.2f}",
            "valor_total_receita_custeio": f"{receita.valor_total_receita_custeio:.2f}",
            "valor_total_receita_capital": f"{receita.valor_total_receita_capital:.2f}",
            "saldo_reprogramar_custeio": f"{receita.saldo_reprogramar_custeio:.2f}",
            "saldo_reprogramar_capital": f"{receita.saldo_reprogramar_capital:.2f}",
            "saldo_devolvido_custeio": f"{receita.saldo_devolvido_custeio:.2f}",
            "saldo_devolvido_capital": f"{receita.saldo_devolvido_capital:.2f}",
            "escolas_atendidas": receita.escolas_atendidas
        }

        return JsonResponse(data)

    except Exception as e:
        return JsonResponse({"erro": f"Erro ao buscar dados de receita: {str(e)}"}, status=500)
###***********************************************************************************************************************

from django.http import JsonResponse
from django.db.models import Q
from .models import EscolaPdde, ReceitaDespesa

def get_dados_escola(request, escola_id):
    """
    Retorna os dados financeiros e o programa da escola selecionada.
    """
    try:
        escola = EscolaPdde.objects.filter(id=escola_id).first()
        if not escola:
            return JsonResponse({"erro": "Escola não encontrada."}, status=404)

        # Busca um único programa vinculado à escola
        receita = ReceitaDespesa.objects.filter(escola=escola).first()
        programa = receita.programa.nome if receita and receita.programa else None

        # Retorna os dados financeiros
        return JsonResponse({
            "programa": programa,
            "saldo_anterior_custeio": str(receita.saldo_anterior_custeio) if receita else "0.00",
            "saldo_anterior_capital": str(receita.saldo_anterior_capital) if receita else "0.00",
            "valor_creditado_custeio": str(receita.valor_creditado_custeio) if receita else "0.00",
            "valor_creditado_capital": str(receita.valor_creditado_capital) if receita else "0.00",
            "saldo_reprogramar_custeio": str(receita.saldo_reprogramar_custeio) if receita else "0.00",
            "saldo_reprogramar_capital": str(receita.saldo_reprogramar_capital) if receita else "0.00",
        })

    except Exception as e:
        return JsonResponse({"erro": f"Erro ao buscar dados da escola: {str(e)}"}, status=500)
###***********************************************************************************************************************

from django.http import JsonResponse
from django.db.models import Sum
from .models import Receita

def get_receita(request, escola_id, programa):
    try:
        # 🔹 Verificar se os parâmetros recebidos estão corretos
        if not escola_id or not programa:
            return JsonResponse({"erro": "Escola e programa são obrigatórios"}, status=400)

        # 🔹 Filtrar corretamente pela escola e pelo programa (ignorando maiúsculas/minúsculas)
        receita = Receita.objects.filter(escola_id=escola_id, programa__iexact=programa).aggregate(
            saldo_anterior_custeio=Sum('saldo_anterior_custeio') or 0,
            saldo_anterior_capital=Sum('saldo_anterior_capital') or 0,
            valor_creditado_custeio=Sum('valor_creditado_custeio') or 0,
            valor_creditado_capital=Sum('valor_creditado_capital') or 0,
            saldo_reprogramar_custeio=Sum('saldo_reprogramar_custeio') or 0,
            saldo_reprogramar_capital=Sum('saldo_reprogramar_capital') or 0
        )

        # 🔹 Formatar valores para evitar exibição incorreta
        def formatar_valor(valor):
            return "{:.2f}".format(float(valor)).replace(".", ",")

        # 🔹 Retornar JSON com valores formatados corretamente
        data = {
            "saldo_anterior_custeio": formatar_valor(receita["saldo_anterior_custeio"]),
            "saldo_anterior_capital": formatar_valor(receita["saldo_anterior_capital"]),
            "valor_creditado_custeio": formatar_valor(receita["valor_creditado_custeio"]),
            "valor_creditado_capital": formatar_valor(receita["valor_creditado_capital"]),
            "saldo_reprogramar_custeio": formatar_valor(receita["saldo_reprogramar_custeio"]),
            "saldo_reprogramar_capital": formatar_valor(receita["saldo_reprogramar_capital"]),
        }

        return JsonResponse(data)

    except Exception as e:
        return JsonResponse({"erro": f"Erro ao buscar dados financeiros: {str(e)}"}, status=500)
###***********************************************************************************************************************

from django.shortcuts import render, get_object_or_404, redirect
from django.contrib import messages
from .models import Pagamento

def editar_pagamento(request, pagamento_id):
    pagamento = get_object_or_404(Pagamento, id=pagamento_id)

    if request.method == "POST":
        pagamento.nome_favorecido = request.POST.get("nome_favorecido")
        pagamento.cnpj_cpf = request.POST.get("cnpj_cpf")
        pagamento.tipo_bem_servico = request.POST.get("tipo_bem_servico")
        pagamento.tipo_pagamento = request.POST.get("tipo_pagamento")
        pagamento.tipo_documento = request.POST.get("tipo_documento")
        pagamento.data_pagamento = request.POST.get("data_pagamento")
        pagamento.valor = request.POST.get("valor")

        pagamento.save()
        messages.success(request, "✅ Pagamento atualizado com sucesso!")
        return redirect("pddelancar_pagamento")

    return render(request, "pdde/editar_pagamento.html", {"pagamento": pagamento})
###***********************************************************************************************************************

from django.shortcuts import render
from .models import Escola, Pagamento

def demonstrativo_execucao_pdde(request):
    escolas = Escola.objects.all()  # Buscar todas as escolas disponíveis
    escola_selecionada = request.GET.get('escola_id', None)  # Verifica se há uma escola selecionada

    pagamentos = Pagamento.objects.filter(escola_id=escola_selecionada) if escola_selecionada else Pagamento.objects.none()

    context = {
        'escolas': escolas,
        'pagamentos': pagamentos,  # Passa os pagamentos para o template
        'escola_selecionada': escola_selecionada,
    }

    return render(request, 'pddereceitadespesa.html', context)
###***********************************************************************************************************************

from django.db.models import Sum
from django.shortcuts import render
from .models import ContaBancaria
from decimal import Decimal

def emitir_conciliacao(request):
    contas = ContaBancaria.objects.all()
    contas_corrigidas = []

    for conta in contas:
        total_credito = conta.lancamentos.filter(tipo='credito').aggregate(total=Sum('valor'))['total'] or Decimal('0.00')
        total_debito = conta.lancamentos.filter(tipo='debito').aggregate(total=Sum('valor'))['total'] or Decimal('0.00')
        saldo_corrigido = total_credito - total_debito

        contas_corrigidas.append({
            'banco': conta.banco,
            'agencia': conta.agencia,
            'conta': conta.conta,
            'saldo_corrigido': saldo_corrigido,
            'escola': conta.escola,  # necessário para acessar escola.nome_conselho no template
        })

    return render(request, 'conciliacao/emitir_conciliacao.html', {
        'contas': contas_corrigidas
    })



###***********************************************************************************************************************

from django.shortcuts import render
from django.http import HttpResponse
from .models import ContaBancaria
from weasyprint import HTML
from django.db.models import Sum
from decimal import Decimal

def gerar_conciliacao_pdf(request):
    """Gera um PDF da conciliação bancária com saldos corrigidos."""
    contas = ContaBancaria.objects.all()
    contas_corrigidas = []

    for conta in contas:
        total_credito = conta.lancamentos.filter(tipo='credito').aggregate(total=Sum('valor'))['total'] or Decimal('0.00')
        total_debito = conta.lancamentos.filter(tipo='debito').aggregate(total=Sum('valor'))['total'] or Decimal('0.00')
        saldo_corrigido = total_credito - total_debito

        contas_corrigidas.append({
            'banco': conta.banco,
            'agencia': conta.agencia,
            'conta': conta.conta,
            'saldo_corrigido': saldo_corrigido,
            'escola': conta.escola,
        })

    html_string = render(
        request,
        "conciliacao/conciliacao_pdf.html",
        {"contas": contas_corrigidas}
    ).content.decode("utf-8")

    pdf_file = HTML(string=html_string).write_pdf()

    response = HttpResponse(pdf_file, content_type="application/pdf")
    response["Content-Disposition"] = 'inline; filename="conciliacao.pdf"'
    return response

###***********************************************************************************************************************
###***********************************************************************************************************************
###****************************************SEPECC PESQUISA DE PREÇOS******************************************************
###***********************************************************************************************************************

from django.shortcuts import render, redirect
from .models import Item, Proponente, Proposta, ApuracaoResultado
from .forms import ItemForm, ProponenteForm, PropostaForm

def pesquisa_precos(request):
    """Tela principal do módulo com submenus"""
    return render(request, "pesquisa_precos/pesquisa_precos.html")
###***********************************************************************************************************************

from django.shortcuts import render, redirect
from .models import Item, Categoria, Subcategoria

def cadastrar_item(request):
    """Cadastro de Itens incluindo Categoria e Subcategoria"""
    categorias = Categoria.objects.all()

    if request.method == "POST":
        nome = request.POST.get("nome")
        descricao = request.POST.get("descricao")
        unidade_medida = request.POST.get("unidade_medida")
        categoria_id = request.POST.get("categoria")
        subcategoria_id = request.POST.get("subcategoria")

        categoria = Categoria.objects.get(id=categoria_id)
        subcategoria = Subcategoria.objects.get(id=subcategoria_id)

        if nome and unidade_medida and categoria and subcategoria:
            Item.objects.create(
                nome=nome,
                descricao=descricao,
                unidade_medida=unidade_medida,
                categoria=categoria,
                subcategoria=subcategoria
            )
            return redirect("listar_item")

    return render(request, "pesquisa_precos/cadastrar_item.html", {
        "categorias": categorias
    })
###***********************************************************************************************************************

from django.shortcuts import render, redirect
from .models import Proponente

def cadastrar_proponente(request):
    """Cadastro de Proponentes diretamente na view"""
    if request.method == "POST":
        nome = request.POST.get("nome")
        cpf_cnpj = request.POST.get("cpf_cnpj")
        email = request.POST.get("email")
        telefone = request.POST.get("telefone")
        endereco = request.POST.get("endereco")
        bairro = request.POST.get("bairro")
        cidade = request.POST.get("cidade")
        estado = request.POST.get("estado")
        cep = request.POST.get("cep")
        tipo_proponente = request.POST.get("tipo_proponente")
        representante_legal = request.POST.get("representante_legal")
        observacoes = request.POST.get("observacoes")

        # Criar o objeto Proponente e salvar no banco de dados
        Proponente.objects.create(
            nome=nome,
            cpf_cnpj=cpf_cnpj,
            email=email,
            telefone=telefone,
            endereco=endereco,
            bairro=bairro,
            cidade=cidade,
            estado=estado,
            cep=cep,
            tipo_proponente=tipo_proponente,
            representante_legal=representante_legal,
            observacoes=observacoes
        )

        return redirect("listar_proponentes")  # Redireciona para a listagem após o cadastro

    return render(request, "pesquisa_precos/cadastrar_proponente.html", {
        "titulo": "Cadastrar Proponente",
        "descricao": "Informe os dados do proponente que fornecerá a proposta de preço."
    })
###***********************************************************************************************************************

from django.shortcuts import render
from django.http import JsonResponse
from django.utils.timezone import now
import json
from .models import Proposta, Item, Proponente, ApuracaoResultado

def cadastrar_propostas(request):
    if request.method == "POST":
        data = json.loads(request.body)
        propostas = data.get("propostas", [])

        for proposta in propostas:
            try:
                item = Item.objects.get(id=proposta["item_id"])
                proponente = Proponente.objects.get(id=proposta["proponente_id"])
                preco_unitario = float(proposta["preco_unitario"])
                quantidade = int(proposta["quantidade"])
                valor_total = preco_unitario * quantidade  # 🔹 Calcula o valor total

                print(f"📌 Salvando proposta: Item={item}, Proponente={proponente}, Preço={preco_unitario}, Quantidade={quantidade}, Valor Total={valor_total}")  # LOG PARA DEPURAÇÃO

                # Criar a proposta com a quantidade e valor total corretos
                nova_proposta = Proposta.objects.create(
                    item=item,
                    proponente=proponente,
                    preco_unitario=preco_unitario,
                    quantidade=quantidade,  # 🔹 Agora está correto
                    valor_total=valor_total,  # 🔹 Agora está correto
                    data_proposta=now(),
                )

            except Exception as e:
                print(f"❌ Erro ao salvar proposta: {e}")  # LOG DE ERRO
                return JsonResponse({"message": "Erro ao salvar proposta!", "error": str(e)}, status=400)

        # **ATUALIZA A TABELA DE APURAÇÃO**
        atualizar_apuracao_resultado()

        return JsonResponse({"message": "Propostas salvas e apuração atualizada com sucesso!"})

    itens = Item.objects.all()
    proponentes = Proponente.objects.all()
    
    return render(request, "pesquisa_precos/cadastrar_propostas.html", {
        "itens": itens,
        "proponentes": proponentes,
        "titulo": "Cadastrar Propostas",
        "descricao": "Escolha os proponentes e itens para registrar os preços."
    })




from django.http import JsonResponse
from django.shortcuts import get_object_or_404, render
from django.utils.timezone import now
from .models import ApuracaoResultado, Proposta, Proponente, Item

def atualizar_apuracao_resultado():
    """
    Atualiza a tabela `ApuracaoResultado` com o menor preço para cada item adjudicado.
    """
    itens_com_propostas = Proposta.objects.values_list('item', flat=True).distinct()

    for item_id in itens_com_propostas:
        melhor_proposta = Proposta.objects.filter(item_id=item_id).order_by("preco_unitario").first()

        if melhor_proposta:
            ApuracaoResultado.objects.update_or_create(
                item=melhor_proposta.item,
                defaults={
                    "proponente_vencedor": melhor_proposta.proponente,
                    "preco_vencedor": melhor_proposta.preco_unitario,
                    "quantidade": melhor_proposta.quantidade,
                    "valor_total": melhor_proposta.valor_total,
                    "data_apuracao": now(),
                    "status": "Pendente"
                }
            )

def apuracao_resultados_view(request):
    """
    Exibe os vencedores e permite adjudicação/revogação com filtro por empresa vencedora.
    """

    # **Chama a função de apuração antes de renderizar os dados**
    atualizar_apuracao_resultado()

    # Obtém todas as empresas vencedoras para o filtro
    empresas_vencedoras = Proponente.objects.filter(apuracaoresultado__status="Adjudicado").distinct()

    # Filtro por empresa selecionada
    empresa_selecionada = request.GET.get('empresa_vencedora', '')

    if empresa_selecionada:
        resultados = ApuracaoResultado.objects.filter(proponente_vencedor__id=empresa_selecionada, status="Adjudicado")
    else:
        resultados = ApuracaoResultado.objects.filter(status="Adjudicado")

    return render(request, "pesquisa_precos/apuracao_resultado.html", {
        "resultados": resultados,
        "empresas_vencedoras": empresas_vencedoras,
        "empresa_selecionada": empresa_selecionada
    })





###***********************************************************************************************************************

def apuracao_resultado(request):
    """Tela de Apuração e Resultado"""
    resultados = ApuracaoResultado.objects.all()

    return render(request, "pesquisa_precos/apuracao_resultado.html", {
        "resultados": resultados
    })
###***********************************************************************************************************************

from django.http import JsonResponse

def get_subcategorias(request, categoria_id):
    """Retorna as subcategorias de uma categoria específica"""
    subcategorias = Subcategoria.objects.filter(categoria_id=categoria_id).values("id", "nome")
    return JsonResponse({"subcategorias": list(subcategorias)})
###***********************************************************************************************************************

from django.shortcuts import render
from .models import Item

def listar_itens(request):
    """Exibe a lista de itens cadastrados"""
    itens = Item.objects.all()
    return render(request, "pesquisa_precos/listar_itens.html", {"itens": itens})
###***********************************************************************************************************************

from django.shortcuts import render
from .models import Proponente

def listar_proponentes(request):
    """Lista todos os proponentes cadastrados"""
    proponentes = Proponente.objects.all()
    return render(request, "pesquisa_precos/listar_proponentes.html", {"proponentes": proponentes})
###***********************************************************************************************************************

from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from django.http import HttpResponse
from .models import Item, Proponente, Proposta  # Certifique-se de importar os modelos corretos

def gerar_pdf_orcamento(request):
    """ Gera um PDF com os orçamentos de preços dos proponentes. """

    # Criar a resposta HTTP com o cabeçalho para PDF
    response = HttpResponse(content_type="application/pdf")
    response["Content-Disposition"] = 'attachment; filename="orcamento.pdf"'

    # Criar um objeto PDF usando ReportLab
    p = canvas.Canvas(response, pagesize=A4)
    largura, altura = A4

    # Definir título
    p.setFont("Helvetica-Bold", 16)
    p.drawString(200, altura - 50, "Orçamento de Propostas")

    # Adicionar informações dos itens e preços por proponente
    p.setFont("Helvetica", 12)
    y_position = altura - 100  # Começar abaixo do título

    itens = Item.objects.all()
    proponentes = Proponente.objects.all()

    # Criar cabeçalho da tabela
    p.drawString(50, y_position, "Item")
    x_position = 200
    for proponente in proponentes:
        p.drawString(x_position, y_position, proponente.nome[:15])  # Limite de caracteres
        x_position += 150

    y_position -= 20  # Ajustar a posição abaixo do cabeçalho

    # Adicionar os itens e seus preços oferecidos pelos proponentes
    for item in itens:
        p.drawString(50, y_position, item.nome)
        x_position = 200
        for proponente in proponentes:
            proposta = Proposta.objects.filter(item=item, proponente=proponente).first()
            if proposta:
                preco = f"R$ {proposta.preco_unitario:.2f}"  # ✅ Corrigido de proposta.preco para proposta.preco_unitario
            else:
                preco = "N/A"
            p.drawString(x_position, y_position, preco)
            x_position += 150
        y_position -= 20  # Mover para a próxima linha

        # Verificar se precisa de uma nova página
        if y_position < 50:
            p.showPage()  # Criar nova página
            y_position = altura - 100  # Reiniciar a posição

    # Finalizar PDF
    p.showPage()
    p.save()

    return response

###***********************************************************************************************************************

import pandas as pd
from django.http import HttpResponse
from .models import Item, Proponente, Proposta  # Certifique-se de importar os modelos corretos

def gerar_excel_orcamento(request):
    """ Gera um arquivo Excel com os orçamentos de preços dos proponentes. """

    # Criar resposta HTTP para arquivo Excel
    response = HttpResponse(content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
    response["Content-Disposition"] = 'attachment; filename="orcamento.xlsx"'

    # Criar um DataFrame para armazenar os dados
    itens = Item.objects.all()
    proponentes = Proponente.objects.all()

    # Criar a estrutura de cabeçalho
    header = ["Item", "Unidade"]
    for proponente in proponentes:
        header.append(proponente.nome)

    # Criar as linhas com os preços
    rows = []
    for item in itens:
        row = [item.nome, item.unidade_medida]
        for proponente in proponentes:
            proposta = Proposta.objects.filter(item=item, proponente=proponente).first()
            if proposta:
                row.append(proposta.preco_unitario)  # ✅ Corrigido de proposta.preco para proposta.preco_unitario
            else:
                row.append("N/A")
        rows.append(row)

    # Criar DataFrame do pandas
    df = pd.DataFrame(rows, columns=header)

    # Criar o arquivo Excel usando `xlsxwriter`
    with pd.ExcelWriter(response, engine="xlsxwriter") as writer:
        df.to_excel(writer, sheet_name="Orçamento", index=False)

    return response

###***********************************************************************************************************************

from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
import json
from .models import Categoria, Subcategoria

@csrf_exempt  # Permite chamadas AJAX sem CSRF Token
def cadastrar_subcategoria(request):
    if request.method == "POST":
        try:
            data = json.loads(request.body)  # Lê os dados do JSON enviado
            categoria_id = data.get("categoria_id")
            nome_subcategoria = data.get("nome")

            if not categoria_id or not nome_subcategoria:
                return JsonResponse({"error": "Todos os campos são obrigatórios!"}, status=400)

            categoria = Categoria.objects.get(id=categoria_id)
            subcategoria = Subcategoria.objects.create(nome=nome_subcategoria, categoria=categoria)

            return JsonResponse({"id": subcategoria.id, "nome": subcategoria.nome}, status=201)

        except Categoria.DoesNotExist:
            return JsonResponse({"error": "Categoria não encontrada!"}, status=404)
        except Exception as e:
            return JsonResponse({"error": str(e)}, status=500)

    return JsonResponse({"error": "Método não permitido"}, status=405)
###***********************************************************************************************************************

from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
import json
from .models import Categoria

@csrf_exempt  # Permite chamadas AJAX sem CSRF Token
def cadastrar_categoria(request):
    """Salva uma nova categoria via requisição AJAX."""
    if request.method == "POST":
        try:
            data = json.loads(request.body)  # Lê os dados do JSON enviado
            nome_categoria = data.get("nome")

            if not nome_categoria:
                return JsonResponse({"error": "O nome da categoria é obrigatório!"}, status=400)

            # Criando a nova categoria no banco de dados
            categoria = Categoria.objects.create(nome=nome_categoria)

            return JsonResponse({"id": categoria.id, "nome": categoria.nome, "message": "Categoria cadastrada com sucesso!"})

        except Exception as e:
            return JsonResponse({"error": str(e)}, status=500)

    return JsonResponse({"error": "Método não permitido"}, status=405)
###***********************************************************************************************************************

from django.http import JsonResponse
from django.shortcuts import get_object_or_404, render
import json
from .models import ApuracaoResultado

def apuracao_resultados(request):
    """ Exibe os vencedores e permite adjudicação/revogação """

    if request.method == "POST":
        try:
            data = json.loads(request.body)  # ✅ Garante que o JSON é lido corretamente
            acao = data.get("acao")
            resultado_id = data.get("resultado_id")

            if not acao or not resultado_id:
                return JsonResponse({"error": "Ação ou ID do resultado ausente!"}, status=400)

            resultado = get_object_or_404(ApuracaoResultado, id=resultado_id)

            if acao == "adjudicar":
                resultado.status = "Adjudicado"
            elif acao == "revogar":
                resultado.status = "Revogado"
            else:
                return JsonResponse({"error": "Ação inválida!"}, status=400)

            resultado.save()

            return JsonResponse({
                "message": f"{acao.capitalize()} com sucesso!",
                "novo_status": resultado.status,
                "quantidade": resultado.quantidade,
                "valor_total": float(resultado.valor_total),
            })

        except json.JSONDecodeError:
            return JsonResponse({"error": "JSON inválido!"}, status=400)
        except Exception as e:
            return JsonResponse({"error": str(e)}, status=500)

    # Se for GET, retorna a página normalmente
    resultados = ApuracaoResultado.objects.all()
    return render(request, "pesquisa_precos/apuracao_resultado.html", {
        "resultados": resultados
    })
###***********************************************************************************************************************

from django.http import JsonResponse
from django.shortcuts import get_object_or_404
from .models import ApuracaoResultado, Proponente

def detalhes_proponente(request, proponente_id):
    """
    Retorna detalhes do proponente, incluindo seus itens adjudicados e o status de cada item.
    """
    proponente = get_object_or_404(Proponente, id=proponente_id)
    resultados = ApuracaoResultado.objects.filter(proponente_vencedor=proponente)

    itens_list = [
        {
            "nome": resultado.item.nome,
            "quantidade": resultado.quantidade,
            "preco_unitario": float(resultado.preco_vencedor),
            "valor_total": float(resultado.valor_total),
            "status": resultado.status  # ✅ Adicionando status do item na API
        } for resultado in resultados
    ]

    return JsonResponse({
        "proponente": {
            "id": proponente.id,
            "nome": proponente.nome
        },
        "itens": itens_list  # ✅ Retornando apenas os itens com status
    })
###***********************************************************************************************************************

from django.shortcuts import render, redirect
from .models import EscolaPdde, Documento

def cadastrar_documento(request):
    escolas = EscolaPdde.objects.all()

    if request.method == "POST":
        escola_id = request.POST.get("escola")
        tipo = request.POST.get("tipo")  # Você precisa ter esse campo no formulário
        numero = request.POST.get("numero")
        data_emissao = request.POST.get("data_emissao")
        descricao = request.POST.get("descricao")
        arquivo = request.FILES.get("arquivo")
        valor_total = request.POST.get("valor_total")

        escola = EscolaPdde.objects.get(id=escola_id)

        Documento.objects.create(
            escola=escola,
            tipo=tipo,
            numero=numero,
            data_emissao=data_emissao,
            descricao=descricao,
            arquivo=arquivo,
            valor_total=valor_total or 0,
        )

        return redirect("pdde")  # Redireciona após o cadastro

    context = {
        "escolas": escolas,
    }
    return render(request, "prestacao_contas/cadastrar_documento.html", context)
###***********************************************************************************************************************

from django.shortcuts import render, redirect
from django.contrib import messages
from django.utils import timezone
from .models import Bem, EscolaPdde, Documento

def cadastrar_bem(request):
    escolas = EscolaPdde.objects.all()
    documentos = Documento.objects.all()

    if request.method == "POST":
        try:
            escola_id = request.POST.get("escola")
            nome_conselho = request.POST.get("nome_conselho")
            nome = request.POST.get("nome")
            documento = request.FILES.get("documento")
            quantidade = int(request.POST.get("quantidade"))
            valor_unitario = float(request.POST.get("valor_unitario").replace(",", "."))

            escola = EscolaPdde.objects.get(id=escola_id)

            bem = Bem.objects.create(
                escola=escola,
                nome_conselho=nome_conselho,
                nome=nome,
                documento=documento,
                quantidade=quantidade,
                valor_unitario=valor_unitario,
                data_cadastro=timezone.now()
            )

            # valor_total será calculado no método `save()` do modelo
            bem.save()

            messages.success(request, "Bem cadastrado com sucesso!")
            return redirect("listar_bens")

        except Exception as e:
            print("Erro ao salvar o bem:", e)
            messages.error(request, "Erro ao cadastrar o bem.")

    return render(request, "prestacao_contas/cadastrar_bem.html", {
        "escolas": escolas,
        "documentos": documentos
    })



###***********************************************************************************************************************

from django.shortcuts import render, redirect
from django.contrib import messages
from .forms import RepresentanteLegalForm

def cadastrar_representante(request):
    form = RepresentanteLegalForm(request.POST or None, request.FILES or None)
    if form.is_valid():
        form.save()
        messages.success(request, "Representante legal cadastrado com sucesso!")
        return redirect('listar_representantes')  # Certifique-se que esta URL existe no seu urls.py
    return render(request, "prestacao_contas/cadastrar_representante.html", {"form": form})

###***********************************************************************************************************************

from django.http import HttpResponse
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from .models import Documento, BemAdquirido
from datetime import datetime

def gerar_prestacao_contas_pdf(request):
    """ Gera um relatório PDF da prestação de contas de bens adquiridos pelas escolas. """
    
    response = HttpResponse(content_type="application/pdf")
    response["Content-Disposition"] = f'attachment; filename="prestacao_contas_{datetime.now().strftime("%Y%m%d")}.pdf"'

    p = canvas.Canvas(response, pagesize=A4)
    largura, altura = A4

    # **Cabeçalho do Documento**
    p.setFont("Helvetica-Bold", 16)
    p.drawString(150, altura - 50, "Prestação de Contas - Bens Adquiridos")

    p.setFont("Helvetica", 12)
    y_position = altura - 80  # Ajustar a posição inicial

    # **Dados dos Documentos**
    documentos = Documento.objects.all()
    
    if not documentos:
        p.drawString(50, y_position, "Nenhum documento cadastrado.")
    else:
        p.drawString(50, y_position, "Lista de Documentos:")
        y_position -= 20

        for doc in documentos:
            p.drawString(50, y_position, f"{doc.get_tipo_display()} Nº {doc.numero} - {doc.escola.nome} ({doc.data_emissao})")
            y_position -= 15

            if y_position < 50:
                p.showPage()
                y_position = altura - 80

    y_position -= 30  # Espaçamento entre seções

    # **Dados dos Bens Adquiridos**
    bens = BemAdquirido.objects.all()

    if not bens:
        p.drawString(50, y_position, "Nenhum bem adquirido cadastrado.")
    else:
        p.drawString(50, y_position, "Lista de Bens Adquiridos:")
        y_position -= 20

        p.drawString(50, y_position, "Item")
        p.drawString(250, y_position, "Quantidade")
        p.drawString(350, y_position, "Valor Unitário")
        p.drawString(450, y_position, "Valor Total")
        y_position -= 20

        total_geral = 0

        for bem in bens:
            total = bem.valor_total
            p.drawString(50, y_position, bem.especificacao)
            p.drawString(250, y_position, str(bem.quantidade))
            p.drawString(350, y_position, f"R$ {bem.valor_unitario:.2f}")
            p.drawString(450, y_position, f"R$ {total:.2f}")
            total_geral += total
            y_position -= 15

            if y_position < 50:
                p.showPage()
                y_position = altura - 80

        y_position -= 20
        p.drawString(50, y_position, f"TOTAL GERAL: R$ {total_geral:.2f}")

    p.showPage()
    p.save()

    return response
###***********************************************************************************************************************

from django.shortcuts import render
from .models import Documento  # ✅ Importando o modelo de Documentos

def listar_documentos(request):
    """
    Exibe todos os documentos cadastrados.
    """
    documentos = Documento.objects.all()
    return render(request, "prestacao_contas/listar_documentos.html", {"documentos": documentos})
###***********************************************************************************************************************

from django.shortcuts import render
from .models import Bem  # Certifique-se de que este modelo existe

def listar_bens(request):
    """
    Lista todos os bens cadastrados.
    """
    bens = Bem.objects.all()
    return render(request, "prestacao_contas/listar_bens.html", {"bens": bens})
###***********************************************************************************************************************

def get_escola_dados(request, escola_id):
    """
    Retorna os dados da escola selecionada via AJAX (JSON).
    """
    escola = get_object_or_404(EscolaPdde, id=escola_id)

    # Obtendo os programas vinculados à escola
    programas_vinculados = list(escola.programas.values_list("nome", flat=True))

    data = {
        "nome": escola.nome,
        "nome_conselho": escola.nome_conselho if escola.nome_conselho else "Não informado",
        "cnpj": escola.cnpj,
        "endereco": escola.endereco if escola.endereco else "Não informado",
        "bairro": escola.bairro if escola.bairro else "Não informado",
        "cidade": escola.cidade if escola.cidade else "Canaã dos Carajás",
        "uf": escola.uf if escola.uf else "PA",
        "programas_vinculados": programas_vinculados if programas_vinculados else ["Nenhum programa vinculado"]
    }

    return JsonResponse(data)
###***********************************************************************************************************************

from django.shortcuts import render
from .models import Documento
from .models import RepresentanteLegal


def listar_documento(request):
    documentos = Documento.objects.all()
    return render(request, "prestacao_contas/listar_documento.html", {"documentos": documentos})
###***********************************************************************************************************************

from django.shortcuts import render
from .models import RepresentanteLegal

# views.py
def listar_representantes(request):
    representantes = RepresentanteLegal.objects.all()
    return render(request, 'prestacao_contas/listar_representantes.html', {'representantes': representantes})
###***********************************************************************************************************************

from django.shortcuts import render, redirect, get_object_or_404
from django.http import HttpResponse
from django.template.loader import render_to_string
from .models import TermoDoacao, BemDoado, EscolaPdde
from .forms import TermoDoacaoForm
from django.utils.dateparse import parse_date
import datetime


def cadastrar_termo_doacao(request):
    if request.method == "POST":
        form = TermoDoacaoForm(request.POST)
        if form.is_valid():
            termo = form.save()
            return redirect('visualizar_termo_doacao', termo_id=termo.id)
    else:
        form = TermoDoacaoForm()

    termos_emitidos = TermoDoacao.objects.all().order_by('-data_emissao')
    termos_com_bens = []

    for termo in termos_emitidos:
        bens = BemDoado.objects.filter(escola=termo.escola)
        termos_com_bens.append({
            "termo": termo,
            "bens": bens
        })

    bens_disponiveis = BemDoado.objects.all().order_by('-data_nota')

    return render(request, "prestacao_contas/cadastrar_termo_doacao.html", {
        "form": form,
        "termos_emitidos": termos_com_bens,
        "bens_disponiveis": bens_disponiveis
    })
###***********************************************************************************************************************

# Visualizar termo de doação
def visualizar_termo_doacao(request, termo_id):
    termo = get_object_or_404(TermoDoacao, id=termo_id)
    bens = BemDoado.objects.filter(escola=termo.escola)
    total_geral = sum(bem.valor_total for bem in bens if bem.valor_total)

    return render(request, "prestacao_contas/visualizar_termo_doacao.html", {
        "termo": termo,
        "bens": bens,
        "total_geral": total_geral
    })
###***********************************************************************************************************************

# Cadastrar bem doado
def cadastrar_bem_doado(request):
    escolas = EscolaPdde.objects.all().order_by('nome')

    if request.method == 'POST':
        escola_id = request.POST.get('escola')
        conselho = request.POST.get('nome_conselho')
        descricao = request.POST.get('descricao')
        numero_nota = request.POST.get('nota_fiscal_numero')
        data_nota = request.POST.get('nota_fiscal_data')
        quantidade = request.POST.get('quantidade')
        valor_unitario = request.POST.get('valor_unitario', '0').replace(",", ".")
        valor_total = request.POST.get('valor_total', '0').replace(",", ".")

        try:
            escola = EscolaPdde.objects.get(id=escola_id)
            BemDoado.objects.create(
                escola=escola,
                conselho=conselho,
                descricao=descricao,
                numero_nota=numero_nota,
                data_nota=parse_date(data_nota) if data_nota else None,
                quantidade=int(quantidade),
                valor_unitario=float(valor_unitario),
                valor_total=float(valor_total),
            )
            return redirect('listar_bens_doados')
        except Exception as e:
            print("Erro ao salvar o bem doado:", e)

    return render(request, 'prestacao_contas/cadastrar_bem_doado.html', {'escolas': escolas})
###***********************************************************************************************************************

# Listar bens doados
def listar_bens_doados(request):
    bens_doados = BemDoado.objects.all().order_by('-data_doacao')
    return render(request, "prestacao_contas/listar_bens_doados.html", {
        "bens_doados": bens_doados
    })
###***********************************************************************************************************************

# Gerar versão HTML do termo (ex: para visualização ou conversão em PDF)
def gerar_pdf_termo_doacao(request, termo_id):
    termo = get_object_or_404(TermoDoacao, id=termo_id)
    bens = BemDoado.objects.filter(escola=termo.escola)
    total_geral = sum(bem.valor_total for bem in bens if bem.valor_total)

    context = {
        'conselho': termo.conselho,
        'escola': termo.escola.nome,
        'programa': 'PDDE',
        'bens': bens,
        'data_atual': datetime.today().strftime('%d/%m/%Y'),
        'total_geral': total_geral
    }

    html = render_to_string("prestacao_contas/termo_doacao_pdf.html", context)
    return HttpResponse(html)  # substitua por geração real de PDF se necessário
###***********************************************************************************************************************

from django.shortcuts import render
from .models import BemDoado
from datetime import datetime


def gerar_termo_doacao(request):
    bens = BemDoado.objects.all()
    data_atual = datetime.date.today()

    return render(request, "prestacao_contas/termo_doacao.html", {
        "bens": bens,
        "data": data_atual
    })
###***********************************************************************************************************************

from django.http import HttpResponse
from django.template.loader import render_to_string
from django.utils import timezone
from weasyprint import HTML
from .models import BemDoado

def gerar_pdf_bem(request, bem_id):
    try:
        bem = BemDoado.objects.get(id=bem_id)
    except BemDoado.DoesNotExist:
        return HttpResponse("Bem não encontrado.", status=404)

    context = {
        'bem': bem,
        'data_atual': timezone.now().date()
    }

    html_string = render_to_string("prestacao_contas/bem_doado_pdf.html", context)

    # Gerar o PDF
    html = HTML(string=html_string)
    pdf_file = html.write_pdf()

    response = HttpResponse(pdf_file, content_type='application/pdf')
    response['Content-Disposition'] = f'attachment; filename="termo_bem_{bem.id}.pdf"'
    return response

########################################################################################################################################################
#################################################PLANO DE GESTOA ESCOLA#################################################################################
########################################################################################################################################################
########################################################################################################################################################

from django.shortcuts import render, redirect
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
from .forms import PlanoGestaoEscolarForm

def enviar_plano_gestao(request):
    if request.method == 'POST':
        form = PlanoGestaoEscolarForm(request.POST, request.FILES)
        if form.is_valid():
            form.save()
            if request.headers.get('x-requested-with') == 'XMLHttpRequest':
                return JsonResponse({'mensagem': 'Enviado com sucesso!'}, status=200)
            return redirect('listar_planos')
        else:
            if request.headers.get('x-requested-with') == 'XMLHttpRequest':
                return JsonResponse({'erros': form.errors}, status=400)
    else:
        form = PlanoGestaoEscolarForm()
    
    return render(request, 'webapp/novo_plano.html', {'form': form})
###***********************************************************************************************************************

from django.shortcuts import render, redirect, get_object_or_404
from .models import PGEPlanoGestaoEscolar
from django.contrib import messages

def listar_pge_planos(request):
    filtro = request.GET.get('filtro_unidade')
    unidades = PGEPlanoGestaoEscolar.objects.values_list('unidade_ensino', flat=True).distinct()

    if filtro:
        planos = PGEPlanoGestaoEscolar.objects.filter(unidade_ensino=filtro)
    else:
        planos = PGEPlanoGestaoEscolar.objects.all()

    return render(request, 'setor_pedagogico/listar_planos.html', {
        'planos': planos,
        'unidades': unidades,
        'filtro': filtro
    })
###***********************************************************************************************************************

def atualizar_status_plano(request, plano_id):
    if request.method == 'POST':
        status = request.POST.get('status')
        try:
            plano = PlanoGestaoEscolar.objects.get(id=plano_id)
            plano.deferido_indeferido = status
            plano.status = status.lower()  # ou "deferido", "indeferido", etc.
            plano.save()
            return JsonResponse({'success': True})
        except PlanoGestaoEscolar.DoesNotExist:
            return JsonResponse({'success': False, 'error': 'Plano não encontrado'})
    return JsonResponse({'success': False, 'error': 'Requisição inválida'})
###***********************************************************************************************************************

from django.http import JsonResponse
from .models import PGEPlanoGestaoEscolar

def obter_dados_unidade(request, unidade_nome):
    try:
        dados = PGEPlanoGestaoEscolar.objects.filter(unidade_ensino=unidade_nome).first()
        if dados:
            return JsonResponse({
                'cargo': dados.cargo,
                'servidor': dados.servidor,
                'telefone': dados.telefone
            })
    except Exception as e:
        return JsonResponse({'erro': str(e)}, status=500)
    
    return JsonResponse({'erro': 'Unidade não encontrada'}, status=404)
###***********************************************************************************************************************

from django.shortcuts import render
from .models import PlanoGestaoEscolar  # <- este é o modelo correto da tabela semedapp_planogestaoescolar

def plano_gestao_escolar_admin(request):
    filtro = request.GET.get('filtro_unidade')
    unidades = PlanoGestaoEscolar.objects.values_list('unidade_ensino', flat=True).distinct()

    if filtro:
        planos = PlanoGestaoEscolar.objects.filter(unidade_ensino=filtro)
    else:
        planos = PlanoGestaoEscolar.objects.all().order_by('-id')

    return render(request, 'setor_pedagogico/pge_admin.html', {
        'planos': planos,
        'unidades': unidades,
        'filtro': filtro
    })
###***********************************************************************************************************************

from django.shortcuts import render
from .forms import PlanoGestaoEscolarForm

def novo_plano_gestao(request):
    form = PlanoGestaoEscolarForm()
    return render(request, 'webapp/novo_plano.html', {'form': form})
###***********************************************************************************************************************

from django.db.utils import OperationalError
from django.contrib.auth import get_user_model
from .models import CadastroEI
from django.http import JsonResponse

User = get_user_model()

def importar_dados(request):
    try:
        # Verifica se tabela auth_user ou customizada existe
        if not User.objects.exists():
            return JsonResponse({"erro": "Tabela de usuários ainda não foi criada. Rode as migrations."}, status=500)

        # Exemplo de criação segura
        coordenador = User.objects.filter(id=86).first()
        professor = User.objects.filter(id=45).first()

        if not coordenador or not professor:
            return JsonResponse({"erro": "Coordenador ou professor não encontrado."}, status=400)

        CadastroEI.objects.create(
            id_matricula=240475,
            unidade_ensino='EMEB LUÍS CARLOS PRESTES',
            formato_letivo='II PERÍODO',
            cpf='9150389270',
            data_nascimento=None,
            ano='2025',
            modalidade='Pré-escolar',
            pessoa_nome='AGATHA HADASSAH SOUSA AMARAL',
            idade=4,
            avaliado='SIM',
            questao_matematica_1='CERTO',
            questao_matematica_2='CERTO',
            questao_matematica_3='CERTO',
            questao_matematica_4='CERTO',
            questao_matematica_5='CERTO',
            questao_matematica_6='CERTO',
            questao_matematica_7='CERTO',
            questao_matematica_8='CERTO',
            questao_matematica_9='CERTO',
            questao_matematica_10='CERTO',
            questao_linguagem_11='CC',
            questao_linguagem_12='PARCIAL',
            questao_linguagem_13='SSVS',
            questao_linguagem_14='CERTO',
            questao_linguagem_15='CERTO',
            questao_linguagem_16='CERTO',
            questao_linguagem_17='CERTO',
            questao_linguagem_18='CERTO',
            questao_linguagem_19='CERTO',
            questao_linguagem_20='ERRADO',
            professor=professor,
            coordenador=coordenador,
            turma='II PERÍODO A-VESP'
        )

        return JsonResponse({"mensagem": "Cadastro inserido com sucesso!"})

    except OperationalError as e:
        return JsonResponse({"erro": f"Erro de banco de dados: {str(e)}"}, status=500)
###***********************************************************************************************************************

from django.shortcuts import render, get_object_or_404
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
from django.views.decorators.http import require_POST
from .models import Certidao
from datetime import date


def emissao_certidoes(request):
    certidoes = Certidao.objects.all().order_by('id')
    return render(request, 'contabilidade/emissao_certidoes.html', {
        'certidoes': certidoes,
    })




###***********************************************************************************************************************

@csrf_exempt
@require_POST
def upload_arquivo_certidao(request):
    certidao_id = request.POST.get('certidao_id')
    arquivo = request.FILES.get('arquivo')

    if not certidao_id or not arquivo:
        return JsonResponse({'status': 'erro', 'mensagem': 'Dados incompletos'}, status=400)

    try:
        certidao = Certidao.objects.get(id=certidao_id)
        certidao.arquivo = arquivo
        certidao.save()
        return JsonResponse({'status': 'ok'})
    except Certidao.DoesNotExist:
        return JsonResponse({'status': 'erro', 'mensagem': 'Certidão não encontrada'}, status=404)
###***********************************************************************************************************************

from django.shortcuts import render
from django.http import HttpResponse
from django.template.loader import get_template
from xhtml2pdf import pisa
from .models import Item  # ajuste conforme seu app/modelo

def gerar_cotacao(request):
    if request.method == 'POST':
        selecionados_ids = request.POST.getlist('selecionados')
        dados = []

        itens = Item.objects.filter(id__in=selecionados_ids)

        for item in itens:
            quantidade = int(request.POST.get(f'quantidade_{item.id}', 1))
            valor_raw = request.POST.get(f'valor_{item.id}', '').replace(',', '.')
            valor_unitario = float(valor_raw) if valor_raw else 0.0
            total = quantidade * valor_unitario

            dados.append({
                'item': item,
                'quantidade': quantidade,
                'valor_unitario': valor_unitario,
                'total': total,
            })

        context = {
            'dados': dados,
        }

        template = get_template('cotacao/pdf_cotacao.html')
        html = template.render(context)
        response = HttpResponse(content_type='application/pdf')
        response['Content-Disposition'] = 'inline; filename="cotacao.pdf"'

        pisa_status = pisa.CreatePDF(html, dest=response)
        if pisa_status.err:
            return HttpResponse('Erro ao gerar PDF', status=500)
        return response

    # GET (exibe os itens)
    itens = Item.objects.all()
    return render(request, 'cotacao/gerar_cotacao.html', {'itens': itens})
###***********************************************************************************************************************

from django.shortcuts import render
from django.db.models import Sum, F, ExpressionWrapper, FloatField
from datetime import datetime
from .models import LancamentoDiario

def livro_diario(request):
    lancamentos = LancamentoDiario.objects.all().order_by('data')
    saldo = 0
    resultados_mensais = {}
    total_semestre1 = 0
    total_semestre2 = 0

    # Atualiza saldo de cada lançamento
    for lanc in lancamentos:
        saldo += lanc.recebimento or 0
        saldo -= lanc.pagamento or 0
        lanc.saldo = saldo

    # Dicionário para nomes dos meses em português
    meses_pt = {
        1: 'Janeiro',
        2: 'Fevereiro',
        3: 'Março',
        4: 'Abril',
        5: 'Maio',
        6: 'Junho',
        7: 'Julho',
        8: 'Agosto',
        9: 'Setembro',
        10: 'Outubro',
        11: 'Novembro',
        12: 'Dezembro'
    }

    # Cálculo mês a mês
    for mes in range(1, 13):
        recebimentos = lancamentos.filter(data__month=mes).aggregate(total=Sum('recebimento'))['total'] or 0
        pagamentos = lancamentos.filter(data__month=mes).aggregate(total=Sum('pagamento'))['total'] or 0
        total = recebimentos - pagamentos
        nome_mes = meses_pt[mes]
        resultados_mensais[nome_mes] = total

    # Cálculo semestral e anual
    total_semestre1 = sum([v for k, v in resultados_mensais.items() if k in ['Janeiro', 'Fevereiro', 'Março', 'Abril', 'Maio', 'Junho']])
    total_semestre2 = sum([v for k, v in resultados_mensais.items() if k in ['Julho', 'Agosto', 'Setembro', 'Outubro', 'Novembro', 'Dezembro']])
    total_ano = total_semestre1 + total_semestre2

    context = {
        'lancamentos': lancamentos,
        'saldo_hoje': saldo,
        'resultados_mensais': resultados_mensais,
        'resultado_semestre1': total_semestre1,
        'resultado_semestre2': total_semestre2,
        'resultado_ano': total_ano,
    }
    return render(request, 'contabilidade/livro_diario/listar.html', context)
###***********************************************************************************************************************

def adicionar_lancamento_diario(request):
    if request.method == "POST":
        try:
            data = request.POST.get('data')
            historico = request.POST.get('historico')
            recebimento = float(request.POST.get('recebimento') or 0)
            pagamento = float(request.POST.get('pagamento') or 0)

            LancamentoDiario.objects.create(
                data=data,
                historico=historico,
                recebimento=recebimento,
                pagamento=pagamento
            )
            messages.success(request, "✅ Lançamento adicionado com sucesso!")
        except Exception as e:
            messages.error(request, f"❌ Erro ao adicionar lançamento: {e}")

        return redirect('livro_diario')

    return render(request, 'contabilidade/livro_diario/adicionar.html')
###***********************************************************************************************************************

def excluir_lancamento_diario(request, id):
    lancamento = get_object_or_404(LancamentoDiario, id=id)
    if request.method == "POST":
        lancamento.delete()
        messages.success(request, "✅ Lançamento excluído com sucesso!")
        return redirect('livro_diario')
    return render(request, 'contabilidade/livro_diario/excluir.html', {'lancamento': lancamento})
###***********************************************************************************************************************

import csv
import pandas as pd
from django.http import HttpResponse
from xhtml2pdf import pisa
from django.template.loader import render_to_string

# Exportar PDF
def exportar_livro_diario_pdf(request):
    lancamentos = LancamentoDiario.objects.all().order_by('data')
    saldo = 0
    for lanc in lancamentos:
        saldo += lanc.recebimento or 0
        saldo -= lanc.pagamento or 0
        lanc.saldo = saldo

    html = render_to_string('contabilidade/livro_diario/livro_diario_pdf.html', {'lancamentos': lancamentos, 'saldo_hoje': saldo})
    response = HttpResponse(content_type='application/pdf')
    response['Content-Disposition'] = 'attachment; filename="livro_diario.pdf"'
    pisa.CreatePDF(html, dest=response)
    return response

# Exportar CSV
def exportar_livro_diario_csv(request):
    lancamentos = LancamentoDiario.objects.all().order_by('data')
    response = HttpResponse(content_type='text/csv')
    response['Content-Disposition'] = 'attachment; filename="livro_diario.csv"'
    writer = csv.writer(response)
    writer.writerow(['Data', 'Histórico', 'Recebimento', 'Pagamento', 'Saldo'])

    saldo = 0
    for lanc in lancamentos:
        saldo += lanc.recebimento or 0
        saldo -= lanc.pagamento or 0
        writer.writerow([lanc.data, lanc.historico, lanc.recebimento, lanc.pagamento, saldo])

    return response
###***********************************************************************************************************************

# Exportar Excel
def exportar_livro_diario_excel(request):
    lancamentos = LancamentoDiario.objects.all().order_by('data')
    data = []
    saldo = 0
    for lanc in lancamentos:
        saldo += lanc.recebimento or 0
        saldo -= lanc.pagamento or 0
        data.append({
            'Data': lanc.data,
            'Histórico': lanc.historico,
            'Recebimento': lanc.recebimento,
            'Pagamento': lanc.pagamento,
            'Saldo': saldo
        })

    df = pd.DataFrame(data)
    response = HttpResponse(content_type='application/vnd.ms-excel')
    response['Content-Disposition'] = 'attachment; filename="livro_diario.xlsx"'
    df.to_excel(response, index=False)
    return response
###***********************************************************************************************************************

from django.http import JsonResponse
from django.shortcuts import get_object_or_404
from .models import Receita, EscolaPdde


def get_info_receita_despesa(request, escola_id):
    try:
        escola = get_object_or_404(EscolaPdde, id=escola_id)
        dados = Receita.objects.filter(escola=escola).aggregate(
            saldo_anterior_custeio=Sum("saldo_anterior_custeio") or 0,
            saldo_anterior_capital=Sum("saldo_anterior_capital") or 0,
            valor_creditado_custeio=Sum("valor_creditado_custeio") or 0,
            valor_creditado_capital=Sum("valor_creditado_capital") or 0,
            recursos_proprios_custeio=Sum("recursos_proprios_custeio") or 0,
            recursos_proprios_capital=Sum("recursos_proprios_capital") or 0,
            rendimento_aplicacao_custeio=Sum("rendimento_aplicacao_custeio") or 0,
            rendimento_aplicacao_capital=Sum("rendimento_aplicacao_capital") or 0,
            valor_despesa_realizada_custeio=Sum("valor_despesa_realizada_custeio") or 0,
            valor_despesa_realizada_capital=Sum("valor_despesa_realizada_capital") or 0,
        )

        saldo_anterior = (
            dados["saldo_anterior_custeio"] + dados["saldo_anterior_capital"] +
            dados["valor_creditado_custeio"] + dados["valor_creditado_capital"] +
            dados["recursos_proprios_custeio"] + dados["recursos_proprios_capital"]
        )
        rendimentos = dados["rendimento_aplicacao_custeio"] + dados["rendimento_aplicacao_capital"]
        receita_total = saldo_anterior + rendimentos
        despesa_total = dados["valor_despesa_realizada_custeio"] + dados["valor_despesa_realizada_capital"]
        superavit_deficit = receita_total - despesa_total

        return JsonResponse({
            "saldo_anterior": float(saldo_anterior),
            "rendimentos": float(rendimentos),
            "receita_total": float(receita_total),
            "despesa_total": float(despesa_total),
            "superavit_deficit": float(superavit_deficit),
        })

    except Exception as e:
        return JsonResponse({"error": f"Erro ao buscar informações: {str(e)}"}, status=500)
###***********************************************************************************************************************

def get_receita_info(request, escola_id):
    try:
        escola = get_object_or_404(EscolaPdde, id=escola_id)
        receita = Receita.objects.filter(escola=escola).first()

        if not receita:
            return JsonResponse({"error": "Nenhuma receita encontrada."})

        return JsonResponse({
            "conselho": escola.nome_conselho or "",
            "cnpj": escola.cnpj or "",
            "programas": [receita.programa.nome],
            "rendimentos": float(receita.rendimento_aplicacao_custeio) + float(receita.rendimento_aplicacao_capital),
            "saldo_anterior": float(receita.saldo_anterior_custeio) + float(receita.saldo_anterior_capital) +
                              float(receita.valor_creditado_custeio) + float(receita.valor_creditado_capital) +
                              float(receita.recursos_proprios_custeio) + float(receita.recursos_proprios_capital),
            "valor_total_receita_custeio": float(receita.valor_total_receita_custeio),
            "valor_total_receita_capital": float(receita.valor_total_receita_capital),
            "valor_despesa_realizada_custeio": float(receita.valor_despesa_realizada_custeio),
            "valor_despesa_realizada_capital": float(receita.valor_despesa_realizada_capital),
            "superavit_deficit": (float(receita.valor_total_receita_custeio) +
                                  float(receita.valor_total_receita_capital) -
                                  float(receita.valor_despesa_realizada_custeio) -
                                  float(receita.valor_despesa_realizada_capital))
        })

    except Exception as e:
        return JsonResponse({"error": f"Erro ao buscar informações: {str(e)}"}, status=500)
###***********************************************************************************************************************

from django.db.models import Sum
from django.http import JsonResponse
from django.shortcuts import get_object_or_404
from .models import Receita, EscolaPdde, Programa, SemedAppEscolaPddeProgramas


def get_info_completa_escola(request, escola_id):
    try:
        escola = get_object_or_404(EscolaPdde, id=escola_id)

        # Programas vinculados
        programas_ids = SemedAppEscolaPddeProgramas.objects.filter(escolapdde_id=escola_id).values_list('programa_id', flat=True)
        programas = Programa.objects.filter(id__in=programas_ids).values_list('nome', flat=True)

        # Consulta dados financeiros
        dados = Receita.objects.filter(escola=escola).aggregate(
            saldo_anterior_custeio=Sum("saldo_anterior_custeio") or 0,
            saldo_anterior_capital=Sum("saldo_anterior_capital") or 0,
            valor_creditado_custeio=Sum("valor_creditado_custeio") or 0,
            valor_creditado_capital=Sum("valor_creditado_capital") or 0,
            recursos_proprios_custeio=Sum("recursos_proprios_custeio") or 0,
            recursos_proprios_capital=Sum("recursos_proprios_capital") or 0,
            rendimento_aplicacao_custeio=Sum("rendimento_aplicacao_custeio") or 0,
            rendimento_aplicacao_capital=Sum("rendimento_aplicacao_capital") or 0,
            valor_despesa_realizada_custeio=Sum("valor_despesa_realizada_custeio") or 0,
            valor_despesa_realizada_capital=Sum("valor_despesa_realizada_capital") or 0,
        )

        saldo_anterior = (
            dados["saldo_anterior_custeio"] + dados["saldo_anterior_capital"] +
            dados["valor_creditado_custeio"] + dados["valor_creditado_capital"] +
            dados["recursos_proprios_custeio"] + dados["recursos_proprios_capital"]
        )
        rendimentos = dados["rendimento_aplicacao_custeio"] + dados["rendimento_aplicacao_capital"]
        receita_total = saldo_anterior + rendimentos
        despesa_total = dados["valor_despesa_realizada_custeio"] + dados["valor_despesa_realizada_capital"]
        superavit_deficit = receita_total - despesa_total

        return JsonResponse({
            # Dados financeiros
            "saldo_anterior": float(saldo_anterior),
            "rendimentos": float(rendimentos),
            "receita_total": float(receita_total),
            "despesa_total": float(despesa_total),
            "superavit_deficit": float(superavit_deficit),

            # Dados da escola
            "conselho": escola.nome_conselho or "",
            "cnpj": escola.cnpj or "",
            "endereco": escola.endereco or "",
            "uf": escola.uf or "",
            "programas": list(programas),
        })

    except Exception as e:
        return JsonResponse({"error": f"Erro ao buscar informações: {str(e)}"}, status=500)
###***********************************************************************************************************************

from django.db.models import Sum
from django.http import JsonResponse
from django.shortcuts import render
from django.utils import timezone
from .models import LancamentoBancario, Receita

# Consolidado Contas Bancárias
def get_consolidado_contas(request):
    try:
        hoje = timezone.now()
        inicio_mes = hoje.replace(day=1)

        saldo_total = LancamentoBancario.objects.aggregate(total=Sum('valor'))['total'] or 0
        despesas_mes = LancamentoBancario.objects.filter(tipo='debito', data__gte=inicio_mes).aggregate(total=Sum('valor'))['total'] or 0
        receitas_mes = LancamentoBancario.objects.filter(tipo='credito', data__gte=inicio_mes).aggregate(total=Sum('valor'))['total'] or 0

        detalhes = LancamentoBancario.objects.all().values(
            'id', 'categoria', 'tipo', 'valor', 'data'
        )

        dados = {
            "saldo_total": round(saldo_total, 2),
            "despesas_mes": round(despesas_mes, 2),
            "receitas_mes": round(receitas_mes, 2),
            "detalhes_contas": list(detalhes)
        }
        return JsonResponse(dados)
    
    except Exception as e:
        return JsonResponse({"error": f"Erro ao buscar dados das contas: {str(e)}"}, status=500)
###***********************************************************************************************************************

# Consolidado Receita PDDE
def get_consolidado_receita(request):
    try:
        dados = Receita.objects.aggregate(
            total_receita_custeio=Sum('valor_total_receita_custeio'),
            total_receita_capital=Sum('valor_total_receita_capital'),
            total_despesa_custeio=Sum('valor_despesa_realizada_custeio'),
            total_despesa_capital=Sum('valor_despesa_realizada_capital')
        )

        receita_custeio = dados['total_receita_custeio'] or 0
        receita_capital = dados['total_receita_capital'] or 0
        despesa_custeio = dados['total_despesa_custeio'] or 0
        despesa_capital = dados['total_despesa_capital'] or 0

        total_entradas = receita_custeio + receita_capital
        total_despesas = despesa_custeio + despesa_capital
        saldo_final = total_entradas - total_despesas

        detalhes = Receita.objects.values(
            'programa', 'data_inicio', 'data_fim',
            'valor_total_receita_custeio', 'valor_total_receita_capital',
            'valor_despesa_realizada_custeio', 'valor_despesa_realizada_capital'
        )

        dados = {
            "saldo_total": round(saldo_final, 2),
            "despesas_total": round(total_despesas, 2),
            "receitas_total": round(total_entradas, 2),
            "detalhes_receita": list(detalhes)
        }

        return JsonResponse(dados)

    except Exception as e:
        return JsonResponse({"error": f"Erro ao buscar dados: {str(e)}"}, status=500)
###***********************************************************************************************************************

# Síntese Receita PDDE
from django.http import JsonResponse
from django.db.models import Sum
from .models import Receita

def get_sintese_receita(request):
    try:
        dados = Receita.objects.aggregate(
            saldo_reprogramado_ea_custeio=Sum('saldo_anterior_custeio') or 0,
            saldo_reprogramado_ea_capital=Sum('saldo_anterior_capital') or 0,
            valor_creditado_fnde_ee_custeio=Sum('valor_creditado_custeio') or 0,
            valor_creditado_fnde_ee_capital=Sum('valor_creditado_capital') or 0,
            recursos_proprios_custeio=Sum('recursos_proprios_custeio') or 0,
            recursos_proprios_capital=Sum('recursos_proprios_capital') or 0,
            rendimento_aplicacao_custeio=Sum('rendimento_aplicacao_custeio') or 0,
            rendimento_aplicacao_capital=Sum('rendimento_aplicacao_capital') or 0,
            devolucao_fnde_custeio=Sum('devolucao_fnde_custeio') or 0,
            devolucao_fnde_capital=Sum('devolucao_fnde_capital') or 0,
            valor_despesa_realizada_custeio=Sum('valor_despesa_realizada_custeio') or 0,
            valor_despesa_realizada_capital=Sum('valor_despesa_realizada_capital') or 0,
            saldo_reprogramado_es_custeio=Sum('saldo_reprogramar_custeio') or 0,
            saldo_reprogramado_es_capital=Sum('saldo_reprogramar_capital') or 0,
            saldo_devolvido_custeio=Sum('saldo_devolvido_custeio') or 0,
            saldo_devolvido_capital=Sum('saldo_devolvido_capital') or 0,
            escolas_atendidas=Sum('escolas_atendidas') or 0,
        )

        total_receita_custeio = (
            (dados['saldo_reprogramado_ea_custeio'] or 0) + 
            (dados['valor_creditado_fnde_ee_custeio'] or 0) + 
            (dados['recursos_proprios_custeio'] or 0) + 
            (dados['rendimento_aplicacao_custeio'] or 0)
        )
        total_receita_capital = (
            (dados['saldo_reprogramado_ea_capital'] or 0) + 
            (dados['valor_creditado_fnde_ee_capital'] or 0) + 
            (dados['recursos_proprios_capital'] or 0) + 
            (dados['rendimento_aplicacao_capital'] or 0)
        )
        total_receita = total_receita_custeio + total_receita_capital

        total_despesa = (dados['valor_despesa_realizada_custeio'] or 0) + (dados['valor_despesa_realizada_capital'] or 0)

        saldo_reprogramado = (dados['saldo_reprogramado_es_custeio'] or 0) + (dados['saldo_reprogramado_es_capital'] or 0)

        percentual_execucao = (total_despesa / total_receita * 100) if total_receita > 0 else 0

        recursos_proprios = (dados['recursos_proprios_custeio'] or 0) + (dados['recursos_proprios_capital'] or 0)
        rendimentos_aplicacao = (dados['rendimento_aplicacao_custeio'] or 0) + (dados['rendimento_aplicacao_capital'] or 0)
        devolucoes_fnde = (dados['devolucao_fnde_custeio'] or 0) + (dados['devolucao_fnde_capital'] or 0)
        escolas_atendidas = dados['escolas_atendidas'] or 0
        

        response = {
            "sintese": {
                "saldo_reprogramado_ea": {
                    "custeio": float(dados['saldo_reprogramado_ea_custeio'] or 0),
                    "capital": float(dados['saldo_reprogramado_ea_capital'] or 0),
                },
                "valor_creditado_fnde_ee": {
                    "custeio": float(dados['valor_creditado_fnde_ee_custeio'] or 0),
                    "capital": float(dados['valor_creditado_fnde_ee_capital'] or 0),
                },
                "recursos_proprios": {
                    "custeio": float(dados['recursos_proprios_custeio'] or 0),
                    "capital": float(dados['recursos_proprios_capital'] or 0),
                },
                "rendimento_aplicacao": {
                    "custeio": float(dados['rendimento_aplicacao_custeio'] or 0),
                    "capital": float(dados['rendimento_aplicacao_capital'] or 0),
                },
                "devolucao_fnde": {
                    "custeio": float(dados['devolucao_fnde_custeio'] or 0),
                    "capital": float(dados['devolucao_fnde_capital'] or 0),
                },
                "valor_total_receita": {
                    "custeio": float(total_receita_custeio),
                    "capital": float(total_receita_capital),
                },
                "valor_despesa_realizada": {
                    "custeio": float(dados['valor_despesa_realizada_custeio'] or 0),
                    "capital": float(dados['valor_despesa_realizada_capital'] or 0),
                },
                "saldo_reprogramado_es": {
                    "custeio": float(dados['saldo_reprogramado_es_custeio'] or 0),
                    "capital": float(dados['saldo_reprogramado_es_capital'] or 0),
                },
                "saldo_devolvido": {
                    "custeio": float(dados['saldo_devolvido_custeio'] or 0),
                    "capital": float(dados['saldo_devolvido_capital'] or 0),
                },
                "escolas_atendidas": int(escolas_atendidas),
                # 👇 Mantendo para Dashboard
                "resumo_cards": {
                    "receita_total": float(total_receita),
                    "despesa_total": float(total_despesa),
                    "saldo_reprogramado": float(saldo_reprogramado),
                    "percentual_execucao": round(percentual_execucao, 2),
                    "recursos_proprios": float(recursos_proprios),
                    "rendimentos_aplicacao": float(rendimentos_aplicacao),
                    "devolucoes_fnde": float(devolucoes_fnde),
                    "escolas_atendidas": int(escolas_atendidas),
                }
            }
        }
        return JsonResponse(response)

    except Exception as e:
        return JsonResponse({"error": f"Erro ao buscar síntese: {str(e)}"}, status=500)
###***********************************************************************************************************************

from django.shortcuts import render
from django.db.models import Sum
from django.http import JsonResponse
from .models import Receita, EscolaPdde, ConselhoMembro, Certidao


def sintese_pdde(request):
    # Filtros recebidos por GET
    filtro_escola = request.GET.get("escola")
    filtro_ano = request.GET.get("ano")
    filtro_programa = request.GET.get("programa")

    receitas = Receita.objects.all()

    if filtro_escola:
        receitas = receitas.filter(escola_id=filtro_escola)
    if filtro_ano:
        receitas = receitas.filter(data_inicio__year=filtro_ano)
    if filtro_programa:
        receitas = receitas.filter(programa=filtro_programa)

    # Dados agregados
    receita_total = receitas.aggregate(
        custeio=Sum('valor_total_receita_custeio') or 0,
        capital=Sum('valor_total_receita_capital') or 0
    )
    despesa_total = receitas.aggregate(
        custeio=Sum('valor_despesa_realizada_custeio') or 0,
        capital=Sum('valor_despesa_realizada_capital') or 0
    )
    saldo_reprogramado = receitas.aggregate(
        custeio=Sum('saldo_reprogramar_custeio') or 0,
        capital=Sum('saldo_reprogramar_capital') or 0
    )
    recursos_proprios = receitas.aggregate(
        custeio=Sum('recursos_proprios_custeio') or 0,
        capital=Sum('recursos_proprios_capital') or 0
    )
    rendimentos_aplicacao = receitas.aggregate(
        custeio=Sum('rendimento_aplicacao_custeio') or 0,
        capital=Sum('rendimento_aplicacao_capital') or 0
    )
    devolucoes_fnde = receitas.aggregate(
        custeio=Sum('devolucao_fnde_custeio') or 0,
        capital=Sum('devolucao_fnde_capital') or 0
    )
    escolas_atendidas = receitas.values('escola').distinct().count()

    total_receita = receita_total["custeio"] + receita_total["capital"]
    total_despesa = despesa_total["custeio"] + despesa_total["capital"]
    total_saldo_reprogramado = saldo_reprogramado["custeio"] + saldo_reprogramado["capital"]

    percentual_execucao = (total_despesa / total_receita * 100) if total_receita > 0 else 0

    resumo_cards = {
        "receita_total": float(total_receita),
        "despesa_total": float(total_despesa),
        "saldo_reprogramado": float(total_saldo_reprogramado),
        "percentual_execucao": round(percentual_execucao, 2),
        "recursos_proprios": float(recursos_proprios["custeio"] + recursos_proprios["capital"]),
        "rendimentos_aplicacao": float(rendimentos_aplicacao["custeio"] + rendimentos_aplicacao["capital"]),
        "devolucoes_fnde": float(devolucoes_fnde["custeio"] + devolucoes_fnde["capital"]),
        "escolas_atendidas": escolas_atendidas,
        "total_escolas": EscolaPdde.objects.count(),
        "total_membros": ConselhoMembro.objects.count(),
        "total_certidoes": Certidao.objects.count(),
    }

    # Se for uma requisição via fetch(), retorna JSON
    if request.headers.get('x-requested-with') == 'XMLHttpRequest':
        return JsonResponse({
            "sintese": {
                "resumo_cards": resumo_cards
            }
        })

    # Se for acesso normal via navegador, renderiza a página
    context = {
        "escolas": EscolaPdde.objects.all(),
        "anos": Receita.objects.dates("data_inicio", "year", order="DESC").distinct(),
        "programas": Receita.objects.values_list("programa", flat=True).distinct(),

        "total_escolas": resumo_cards["total_escolas"],
        "total_membros": resumo_cards["total_membros"],
        "total_certidoes": resumo_cards["total_certidoes"],

        "dados": resumo_cards,
        "filtros_aplicados": {
            "escola": filtro_escola,
            "ano": filtro_ano,
            "programa": filtro_programa,
        }
    }

    return render(request, "pdde/sintese_pdde.html", context)


###***********************************************************************************************************************

def cadastrar_info_pdde(request):
    return render(request, 'pdde/cadastrar_info_pdde.html')

###***********************************************************************************************************************

from django.template.loader import render_to_string
from django.http import HttpResponse
from weasyprint import HTML
import json

def gerar_pdf_sintese(request):
    try:
        # Busca os dados da API interna
        from .views import get_sintese_receita
        response = get_sintese_receita(request)

        # Converte o JSON para dicionário
        data = json.loads(response.content.decode('utf-8'))

        # Renderiza o template para HTML com os dados
        html_string = render_to_string('contabilidade/sintese_pdf.html', {'dados': data})

        # Gera o PDF
        html = HTML(string=html_string)
        pdf = html.write_pdf()

        # Retorna o PDF como resposta HTTP
        response = HttpResponse(pdf, content_type='application/pdf')
        response['Content-Disposition'] = 'inline; filename="sintese_execucao_pdde.pdf"'
        return response

    except Exception as e:
        return HttpResponse(f"Erro ao gerar PDF: {str(e)}", status=500)
###***********************************************************************************************************************

from django.shortcuts import render
from django.http import FileResponse, Http404
import os
from django.conf import settings

def lista_manuais(request):
    manuais = [
        {"nome": "Constituição da UEX 2024", "arquivo": "CONSTITUIÇÃO DA UEX 2024.pdf"},
        {"nome": "Guia de Execução dos Recursos do PDDE", "arquivo": "GUIADEEXECUODOSRECURSOSDOPDDEv4FINAL novo.pdf"},
        {"nome": "Novidades Resolução 15 PDDE", "arquivo": "Novidades_Resoluo15PDDE.pdf"},
        {"nome": "Resolução nº 6/2023 - Utilização de Saldos Financeiros", "arquivo": "RESOLUÇÃO 6, de 4-5-23 Autoriza utilização de saldos financeiros.pdf"},
        {"nome": "Resolução nº 15/2021", "arquivo": "RESOLUÇÃO Nº 15, DE 16 DE SETEMBRO DE 2021.pdf"},
    ]
    return render(request, 'contabilidade/manuais.html', {'manuais': manuais})
###***********************************************************************************************************************

def download_manual(request, arquivo):
    caminho = os.path.join(settings.MEDIA_ROOT, 'download_manual', arquivo)
    if os.path.exists(caminho):
        return FileResponse(open(caminho, 'rb'), as_attachment=True, filename=arquivo)
    else:
        raise Http404("Arquivo não encontrado.")

###***********************************************************************************************************************

from django.http import JsonResponse
from .models import PlanoGestaoEscolar, MotivoIndeferimento

def salvar_motivo_indeferimento(request):
    if request.method == 'POST':
        plano_id = request.POST.get('plano_id')
        motivo = request.POST.get('motivo')
        prazo = request.POST.get('prazo')
        parecer_direcao = request.POST.get('parecer_direcao')
        orientacoes = request.POST.get('orientacoes')

        if not plano_id or not motivo or not prazo or not parecer_direcao:
            return JsonResponse({'success': False, 'error': 'Todos os campos obrigatórios devem ser preenchidos.'})

        try:
            plano = PlanoGestaoEscolar.objects.get(id=plano_id)
            indeferimento = MotivoIndeferimento.objects.create(
                plano=plano,
                motivo=motivo,
                prazo_reenvio=prazo,
                parecer_direcao=parecer_direcao,
                orientacoes=orientacoes
            )

            if request.FILES.get('anexo_parecer'):
                indeferimento.anexo_parecer = request.FILES['anexo_parecer']
                indeferimento.save()

            return JsonResponse({'success': True})
        except PlanoGestaoEscolar.DoesNotExist:
            return JsonResponse({'success': False, 'error': 'Plano de gestão não encontrado.'})
        except Exception as e:
            return JsonResponse({'success': False, 'error': str(e)})
    else:
        return JsonResponse({'success': False, 'error': 'Método inválido.'})
###***********************************************************************************************************************

from django.templatetags.static import static
from django.shortcuts import get_object_or_404
from weasyprint import HTML, CSS
from django.template.loader import render_to_string
from django.http import HttpResponse
from .models import PlanoGestaoEscolar, MotivoIndeferimento

def gerar_pdf_indeferimento(request, plano_id):
    plano = get_object_or_404(PlanoGestaoEscolar, id=plano_id)
    indeferimento = MotivoIndeferimento.objects.filter(plano=plano).first()

    # URLs absolutas das imagens
    logo_url = request.build_absolute_uri(static('assets/dist/img/logo_prefeitura.png'))
    footer_logo_url = request.build_absolute_uri(static('assets/dist/img/logoFooter.png'))

    # Renderização do template
    html_string = render_to_string(
        'setor_pedagogico/pge_indeferimento_pdf.html',
        {
            'plano': plano,
            'indeferimento': indeferimento,
            'logo_url': logo_url,
            'footer_logo_url': footer_logo_url,
        }
    )

    # Geração do PDF com base_url e página A4 configurada
    pdf = HTML(string=html_string, base_url=request.build_absolute_uri()).write_pdf()

    # Retorno do PDF na resposta
    response = HttpResponse(pdf, content_type='application/pdf')
    response['Content-Disposition'] = f'inline; filename="indeferimento_plano_{plano.id}.pdf"'
    return response
###***********************************************************************************************************************

from django.core.mail import EmailMessage

def enviar_email_indeferimento(request):
    if request.method == "POST":
        plano_id = request.POST.get('plano_id')
        email_destino = request.POST.get('email_destino')
        mensagem = request.POST.get('mensagem')

        plano = get_object_or_404(PlanoGestaoEscolar, id=plano_id)
        indeferimento = MotivoIndeferimento.objects.filter(plano=plano).first()

        # Gera o PDF
        from .views import gerar_pdf_indeferimento
        response = gerar_pdf_indeferimento(request, plano_id)
        pdf_content = response.content

        # Enviar e-mail
        email = EmailMessage(
            subject=f"Indeferimento do Plano de Gestão - {plano.unidade_ensino}",
            body=mensagem,
            to=[email_destino],
        )
        email.attach(f"Indeferimento_{plano.unidade_ensino}.pdf", pdf_content, 'application/pdf')
        email.send()

        return JsonResponse({'success': True})
    return JsonResponse({'error': 'Requisição inválida'})
###***********************************************************************************************************************
###***********************************************************************************************************************
###***********************************************************************************************************************
###**********************************RELATORIOS SAMACC********************************************************************
###***********************************************************************************************************************
###***********************************************************************************************************************

from django.http import HttpResponse
from django.template.loader import get_template
from django.templatetags.static import static
from xhtml2pdf import pisa
import os
from .models import CadastroEI

def gerar_pdf_mapeamento(request):
    nome_escola = request.GET.get('nome_escola')
    turma = request.GET.get('turma')
    ano = request.GET.get('ano')
    modalidade = request.GET.get('modalidade')

    alunos = CadastroEI.objects.all()
    if nome_escola:
        alunos = alunos.filter(unidade_ensino__icontains=nome_escola)
    if turma:
        alunos = alunos.filter(turma__nome__icontains=turma)
    if ano:
        alunos = alunos.filter(ano=ano)
    if modalidade:
        alunos = alunos.filter(modalidade__icontains=modalidade)

    total_escolas = alunos.values('unidade_ensino').distinct().count()
    total_alunos = alunos.count()
    total_educadores = alunos.values('professor').distinct().exclude(professor=None).count()

    # --- Formulário 1: Linguagem (Q11 a Q20)
    campos_linguagem = [f"questao_linguagem_{i}" for i in range(11, 21)]
    formulario_1 = 0
    for aluno in alunos:
        respostas = [getattr(aluno, campo) for campo in campos_linguagem]
        if any(resp not in [None, '', 'FALTOU'] for resp in respostas):
            formulario_1 += 1

    # --- Formulário 2: Matemática (Q1 a Q10)
    campos_matematica = [f"questao_matematica_{i}" for i in range(1, 11)]
    formulario_2 = 0
    for aluno in alunos:
        respostas = [getattr(aluno, campo) for campo in campos_matematica]
        if any(resp not in [None, '', 'FALTOU'] for resp in respostas):
            formulario_2 += 1

    context = {
        'total_escolas': total_escolas,
        'total_alunos': total_alunos,
        'total_educadores': total_educadores,
        'formulario_1': formulario_1,
        'formulario_2': formulario_2,
        'cabecalho_path': request.build_absolute_uri(static('assets/dist/img/logocabecalho.png')),
        'rodape_path': request.build_absolute_uri(static('assets/dist/img/logorodape.png')),
    }

    template = get_template('relatorios/mapeamento_aplicacao.html')
    html = template.render(context)

    response = HttpResponse(content_type='application/pdf')
    response['Content-Disposition'] = 'inline; filename="mapeamento_aplicacao.pdf"'
    pisa_status = pisa.CreatePDF(html, dest=response)

    if pisa_status.err:
        return HttpResponse('Erro ao gerar PDF <pre>' + html + '</pre>')

    return response







###***********************************************************************************************************************

from django.template.loader import render_to_string
from django.http import HttpResponse
from weasyprint import HTML
import tempfile, os
from semedapp.models import CadastroEI, Turma
from django.db.models import Q

def gerar_pdf_detalhado(request):
    # Filtros da URL
    nome_escola = request.GET.get('nome_escola', '')
    turma = request.GET.get('turma', '')
    ano = request.GET.get('ano', '')
    modalidade = request.GET.get('modalidade', '')

    # Alunos avaliados
    alunos = CadastroEI.objects.filter(avaliado='SIM')
    if nome_escola:
        alunos = alunos.filter(unidade_ensino__icontains=nome_escola)
    if turma:
        alunos = alunos.filter(turma__nome__icontains=turma)
    if ano:
        alunos = alunos.filter(ano=ano)
    if modalidade:
        alunos = alunos.filter(modalidade__icontains=modalidade)

    total_alunos = alunos.count()
    total_turmas = Turma.objects.count()

    # Campos das questões
    campos_matematica = [f"questao_matematica_{i}" for i in range(1, 11)]
    campos_linguagem = [f"questao_linguagem_{i}" for i in range(12, 21) if i != 13]
    campos_total = campos_matematica + campos_linguagem

    # Identificar alunos com laudo
    filtro_laudo = Q()
    for campo in campos_total:
        filtro_laudo |= Q(**{campo: 'CRIANÇA COM LAUDO'})

    ids_com_laudo = alunos.filter(filtro_laudo).values_list('id_matricula', flat=True).distinct()
    total_laudo_rede = ids_com_laudo.count()

    # Alunos que devem de fato participar (sem laudo)
    alunos_validos = alunos.exclude(id_matricula__in=ids_com_laudo)

    # Contadores de faltantes por caderno
    faltantes_matematica = 0
    faltantes_linguagem = 0
    total_faltantes_geral = 0

    for aluno in alunos_validos:
        q_mat = [getattr(aluno, campo) for campo in campos_matematica]
        q_lin = [getattr(aluno, campo) for campo in campos_linguagem]

        if all(q in [None, '', 'FALTOU'] for q in q_mat):
            faltantes_matematica += 1
        if all(q in [None, '', 'FALTOU'] for q in q_lin):
            faltantes_linguagem += 1
        if all(q in [None, '', 'FALTOU'] for q in q_mat + q_lin):
            total_faltantes_geral += 1

    total_validos = alunos_validos.count()

    # Cálculo da participação correto
    def calcular_participacao(validos, faltantes):
        if validos <= 0:
            return "0%"
        participacao = ((validos - min(faltantes, validos)) / validos) * 100
        return f"{round(participacao)}%"



    participacao_matematica = calcular_participacao(total_validos, faltantes_matematica)
    participacao_linguagem = calcular_participacao(total_validos, faltantes_linguagem)

    # Dados finais para template
    resultado_por_unidade = [{
        'unidade': nome_escola or 'Todas',
        'total_nei': total_alunos,
        'total_turmas': total_turmas,
        'participacao_matematica': participacao_matematica,
        'participacao_linguagem': participacao_linguagem,
        'total_faltantes_matematica': faltantes_matematica,
        'total_faltantes_linguagem': faltantes_linguagem,
        'total_faltantes_geral': total_faltantes_geral,
        'total_laudo': total_laudo_rede,
        'modalidades_q11': {'CC': 0, 'SC': 0},
        'modalidades_q13': {'PSI': 0, 'PSII': 0, 'SSVS': 0, 'SCVS': 0, 'SA': 0, 'ALF': 0},
        'dados_matematica': {},  # pode preencher se desejar mostrar por questão
        'dados_linguagem': {},   # idem
    }]

    html_string = render_to_string('relatorios/relatorio_detalhado.html', {
        'resultado_por_unidade': resultado_por_unidade
    })

    # Geração do PDF
    fd, path = tempfile.mkstemp(suffix='.pdf')
    os.close(fd)
    try:
        HTML(string=html_string).write_pdf(path)
        with open(path, 'rb') as pdf:
            response = HttpResponse(pdf.read(), content_type='application/pdf')
            response['Content-Disposition'] = 'inline; filename="relatorio_detalhado.pdf"'
        return response
    finally:
        os.remove(path)

###***********************************************************************************************************************

from django.shortcuts import render
from semedapp.models import CadastroEI, Turma
from django.db.models import Q

def relatorio_saida_rede(request):
    nome_escola = request.GET.get('nome_escola', '')
    ano = request.GET.get('ano', '')
    modalidade = request.GET.get('modalidade', '')

    alunos = CadastroEI.objects.filter(avaliado='SIM')
    if nome_escola:
        alunos = alunos.filter(unidade_ensino__icontains=nome_escola)
    if ano:
        alunos = alunos.filter(ano=ano)
    if modalidade:
        alunos = alunos.filter(modalidade__icontains=modalidade)

    unidades = alunos.values_list('unidade_ensino', flat=True).distinct()
    resultado_por_unidade = []

    for unidade in unidades:
        alunos_nei = alunos.filter(unidade_ensino=unidade)
        total_nei = alunos_nei.count()

        dados_matematica = {}
        dados_linguagem = {}
        faltantes_matematica = 0
        faltantes_linguagem = 0
        total_faltantes_geral = 0

        # Modalidades
        modalidades_q11 = {label: alunos_nei.filter(questao_linguagem_11=label).values('id_matricula').distinct().count()
                           for label in ['CC', 'SC']}
        modalidades_q13 = {label: alunos_nei.filter(questao_linguagem_13=label).values('id_matricula').distinct().count()
                           for label in ['PSI', 'PSII', 'SSVS', 'SCVS', 'SA', 'ALF']}

        # Campos das questões
        campos_matematica = [f"questao_matematica_{i}" for i in range(1, 11)]
        campos_linguagem = [f"questao_linguagem_{i}" for i in range(12, 21) if i != 13]

        # --- Matemática
        for campo in campos_matematica:
            dados_matematica[int(campo.split('_')[-1])] = {
                'CERTO': alunos_nei.filter(**{campo: 'CERTO'}).count(),
                'ERRADO': alunos_nei.filter(**{campo: 'ERRADO'}).count(),
                'PARCIAL': alunos_nei.filter(**{campo: 'PARCIAL'}).count(),
                'BRANCO': alunos_nei.filter(**{campo: ''}).count(),
                'LAUDO': alunos_nei.filter(**{campo: 'CRIANÇA COM LAUDO'}).count(),
                'FALTOU': alunos_nei.filter(**{f"{campo}__isnull": True}).count(),
            }
            faltantes_matematica += dados_matematica[int(campo.split('_')[-1])]['FALTOU']

        # --- Linguagem
        for campo in campos_linguagem:
            dados_linguagem[int(campo.split('_')[-1])] = {
                'CERTO': alunos_nei.filter(**{campo: 'CERTO'}).count(),
                'ERRADO': alunos_nei.filter(**{campo: 'ERRADO'}).count(),
                'PARCIAL': alunos_nei.filter(**{campo: 'PARCIAL'}).count(),
                'BRANCO': alunos_nei.filter(**{campo: ''}).count(),
                'LAUDO': alunos_nei.filter(**{campo: 'CRIANÇA COM LAUDO'}).count(),
                'FALTOU': alunos_nei.filter(**{f"{campo}__isnull": True}).count(),
            }

        # --- Identificar alunos com LAUDO em linguagem
        filtro_laudo_linguagem = Q()
        for campo in campos_linguagem:
            filtro_laudo_linguagem |= Q(**{campo: 'CRIANÇA COM LAUDO'})

        ids_com_laudo = alunos_nei.filter(filtro_laudo_linguagem).values_list('id_matricula', flat=True).distinct()
        alunos_validos_linguagem = alunos_nei.exclude(id_matricula__in=ids_com_laudo)

        # --- Contar faltantes gerais e de linguagem
        for aluno in alunos_validos_linguagem:
            respostas_ling = [getattr(aluno, campo) for campo in campos_linguagem]
            respostas_mat = [getattr(aluno, campo) for campo in campos_matematica]

            if all(r in [None, '', 'FALTOU'] for r in respostas_ling):
                faltantes_linguagem += 1

            if all(r in [None, '', 'FALTOU'] for r in respostas_mat + respostas_ling):
                total_faltantes_geral += 1

        # --- Participações
        def calcular_participacao(total, faltantes):
            if total == 0:
                return "0%"
            return f"{round(((total - faltantes) / total) * 100)}%"

        participacao_matematica = calcular_participacao(total_nei, faltantes_matematica)
        participacao_linguagem = calcular_participacao(alunos_validos_linguagem.count(), faltantes_linguagem)

        resultado_por_unidade.append({
            'unidade': unidade,
            'total_nei': total_nei,
            'dados_matematica': dados_matematica,
            'dados_linguagem': dados_linguagem,
            'participacao_matematica': participacao_matematica,
            'participacao_linguagem': participacao_linguagem,
            'total_faltantes_geral': total_faltantes_geral,
            'modalidades_q11': modalidades_q11,
            'modalidades_q13': modalidades_q13,
        })

    contexto = {
        'resultado_por_unidade': resultado_por_unidade,
        'filtros': {
            'nome_escola': nome_escola,
            'ano': ano,
            'modalidade': modalidade,
        }
    }

    return render(request, 'relatorios/relatorio_saida_rede.html', contexto)

###***********************************************************************************************************************

from django.http import HttpResponse
from django.template.loader import get_template
from django.templatetags.static import static
from xhtml2pdf import pisa
from semedapp.models import CadastroEI

def gerar_pdf_saida_rede(request):
    nome_escola = request.GET.get('nome_escola', '')
    ano = request.GET.get('ano', '')
    modalidade = request.GET.get('modalidade', '')

    alunos = CadastroEI.objects.filter(avaliado='SIM')
    if nome_escola:
        alunos = alunos.filter(unidade_ensino__icontains=nome_escola)
    if ano:
        alunos = alunos.filter(ano=ano)
    if modalidade:
        alunos = alunos.filter(modalidade__icontains=modalidade)

    unidades = alunos.values_list('unidade_ensino', flat=True).distinct()
    resultado_por_unidade = []

    for unidade in unidades:
        alunos_nei = alunos.filter(unidade_ensino=unidade)
        total_nei = alunos_nei.count()

        campos_matematica = [f"questao_matematica_{i}" for i in range(1, 11)]
        campos_linguagem = [f"questao_linguagem_{i}" for i in range(12, 21) if i != 13]

        dados_matematica = {
            i: {
                'CERTO': alunos_nei.filter(**{campo: 'CERTO'}).count(),
                'ERRADO': alunos_nei.filter(**{campo: 'ERRADO'}).count(),
                'PARCIAL': alunos_nei.filter(**{campo: 'PARCIAL'}).count(),
                'BRANCO': alunos_nei.filter(**{campo: 'BRANCO'}).count(),
                'LAUDO': alunos_nei.filter(**{campo: 'CRIANÇA COM LAUDO'}).count(),
                'FALTOU': alunos_nei.filter(**{campo: 'FALTOU'}).count(),
            } for i, campo in enumerate(campos_matematica, 1)
        }

        dados_linguagem = {
            i: {
                'CERTO': alunos_nei.filter(**{campo: 'CERTO'}).count(),
                'ERRADO': alunos_nei.filter(**{campo: 'ERRADO'}).count(),
                'PARCIAL': alunos_nei.filter(**{campo: 'PARCIAL'}).count(),
                'BRANCO': alunos_nei.filter(**{campo: 'BRANCO'}).count(),
                'LAUDO': alunos_nei.filter(**{campo: 'CRIANÇA COM LAUDO'}).count(),
                'FALTOU': alunos_nei.filter(**{campo: 'FALTOU'}).count(),
            } for i, campo in enumerate(campos_linguagem, 12)
        }

        modalidades_q11 = {m: alunos_nei.filter(questao_linguagem_11=m).count() for m in ['CC', 'SC']}
        modalidades_q13 = {m: alunos_nei.filter(questao_linguagem_13=m).count() for m in ['PSI', 'PSII', 'SSVS', 'SCVS', 'SA', 'ALF']}

        faltantes_matematica = sum(
            1 for a in alunos_nei if all(getattr(a, campo) in [None, '', 'FALTOU'] for campo in campos_matematica)
        )
        faltantes_linguagem = sum(
            1 for a in alunos_nei if all(getattr(a, campo) in [None, '', 'FALTOU'] for campo in campos_linguagem)
        )
        total_faltantes_geral = sum(
            1 for a in alunos_nei if all(getattr(a, campo) in [None, '', 'FALTOU'] for campo in campos_matematica + campos_linguagem)
        )

        def calc_part(total, falt):
            return f"{round(((total - min(falt, total)) / total) * 100)}%" if total else "0%"

        resultado_por_unidade.append({
            'unidade': unidade,
            'total_nei': total_nei,
            'dados_matematica': dados_matematica,
            'dados_linguagem': dados_linguagem,
            'participacao_matematica': calc_part(total_nei, faltantes_matematica),
            'participacao_linguagem': calc_part(total_nei, faltantes_linguagem),
            'modalidades_q11': modalidades_q11,
            'modalidades_q13': modalidades_q13,
            'total_faltantes_geral': total_faltantes_geral,
        })

    context = {
        'resultado_por_unidade': resultado_por_unidade,
        'cabecalho_superior': request.build_absolute_uri(static('assets/dist/img/logocabecalho.png')),
        'cabecalho_inferior': request.build_absolute_uri(static('assets/dist/img/logorodape.png')),
    }

    template = get_template('relatorios/relatorio_saida_rede.html')
    html = template.render(context)

    response = HttpResponse(content_type='application/pdf')
    response['Content-Disposition'] = 'inline; filename="relatorio_saida_rede.pdf"'

    pisa_status = pisa.CreatePDF(html, dest=response)
    if pisa_status.err:
        return HttpResponse("Erro ao gerar PDF")
    return response


###***********************************************************************************************************************

from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
from django.http import HttpResponse
from .models import CadastroEI
import matplotlib.pyplot as plt
import base64


# 🔧 Função para gerar gráficos em base64
def gerar_grafico_base64(campo, categorias, alunos, titulo):
    escolas = alunos.values_list('unidade_ensino', flat=True).distinct()
    fig, ax = plt.subplots(figsize=(10, 6))

    width = 0.8 / len(escolas)  # largura de cada barra ajustada ao número de escolas
    x = range(len(categorias))

    for idx, escola in enumerate(escolas):
        dados = [alunos.filter(unidade_ensino=escola, **{campo: c}).count() for c in categorias]
        positions = [i + (idx * width) for i in x]
        bars = ax.bar(positions, dados, width=width, label=escola)
        ax.bar_label(bars, labels=[str(d) for d in dados], fontsize=8, padding=2)

    ax.set_title(titulo)
    ax.set_ylabel("Total")
    ax.set_xticks([i + width*(len(escolas)-1)/2 for i in x])
    ax.set_xticklabels(categorias, rotation=45)
    ax.legend(fontsize=6)
    plt.tight_layout()

    buffer = BytesIO()
    fig.savefig(buffer, format='png')
    plt.close(fig)
    return base64.b64encode(buffer.getvalue()).decode()


# 🔧 Função para adicionar slide com imagem (gráfico)
def adicionar_slide_com_imagem(prs, titulo, imagem_base64):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = titulo
    imagem_bytes = BytesIO(base64.b64decode(imagem_base64))
    slide.shapes.add_picture(imagem_bytes, Inches(0.5), Inches(1.5), height=Inches(4.5))


def exportar_pptx_saida_rede(request):
    nome_escola = request.GET.get('nome_escola', '')
    ano = request.GET.get('ano', '')
    modalidade = request.GET.get('modalidade', '')

    alunos = CadastroEI.objects.filter(avaliado='SIM')
    if nome_escola:
        alunos = alunos.filter(unidade_ensino__icontains=nome_escola)
    if ano:
        alunos = alunos.filter(ano=ano)
    if modalidade:
        alunos = alunos.filter(modalidade__icontains=modalidade)

    unidades = alunos.values_list('unidade_ensino', flat=True).distinct()
    prs = Presentation()

    # Slide de título
    title_slide_layout = prs.slide_layouts[0]
    prs.slides.add_slide(title_slide_layout).shapes.title.text = "RELATÓRIO SAMACC 2025"

    # Slide de Mapeamento
    slide_mapeamento = prs.slides.add_slide(prs.slide_layouts[5])
    slide_mapeamento.shapes.title.text = "Mapeamento de Aplicação"
    total_escolas = alunos.values('unidade_ensino').distinct().count()
    total_alunos = alunos.count()
    total_educadores = alunos.values('professor').distinct().exclude(professor=None).count()
    formulario_1 = alunos.exclude(questao_linguagem_11=None).exclude(questao_linguagem_11='').count()
    formulario_2 = alunos.exclude(questao_matematica_1=None).exclude(questao_matematica_1='').count()

    table = slide_mapeamento.shapes.add_table(5, 2, Inches(0.5), Inches(1.5), Inches(9), Inches(3)).table
    dados_mapeamento = [
        ("Escolas Atendidas", total_escolas),
        ("Total de Alunos", f"{total_alunos} alunos"),
        ("Educadores", total_educadores),
        ("Formulários - 1º Caderno (Linguagem)", formulario_1),
        ("Formulários - 2º Caderno (Matemática)", formulario_2),
    ]
    for i, (label, value) in enumerate(dados_mapeamento):
        table.cell(i, 0).text = label
        table.cell(i, 1).text = str(value)
        for paragraph in table.cell(i, 0).text_frame.paragraphs + table.cell(i, 1).text_frame.paragraphs:
            paragraph.font.size = Pt(14)

    # Tabelas por unidade (questões)
    def add_table(slide, title_text, dados, top):
        rows = len(dados) + 1
        cols = 7
        table_shape = slide.shapes.add_table(rows, cols, Inches(0.3), top, Inches(9), Inches(0.4 + 0.2 * rows))
        table = table_shape.table
        headers = ['Q', 'CERTO', 'ERRADO', 'PARCIAL', 'BRANCO', 'LAUDO', 'FALTOU']
        for j, header in enumerate(headers):
            table.cell(0, j).text = header
            table.cell(0, j).text_frame.paragraphs[0].font.bold = True
        for i, (q, resultado) in enumerate(dados.items(), start=1):
            table.cell(i, 0).text = str(q)
            table.cell(i, 1).text = str(resultado['CERTO'])
            table.cell(i, 2).text = str(resultado['ERRADO'])
            table.cell(i, 3).text = str(resultado['PARCIAL'])
            table.cell(i, 4).text = str(resultado['BRANCO'])
            table.cell(i, 5).text = str(resultado['LAUDO'])
            table.cell(i, 6).text = str(resultado['FALTOU'])

    # Loop das escolas (por unidade)
    for unidade in unidades:
        alunos_nei = alunos.filter(unidade_ensino=unidade)
        total_nei = alunos_nei.count()

        # Matemática
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"{unidade} - {total_nei} alunos"
        dados_matematica = {
            i: {
                'CERTO': alunos_nei.filter(**{f"questao_matematica_{i}": 'CERTO'}).count(),
                'ERRADO': alunos_nei.filter(**{f"questao_matematica_{i}": 'ERRADO'}).count(),
                'PARCIAL': alunos_nei.filter(**{f"questao_matematica_{i}": 'PARCIAL'}).count(),
                'BRANCO': alunos_nei.filter(**{f"questao_matematica_{i}": 'BRANCO'}).count(),
                'LAUDO': alunos_nei.filter(**{f"questao_matematica_{i}": 'CRIANÇA COM LAUDO'}).count(),
                'FALTOU': alunos_nei.filter(**{f"questao_matematica_{i}": 'FALTOU'}).count(),
            }
            for i in range(1, 11)
        }
        add_table(slide, "Matemática", dados_matematica, top=Inches(1.8))

        # Linguagem
        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        slide2.shapes.title.text = f"{unidade} - Linguagem"
        dados_linguagem = {
            i - 11: {
                'CERTO': alunos_nei.filter(**{f"questao_linguagem_{i}": 'CERTO'}).count(),
                'ERRADO': alunos_nei.filter(**{f"questao_linguagem_{i}": 'ERRADO'}).count(),
                'PARCIAL': alunos_nei.filter(**{f"questao_linguagem_{i}": 'PARCIAL'}).count(),
                'BRANCO': alunos_nei.filter(**{f"questao_linguagem_{i}": 'BRANCO'}).count(),
                'LAUDO': alunos_nei.filter(**{f"questao_linguagem_{i}": 'CRIANÇA COM LAUDO'}).count(),
                'FALTOU': alunos_nei.filter(**{f"questao_linguagem_{i}": 'FALTOU'}).count(),
            }
            for i in [i for i in range(12, 21) if i != 13]
        }
        add_table(slide2, "Linguagem", dados_linguagem, top=Inches(1.8))

        # Modalidades (por unidade)
        slide3 = prs.slides.add_slide(prs.slide_layouts[5])
        slide3.shapes.title.text = f"{unidade} - Modalidades"

        q11 = {'CC': alunos_nei.filter(questao_linguagem_11='CC').count(),
               'SC': alunos_nei.filter(questao_linguagem_11='SC').count()}
        q13 = {m: alunos_nei.filter(questao_linguagem_13=m).count()
               for m in ['PSI', 'PSII', 'SSVS', 'SCVS', 'SA', 'ALF']}

        table_q11 = slide3.shapes.add_table(2, 2, Inches(0.5), Inches(1.8), Inches(4), Inches(1)).table
        table_q11.cell(0, 0).text = "CC"
        table_q11.cell(0, 1).text = "SC"
        table_q11.cell(1, 0).text = str(q11['CC'])
        table_q11.cell(1, 1).text = str(q11['SC'])

        table_q13 = slide3.shapes.add_table(2, 6, Inches(0.5), Inches(2.9), Inches(9), Inches(1.2)).table
        for idx, key in enumerate(q13.keys()):
            table_q13.cell(0, idx).text = key
            table_q13.cell(1, idx).text = str(q13[key])

    # ✅ SLIDE FINAL 1: RELATÓRIO GERAL (compactado)
    slide_final = prs.slides.add_slide(prs.slide_layouts[5])
    slide_final.shapes.title.text = "Relatório Geral - Efetivados x Pretendidos"

    escolas = alunos.values_list('unidade_ensino', flat=True).distinct()
    num_escolas = len(escolas)

    altura_por_linha = 0.25
    altura_total = altura_por_linha * (num_escolas + 2)
    table_final = slide_final.shapes.add_table(num_escolas + 2, 5, Inches(0.3), Inches(1.9), Inches(9), Inches(altura_total)).table

    headers = ["Escola", "Pretendidos C1", "Efetivados C1", "Pretendidos C2", "Efetivados C2"]
    for j, h in enumerate(headers):
        cell = table_final.cell(0, j)
        cell.text = h
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        cell.text_frame.paragraphs[0].font.bold = True

    total_c1_pret = total_c1_efet = total_c2_pret = total_c2_efet = 0
    for i, escola in enumerate(escolas, start=1):
        registros = alunos.filter(unidade_ensino=escola)
        # Caderno 1 = Matemática (Q1 a Q10)
        campos_caderno1 = [f"questao_linguagem_{i}" for i in range(12, 21) if i != 13]
        # Caderno 2 = Linguagem (Q12 a Q20, exceto Q13)
        campos_caderno2 = [f"questao_matematica_{i}" for i in range(1, 11)]

        # Pretendidos = total de alunos avaliados
        c1_pret = c2_pret = registros.count()

        # Efetivados = pelo menos 1 questão respondida (com ou sem laudo)
        def respondeu_em(campo_list, aluno):
            return any(getattr(aluno, campo) not in [None, '', 'FALTOU'] for campo in campo_list)

        c1_efet = c2_efet = 0
        for aluno in registros:
            if respondeu_em(campos_caderno1, aluno):
                c1_efet += 1
            if respondeu_em(campos_caderno2, aluno):
                c2_efet += 1


        total_c1_pret += c1_pret
        total_c1_efet += c1_efet
        total_c2_pret += c1_pret
        total_c2_efet += c2_efet

        valores = [escola, c2_pret, c2_efet, c1_pret, c1_efet]
        for j, valor in enumerate(valores):
            cell = table_final.cell(i, j)
            cell.text = str(valor)
            cell.text_frame.paragraphs[0].font.size = Pt(9)

    totais = ["TOTAL", total_c2_pret, total_c2_efet, total_c1_pret, total_c1_efet]
    for j, valor in enumerate(totais):
        cell = table_final.cell(num_escolas + 1, j)
        cell.text = str(valor)
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        cell.text_frame.paragraphs[0].font.bold = True


    # ✅ SLIDE FINAL 2: AVALIAÇÃO DAS HABILIDADES (CERTO, ERRADO, etc.)
    slide_aval = prs.slides.add_slide(prs.slide_layouts[5])
    slide_aval.shapes.title.text = "Avaliação das Habilidades - Linguagem"

    habilidades_config = {
        'EI05LGH1: GRAFAR O NOME PRÓPRIO COMPLETO COM AUXÍLIO DO CRACHÁ': 'questao_linguagem_11',
        'EI05LGH2: IDENTIFICAR AS LETRAS DO NOME NO ALFABETO': 'questao_linguagem_12',
        'EI05LGH3: ESCREVER PALAVRAS FAZENDO AJUSTE DO FALADO COM O ESCRITO': 'questao_linguagem_14',
        'EI05LGH4: DIFERENCIAR LETRAS DE OUTROS SINAIS GRÁFICOS COMO NÚMEROS, DESENHOS E SINAIS': 'questao_linguagem_15',
        'EI05LGH5: IDENTIFICAR E ESCREVER NA TABELA AS LETRAS DO ALFABETO QUE ESTÃO FALTANDO': 'questao_linguagem_16',
        'EI05LGH6: RECONHECER, A PARTIR DA PALAVRA OUVIDA E DA IMAGEM, O VALOR SONORO DE UMA LETRA': 'questao_linguagem_17',
        'EI05LGH7: RECONHECER, A PARTIR DA PALAVRA OUVIDA E DA IMAGEM, O VALOR SONORO DE UMA LETRA': 'questao_linguagem_18',
        'EI05LGH8: RELACIONAR O SIGNIFICANTE AO SIGNIFICADO COM A LETRA INICIAL VARIADA': 'questao_linguagem_19',
        'EI05LGH9: IDENTIFICAR O NÚMERO DE LETRAS DE UMA PALAVRA, TENDO EM VISTA A ESCRITA PARA CONTAGEM': 'questao_linguagem_20',
        'EI05LGH10: PERSONAGENS E SUAS AÇÕES EM CONTO CONSIDERANDO CONHECIDO': 'questao_linguagem_13',
    }

    headers_avaliacao = ['Habilidade', 'CERTO', 'ERRADO', 'PARCIAL', 'BRANCO', 'CCL']
    num_rows = len(habilidades_config) + 1
    num_cols = len(headers_avaliacao)

    # Criação da tabela
    table_aval = slide_aval.shapes.add_table(
        num_rows, num_cols,
        Inches(0.3), Inches(1.9), Inches(9), Inches(4.5)
    ).table

    # Ajuste de largura das colunas
    table_aval.columns[0].width = Inches(4.0)  # Descrição maior
    for col_idx in range(1, num_cols):
        table_aval.columns[col_idx].width = Inches((9 - 4.0) / (num_cols - 1))  # Divide o restante igualmente

    # Cabeçalhos
    for j, header in enumerate(headers_avaliacao):
        cell = table_aval.cell(0, j)
        cell.text = header
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(10)

    # Dados
    for i, (habilidade, campo) in enumerate(habilidades_config.items(), start=1):
        # Coluna 0: Habilidade
        cell_desc = table_aval.cell(i, 0)
        cell_desc.text = habilidade
        cell_desc.text_frame.paragraphs[0].font.size = Pt(8)

        # Demais colunas
        valores = [
            alunos.filter(**{campo: 'CERTO'}).count(),
            alunos.filter(**{campo: 'ERRADO'}).count(),
            alunos.filter(**{campo: 'PARCIAL'}).count(),
            alunos.filter(**{campo: 'BRANCO'}).count(),
            alunos.filter(**{campo: 'CRIANÇA COM LAUDO'}).count()
        ]
        for j, valor in enumerate(valores, start=1):
            table_aval.cell(i, j).text = str(valor)
            table_aval.cell(i, j).text_frame.paragraphs[0].font.size = Pt(9)




    # ✅ SLIDE FINAL 3: HABILIDADES POR MODALIDADE (CC, SC, PSI...)
    slide_modal = prs.slides.add_slide(prs.slide_layouts[5])
    slide_modal.shapes.title.text = "Habilidades Específicas por Modalidade"

    modalidades = ['CC', 'SC', 'PSI', 'PSII', 'SSVS', 'SCVS', 'SA', 'ALF']
    num_rows = len(habilidades_config) + 1
    num_cols = len(modalidades) + 1  # +1 para coluna "Hab."

    # Tabela
    table_mod = slide_modal.shapes.add_table(
        num_rows, num_cols,
        Inches(0.5), Inches(1.9), Inches(8.5), Inches(4.5)
    ).table

    # Ajuste de largura de colunas
    table_mod.columns[0].width = Inches(4.0)  # Coluna de descrição
    for j in range(1, num_cols):
        table_mod.columns[j].width = Inches((8.5 - 4.0) / (num_cols - 1))  # Restante dividido igualmente

    # Cabeçalho
    table_mod.cell(0, 0).text = "Hab."
    table_mod.cell(0, 0).text_frame.paragraphs[0].font.bold = True
    table_mod.cell(0, 0).text_frame.paragraphs[0].font.size = Pt(10)

    for j, mod in enumerate(modalidades, start=1):
        cell = table_mod.cell(0, j)
        cell.text = mod
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(9)

    # Preenchimento dos dados
    for i, (descricao, campo) in enumerate(habilidades_config.items(), start=1):
        # Coluna 0: descrição da habilidade
        cell_desc = table_mod.cell(i, 0)
        cell_desc.text = descricao
        cell_desc.text_frame.paragraphs[0].font.size = Pt(8)

        for j, mod in enumerate(modalidades, start=1):
            count = alunos.filter(**{campo: mod}).count()
            cell = table_mod.cell(i, j)
            cell.text = str(count)
            cell.text_frame.paragraphs[0].font.size = Pt(9)



        # ✅ Geração dos gráficos (fora do loop)
        grafico_q11 = gerar_grafico_base64(
            campo='questao_linguagem_11',
            categorias=['CC', 'SC', 'BRANCO', 'CRIANÇA COM LAUDO', 'FALTOU'],
            alunos=alunos,
            titulo='Q11 por Escola'
        )

        grafico_q13 = gerar_grafico_base64(
            campo='questao_linguagem_13',
            categorias=['PSI', 'PSII', 'SSVS', 'SCVS', 'SA', 'ALF', 'BRANCO', 'CRIANÇA COM LAUDO', 'FALTOU'],
            alunos=alunos,
            titulo='Q13 por Escola'
        )

        # ✅ Inserção dos gráficos no PowerPoint
        def adicionar_slide_com_imagem(prs, titulo, imagem_base64):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = titulo
            imagem_bytes = BytesIO(base64.b64decode(imagem_base64))
            slide.shapes.add_picture(imagem_bytes, Inches(0.5), Inches(1.5), height=Inches(4.5))

        adicionar_slide_com_imagem(prs, "📊 Gráfico Q11 por Escola", grafico_q11)
        adicionar_slide_com_imagem(prs, "📊 Gráfico Q13 por Escola", grafico_q13)


        # ✅ SLIDE FINAL 4: AVALIAÇÃO DAS HABILIDADES - MATEMÁTICA
    

    slide_aval_mat = prs.slides.add_slide(prs.slide_layouts[5])
    slide_aval_mat.shapes.title.text = "Avaliação das Habilidades - Matemática"

    habilidades_matematica = {
        'EI05CMH1: IDENTIFICAR OS NUMERAIS E GRAFAR ATÉ 10': 'questao_matematica_1',
        'EI05CMH2: ASSOCIAR UMA QUANTIDADE DE OBJETOS A UM NÚMERO NATURAL': 'questao_matematica_2',
        'EI05CMH3: ESTABELECER NOÇÕES DE ADIÇÃO': 'questao_matematica_3',
        'EI05CMH4: AGRUPAR ELEMENTOS PELAS SUAS CARACTERÍSTICAS SIMILARES': 'questao_matematica_4',
        'EI05CMH5: NOÇÕES DE MEDIDAS DE TEMPO (DIA E NOITE)': 'questao_matematica_5',
        'EI05CMH6: CONSTRUIR CONCEITOS MATEMÁTICOS DE COMPRIMENTO: ALTO, BAIXO': 'questao_matematica_6',
        'EI05CMH7: CONSTRUIR CONCEITOS MATEMÁTICOS DE GRANDEZA: (MAIOR, MENOR, FINO, GROSSO)': 'questao_matematica_7',
        'EI05CMH8: NOÇÕES DE POSIÇÃO (FRENTE, ATRÁS, DIREITA, ESQUERDA, PRIMEIRO, ÚLTIMO)': 'questao_matematica_8',
        'EI05CMH9: CLASSIFICAÇÃO DAS FORMAS GEOMÉTRICAS PLANAS (CÍRCULO, QUADRADO, RETÂNGULO E TRIÂNGULO)': 'questao_matematica_9',
        'EI05CMH10: NOÇÕES DE ESTATÍSTICA: LEITURA DE GRÁFICOS SIMPLES': 'questao_matematica_10',
    }

    headers_mat = ['Habilidade', 'CERTO', 'ERRADO', 'PARCIAL', 'BRANCO', 'CCL']
    table_mat = slide_aval_mat.shapes.add_table(
        len(habilidades_matematica)+1, len(headers_mat),
        Inches(0.3), Inches(1.9), Inches(9), Inches(4.5)
    ).table

    # ✅ Ajustar larguras das colunas
    table_mat.columns[0].width = Inches(4.0)  # Habilidade
    for j in range(1, len(headers_mat)):
        table_mat.columns[j].width = Inches(1.0)  # CERTO a CCL (5 colunas restantes)

    # Cabeçalhos
    for j, header in enumerate(headers_mat):
        cell = table_mat.cell(0, j)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)

    # Linhas de dados
    for i, (habilidade, campo) in enumerate(habilidades_matematica.items(), start=1):
        cell_hab = table_mat.cell(i, 0)
        cell_hab.text = habilidade
        cell_hab.text_frame.paragraphs[0].font.size = Pt(8)  # Fonte menor para habilidade longa

        # Colunas de resposta
        table_mat.cell(i, 1).text = str(alunos.filter(**{campo: 'CERTO'}).count())
        table_mat.cell(i, 2).text = str(alunos.filter(**{campo: 'ERRADO'}).count())
        table_mat.cell(i, 3).text = str(alunos.filter(**{campo: 'PARCIAL'}).count())
        table_mat.cell(i, 4).text = str(alunos.filter(**{campo: 'BRANCO'}).count())
        table_mat.cell(i, 5).text = str(alunos.filter(**{campo: 'CRIANÇA COM LAUDO'}).count())




        # ✅ GRÁFICO DE MATEMÁTICA - POR HABILIDADE
        # ✅ GRÁFICO DE MATEMÁTICA - POR HABILIDADE COM RÓTULOS
    habilidades_codigos = [key.split(':')[0].strip() for key in habilidades_matematica.keys()]
    categorias = ['CERTO', 'ERRADO', 'PARCIAL', 'BRANCO', 'CRIANÇA COM LAUDO']
    dados = {categoria: [] for categoria in categorias}

    for cod, campo in habilidades_matematica.items():
        for categoria in categorias:
            dados[categoria].append(alunos.filter(**{campo: categoria}).count())

    fig, ax = plt.subplots(figsize=(12, 6))
    x = list(range(len(habilidades_codigos)))
    bottom = [0] * len(habilidades_codigos)

    for cat in categorias:
        valores = dados[cat]
        bars = ax.bar(x, valores, bottom=bottom, label=cat)
        ax.bar_label(bars, labels=[str(v) for v in valores], fontsize=8, padding=2)
        bottom = [i + j for i, j in zip(bottom, valores)]

    ax.set_xticks(x)
    ax.set_xticklabels(habilidades_codigos, rotation=45, fontsize=9)
    ax.set_title("📊 Respostas por Habilidade - Matemática")
    ax.set_ylabel("Total")
    ax.legend()
    plt.tight_layout()

    buffer = BytesIO()
    fig.savefig(buffer, format='png')
    plt.close(fig)
    img_base64 = base64.b64encode(buffer.getvalue()).decode()

    adicionar_slide_com_imagem(prs, "📊 Gráfico Geral - Habilidades de Matemática", img_base64)



    # EXPORTAÇÃO DO ARQUIVO .pptx
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    response = HttpResponse(
        buffer.read(),
        content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )
    response['Content-Disposition'] = 'attachment; filename=relatorio_samacc_2025.pptx'
    return response

###***********************************************************************************************************************

from django.http import HttpResponse
from django.template.loader import get_template
from django.templatetags.static import static
from xhtml2pdf import pisa
from .models import CadastroEI

def gerar_pdf_relatorio_geral(request):
    nome_escola = request.GET.get('nome_escola')
    ano = request.GET.get('ano')
    modalidade = request.GET.get('modalidade')

    alunos = CadastroEI.objects.filter(avaliado='SIM')
    if nome_escola:
        alunos = alunos.filter(unidade_ensino__icontains=nome_escola)
    if ano:
        alunos = alunos.filter(ano=ano)
    if modalidade:
        alunos = alunos.filter(modalidade__icontains=modalidade)

    unidades = alunos.values_list('unidade_ensino', flat=True).distinct()
    campos_linguagem = [f"questao_linguagem_{i}" for i in range(12, 21) if i != 13]
    campos_matematica = [f"questao_matematica_{i}" for i in range(1, 11)]

    tabela_dados = []
    for unidade in unidades:
        registros = alunos.filter(unidade_ensino=unidade)

        registros_matematica_validos = [
            aluno for aluno in registros
            if not all(getattr(aluno, campo) in [None, '', 'FALTOU'] for campo in campos_matematica)
        ]
        registros_linguagem_validos = [
            aluno for aluno in registros
            if not all(getattr(aluno, campo) in [None, '', 'FALTOU'] for campo in campos_linguagem)
        ]

        total_pretendidos = registros.count()
        tabela_dados.append({
            'nome': unidade,
            'pretendidos_1': total_pretendidos,
            'efetivados_1': len(registros_matematica_validos),
            'pretendidos_2': total_pretendidos,
            'efetivados_2': len(registros_linguagem_validos),
        })

    context = {
        'tabela_dados': tabela_dados,
        'cabecalho_superior': request.build_absolute_uri(static('assets/dist/img/logocabecalho.png')),
        'cabecalho_inferior': request.build_absolute_uri(static('assets/dist/img/logorodape.png')),
    }

    template = get_template('relatorios/relatorio_geral.html')
    html = template.render(context)

    response = HttpResponse(content_type='application/pdf')
    response['Content-Disposition'] = 'inline; filename="relatorio_geral.pdf"'

    pisa_status = pisa.CreatePDF(html, dest=response)
    if pisa_status.err:
        return HttpResponse(f'Erro ao gerar PDF: <pre>{html}</pre>')

    return response


###***********************************************************************************************************************

from django.http import HttpResponse
from django.template.loader import get_template
from django.templatetags.static import static
from xhtml2pdf import pisa
from .models import CadastroEI

def gerar_pdf_habilidades_especificas(request):
    habilidades = {
        'EI05LGH1': {'descricao': 'Grafar o nome próprio com auxílio do crachá.', 'campo': 'questao_linguagem_11'},
        'EI05LGH2': {'descricao': 'Identificar as letras do nome no alfabeto.', 'campo': 'questao_linguagem_12'},
        'EI05LGH3': {'descricao': 'Escrever palavras fazendo ajuste do falado com o escrito.', 'campo': 'questao_linguagem_14'},
        'EI05LGH4': {'descricao': 'Diferenciar letras de outros sinais gráficos como números, desenhos e sinais.', 'campo': 'questao_linguagem_15'},
        'EI05LGH5': {'descricao': 'Identificar e escrever na tabela as letras do alfabeto que estão faltando.', 'campo': 'questao_linguagem_16'},
        'EI05LGH6': {'descricao': 'Reconhecer o valor sonoro da letra inicial da palavra.', 'campo': 'questao_linguagem_17'},
        'EI05LGH7': {'descricao': 'Reconhecer o valor sonoro da letra final da palavra.', 'campo': 'questao_linguagem_18'},
        'EI05LGH8': {'descricao': 'Relacionar o significante ao significado (imagem/palavra).', 'campo': 'questao_linguagem_19'},
        'EI05LGH9': {'descricao': 'Identificar o número de letras de uma palavra para contagem.', 'campo': 'questao_linguagem_20'},
        'EI05LGH10': {'descricao': 'Identificar personagens e suas ações em conto conhecido.', 'campo': 'questao_linguagem_13'},
    }

    modalidades = ['CC', 'SC', 'PSI', 'PSII', 'SSVS', 'SCVS', 'SA', 'ALF', 'BRANCO', 'CCL']
    avaliacoes = ['CERTO', 'ERRADO', 'PARCIAL', 'BRANCO', 'CCL']
    mapeamento_valores = {'CCL': 'CRIANÇA COM LAUDO'}

    dados = {}
    for cod, info in habilidades.items():
        campo = info['campo']

        totais_modalidade = {
            modalidade: CadastroEI.objects.filter(**{campo: mapeamento_valores.get(modalidade, modalidade)}).count()
            for modalidade in modalidades
        }

        totais_avaliacao = {
            avaliacao: CadastroEI.objects.filter(**{campo: mapeamento_valores.get(avaliacao, avaliacao)}).count()
            for avaliacao in avaliacoes
        }

        if any(totais_modalidade.values()) or any(totais_avaliacao.values()):
            dados[cod] = {
                'descricao': info['descricao'],
                'totais_modalidade': totais_modalidade,
                'totais_avaliacao': totais_avaliacao,
            }

    # Caminhos públicos usando static
    cabecalho_superior = request.build_absolute_uri(static('assets/dist/img/cabecalho_superior.png'))
    cabecalho_inferior = request.build_absolute_uri(static('assets/dist/img/cabecalho_inferior.png'))

    context = {
        'dados': dados,
        'modalidades': modalidades,
        'avaliacoes': avaliacoes,
        'cabecalho_superior': cabecalho_superior,
        'cabecalho_inferior': cabecalho_inferior,
    }

    template_path = 'relatorios/relatorio_habilidades_especificas.html'
    template = get_template(template_path)
    html = template.render(context)

    response = HttpResponse(content_type='application/pdf')
    response['Content-Disposition'] = 'inline; filename="relatorio_habilidades_especificas.pdf"'

    pisa_status = pisa.CreatePDF(html, dest=response)
    if pisa_status.err:
        return HttpResponse(f'Erro ao gerar PDF: <pre>{html}</pre>')

    return response

###***********************************************************************************************************************

import matplotlib.pyplot as plt
import tempfile
import os
from .models import CadastroEI
from weasyprint import HTML
from django.template.loader import render_to_string
from django.http import HttpResponse


def gerar_grafico_habilidade_resposta():
    habilidades = {
        'EI05LGH1': 'questao_linguagem_11',
        'EI05LGH2': 'questao_linguagem_12',
        'EI05LGH3': 'questao_linguagem_14',
        'EI05LGH4': 'questao_linguagem_15',
        'EI05LGH5': 'questao_linguagem_16',
    }

    # Nomes para legenda
    avaliacoes_exibicao = ['CERTO', 'ERRADO', 'PARCIAL', 'BRANCO', 'CCL']

    # Mapeamento interno do valor para consulta no banco
    mapeamento = {
        'CCL': 'CRIANÇA COM LAUDO'
    }

    labels = list(habilidades.keys())
    data_por_avaliacao = {aval: [] for aval in avaliacoes_exibicao}

    for cod, campo in habilidades.items():
        for aval in avaliacoes_exibicao:
            valor_banco = mapeamento.get(aval, aval)
            total = CadastroEI.objects.filter(**{campo: valor_banco}).count()
            data_por_avaliacao[aval].append(total)

    # Gerar gráfico
    fig, ax = plt.subplots(figsize=(10, 6))
    bottom = [0] * len(labels)

    for aval in avaliacoes_exibicao:
        valores = data_por_avaliacao[aval]
        bars = ax.bar(labels, valores, bottom=bottom, label=aval)
        bottom = [sum(x) for x in zip(bottom, valores)]

    ax.set_title("Respostas por Habilidade", fontsize=14)
    ax.set_ylabel("Total")
    ax.legend()

    temp_img = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    plt.tight_layout()
    plt.savefig(temp_img.name, dpi=150)
    plt.close(fig)

    return temp_img.name  # retorna o caminho da imagem


from django.http import HttpResponse
from django.template.loader import render_to_string
from weasyprint import HTML
import tempfile, os

def gerar_pdf_grafico_habilidade_resposta(request):
    from .views import gerar_grafico_habilidade_resposta

    imagem_grafico = gerar_grafico_habilidade_resposta()

    contexto = {
        'imagem_grafico': f"file://{imagem_grafico}",
        'cabecalho_path': 'file:///C:/Users/sysco/Desktop/SIEDGE/semedweb/static/assets/dist/img/logocabecalho.png',
        'rodape_path': 'file:///C:/Users/sysco/Desktop/SIEDGE/semedweb/static/assets/dist/img/logorodape.png',

    }

    html_string = render_to_string('relatorios/relatorio_grafico_habilidades.html', contexto)

    with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as output:
        HTML(string=html_string, base_url="file:///C:/Users/sysco/Desktop/SIEDGE/semedweb/").write_pdf(target=output.name)
        output.seek(0)
        response = HttpResponse(output.read(), content_type='application/pdf')
        response['Content-Disposition'] = 'inline; filename="grafico_habilidades.pdf"'

    os.remove(imagem_grafico)
    os.remove(output.name)
    return response

###***********************************************************************************************************************

def gerar_pdf_grafico_habilidades(request):
    caminho_imagem = gerar_grafico_habilidade_resposta()

    contexto = {
        'imagem_grafico': caminho_imagem,
    }

    html_string = render_to_string('relatorios/relatorio_com_grafico.html', contexto)
    fd, path = tempfile.mkstemp(suffix='.pdf')
    os.close(fd)

    try:
        HTML(string=html_string, base_url='.').write_pdf(path)
        with open(path, 'rb') as pdf:
            response = HttpResponse(pdf.read(), content_type='application/pdf')
            response['Content-Disposition'] = 'inline; filename="grafico_habilidades.pdf"'
            return response
    finally:
        os.remove(path)
        os.remove(caminho_imagem)  # remove a imagem temporária
###***********************************************************************************************************************

from io import BytesIO
from django.http import HttpResponse
from django.template.loader import get_template
from django.templatetags.static import static
from xhtml2pdf import pisa
import matplotlib.pyplot as plt
import base64
import tempfile
from .models import CadastroEI

def gerar_graficos_pdf(request):
    nome_escola = request.GET.get('nome_escola', '')
    ano = request.GET.get('ano', '')
    modalidade = request.GET.get('modalidade', '')

    alunos = CadastroEI.objects.filter(avaliado='SIM')
    if nome_escola:
        alunos = alunos.filter(unidade_ensino__icontains=nome_escola)
    if ano:
        alunos = alunos.filter(ano=ano)
    if modalidade:
        alunos = alunos.filter(modalidade__icontains=modalidade)

    escolas = alunos.values_list('unidade_ensino', flat=True).distinct()
    respostas = ['CERTO', 'ERRADO', 'PARCIAL', 'BRANCO', 'CCL']
    mapeamento = {'CCL': 'CRIANÇA COM LAUDO'}

    graficos = []

    for escola in escolas:
        dados = {}
        for r in respostas:
            valor_db = mapeamento.get(r, r)
            dados[r] = alunos.filter(unidade_ensino=escola, questao_linguagem_11=valor_db).count()

        fig, ax = plt.subplots()
        ax.bar(dados.keys(), dados.values(), color='skyblue')
        ax.set_title(f'{escola} - Linguagem Q11')
        ax.set_ylim(0, max(dados.values()) + 5)
        for i, v in enumerate(dados.values()):
            ax.text(i, v + 0.5, str(v), ha='center', fontsize=8)

        buffer = BytesIO()
        plt.tight_layout()
        fig.savefig(buffer, format='png')
        plt.close(fig)

        imagem_base64 = base64.b64encode(buffer.getvalue()).decode()
        graficos.append({'escola': escola, 'grafico': imagem_base64})

    # Usa caminho absoluto para imagens como no exemplo
    cabecalho_superior = request.build_absolute_uri(static('assets/dist/img/logocabecalho.png'))
    cabecalho_inferior = request.build_absolute_uri(static('assets/dist/img/logorodape.png'))

    context = {
        'graficos': graficos,
        'cabecalho_superior': cabecalho_superior,
        'cabecalho_inferior': cabecalho_inferior,
    }

    template = get_template('relatorios/relatorio_graficos_pisa.html')
    html_string = template.render(context)

    response = HttpResponse(content_type='application/pdf')
    response['Content-Disposition'] = 'inline; filename="graficos_linguagem_q11.pdf"'

    pisa_status = pisa.CreatePDF(html_string, dest=response)
    if pisa_status.err:
        return HttpResponse('Erro ao gerar PDF <pre>' + html_string + '</pre>')

    return response

###***********************************************************************************************************************

from django.contrib.staticfiles.storage import staticfiles_storage
from django.template.loader import render_to_string
from django.http import HttpResponse
from weasyprint import HTML
from .models import CadastroEI
import matplotlib.pyplot as plt
import base64
from io import BytesIO
import tempfile, os

def gerar_graficos_q11_q13(request):
    nome_escola = request.GET.get('nome_escola', '')
    ano = request.GET.get('ano', '')
    modalidade = request.GET.get('modalidade', '')

    alunos_base = CadastroEI.objects.filter(avaliado='SIM')
    if nome_escola:
        alunos_base = alunos_base.filter(unidade_ensino__icontains=nome_escola)
    if ano:
        alunos_base = alunos_base.filter(ano=ano)
    if modalidade:
        alunos_base = alunos_base.filter(modalidade__icontains=modalidade)

    escolas = alunos_base.values_list('unidade_ensino', flat=True).distinct()

    respostas_q11_display = ['CC', 'SC', 'BRANCO', 'CCL', 'FALTOU']
    respostas_q13_display = ['PSI', 'PSII', 'SSVS', 'SCVS', 'SA', 'ALF', 'BRANCO', 'CCL', 'FALTOU']
    mapeamento_valores = {'CCL': 'CRIANÇA COM LAUDO'}

    graficos = []

    for escola in escolas:
        alunos = alunos_base.filter(unidade_ensino=escola)

        # Gráfico Q11
        dados_q11 = [alunos.filter(questao_linguagem_11=mapeamento_valores.get(r, r)).count()
                     for r in respostas_q11_display]
        fig1, ax1 = plt.subplots(figsize=(10, 4.5))
        ax1.bar(respostas_q11_display, dados_q11, color='skyblue')
        ax1.set_title(f'{escola} - Q1 (Com/sem crachá)', fontsize=12)
        ax1.set_ylim(0, max(dados_q11 + [1]) + 5)
        ax1.tick_params(axis='x', labelrotation=30, labelsize=9)
        for i, v in enumerate(dados_q11):
            ax1.text(i, v + 0.5, str(v), ha='center', fontsize=9)
        buffer1 = BytesIO()
        plt.tight_layout()
        fig1.savefig(buffer1, format='png')
        plt.close(fig1)
        imagem_q11 = base64.b64encode(buffer1.getvalue()).decode()

        # Gráfico Q13
        dados_q13 = [alunos.filter(questao_linguagem_13=mapeamento_valores.get(r, r)).count()
                     for r in respostas_q13_display]
        fig2, ax2 = plt.subplots(figsize=(10, 4.5))
        ax2.bar(respostas_q13_display, dados_q13, color='salmon')
        ax2.set_title(f'{escola} - Q3 (Níveis de Escrita)', fontsize=12)
        ax2.set_ylim(0, max(dados_q13 + [1]) + 5)
        ax2.tick_params(axis='x', labelrotation=30, labelsize=9)
        for i, v in enumerate(dados_q13):
            ax2.text(i, v + 0.5, str(v), ha='center', fontsize=9)
        buffer2 = BytesIO()
        plt.tight_layout()
        fig2.savefig(buffer2, format='png')
        plt.close(fig2)
        imagem_q13 = base64.b64encode(buffer2.getvalue()).decode()

        graficos.append({
            'escola': escola,
            'grafico_q11': imagem_q11,
            'grafico_q13': imagem_q13,
        })

    context = {
        'graficos': graficos,
        'cabecalho_path': staticfiles_storage.url('assets/dist/img/logocabecalho.png'),
        'rodape_path': staticfiles_storage.url('assets/dist/img/logorodape.png'),
    }

    html_string = render_to_string('relatorios/relatorio_graficos_q11_q13.html', context)

    fd, path = tempfile.mkstemp(suffix='.pdf')
    os.close(fd)

    try:
        HTML(string=html_string, base_url=request.build_absolute_uri('/')).write_pdf(path)
        with open(path, 'rb') as pdf:
            return HttpResponse(pdf.read(), content_type='application/pdf')
    finally:
        os.remove(path)

###***********************************************************************************************************************

from django.http import HttpResponse
from django.template.loader import render_to_string
from django.templatetags.static import static
from weasyprint import HTML
from .models import CadastroEI
from django.db.models import Count
import tempfile, os

def relatorio_saida_matematica_view(request):
    nome_escola = request.GET.get('nome_escola', '')
    ano = request.GET.get('ano', '')
    modalidade = request.GET.get('modalidade', '')

    alunos = CadastroEI.objects.filter(avaliado="SIM")
    if nome_escola:
        alunos = alunos.filter(unidade_ensino__icontains=nome_escola)
    if ano:
        alunos = alunos.filter(ano=ano)
    if modalidade:
        alunos = alunos.filter(modalidade__icontains=modalidade)

    escolas = alunos.values_list("unidade_ensino", flat=True).distinct()

    resultado_matematica = {}
    for escola in escolas:
        alunos_escola = alunos.filter(unidade_ensino=escola)
        questoes = {}
        for i in range(1, 11):
            campo = f'questao_matematica_{i}'
            contagem = alunos_escola.values(campo).annotate(total=Count(campo))
            questoes[i] = {
                'CERTO': sum(c['total'] for c in contagem if c[campo] == 'CERTO'),
                'ERRADO': sum(c['total'] for c in contagem if c[campo] == 'ERRADO'),
                'PARCIAL': sum(c['total'] for c in contagem if c[campo] == 'PARCIAL'),
                'BRANCO': sum(c['total'] for c in contagem if c[campo] == 'BRANCO'),
                'LAUDO': sum(c['total'] for c in contagem if c[campo] == 'CRIANÇA COM LAUDO'),
                'FALTOU': sum(c['total'] for c in contagem if c[campo] == 'FALTOU'),
            }
        resultado_matematica[escola] = questoes

    # Gera URLs públicas a partir do domínio da aplicação
    cabecalho_url = request.build_absolute_uri(static('assets/dist/img/logocabecalho.png'))
    rodape_url = request.build_absolute_uri(static('assets/dist/img/logorodape.png'))

    context = {
        'titulo': "Relatório de Saída - SAMACC 2025 - Matemática",
        'resultado_matematica': resultado_matematica,
        'cabecalho_path': cabecalho_url,
        'rodape_path': rodape_url,
    }

    html_string = render_to_string('relatorios/relatorio_saida_matematica.html', context)

    fd, path = tempfile.mkstemp(suffix='.pdf')
    os.close(fd)

    try:
        HTML(string=html_string, base_url=request.build_absolute_uri('/')).write_pdf(path)
        with open(path, 'rb') as pdf_file:
            response = HttpResponse(pdf_file.read(), content_type='application/pdf')
            response['Content-Disposition'] = 'inline; filename="relatorio_saida_matematica.pdf"'
            return response
    finally:
        os.remove(path)

from django.http import JsonResponse

def get_sintese_dinamica(request):
    # mesmo filtro da view acima
    receitas = Receita.objects.all()
    if request.GET.get("escola"):
        receitas = receitas.filter(escola_id=request.GET.get("escola"))
    if request.GET.get("ano"):
        receitas = receitas.filter(data_inicio__year=request.GET.get("ano"))
    if request.GET.get("programa"):
        receitas = receitas.filter(programa=request.GET.get("programa"))

    dados = receitas.aggregate(
        receita_total=Sum('valor_total_receita_custeio') + Sum('valor_total_receita_capital'),
        despesa_total=Sum('valor_despesa_realizada_custeio') + Sum('valor_despesa_realizada_capital'),
        saldo_reprogramado=Sum('saldo_reprogramar_custeio') + Sum('saldo_reprogramar_capital'),
        recursos_proprios=Sum('recursos_proprios_custeio') + Sum('recursos_proprios_capital'),
        rendimentos_aplicacao=Sum('rendimento_aplicacao_custeio') + Sum('rendimento_aplicacao_capital'),
        devolucoes_fnde=Sum('devolucao_fnde_custeio') + Sum('devolucao_fnde_capital'),
        escolas_atendidas=receitas.values('escola').distinct().count(),
    )

    return JsonResponse(dados)



def sintese_pdde_view(request):
    return render(request, 'pdde/sintese_pdde.html')
